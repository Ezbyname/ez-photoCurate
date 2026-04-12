"""
E-z Photo Organizer Web UI — Step-by-step wizard for building photo collections.
Runs locally, opens in browser.

Usage:
    python app.py
    python app.py --port 5050
"""

import os
import sys
import json
import threading
import webbrowser
import argparse
from datetime import datetime
from collections import defaultdict

import time
import base64
import shutil
import numpy as np


try:
    import clip_engine
except ImportError:
    clip_engine = None


class NumpyEncoder(json.JSONEncoder):
    """Handle numpy types when serializing to JSON."""
    def default(self, obj):
        if isinstance(obj, (np.integer,)):
            return int(obj)
        if isinstance(obj, (np.floating,)):
            return float(obj)
        if isinstance(obj, np.ndarray):
            return obj.tolist()
        return super().default(obj)

from flask import Flask, request, jsonify, send_file, session, redirect

sys.stdout.reconfigure(line_buffering=True)

PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECTS_DIR = os.path.join(PROJECT_DIR, "projects")
SCAN_DB_PATH = os.path.join(PROJECT_DIR, "scan_db.json")
CLIP_VECTORS_PATH = os.path.join(PROJECT_DIR, "clip_vectors.npz")
IMAGE_VECTORS_PATH = os.path.join(PROJECT_DIR, "image_vectors.npz")
FACE_ENCODINGS_PATH = os.path.join(PROJECT_DIR, "face_encodings.npz")
CONFIG_PATH = os.path.join(PROJECT_DIR, "curate_config.json")

app = Flask(__name__, static_folder=None)
app.secret_key = os.environ.get("SECRET_KEY", "ez-photo-organizer-dev-key-change-in-prod")
app.config["PERMANENT_SESSION_LIFETIME"] = 60 * 60 * 24 * 30  # 30 days

# ── Auth ─────────────────────────────────────────────────────────────────────
from auth import init_db, register_auth_routes, login_required
init_db()
register_auth_routes(app)

@app.before_request
def require_auth():
    """Protect all routes except login/auth endpoints."""
    open_paths = ("/login", "/api/auth/")
    if any(request.path == p or request.path.startswith(p) for p in open_paths):
        return None
    if not session.get("user"):
        if request.path.startswith("/api/"):
            return jsonify({"error": "Not authenticated"}), 401
        return redirect("/login")
    return None

# Lock for scan_db.json reads/writes
_db_lock = threading.Lock()

# ── Background task state (multi-job) ─────────────────────────────────────────
import uuid as _uuid

_jobs = {}  # job_id -> job dict
_jobs_lock = threading.Lock()
# _current_job_id tracks the "active" job for legacy single-task API compat
_current_job_id = None

# Legacy single-task shim — points to current job or a dummy
_task = {"running": False, "type": None, "progress": "", "lines": [], "done": False, "error": None, "cancelled": False}
_task_lock = threading.Lock()


def _create_job(task_type, project_name=None):
    """Create a new background job and return its ID."""
    global _current_job_id
    job_id = _uuid.uuid4().hex[:8]
    job = {
        "id": job_id,
        "running": True,
        "type": task_type,
        "project_name": project_name or "",
        "progress": "Starting...",
        "percent": 0,
        "lines": [],
        "done": False,
        "error": None,
        "cancelled": False,
        "started_at": datetime.now().isoformat(),
    }
    with _jobs_lock:
        _jobs[job_id] = job
        _current_job_id = job_id
    # Sync legacy _task
    _sync_legacy_task(job)
    return job_id


def _get_job(job_id):
    with _jobs_lock:
        return _jobs.get(job_id)


def _job_is_cancelled(job_id):
    with _jobs_lock:
        job = _jobs.get(job_id)
        return job["cancelled"] if job else True


def _update_job(job_id, line):
    with _jobs_lock:
        job = _jobs.get(job_id)
        if not job:
            return
        job["progress"] = line
        job["lines"].append(line)
        if len(job["lines"]) > 200:
            job["lines"] = job["lines"][-100:]
    _sync_legacy_task(job)


def _update_job_percent(job_id, percent):
    with _jobs_lock:
        job = _jobs.get(job_id)
        if job:
            job["percent"] = min(100, max(0, int(percent)))


def _finish_job(job_id, error=None):
    with _jobs_lock:
        job = _jobs.get(job_id)
        if not job:
            return
        job["running"] = False
        job["done"] = True
        job["error"] = error
        if not error or error == "Cancelled":
            job["percent"] = 100
    _sync_legacy_task(job)


def _sync_legacy_task(job):
    """Keep legacy _task in sync with the current job for backward compat."""
    with _task_lock:
        _task["running"] = job["running"]
        _task["type"] = job["type"]
        _task["progress"] = job["progress"]
        _task["lines"] = job["lines"]
        _task["done"] = job["done"]
        _task["error"] = job["error"]
        _task["cancelled"] = job["cancelled"]


def _any_job_running():
    with _jobs_lock:
        return any(j["running"] for j in _jobs.values())


# Legacy compat wrappers — used by existing scan/select/export code
def _reset_task(task_type):
    global _current_job_id
    # If called directly (old-style), create a job
    job_id = _current_job_id
    if job_id and _get_job(job_id) and _get_job(job_id)["running"]:
        # Already created via _create_job, just sync
        return
    job_id = _create_job(task_type)


def _is_cancelled():
    job_id = _current_job_id
    if job_id:
        return _job_is_cancelled(job_id)
    with _task_lock:
        return _task["cancelled"]


def _update_task(line):
    job_id = _current_job_id
    if job_id:
        _update_job(job_id, line)
        return
    with _task_lock:
        _task["progress"] = line
        _task["lines"].append(line)
        if len(_task["lines"]) > 200:
            _task["lines"] = _task["lines"][-100:]


def _update_task_percent(percent):
    job_id = _current_job_id
    if job_id:
        _update_job_percent(job_id, percent)


def _finish_task(error=None):
    job_id = _current_job_id
    if job_id:
        _finish_job(job_id, error)
        return
    with _task_lock:
        _task["running"] = False
        _task["done"] = True
        _task["error"] = error


# ── Helpers ───────────────────────────────────────────────────────────────────

def load_config():
    if os.path.isfile(CONFIG_PATH):
        with open(CONFIG_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return None


_active_project_dir = None   # cached path to active project folder
_last_project_sync = 0       # timestamp of last scan_db sync


def _set_active_project(config):
    """Update cached project dir from config. Called only when config changes."""
    global _active_project_dir
    import re
    name = (config or {}).get("project_name", "")
    if not name:
        _active_project_dir = None
        return
    safe = re.sub(r'[<>:"/\\|?*]', '_', name.strip())[:100]
    pdir = os.path.join(PROJECTS_DIR, safe)
    _active_project_dir = pdir if os.path.isdir(pdir) else None


def _auto_save_to_project(config_only=False, scan_db_only=False):
    """Sync working files to active project folder.
    Hot path (scan_db): one timestamp check + throttled copy, no file reads.
    Config saves always sync immediately (infrequent)."""
    global _last_project_sync
    pdir = _active_project_dir
    if not pdir:
        return
    try:
        if not scan_db_only:
            shutil.copy2(CONFIG_PATH, os.path.join(pdir, "curate_config.json"))

        if not config_only:
            now = time.time()
            if now - _last_project_sync < 30:
                return
            # Copy sidecars BEFORE scan_db — on crash, project has
            # newer sidecars + older scan_db (safe: extra vectors harmless)
            if os.path.isfile(IMAGE_VECTORS_PATH):
                shutil.copy2(IMAGE_VECTORS_PATH,
                             os.path.join(pdir, "image_vectors.npz"))
            if os.path.isfile(CLIP_VECTORS_PATH):
                shutil.copy2(CLIP_VECTORS_PATH,
                             os.path.join(pdir, "clip_vectors.npz"))
            if os.path.isfile(FACE_ENCODINGS_PATH):
                shutil.copy2(FACE_ENCODINGS_PATH,
                             os.path.join(pdir, "face_encodings.npz"))
            shutil.copy2(SCAN_DB_PATH, os.path.join(pdir, "scan_db.json"))
            _last_project_sync = now
    except Exception:
        pass


def save_config(config):
    with open(CONFIG_PATH, "w", encoding="utf-8") as f:
        json.dump(config, f, ensure_ascii=False, indent=2)
    _set_active_project(config)
    _auto_save_to_project(config_only=True)


def load_scan_db():
    if os.path.isfile(SCAN_DB_PATH):
        with _db_lock:
            try:
                with open(SCAN_DB_PATH, "r", encoding="utf-8") as f:
                    return json.load(f)
            except (json.JSONDecodeError, ValueError):
                return None
    return None


def _save_image_vectors(images, path=None):
    """Extract image_vector from image entries and save to npz sidecar.
    Returns the number of vectors saved.  Never raises — logs errors so
    callers (save_scan_db) can always proceed to write JSON."""
    path = path or IMAGE_VECTORS_PATH
    vectors = {}
    for img in images:
        v = img.get("image_vector")
        h = img.get("hash")
        if v is not None and h:
            if isinstance(v, list):
                vectors[h] = np.array(v, dtype=np.float32)
            elif isinstance(v, np.ndarray):
                vectors[h] = v.astype(np.float32)
    if not vectors:
        return 0
    # np.savez_compressed appends .npz if filename doesn't end with it,
    # so use a tmp name that already ends with .npz for predictable paths.
    tmp = path + ".tmp.npz"
    try:
        # Merge with existing vectors (don't lose vectors from images not in this save)
        if os.path.isfile(path):
            try:
                existing = dict(np.load(path, allow_pickle=False))
                existing.update(vectors)
                vectors = existing
            except Exception as e:
                print(f"[WARN] Could not merge existing {path} (corrupt?): {e} — "
                      f"saving {len(vectors)} new vectors only", flush=True)
        np.savez_compressed(tmp, **vectors)
        for attempt in range(5):
            try:
                os.replace(tmp, path)
                break
            except PermissionError:
                if attempt < 4:
                    time.sleep(1 + attempt)
                    continue
                print(f"[WARN] Could not replace {path} after 5 retries — "
                      f"vectors not persisted this cycle", flush=True)
                break
    except Exception as e:
        print(f"[WARN] Failed to save image vectors to {path}: {e}", flush=True)
    finally:
        # Clean up stale tmp if replace failed or we crashed mid-write
        try:
            if os.path.isfile(tmp):
                os.remove(tmp)
        except OSError:
            pass
    return len(vectors)


def load_image_vectors(path=None):
    """Load image vectors from npz sidecar. Returns dict {hash: np.array}.
    Returns empty dict (never raises) if file missing or corrupt."""
    path = path or IMAGE_VECTORS_PATH
    if os.path.isfile(path):
        try:
            data = np.load(path, allow_pickle=False)
            return dict(data)
        except Exception as e:
            print(f"[WARN] Could not load {path} (corrupt?): {e} — "
                  f"vectors will be recomputed from files where possible", flush=True)
    return {}


def save_scan_db(db):
    images = db.get("images", [])
    # Save vectors to npz sidecar (before stripping from entries).
    # _save_image_vectors never raises — if npz write fails, JSON still gets saved.
    _save_image_vectors(images)
    tmp_path = SCAN_DB_PATH + ".tmp"
    try:
        with _db_lock:
            with open(tmp_path, "w", encoding="utf-8") as f:
                # Strip image_vector from entries to keep JSON small
                def _strip_vectors(obj):
                    if isinstance(obj, dict) and "images" in obj:
                        stripped = dict(obj)
                        stripped["images"] = [
                            {k: v for k, v in img.items() if k != "image_vector"}
                            for img in obj["images"]
                        ]
                        return stripped
                    return obj
                json.dump(_strip_vectors(db), f, ensure_ascii=False, cls=NumpyEncoder)
            for attempt in range(5):
                try:
                    os.replace(tmp_path, SCAN_DB_PATH)
                    break
                except PermissionError:
                    if attempt < 4:
                        time.sleep(1 + attempt)
                        continue
                    raise
    except Exception:
        # Clean up stale tmp on failure
        try:
            if os.path.isfile(tmp_path):
                os.remove(tmp_path)
        except OSError:
            pass
        raise
    _auto_save_to_project(scan_db_only=True)


def verify_vector_consistency(db=None, repair=False):
    """Check that image_vectors.npz covers all hashes in scan_db.
    Returns dict with diagnostics.  If repair=True, recomputes missing
    vectors for images (videos cannot be recomputed cheaply)."""
    if db is None:
        db = load_scan_db() or {}
    images = db.get("images", [])
    npz_vecs = load_image_vectors()

    total = 0
    covered = 0
    missing_images = []  # can recompute
    missing_videos = []  # cannot recompute
    orphan_keys = set(npz_vecs.keys())

    for img in images:
        h = img.get("hash")
        if not h:
            continue
        total += 1
        if h in npz_vecs:
            covered += 1
            orphan_keys.discard(h)
        else:
            if img.get("media_type") == "video":
                missing_videos.append(h)
            else:
                missing_images.append(h)

    repaired = 0
    if repair and missing_images:
        from curate import compute_image_vector
        new_vecs = {}
        for img in images:
            h = img.get("hash")
            if h and h in missing_images:
                fpath = img.get("path", "").replace("/", os.sep)
                if fpath and os.path.isfile(fpath):
                    try:
                        v = compute_image_vector(fpath)
                        if v is not None:
                            new_vecs[h] = v.astype(np.float32)
                            repaired += 1
                    except Exception:
                        pass
        if new_vecs:
            # Merge repaired vectors into npz
            npz_vecs.update(new_vecs)
            tmp = IMAGE_VECTORS_PATH + ".tmp.npz"
            try:
                np.savez_compressed(tmp, **npz_vecs)
                os.replace(tmp, IMAGE_VECTORS_PATH)
            except Exception as e:
                print(f"[WARN] Could not save repaired vectors: {e}", flush=True)
            finally:
                try:
                    if os.path.isfile(tmp):
                        os.remove(tmp)
                except OSError:
                    pass

    result = {
        "total_entries": total,
        "covered": covered,
        "missing_images": len(missing_images),
        "missing_videos": len(missing_videos),
        "orphan_vectors": len(orphan_keys),
        "repaired": repaired,
        "consistent": len(missing_images) == 0 and len(missing_videos) == 0,
    }
    if missing_images or missing_videos:
        print(f"[WARN] Vector consistency: {len(missing_images)} images + "
              f"{len(missing_videos)} videos missing from npz "
              f"(out of {total} entries)", flush=True)
    return result


def _check_nsfw(fpath, classifier):
    """Check if image contains NSFW content using NudeDetector.
    Returns (is_nsfw: bool, labels: list of detected NSFW class names)."""
    # NudeDetector labels considered inappropriate
    NSFW_LABELS = {
        "FEMALE_BREAST_EXPOSED", "FEMALE_GENITALIA_EXPOSED",
        "MALE_GENITALIA_EXPOSED", "BUTTOCKS_EXPOSED",
        "ANUS_EXPOSED",
    }
    try:
        detections = classifier.detect(fpath)
        found = [d["class"] for d in detections
                 if d["class"] in NSFW_LABELS and d["score"] >= 0.6]
        return bool(found), found
    except Exception:
        return False, []


def _estimate_age(fpath):
    """Estimate the age of the primary face in an image using DeepFace.
    Returns estimated age (int) or None on failure."""
    try:
        from deepface import DeepFace
        results = DeepFace.analyze(
            img_path=fpath, actions=["age"],
            enforce_detection=False, silent=True,
            detector_backend="opencv",
        )
        if isinstance(results, list) and results:
            return results[0].get("age")
        elif isinstance(results, dict):
            return results.get("age")
    except Exception:
        pass
    return None


def _fast_face_detect(fpath, ref_encodings, face_names, tolerance=0.6):
    """Two-stage face detection: thumbnail pre-screen + full-res detection.
    Returns (face_count, faces_found, ok, best_distance, best_encoding)."""
    import face_recognition as fr
    import numpy as np
    from PIL import Image, ImageOps
    try:
        pil_img = ImageOps.exif_transpose(Image.open(fpath)).convert("RGB")
        w, h = pil_img.size

        # Stage A: Quick face check on 256px thumbnail
        thumb = pil_img.copy()
        thumb.thumbnail((256, 256), Image.LANCZOS)
        thumb_arr = np.array(thumb)
        thumb_locs = fr.face_locations(thumb_arr, model="hog")
        if not thumb_locs:
            return 0, [], True, None, None

        # Stage B: Full detection on larger image (only resize if > 1200px)
        max_dim = 1200
        if w > max_dim or h > max_dim:
            pil_img.thumbnail((max_dim, max_dim), Image.LANCZOS)
        arr = np.array(pil_img)
        locations = fr.face_locations(arr, model="hog")
        if not locations:
            return 0, [], True, None, None
        encodings = fr.face_encodings(arr, locations)
        if not encodings:
            return len(locations), [], True, None, None
        found = set()
        best_dist = 999.0
        best_enc = encodings[0]  # largest face (first detected)
        for person, ref_encs in ref_encodings.items():
            for ref_enc in ref_encs:
                dists = fr.face_distance(encodings, ref_enc)
                min_d = float(np.min(dists))
                if min_d < best_dist:
                    best_dist = min_d
                if min_d <= tolerance:
                    found.add(person)
        return len(locations), sorted(found), True, round(best_dist, 3) if best_dist < 999 else None, best_enc
    except Exception:
        return 0, [], False, None, None


def _should_skip_face_detect(entry):
    """Pre-filter: skip face detection on images unlikely to contain usable faces."""
    w, h = entry.get("width", 0), entry.get("height", 0)
    if w < 200 or h < 200:
        return True
    if entry.get("is_screenshot"):
        return True
    if entry.get("size_kb", 0) < 50:
        return True
    return False


def _verify_single_photo(fpath):
    """Verify a single photo for face detection. Returns dict with status/message."""
    try:
        import face_recognition as fr
        import numpy as np
        from PIL import Image
        pil_img = Image.open(fpath).convert("RGB")
        w, h = pil_img.size
        max_dim = 1600
        if w > max_dim or h > max_dim:
            pil_img.thumbnail((max_dim, max_dim), Image.LANCZOS)
        arr = np.array(pil_img)
        locations = fr.face_locations(arr, model="hog")
        if not locations:
            return {"status": "no_face", "message": "No face detected", "dimensions": f"{w}x{h}"}
        elif len(locations) > 1:
            areas = [(b-t)*(r-l) for t, r, b, l in locations]
            best = areas.index(max(areas))
            enc = fr.face_encodings(arr, [locations[best]])
            if enc:
                return {"status": "ok_multi", "message": f"{len(locations)} faces, using largest", "dimensions": f"{w}x{h}"}
            return {"status": "encode_fail", "message": "Face found but encoding failed", "dimensions": f"{w}x{h}"}
        else:
            enc = fr.face_encodings(arr, locations)
            if enc:
                return {"status": "ok", "message": "Face detected and encoded", "dimensions": f"{w}x{h}"}
            return {"status": "encode_fail", "message": "Encoding failed", "dimensions": f"{w}x{h}"}
    except Exception as e:
        return {"status": "error", "message": str(e)}


def auto_rotate_image(filepath):
    """Apply EXIF orientation and save back. Returns True if rotated."""
    from PIL import Image, ImageOps
    try:
        img = Image.open(filepath)
        rotated = ImageOps.exif_transpose(img)
        if rotated is not img:
            rotated.save(filepath, quality=95)
            return True
    except Exception:
        pass
    return False


def list_templates():
    templates_dir = os.path.join(PROJECT_DIR, "templates")
    result = []
    if not os.path.isdir(templates_dir):
        return result
    for fname in sorted(os.listdir(templates_dir)):
        if not fname.endswith(".json"):
            continue
        with open(os.path.join(templates_dir, fname), "r", encoding="utf-8") as f:
            t = json.load(f)
        result.append({
            "event_type": t["event_type"],
            "display_name": t["display_name"],
            "description": t["description"],
            "categorization": t["categorization"],
            "num_categories": len(t.get("categories", [])),
            "required_fields": t.get("required_fields", {}),
            "tips": t.get("tips", []),
            "extra": t.get("extra", False),
            "categories": [{"id": c["id"], "display": c["display"], "target": c.get("target")} for c in t.get("categories", [])],
        })
    return result


def load_template(event_type):
    tpath = os.path.join(PROJECT_DIR, "templates", f"{event_type}.json")
    if os.path.isfile(tpath):
        with open(tpath, "r", encoding="utf-8") as f:
            return json.load(f)
    return None


# ── API Routes ────────────────────────────────────────────────────────────────

@app.route("/api/templates")
def api_templates():
    return jsonify(list_templates())


@app.route("/api/config", methods=["GET"])
def api_get_config():
    config = load_config()
    if config:
        return jsonify(config)
    return jsonify(None)


@app.route("/api/config", methods=["POST"])
def api_save_config():
    config = request.json
    save_config(config)
    return jsonify({"ok": True})


@app.route("/api/config/sources", methods=["POST"])
def api_save_sources():
    """Update only sources in the config, preserving everything else."""
    data = request.json
    config = load_config() or {}
    config["sources"] = data.get("sources", [])
    save_config(config)
    return jsonify({"ok": True})


@app.route("/api/browse")
def api_browse():
    """List subdirectories at a given path for the folder picker.
    GET /api/browse              → list drive letters (Windows) or /
    GET /api/browse?path=C:\\    → list folders in C:\\
    Returns {path, parent, folders: [{name, path, has_children}]}"""
    import string
    raw_path = request.args.get("path", "").strip()

    # ── Root: list available drives (Windows) or filesystem root ──
    if not raw_path:
        if sys.platform == "win32":
            drives = []
            for letter in string.ascii_uppercase:
                dp = f"{letter}:\\"
                if os.path.isdir(dp):
                    drives.append({"name": f"{letter}:", "path": dp, "has_children": True})
            return jsonify({"path": "", "parent": None, "folders": drives})
        else:
            raw_path = "/"

    # Normalize
    browse_path = os.path.normpath(raw_path)
    if not os.path.isdir(browse_path):
        return jsonify({"error": f"Not a directory: {raw_path}"}), 400

    # Parent path (None at drive root / filesystem root)
    parent = os.path.dirname(browse_path)
    if parent == browse_path:
        parent = ""  # at root — signal to show drives list

    folders = []
    try:
        with os.scandir(browse_path) as it:
            for entry in sorted(it, key=lambda e: e.name.lower()):
                if not entry.is_dir(follow_symlinks=False):
                    continue
                # Skip hidden/system folders
                name = entry.name
                if name.startswith(".") or name.startswith("$"):
                    continue
                # Check if folder has subdirectories (for expand arrow hint)
                has_sub = False
                try:
                    with os.scandir(entry.path) as sub:
                        for s in sub:
                            if s.is_dir(follow_symlinks=False):
                                has_sub = True
                                break
                except PermissionError:
                    pass
                folders.append({
                    "name": name,
                    "path": entry.path,
                    "has_children": has_sub,
                })
    except PermissionError:
        return jsonify({"path": browse_path, "parent": parent,
                        "folders": [], "error": "Access denied"})
    except OSError as e:
        return jsonify({"error": str(e)}), 400

    return jsonify({"path": browse_path, "parent": parent, "folders": folders})


@app.route("/api/browse/validate")
def api_browse_validate():
    """Validate a source folder path.
    GET /api/browse/validate?path=C:\\Photos
    Returns {path, valid, accessible, media_count, subfolder_count}"""
    raw_path = request.args.get("path", "").strip()
    if not raw_path:
        return jsonify({"path": "", "valid": False, "accessible": False,
                        "media_count": 0, "subfolder_count": 0})

    norm = os.path.normpath(raw_path)
    exists = os.path.isdir(norm)
    if not exists:
        return jsonify({"path": raw_path, "valid": False, "accessible": False,
                        "media_count": 0, "subfolder_count": 0})

    # Count media files (walk up to a cap for responsiveness)
    MEDIA_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".heic",
                  ".webp", ".mp4", ".mov", ".avi", ".mkv", ".wmv", ".webm",
                  ".m4v", ".mpg", ".mpeg", ".3gp"}
    media_count = 0
    subfolder_count = 0
    cap = 50000  # stop counting after this to stay fast
    accessible = True
    try:
        for dirpath, dirnames, filenames in os.walk(norm):
            subfolder_count += len(dirnames)
            for fn in filenames:
                if os.path.splitext(fn)[1].lower() in MEDIA_EXTS:
                    media_count += 1
                    if media_count >= cap:
                        break
            if media_count >= cap:
                break
    except PermissionError:
        accessible = False

    return jsonify({
        "path": raw_path,
        "valid": True,
        "accessible": accessible,
        "media_count": media_count,
        "capped": media_count >= cap,
        "subfolder_count": subfolder_count,
    })


_native_picker = {"active": False, "result": None}  # guarded by GIL

@app.route("/api/browse/native", methods=["POST"])
def api_browse_native():
    """Open the native Windows folder picker (non-blocking).
    POST {initialdir?, poll?}
    First call (poll absent): launches dialog, returns {status:'open'}.
    Subsequent calls (poll=true): returns {status:'open'|'done', paths:[]}."""
    import subprocess as _sp
    data = request.json or {}

    # Poll mode — check if dialog finished
    if data.get("poll"):
        if not _native_picker["active"]:
            res = _native_picker["result"]
            _native_picker["result"] = None
            return jsonify({"status": "done", "paths": res or []})
        return jsonify({"status": "open"})

    # Already open — don't launch another
    if _native_picker["active"]:
        return jsonify({"status": "open"})

    initial = data.get("initialdir", "").strip()
    if initial and not os.path.isdir(initial):
        initial = ""

    ps = (
        'Add-Type -AssemblyName System.Windows.Forms;'
        '$owner = New-Object System.Windows.Forms.Form;'
        '$owner.TopMost = $true;'
        '$owner.ShowInTaskbar = $false;'
        '$owner.WindowState = "Minimized";'
        '$owner.Show();'
        '$owner.Hide();'
        '$f = New-Object System.Windows.Forms.OpenFileDialog;'
        '$f.ValidateNames = $false;'
        '$f.CheckFileExists = $false;'
        '$f.CheckPathExists = $true;'
        "$f.FileName = 'Select Folder';"
        "$f.Title = 'Select Source Folder';"
        "$f.Filter = 'Folders|no.files';"
    )
    if initial:
        ps += f"$f.InitialDirectory = '{initial.replace(chr(39), chr(39)+chr(39))}';"
    ps += (
        "if ($f.ShowDialog($owner) -eq 'OK') {"
        "  Write-Output ([System.IO.Path]::GetDirectoryName($f.FileName))"
        "}"
        "$owner.Dispose();"
    )

    _native_picker["active"] = True
    _native_picker["result"] = None

    def _run():
        try:
            r = _sp.run(
                ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", ps],
                capture_output=True, text=True, timeout=120,
            )
            path = r.stdout.strip()
            _native_picker["result"] = [path] if path else []
        except Exception:
            _native_picker["result"] = []
        finally:
            _native_picker["active"] = False

    threading.Thread(target=_run, daemon=True).start()
    return jsonify({"status": "open"})


@app.route("/api/init", methods=["POST"])
def api_init():
    data = request.json
    event_type = data.get("event_type")
    template = load_template(event_type)
    if not template:
        return jsonify({"error": f"Unknown event type: {event_type}"}), 400

    defaults = template.get("defaults", {})
    config = {
        "event_type": event_type,
        "template": template["display_name"],
        "categorization": template["categorization"],
        "ref_faces_dir": "./ref_faces",
        "face_names": data.get("face_names", []),
        "face_tolerance": defaults.get("face_tolerance", 0.6),
        "sources": data.get("sources", []),
        "target_per_category": defaults.get("target_per_category", 75),
        "min_size_kb": defaults.get("min_size_kb", 80),
        "min_dim": defaults.get("min_dim", 600),
        "thumb_size": defaults.get("thumb_size", 120),
        "categories": template["categories"],
    }

    # Fill required fields
    if data.get("subject_birthday"):
        config["subject_birthday"] = data["subject_birthday"]
    if data.get("event_date"):
        config["event_date"] = data["event_date"]
    if data.get("end_date"):
        config["end_date"] = data["end_date"]
    if data.get("year"):
        config["year"] = data["year"]

    save_config(config)
    return jsonify({"ok": True, "config": config})


@app.route("/api/scan/start", methods=["POST"])
def api_scan_start():
    if _any_job_running():
        return jsonify({"error": "A task is already running"}), 409

    full = request.json.get("full", False) if request.json else False
    nsfw_filter = request.json.get("nsfw_filter", False) if request.json else False
    age_estimation = request.json.get("age_estimation", None) if request.json else None

    config = load_config()
    proj_name = (config.get("project_name") or config.get("event_name") or "") if config else ""
    job_id = _create_job("scan", proj_name)

    def run_scan():
        try:
            import time as _time
            from concurrent.futures import ThreadPoolExecutor, as_completed
            _reset_task("scan")
            _update_task("Loading config...")

            config = load_config()
            if not config:
                _finish_task("No config found. Complete setup first.")
                return

            # Import scan logic
            from curate import (
                load_reference_faces, detect_faces_in_image,
                get_image_date, is_screenshot, file_hash,
                guess_device_source, make_thumbnail_b64,
                categorize_by_template, load_template as load_tpl,
                age_days_to_bracket, IMAGE_EXTS, VIDEO_EXTS, MEDIA_EXTS, REEF_BIRTHDAY,
                get_video_date, get_video_info, make_video_thumbnail_b64,
                get_exif_gps, reverse_geocode_batch, infer_location_from_path,
                compute_image_vector, compute_dhash, cluster_similar_images,
                analyze_video_frames,
            )
            from PIL import Image as PILImage

            sources = config.get("sources", [])
            face_names = config.get("face_names", [])
            face_match_mode = config.get("face_match_mode", "any")
            face_dir = config.get("ref_faces_dir", "")
            tolerance = config.get("face_tolerance", 0.6)
            min_size_kb = config.get("min_size_kb", 80)
            min_dim = config.get("min_dim", 600)
            thumb_size = config.get("thumb_size", 120)

            # Template
            event_type = config.get("event_type")
            template = None
            if event_type:
                template = load_tpl(event_type)
                if not template and "categories" in config and "categorization" in config:
                    template = {"categorization": config["categorization"], "categories": config["categories"]}
            use_template = template is not None

            # Faces
            ref_encodings = {}
            use_faces = False
            if face_names and face_dir and os.path.isdir(face_dir):
                _update_task("Loading face references...")
                if _is_cancelled():
                    _update_task("Stopped by user.")
                    _finish_task("Cancelled")
                    return
                ref_encodings = load_reference_faces(face_dir)
                use_faces = bool(ref_encodings)

            if _is_cancelled():
                _update_task("Stopped by user.")
                _finish_task("Cancelled")
                return

            # NSFW detection
            nsfw_classifier = None
            if nsfw_filter:
                try:
                    from nudenet import NudeDetector
                    _update_task("Loading NSFW detection model (first time may download ~25 MB)...")
                    nsfw_classifier = NudeDetector()
                    _update_task("NSFW filter ready.")
                except ImportError:
                    _update_task("nudenet not installed — installing...")
                    import subprocess
                    subprocess.check_call([sys.executable, "-m", "pip", "install", "nudenet"], stdout=subprocess.DEVNULL)
                    from nudenet import NudeDetector
                    nsfw_classifier = NudeDetector()
                    _update_task("NSFW filter installed and ready.")
                except Exception as e:
                    _update_task(f"NSFW filter failed to load: {e}. Continuing without it.")

            if _is_cancelled():
                _update_task("Stopped by user.")
                _finish_task("Cancelled")
                return

            # Age estimation setup
            age_est_enabled = False
            age_est_scope = "all"
            age_est_folders = []
            if age_estimation and age_estimation.get("enabled"):
                try:
                    from deepface import DeepFace
                    _update_task("Loading age estimation model (first run may download ~500 MB)...")
                    # Warm up the model with a dummy analysis
                    import numpy as np
                    _dummy = np.zeros((100, 100, 3), dtype=np.uint8)
                    try:
                        DeepFace.analyze(_dummy, actions=["age"], enforce_detection=False, silent=True)
                    except Exception:
                        pass
                    age_est_enabled = True
                    age_est_scope = age_estimation.get("scope", "all")
                    age_est_folders = [f.replace("\\", "/") for f in age_estimation.get("folders", [])]
                    _update_task("Age estimation model ready.")
                except ImportError:
                    _update_task("deepface not installed — installing...")
                    import subprocess
                    subprocess.check_call([sys.executable, "-m", "pip", "install", "deepface", "tf-keras"], stdout=subprocess.DEVNULL)
                    from deepface import DeepFace
                    age_est_enabled = True
                    age_est_scope = age_estimation.get("scope", "all")
                    age_est_folders = [f.replace("\\", "/") for f in age_estimation.get("folders", [])]
                    _update_task("Age estimation installed and ready.")
                except Exception as e:
                    _update_task(f"Age estimation failed to load: {e}. Continuing without it.")

            if _is_cancelled():
                _update_task("Stopped by user.")
                _finish_task("Cancelled")
                return

            # Existing DB for incremental — only reuse if same project config
            existing_db = {}
            if os.path.isfile(SCAN_DB_PATH) and not full:
                old = load_scan_db()
                if old:
                    old_cfg = old.get("config", {})
                    old_sources = set()
                    for s in old_cfg.get("sources", []):
                        old_sources.add(s.get("path", s) if isinstance(s, dict) else s)
                    new_sources = set()
                    for s in sources:
                        new_sources.add(s.get("path", s) if isinstance(s, dict) else s)
                    old_faces = set(old_cfg.get("face_names", []))
                    new_faces = set(face_names)
                    if old_sources == new_sources and old_faces == new_faces:
                        for img in old.get("images", []):
                            existing_db[img["hash"]] = img
                        _update_task(f"Incremental mode: {len(existing_db)} cached images")
                    else:
                        _update_task("Config changed — running fresh scan...")

            if _is_cancelled():
                _update_task("Stopped by user.")
                _finish_task("Cancelled")
                return

            # === PASS 1: Fast metadata extraction with parallel I/O ===
            _update_task("Pass 1: Collecting file list...")
            pass1_start = _time.monotonic()

            def _compute_photo_grade(fpath, w, h):
                """
                Comprehensive photo quality grading (0-100 each dimension).
                Returns dict with sub-scores and composite grade.
                Computes: resolution, sharpness, noise, compression, color,
                exposure, focus, distortion + composite.
                """
                try:
                    import cv2
                    import numpy as np
                    from PIL import ImageStat

                    pil_img = PILImage.open(fpath)
                    pil_rgb = pil_img.convert("RGB")
                    file_size_kb = os.path.getsize(fpath) / 1024
                    megapixels = (w * h) / 1_000_000

                    # Resize for analysis (consistent measurement)
                    measure_size = min(800, max(w, h))
                    scale = measure_size / max(w, h)
                    if scale < 1:
                        pil_small = pil_rgb.resize((int(w * scale), int(h * scale)), PILImage.LANCZOS)
                    else:
                        pil_small = pil_rgb
                    rgb_arr = np.array(pil_small)
                    gray_arr = cv2.cvtColor(rgb_arr, cv2.COLOR_RGB2GRAY)
                    sh, sw = gray_arr.shape

                    # ── 1. Resolution (0-100): megapixels scaled ──
                    # 0.5MP=20, 2MP=50, 5MP=70, 12MP=90, 20MP+=100
                    resolution = min(100, max(0, 20 + megapixels * 10))

                    # ── 2. Sharpness (0-100): Laplacian variance ──
                    laplacian = cv2.Laplacian(gray_arr, cv2.CV_64F)
                    sharpness_var = laplacian.var()
                    sharpness = min(100, max(0, sharpness_var / 8))

                    # ── 3. Noise/Grain (0-100, higher=less noise=better) ──
                    # Estimate noise via high-pass filter standard deviation
                    # Median filter removes signal, difference = noise estimate
                    median_filtered = cv2.medianBlur(gray_arr, 5)
                    noise_diff = gray_arr.astype(np.float32) - median_filtered.astype(np.float32)
                    noise_std = noise_diff.std()
                    # noise_std: 0-3 = very clean, 5-10 = some noise, 15+ = noisy
                    noise = min(100, max(0, 100 - noise_std * 5))

                    # ── 4. Compression artifacts (0-100, higher=less artifacts) ──
                    # KB per megapixel: higher = less compressed = better
                    kb_per_mp = file_size_kb / max(megapixels, 0.01)
                    # 200 KB/MP = heavy JPEG, 500 = normal, 1000+ = high quality
                    compression = min(100, max(0, kb_per_mp / 15))

                    # ── 5. Color accuracy & dynamic range (0-100) ──
                    stat = ImageStat.Stat(pil_small)
                    # Color saturation: convert to HSV, measure saturation channel
                    hsv_arr = cv2.cvtColor(rgb_arr, cv2.COLOR_RGB2HSV)
                    mean_sat = hsv_arr[:, :, 1].mean()  # 0-255
                    sat_score = min(100, mean_sat / 1.5)  # ~170 sat = 100

                    # Dynamic range: difference between darkest and brightest regions
                    p5 = np.percentile(gray_arr, 5)
                    p95 = np.percentile(gray_arr, 95)
                    dyn_range = p95 - p5  # 0-255
                    range_score = min(100, max(0, dyn_range / 2.0))  # 200+ range = 100

                    color = sat_score * 0.4 + range_score * 0.6

                    # ── 6. Exposure & lighting (0-100) ──
                    mean_brightness = sum(stat.mean[:3]) / 3  # 0-255
                    if 80 <= mean_brightness <= 180:
                        exposure = 100
                    elif mean_brightness < 40 or mean_brightness > 230:
                        exposure = 20
                    elif mean_brightness < 80:
                        exposure = 20 + (mean_brightness - 40) * 2
                    else:  # > 180
                        exposure = 20 + (230 - mean_brightness) * 1.6

                    # Bonus/penalty for contrast (std dev of brightness)
                    brightness_std = np.std(gray_arr)
                    # Good contrast: std 40-80, too flat: <20, too harsh: >100
                    if 30 <= brightness_std <= 80:
                        contrast_bonus = 0
                    elif brightness_std < 15:
                        contrast_bonus = -15  # flat/washed out
                    elif brightness_std > 100:
                        contrast_bonus = -10  # harsh
                    else:
                        contrast_bonus = -5
                    exposure = max(0, min(100, exposure + contrast_bonus))

                    # ── 7. Focus & depth of field (0-100) ──
                    # Compare sharpness in center vs edges
                    ch, cw = sh // 4, sw // 4
                    center = gray_arr[ch:sh-ch, cw:sw-cw]
                    center_lap = cv2.Laplacian(center, cv2.CV_64F).var()

                    # Edge regions
                    top = gray_arr[:ch, :]
                    bottom = gray_arr[sh-ch:, :]
                    edge_lap = (cv2.Laplacian(top, cv2.CV_64F).var() +
                                cv2.Laplacian(bottom, cv2.CV_64F).var()) / 2

                    # Good focus: center sharper than edges (subject in focus, bg bokeh ok)
                    if center_lap > 20:
                        focus_sharpness = min(100, center_lap / 8)
                        # Slight bonus if center is sharper (good depth of field control)
                        if edge_lap > 0 and center_lap / max(edge_lap, 1) > 1.5:
                            focus_bonus = 5
                        else:
                            focus_bonus = 0
                        focus = min(100, focus_sharpness + focus_bonus)
                    else:
                        focus = max(0, center_lap / 8 * 100)

                    # ── 8. Lack of distortion (0-100) ──
                    # Penalize extreme aspect ratios and very small images
                    ratio = w / h if h > 0 else 1
                    if 0.6 <= ratio <= 1.8:
                        distortion = 100
                    elif ratio < 0.4 or ratio > 3:
                        distortion = 30
                    else:
                        distortion = 65
                    # Small dimension penalty
                    min_dim_val = min(w, h)
                    if min_dim_val < 500:
                        distortion = max(0, distortion - 20)

                    pil_img.close()

                    # ── Composite grade (weighted) ──
                    composite = (
                        resolution * 0.10 +
                        sharpness * 0.20 +
                        noise * 0.10 +
                        compression * 0.05 +
                        color * 0.10 +
                        exposure * 0.15 +
                        focus * 0.20 +
                        distortion * 0.10
                    )

                    return {
                        "resolution": round(resolution, 1),
                        "sharpness": round(sharpness, 1),
                        "noise": round(noise, 1),
                        "compression": round(compression, 1),
                        "color": round(color, 1),
                        "exposure": round(exposure, 1),
                        "focus": round(focus, 1),
                        "distortion": round(distortion, 1),
                        "composite": round(composite, 1),
                        "blur_score": round(sharpness_var, 1),
                    }
                except Exception:
                    return None

            def _extract_metadata(fpath, fname, fhash, is_video, src_label, rel_dir):
                """Thread-safe metadata extraction (no face detection)."""
                try:
                    file_size = os.path.getsize(fpath)
                    if not is_video and file_size < min_size_kb * 1024:
                        return None, "too_small", fhash
                    if is_video:
                        w, h, duration = get_video_info(fpath)
                    else:
                        img = PILImage.open(fpath)
                        w, h = img.size
                        img.close()
                        duration = 0
                    if not is_video and w < min_dim and h < min_dim:
                        return None, "low_res", fhash

                    # Photo quality grading
                    photo_grade = None
                    blur_score = None
                    video_analysis = None
                    if is_video:
                        # Analyze video frames for quality, vector, dHash
                        video_analysis = analyze_video_frames(
                            fpath, ref_encodings=None, face_names=None,
                            n_frames=5,
                        )
                        if video_analysis:
                            photo_grade = video_analysis.get("photo_grade")
                            if photo_grade:
                                blur_score = photo_grade.get("blur_score")
                    else:
                        photo_grade = _compute_photo_grade(fpath, w, h)
                        if photo_grade:
                            blur_score = photo_grade.get("blur_score")

                    if is_video:
                        img_date = get_video_date(fpath)
                    else:
                        img_date = get_image_date(fpath, rel_dir)

                    age_days = None
                    bracket = None
                    if img_date:
                        if use_template:
                            bracket = categorize_by_template(template, config, img_date, fpath=fpath)
                            bday_str = config.get("subject_birthday")
                            if bday_str:
                                from datetime import datetime as dt
                                bday = dt.strptime(bday_str, "%Y-%m-%d")
                                age_days = (img_date - bday).days
                        else:
                            age_days = (img_date - REEF_BIRTHDAY).days
                            bracket = age_days_to_bracket(age_days)

                    screenshot = False
                    if not is_video and file_size < 500 * 1024:
                        screenshot = is_screenshot(fpath)

                    # GPS extraction (images only)
                    gps_coords = None
                    if not is_video:
                        gps_coords = get_exif_gps(fpath)

                    if is_video:
                        # Prefer sharpest-frame thumbnail from analysis
                        thumb = (video_analysis or {}).get("best_thumb") or make_video_thumbnail_b64(fpath, thumb_size)
                    else:
                        thumb = make_thumbnail_b64(fpath, thumb_size)
                    device = guess_device_source(fname)

                    # Compute perceptual vector + dHash for similarity clustering
                    image_vector = None
                    dhash_val = None
                    if is_video and video_analysis:
                        # Use pre-computed values from video frame analysis
                        image_vector = video_analysis.get("image_vector")
                        dhash_val = video_analysis.get("dhash")
                    elif not is_video:
                        try:
                            image_vector = compute_image_vector(fpath)
                        except Exception:
                            pass
                        try:
                            dhash_val = compute_dhash(fpath)
                        except Exception:
                            pass

                    # CLIP semantic embedding + auto-tags.
                    #
                    # State machine (new images enter as NEVER PROCESSED):
                    #   NEVER PROCESSED  → has_clip=False, tag_meta=None
                    #   PROCESSED OK     → has_clip=True,  tag_meta has both versions
                    #   FAILED           → has_clip=False,  tag_meta=None (same as never)
                    #
                    # Atomicity: _tags/_tag_meta stay empty unless BOTH embedding
                    # AND tag generation succeed and the vector is persisted.
                    # On failure the entry lands in NEVER PROCESSED / FAILED state
                    # and will be retried on next scan.
                    _tags = []
                    _tag_meta = None
                    if clip_ready:
                        try:
                            if is_video and thumb:
                                from io import BytesIO as _BytesIO
                                _pil = PILImage.open(
                                    _BytesIO(base64.b64decode(thumb)))
                                _cv = clip_engine.compute_embedding_pil(_pil)
                            elif not is_video:
                                _cv = clip_engine.compute_embedding(fpath)
                            else:
                                _cv = None
                            if _cv is not None:
                                _t, _m = clip_engine.generate_tags(_cv)
                                # All succeeded — commit results atomically
                                with _clip_lock_local:
                                    _clip_vectors_local[fhash] = _cv
                                _tags = _t
                                _tag_meta = _m
                        except Exception:
                            # Embedding or tagging failed — leave defaults
                            _tags = []
                            _tag_meta = None

                    entry = {
                        "hash": fhash,
                        "path": fpath.replace("\\", "/"),
                        "filename": fname,
                        "source_label": src_label,
                        "device": device,
                        "media_type": "video" if is_video else "image",
                        "date": img_date.strftime("%Y-%m-%d") if img_date else None,
                        "age_days": age_days,
                        "category": bracket,
                        "face_count": 0,
                        "faces_found": [],
                        "has_target_face": False,
                        "face_distance": None,
                        "width": w, "height": h,
                        "duration": round(duration, 1) if is_video else 0,
                        "size_kb": round(file_size / 1024),
                        "is_screenshot": screenshot,
                        "thumb": thumb,
                        "gps_lat": gps_coords[0] if gps_coords else None,
                        "gps_lon": gps_coords[1] if gps_coords else None,
                        "location": None,  # resolved in batch after Pass 1
                        "blur_score": blur_score,
                        "photo_grade": photo_grade,
                        "image_vector": image_vector,
                        "dhash": dhash_val,
                        "tags": _tags,
                        "tag_meta": _tag_meta,
                        "has_clip": bool(_tag_meta),
                    }

                    # For manual categorization: assign temporary category (will be refined after Pass 2)
                    if not bracket and config.get("categorization") == "manual" and config.get("categories"):
                        bracket = config["categories"][0]["id"]
                        entry["category"] = bracket
                        entry["_needs_recategorize"] = True

                    # Mark thematic fallback assignments for refinement after Pass 2
                    if bracket and use_template and template and config.get("categorization") in ("date_time_ranges", "date_ranges"):
                        for _tc in template.get("categories", []):
                            if _tc["id"] == bracket and _tc.get("day_offset") == -1:
                                entry["_needs_recategorize"] = True
                                break
                            # Full-year catch-all in date_ranges
                            if _tc["id"] == bracket and _tc.get("month_from") == 1 and _tc.get("month_to") == 12:
                                entry["_needs_recategorize"] = True
                                break

                    # Preliminary status (will be updated after face detection in Pass 2)
                    blur_threshold = config.get("blur_threshold", 50)
                    if screenshot:
                        entry["status"] = "rejected"
                        entry["reject_reason"] = "screenshot"
                    elif blur_score is not None and blur_score < blur_threshold:
                        entry["status"] = "rejected"
                        entry["reject_reason"] = "blurry"
                    elif not bracket:
                        entry["status"] = "pool"
                        entry["reject_reason"] = "no_date"
                    else:
                        entry["status"] = "qualified"
                        entry["reject_reason"] = None

                    return entry, None, fhash
                except Exception:
                    return None, "unreadable", fhash

            # Collect all file paths first
            all_file_info = []  # (fpath, fname, is_video, src_label, rel_dir, src_path)
            for source in sources:
                if isinstance(source, str):
                    src_path = source
                    src_label = os.path.basename(source) or source
                else:
                    src_path = source.get("path", "")
                    src_label = source.get("label", "Unknown")
                if not os.path.isdir(src_path):
                    _update_task(f"Source not found: {src_label}")
                    continue
                for dirpath, dirnames, filenames in os.walk(src_path):
                    rel_dir = os.path.relpath(dirpath, src_path)
                    for fname in filenames:
                        ext = os.path.splitext(fname)[1].lower()
                        if ext not in MEDIA_EXTS:
                            continue
                        fpath = os.path.join(dirpath, fname)
                        is_video = ext in VIDEO_EXTS
                        all_file_info.append((fpath, fname, is_video, src_label, rel_dir, src_path))

            total_media_files = len(all_file_info)
            _update_task(f"Found {total_media_files} media files. Starting Pass 1 (metadata)...")

            # Pre-flight: CLIP semantic embedding support
            clip_ready = False
            _clip_vectors_local = {}
            _clip_lock_local = threading.Lock()
            _clip_embed_ver = None   # e.g. "clip-vit-b32-onnx:1.0"
            _clip_tagger_ver = None  # e.g. "a3f2c8e91b04"
            if clip_engine is not None:
                try:
                    if not clip_engine.is_available():
                        _update_task("Downloading CLIP model for image understanding...")
                        clip_engine.ensure_models(
                            progress_fn=lambda msg: _update_task(str(msg)))
                    if clip_engine.is_available():
                        _update_task("Loading CLIP model...")
                        clip_engine._get_vision_session()
                        clip_engine._get_tag_embeddings()
                        clip_ready = True
                        _clip_embed_ver = clip_engine.EMBED_VERSION
                        _clip_tagger_ver = clip_engine.TAGGER_VERSION
                        # Pre-load existing vectors so incremental backfill
                        # merges with them (skip on full rescan to drop stale data)
                        if not full and os.path.isfile(CLIP_VECTORS_PATH):
                            _clip_vectors_local = clip_engine.load_vectors(
                                CLIP_VECTORS_PATH)
                            _update_task(
                                f"CLIP ready (embed={_clip_embed_ver}, "
                                f"tagger={_clip_tagger_ver}), "
                                f"{len(_clip_vectors_local)} cached vectors loaded")
                        else:
                            _update_task(
                                f"CLIP ready (embed={_clip_embed_ver}, "
                                f"tagger={_clip_tagger_ver})")
                except Exception as e:
                    print(f"CLIP not available: {e}")
                    import traceback; traceback.print_exc()

            all_images = []
            seen_hashes = set()
            skipped = defaultdict(int)
            scanned = 0
            pass1_times = []

            # Process in batches with ThreadPoolExecutor
            BATCH_SIZE = 50
            with ThreadPoolExecutor(max_workers=4) as executor:
                for batch_start in range(0, total_media_files, BATCH_SIZE):
                    if _is_cancelled():
                        if all_images:
                            save_scan_db({"scan_date": datetime.now().isoformat(), "config": config,
                                "stats": {"total_scanned": scanned, "total_kept": len(all_images), "skipped": dict(skipped), "partial": True},
                                "images": all_images})
                        _update_task(f"Stopped. Saved {len(all_images)} images.")
                        _finish_task("Cancelled")
                        return

                    batch = all_file_info[batch_start:batch_start + BATCH_SIZE]
                    batch_t0 = _time.monotonic()

                    futures = {}
                    for fpath, fname, is_video, src_label, rel_dir, src_path in batch:
                        # Hash + dedup on main thread (shared state)
                        try:
                            fhash = file_hash(fpath)
                        except Exception:
                            skipped["unreadable"] += 1
                            scanned += 1
                            continue

                        if fhash in seen_hashes:
                            skipped["duplicate"] += 1
                            scanned += 1
                            continue
                        seen_hashes.add(fhash)

                        # Reuse from existing DB
                        if fhash in existing_db:
                            entry = existing_db[fhash]
                            entry["path"] = fpath.replace("\\", "/")
                            entry["source_label"] = src_label

                            # NSFW check on cached entries if filter enabled and not yet checked
                            if nsfw_classifier and "nsfw" not in entry and entry.get("status") != "rejected":
                                is_nsfw, nsfw_labels = _check_nsfw(fpath, nsfw_classifier)
                                entry["nsfw"] = is_nsfw
                                if is_nsfw:
                                    entry["nsfw_labels"] = nsfw_labels
                                    entry["status"] = "rejected"
                                    entry["reject_reason"] = "nsfw"

                            # GPS extraction on cached entries if not yet done
                            if "gps_lat" not in entry and entry.get("media_type") != "video":
                                gps = get_exif_gps(fpath)
                                entry["gps_lat"] = gps[0] if gps else None
                                entry["gps_lon"] = gps[1] if gps else None
                                # location resolved in batch after Pass 1

                            # Age estimation on cached entries if enabled and not yet done
                            if (age_est_enabled
                                and "estimated_age" not in entry
                                and entry.get("has_target_face")
                                and entry.get("media_type") != "video"
                                and entry.get("status") != "rejected"):
                                run_age = False
                                if age_est_scope == "all":
                                    run_age = True
                                elif age_est_scope == "folders":
                                    src_norm = src_path.replace("\\", "/")
                                    run_age = src_norm in age_est_folders
                                if run_age:
                                    est = _estimate_age(fpath)
                                    if est is not None:
                                        entry["estimated_age"] = est

                            # CLIP backfill on cached entries.
                            #
                            # Five entry states and their transitions:
                            #
                            #  NEVER PROCESSED   has_clip=False
                            #    → needs_embed=True  → full recompute
                            #
                            #  PROCESSED CURRENT  has_clip=True, both versions match,
                            #                     vector in sidecar
                            #    → skip (no work)
                            #
                            #  OUTDATED EMBED     has_clip=True, embed_version !=
                            #    → needs_embed=True  → full recompute
                            #
                            #  OUTDATED TAGGER    has_clip=True, embed_version ok,
                            #                     tagger_version !=
                            #    → needs_embed=False, needs_tags=True
                            #    → re-tag from existing vector (no image I/O)
                            #
                            #  SANITY RECOVERY    has_clip=True, versions ok,
                            #                     but vector missing from sidecar
                            #                     (crash between periodic scan_db
                            #                     save and sidecar write)
                            #    → needs_embed=True  → full recompute, self-heals
                            #
                            # On any failure the except leaves the entry unchanged,
                            # so it stays in its current state and will be retried.
                            if clip_ready and entry.get("status") != "rejected":
                                _tm = entry.get("tag_meta") or {}
                                _has = entry.get("has_clip", False)
                                _needs_embed = (
                                    not _has
                                    or _tm.get("embed_version") != _clip_embed_ver
                                    or fhash not in _clip_vectors_local
                                )
                                _needs_tags = (
                                    _needs_embed
                                    or _tm.get("tagger_version") != _clip_tagger_ver
                                )

                                if _needs_tags:
                                    try:
                                        _cv = None
                                        if _needs_embed:
                                            _is_vid = entry.get("media_type") == "video"
                                            if _is_vid and entry.get("thumbnail"):
                                                from io import BytesIO as _BytesIO
                                                _pil = PILImage.open(
                                                    _BytesIO(base64.b64decode(
                                                        entry["thumbnail"])))
                                                _cv = clip_engine.compute_embedding_pil(
                                                    _pil)
                                            elif not _is_vid:
                                                _cv = clip_engine.compute_embedding(
                                                    fpath)
                                        else:
                                            # Embedding current — reuse from sidecar
                                            _cv = _clip_vectors_local.get(fhash)

                                        if _cv is not None:
                                            _bt, _bm = clip_engine.generate_tags(_cv)
                                            # All succeeded — commit atomically:
                                            # vector first, then entry fields.
                                            with _clip_lock_local:
                                                _clip_vectors_local[fhash] = _cv
                                            entry["tags"] = _bt
                                            entry["tag_meta"] = _bm
                                            entry["has_clip"] = True
                                        # else: couldn't get embedding, leave as-is
                                    except Exception:
                                        pass  # failed — entry unchanged, retried next scan

                            all_images.append(entry)
                            scanned += 1
                            continue

                        # Submit for parallel metadata extraction (pass pre-computed hash)
                        fut = executor.submit(_extract_metadata, fpath, fname, fhash, is_video, src_label, rel_dir)
                        futures[fut] = (fpath, fname)

                    # Collect results
                    for fut in as_completed(futures):
                        entry, skip_reason, fhash = fut.result()
                        scanned += 1
                        if skip_reason:
                            skipped[skip_reason] += 1
                        elif entry:
                            all_images.append(entry)

                    batch_elapsed = _time.monotonic() - batch_t0
                    pass1_times.append(batch_elapsed / max(len(batch), 1))

                    # Progress + ETA
                    pct = int(scanned * 50 / total_media_files) if total_media_files else 0  # Pass 1 = 0-50%
                    _update_task_percent(pct)
                    if len(pass1_times) >= 3:
                        avg_per_file = sum(pass1_times[-20:]) / len(pass1_times[-20:])
                        remaining = total_media_files - scanned
                        eta_sec = int(avg_per_file * remaining)
                        eta_str = f"{eta_sec // 60}m {eta_sec % 60}s" if eta_sec >= 60 else f"{eta_sec}s"
                        _update_task(f"Pass 1: {scanned}/{total_media_files} files ({len(all_images)} kept) — ETA: {eta_str}")
                    else:
                        _update_task(f"Pass 1: {scanned}/{total_media_files} files ({len(all_images)} kept)...")

                    # Periodic save every ~2% of total files
                    save_interval = max(BATCH_SIZE, total_media_files // 50)
                    if scanned % save_interval < BATCH_SIZE and all_images:
                        save_scan_db({"scan_date": datetime.now().isoformat(), "config": config,
                            "stats": {"total_scanned": scanned, "total_kept": len(all_images), "skipped": dict(skipped), "partial": True},
                            "images": all_images})

            pass1_elapsed = _time.monotonic() - pass1_start
            _update_task(f"Pass 1 complete: {len(all_images)} images in {pass1_elapsed:.0f}s. Resolving locations...")

            # Batch reverse-geocode GPS coordinates
            coords_to_resolve = []
            for i, img in enumerate(all_images):
                if img.get("gps_lat") is not None and not img.get("location"):
                    coords_to_resolve.append((i, (img["gps_lat"], img["gps_lon"])))
            if coords_to_resolve:
                _update_task(f"Reverse geocoding {len(coords_to_resolve)} GPS locations...")
                try:
                    loc_results = reverse_geocode_batch([c for _, c in coords_to_resolve])
                    for (idx, _), loc in zip(coords_to_resolve, loc_results):
                        if loc:
                            all_images[idx]["location"] = loc
                except Exception:
                    pass

            # Folder name fallback for images without GPS location
            for img in all_images:
                if not img.get("location"):
                    loc = infer_location_from_path(img.get("path", ""))
                    if loc:
                        img["location"] = loc

            gps_count = sum(1 for img in all_images if img.get("location"))
            if gps_count:
                _update_task(f"Located {gps_count} images. Starting face detection...")
            else:
                _update_task(f"Starting face detection...")

            # NSFW check on new (non-cached) images
            if nsfw_classifier:
                nsfw_count = 0
                for entry in all_images:
                    if entry.get("status") == "rejected":
                        continue
                    if "nsfw" in entry:
                        continue
                    is_nsfw, nsfw_labels = _check_nsfw(entry["path"].replace("/", os.sep), nsfw_classifier)
                    entry["nsfw"] = is_nsfw
                    if is_nsfw:
                        entry["nsfw_labels"] = nsfw_labels
                        entry["status"] = "rejected"
                        entry["reject_reason"] = "nsfw"
                        nsfw_count += 1
                if nsfw_count:
                    _update_task(f"NSFW filter: {nsfw_count} images rejected")

            # Save CLIP vectors to sidecar file (not in scan_db — too large).
            # Crash safety: save_vectors writes to .tmp then os.replace
            # (atomic on NTFS/POSIX).  If crash happens between a periodic
            # scan_db save (which may have has_clip=True) and this sidecar
            # write, the sanity check (fhash not in _clip_vectors_local)
            # in the backfill block will detect the missing vector on the
            # next incremental scan and recompute it.
            if _clip_vectors_local:
                _clip_tagged = sum(1 for img in all_images if img.get("has_clip"))
                try:
                    clip_engine.save_vectors(
                        _clip_vectors_local, CLIP_VECTORS_PATH)
                    _update_task(
                        f"CLIP: {_clip_tagged} images tagged, "
                        f"{len(_clip_vectors_local)} vectors saved. "
                        f"Starting face detection...")
                except Exception as e:
                    print(f"Failed to save CLIP vectors: {e}")

            # === PASS 2: Face detection (single-threaded, selective) ===
            if use_faces:
                need_face = [e for e in all_images
                             if not e.get("_face_checked")
                             and e.get("status") != "rejected"
                             and not _should_skip_face_detect(e)]

                # Also include cached entries needing face-check or distance backfill
                for e in all_images:
                    if e in need_face:
                        continue
                    if e.get("status") == "rejected":
                        continue
                    needs_face_check = not e.get("_face_checked")
                    needs_distance_backfill = (e.get("has_target_face")
                        and e.get("face_distance") is None)
                    if (needs_face_check or needs_distance_backfill) and not _should_skip_face_detect(e):
                        need_face.append(e)

                total_face = len(need_face)
                _update_task(f"Pass 2: Face detection on {total_face} images...")
                pass2_start = _time.monotonic()
                face_checked = 0
                pass2_times = []
                _face_encs_local = {}  # hash → 128-dim encoding (best face per image)

                for entry in need_face:
                    if _is_cancelled():
                        if all_images:
                            save_scan_db({"scan_date": datetime.now().isoformat(), "config": config,
                                "stats": {"total_scanned": scanned, "total_kept": len(all_images), "skipped": dict(skipped), "partial": True},
                                "images": all_images})
                        _update_task(f"Stopped. Saved {len(all_images)} images ({face_checked} face-checked).")
                        _finish_task("Cancelled")
                        return

                    t0 = _time.monotonic()
                    fpath = entry["path"].replace("/", os.sep)
                    try:
                        is_vid = entry.get("media_type") == "video"
                        if is_vid:
                            # Video: analyze frames for face detection
                            va = analyze_video_frames(
                                fpath, ref_encodings=ref_encodings,
                                face_names=face_names, tolerance=tolerance,
                                face_match_mode=face_match_mode, n_frames=5,
                            )
                            if va:
                                fc = va["face_count"]
                                ff = va["faces_found"]
                                best_d = va["face_distance"]
                                entry["face_count"] = fc
                                entry["faces_found"] = ff
                                entry["has_target_face"] = va["has_target_face"]
                                if best_d is not None:
                                    entry["face_distance"] = best_d
                            else:
                                fc, ff = 0, []
                                entry["face_count"] = 0
                                entry["has_target_face"] = False
                        else:
                            fc, ff, ok, best_d, best_enc = _fast_face_detect(fpath, ref_encodings, face_names, tolerance)
                            entry["face_count"] = fc
                            entry["faces_found"] = ff
                            if best_enc is not None:
                                _face_encs_local[entry["hash"]] = best_enc
                            entry["has_target_face"] = (all(n in ff for n in face_names) if face_match_mode == "all" else any(n in ff for n in face_names)) if face_names else (fc > 0)
                            if best_d is not None:
                                entry["face_distance"] = best_d
                        entry["_face_checked"] = True

                        # Update status based on face results
                        if face_names and not entry["has_target_face"]:
                            if entry.get("status") in ("qualified", "selected"):
                                entry["status"] = "pool"
                                entry["reject_reason"] = "no_faces" if fc == 0 else "wrong_person"
                        elif entry.get("status") == "pool" and entry.get("reject_reason") in ("no_faces", "wrong_person"):
                            if entry.get("category"):
                                entry["status"] = "qualified"
                                entry["reject_reason"] = None
                    except Exception:
                        entry["_face_checked"] = True

                    face_checked += 1
                    elapsed = _time.monotonic() - t0
                    pass2_times.append(elapsed)

                    # Progress + ETA
                    if face_checked % 10 == 0 or face_checked == total_face:
                        pct = 50 + int(face_checked * 50 / total_face) if total_face else 100
                        _update_task_percent(pct)
                        if len(pass2_times) >= 5:
                            avg = sum(pass2_times[-50:]) / len(pass2_times[-50:])
                            remaining = total_face - face_checked
                            eta_sec = int(avg * remaining)
                            eta_str = f"{eta_sec // 60}m {eta_sec % 60}s" if eta_sec >= 60 else f"{eta_sec}s"
                            found_faces = sum(1 for e in all_images if e.get("has_target_face"))
                            _update_task(f"Pass 2: {face_checked}/{total_face} face-checked, {found_faces} matched — ETA: {eta_str}")
                        else:
                            _update_task(f"Pass 2: {face_checked}/{total_face} face-checked...")

                    # Periodic save every ~5% of face-checks
                    face_save_interval = max(10, total_face // 20)
                    if face_checked % face_save_interval == 0:
                        save_scan_db({"scan_date": datetime.now().isoformat(), "config": config,
                            "stats": {"total_scanned": scanned, "total_kept": len(all_images), "skipped": dict(skipped), "partial": True},
                            "images": all_images})

                # Mark skipped images as face-checked too
                for entry in all_images:
                    if not entry.get("_face_checked") and entry.get("media_type") != "video" and entry.get("status") != "rejected":
                        entry["_face_checked"] = True
                        if face_names:
                            entry["status"] = "pool"
                            entry["reject_reason"] = "no_faces"

                pass2_elapsed = _time.monotonic() - pass2_start
                _update_task(f"Face detection complete: {face_checked} images in {pass2_elapsed:.0f}s")

                # Persist face encodings for reuse
                if _face_encs_local:
                    try:
                        np.savez_compressed(
                            FACE_ENCODINGS_PATH + ".tmp",
                            **{k: v for k, v in _face_encs_local.items()})
                        os.replace(
                            FACE_ENCODINGS_PATH + ".tmp",
                            FACE_ENCODINGS_PATH)
                    except Exception:
                        pass

            # Age estimation (if enabled, single-threaded)
            if age_est_enabled:
                age_count = 0
                for entry in all_images:
                    if (entry.get("has_target_face") and entry.get("media_type") != "video"
                        and entry.get("status") != "rejected" and "estimated_age" not in entry):
                        run_age = False
                        if age_est_scope == "all":
                            run_age = True
                        elif age_est_scope == "folders":
                            run_age = True  # simplified — folder check done during Pass 1 for cached
                        if run_age:
                            est = _estimate_age(entry["path"].replace("/", os.sep))
                            if est is not None:
                                entry["estimated_age"] = est
                                age_count += 1
                if age_count:
                    _update_task(f"Age estimation: {age_count} images processed")

            # Re-categorize images that need refinement (manual templates + thematic fallbacks)
            # Now we have face_count from Pass 2, so heuristics can use it
            if config.get("categories"):
                cats_list = config["categories"]
                cat_type = config.get("categorization", "")
                recat_count = 0

                if cat_type == "manual":
                    # Manual templates: use generalized heuristic
                    from curate import categorize_heuristic
                    for entry in all_images:
                        if entry.get("_needs_recategorize") and entry.get("status") != "rejected":
                            new_cat = categorize_heuristic(entry, cats_list, template)
                            if new_cat:
                                entry["category"] = new_cat
                            recat_count += 1
                        entry.pop("_needs_recategorize", None)

                elif cat_type in ("date_time_ranges", "date_ranges"):
                    # Thematic fallback refinement using face count
                    from curate import refine_thematic_category
                    thematic_cats = [c for c in cats_list
                                     if c.get("day_offset") == -1
                                     or (c.get("month_from") == 1 and c.get("month_to") == 12)]
                    if thematic_cats:
                        for entry in all_images:
                            if entry.get("_needs_recategorize") and entry.get("status") != "rejected":
                                new_cat = refine_thematic_category(entry, thematic_cats)
                                if new_cat:
                                    entry["category"] = new_cat
                                recat_count += 1
                            entry.pop("_needs_recategorize", None)

                if recat_count:
                    dist = {}
                    for entry in all_images:
                        c = entry.get("category")
                        if c:
                            dist[c] = dist.get(c, 0) + 1
                    dist_str = ", ".join(f"{k}: {v}" for k, v in sorted(dist.items()))
                    _update_task(f"Smart categorization: {recat_count} images sorted ({dist_str})")

            # === PASS 3: Similarity clustering ===
            _update_task("Pass 3: Clustering similar images...")
            # Load CLIP vectors for semantic clustering (if available)
            _cluster_clip_vecs = None
            if clip_engine is not None and os.path.isfile(CLIP_VECTORS_PATH):
                try:
                    _cluster_clip_vecs = clip_engine.load_vectors(
                        CLIP_VECTORS_PATH)
                except Exception:
                    pass
            _cluster_img_vecs = load_image_vectors()
            cluster_result = cluster_similar_images(
                all_images,
                vector_threshold=config.get("cluster_vector_threshold", 0.92),
                dhash_threshold=config.get("cluster_dhash_threshold", 5),
                clip_vectors=_cluster_clip_vecs,
                image_vectors=_cluster_img_vecs,
                progress_cb=_update_task,
            )
            if cluster_result["clusters"] > 0:
                _update_task(
                    f"Clustering: {cluster_result['clusters']} groups found, "
                    f"{cluster_result['suppressed']} duplicates suppressed "
                    f"(largest group: {cluster_result['largest_cluster']}, "
                    f"{cluster_result['elapsed']:.1f}s)"
                )

            # Final save
            total_elapsed = _time.monotonic() - pass1_start
            db = {
                "scan_date": datetime.now().isoformat(),
                "config": config,
                "stats": {
                    "total_scanned": scanned, "total_kept": len(all_images),
                    "skipped": dict(skipped),
                    "clusters": cluster_result["clusters"],
                    "suppressed_duplicates": cluster_result["suppressed"],
                },
                "images": all_images,
            }
            save_scan_db(db)

            found_target = sum(1 for e in all_images if e.get("has_target_face"))
            blurry_count = sum(1 for e in all_images if e.get("reject_reason") == "blurry")
            graded = [e["photo_grade"]["composite"] for e in all_images if e.get("photo_grade")]
            grade_avg = round(sum(graded) / len(graded), 1) if graded else 0
            sup_msg = f", {cluster_result['suppressed']} near-duplicates suppressed" if cluster_result["suppressed"] else ""
            _update_task(f"Done! {len(all_images)} images, {found_target} with target face, "
                         f"{blurry_count} blurry rejected, avg grade {grade_avg}/100{sup_msg} ({total_elapsed:.0f}s)")
            _finish_task()

        except Exception as e:
            import traceback
            tb_str = traceback.format_exc()
            traceback.print_exc()
            # Also write to scan_error.log for reliable diagnosis
            try:
                with open(os.path.join(PROJECT_DIR, "scan_error.log"), "w", encoding="utf-8") as ef:
                    ef.write(f"Scan error at {datetime.now().isoformat()}\n")
                    ef.write(tb_str)
            except Exception:
                pass
            # Save whatever we have on crash
            try:
                if all_images:
                    save_scan_db({"scan_date": datetime.now().isoformat(), "config": config,
                        "stats": {"total_scanned": scanned, "total_kept": len(all_images), "skipped": dict(skipped), "partial": True},
                        "images": all_images})
            except Exception:
                pass
            _finish_task(f"{type(e).__name__}: {e}")

    threading.Thread(target=run_scan, daemon=True).start()
    return jsonify({"ok": True, "job_id": job_id})


@app.route("/api/scan/status")
def api_scan_status():
    with _task_lock:
        return jsonify({
            "running": _task["running"],
            "type": _task["type"],
            "progress": _task["progress"],
            "lines": _task["lines"][-20:],
            "done": _task["done"],
            "error": _task["error"],
            "cancelled": _task["cancelled"],
        })


@app.route("/api/vectors/verify")
def api_vectors_verify():
    """Check consistency between scan_db.json and image_vectors.npz."""
    repair = request.args.get("repair", "").lower() in ("1", "true", "yes")
    result = verify_vector_consistency(repair=repair)
    return jsonify(result)


@app.route("/api/task/stop", methods=["POST"])
def api_task_stop():
    """Stop by job_id or legacy (current job)."""
    job_id = (request.json or {}).get("job_id") if request.is_json else None
    if not job_id:
        job_id = _current_job_id
    if job_id:
        with _jobs_lock:
            job = _jobs.get(job_id)
            if job and job["running"]:
                job["cancelled"] = True
                _sync_legacy_task(job)
                return jsonify({"ok": True})
    # Legacy fallback
    with _task_lock:
        if _task["running"]:
            _task["cancelled"] = True
            return jsonify({"ok": True})
    return jsonify({"ok": False, "msg": "No task running"})


@app.route("/api/jobs")
def api_jobs():
    """Return all jobs (running + recently finished)."""
    with _jobs_lock:
        jobs_list = []
        for j in _jobs.values():
            jobs_list.append({
                "id": j["id"],
                "type": j["type"],
                "project_name": j.get("project_name", ""),
                "running": j["running"],
                "done": j["done"],
                "error": j["error"],
                "cancelled": j["cancelled"],
                "progress": j["progress"],
                "percent": j.get("percent", 0),
                "started_at": j.get("started_at", ""),
            })
        # Sort: running first, then by start time desc
        jobs_list.sort(key=lambda x: (not x["running"], x["started_at"]), reverse=False)
        return jsonify({"jobs": jobs_list})


@app.route("/api/jobs/<job_id>")
def api_job_detail(job_id):
    """Get detailed status of a specific job."""
    with _jobs_lock:
        job = _jobs.get(job_id)
        if not job:
            return jsonify({"error": "Job not found"}), 404
        return jsonify({
            "id": job["id"],
            "type": job["type"],
            "project_name": job.get("project_name", ""),
            "running": job["running"],
            "done": job["done"],
            "error": job["error"],
            "cancelled": job["cancelled"],
            "progress": job["progress"],
            "percent": job.get("percent", 0),
            "lines": job["lines"][-20:],
            "started_at": job.get("started_at", ""),
        })


@app.route("/api/jobs/<job_id>/stop", methods=["POST"])
def api_job_stop(job_id):
    """Stop a specific job."""
    with _jobs_lock:
        job = _jobs.get(job_id)
        if not job:
            return jsonify({"error": "Job not found"}), 404
        if job["running"]:
            job["cancelled"] = True
            return jsonify({"ok": True})
        return jsonify({"ok": False, "msg": "Job not running"})


@app.route("/api/jobs/<job_id>/dismiss", methods=["POST"])
def api_job_dismiss(job_id):
    """Remove a finished job from the list."""
    with _jobs_lock:
        job = _jobs.get(job_id)
        if not job:
            return jsonify({"error": "Job not found"}), 404
        if job["running"]:
            return jsonify({"error": "Cannot dismiss a running job"}), 400
        del _jobs[job_id]
        return jsonify({"ok": True})


@app.route("/api/analyze")
def api_analyze():
    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data. Run scan first."}), 404
    config = load_config()
    # Sync scan_db config with current config categories
    if config and config.get("categories"):
        db["config"] = config
    from event_agent import analyze_collection
    analysis = analyze_collection(db)
    # Filter to only categories in current config
    if config and isinstance(config.get("categories"), list):
        valid_ids = {c["id"] for c in config["categories"]}
        analysis["categories"] = [c for c in analysis.get("categories", []) if c["id"] in valid_ids]
        analysis["total_target"] = sum(c.get("target", 0) for c in analysis["categories"])
        analysis["total_qualified"] = sum(c.get("count", 0) for c in analysis["categories"])
    return jsonify(analysis)


@app.route("/api/auto-select", methods=["POST"])
def api_auto_select():
    if _any_job_running():
        return jsonify({"error": "A task is already running"}), 409

    data = request.json or {}
    strategy = data.get("strategy", "balanced")
    threshold = data.get("sim_threshold", 0.85)
    config = load_config()
    proj_name = (config.get("project_name") or config.get("event_name") or "") if config else ""
    job_id = _create_job("auto-select", proj_name)

    def run_select():
        try:
            _reset_task("auto-select")
            db = load_scan_db()
            if not db:
                _finish_task("No scan data.")
                return

            # Inject current config so auto_select sees up-to-date categories
            current_config = load_config()
            if current_config:
                db["config"] = current_config

            # Monkey-patch print to capture progress
            import builtins
            _orig_print = builtins.print
            def _capture_print(*args, **kwargs):
                if _is_cancelled():
                    raise InterruptedError("Cancelled by user")
                line = " ".join(str(a) for a in args)
                # Skip internal/traceback lines
                skip = line.lstrip().startswith(("File \"", "Traceback (", "json.decoder"))
                if not skip:
                    _update_task(line.strip())
                _orig_print(*args, **kwargs)
            builtins.print = _capture_print

            from event_agent import auto_select
            db, report = auto_select(db, strategy=strategy, sim_threshold=threshold)

            builtins.print = _orig_print

            save_scan_db(db)

            _update_task(f"Done! Selected {report['total_selected']} images.")
            _finish_task()
        except InterruptedError:
            builtins.print = _orig_print
            _update_task("Stopped by user.")
            _finish_task("Cancelled")
        except Exception as e:
            try:
                builtins.print = _orig_print
            except Exception:
                pass
            _finish_task(str(e))

    threading.Thread(target=run_select, daemon=True).start()
    return jsonify({"ok": True, "job_id": job_id})


@app.route("/api/quick-fill", methods=["POST"])
def api_quick_fill():
    """Quick fill: select top images by quality score without vector diversity (fast)."""
    if _any_job_running():
        return jsonify({"error": "A task is already running"}), 409

    config = load_config()
    proj_name = (config.get("project_name") or config.get("event_name") or "") if config else ""
    job_id = _create_job("auto-select", proj_name)

    def run_quick():
        try:
            _reset_task("auto-select")
            db = load_scan_db()
            if not db:
                _finish_task("No scan data.")
                return

            current_config = load_config()
            if current_config:
                db["config"] = current_config

            config = db.get("config", {})
            images = db.get("images", [])
            target_per_cat = config.get("target_per_category", 75)

            # Get categories
            from event_agent import get_categories_from_config, compute_quality_score, get_event_type, EVENT_KNOWLEDGE
            from ranking_engine import RankingEngine
            categories = get_categories_from_config(config, images)
            event_type = get_event_type(config)
            knowledge = EVENT_KNOWLEDGE.get(event_type, EVENT_KNOWLEDGE.get("photo_book"))
            weights = knowledge.get("quality_weights", {})
            # Inject taste quiz so scoring can use it
            taste_quiz = config.get("taste_quiz")
            if taste_quiz:
                weights["_taste_quiz"] = taste_quiz

            # Group by category (only face-matched when face names are configured)
            # Skip suppressed duplicates — only representatives and unclustered images compete
            face_names = config.get("face_names", [])
            unlimited = config.get("unlimited_mode", False)
            is_manual_cat = config.get("categorization") == "manual"
            suppressed_count = 0
            by_cat = {}
            for img in images:
                if img.get("status") == "rejected":
                    continue
                # Skip suppressed duplicates (non-representative cluster members)
                if img.get("suppressed_by"):
                    suppressed_count += 1
                    continue
                has_face = img.get("has_target_face")
                if face_names and not has_face:
                    if is_manual_cat:
                        # Manual templates: include non-face media with penalty
                        # (action shots / videos may lack recognizable faces)
                        img["_no_face_penalty"] = True
                    else:
                        # Auto-categorized: require face match (images and videos)
                        continue
                cat = img.get("category")
                if cat:
                    by_cat.setdefault(cat, []).append(img)
            if suppressed_count:
                _update_task(f"Skipped {suppressed_count} suppressed duplicates")

            # ── Ranking Engine setup ──
            from curate import compute_image_vector, compute_dhash

            ranker = RankingEngine(
                face_names=face_names,
                is_manual_cat=is_manual_cat,
            )

            # Vector/dHash accessors: load from npz sidecar, fallback to computation
            _vec_cache = load_image_vectors()  # pre-load all vectors from npz
            _dhash_cache = {}

            def _get_vector(img):
                h = img.get("hash")
                if h and h in _vec_cache:
                    v = _vec_cache[h]
                    return np.array(v, dtype=np.float32) if isinstance(v, list) else v
                # Legacy: try persisted vector in entry (pre-migration data)
                v = img.get("image_vector")
                if v is not None:
                    v = np.array(v, dtype=np.float32) if isinstance(v, list) else v
                    if h:
                        _vec_cache[h] = v
                    return v
                # Fallback: compute from file (images only — videos get vectors during scan)
                fpath = img.get("path", "").replace("/", os.sep)
                if not fpath or not os.path.isfile(fpath) or img.get("media_type") == "video":
                    return None
                v = compute_image_vector(fpath)
                if h and v is not None:
                    _vec_cache[h] = v
                return v

            def _get_dhash(img):
                h = img.get("hash")
                if h and h in _dhash_cache:
                    return _dhash_cache[h]
                # Try persisted dHash first
                dh = img.get("dhash")
                if dh is not None:
                    if h:
                        _dhash_cache[h] = dh
                    return dh
                # Fallback: compute from file (images only — videos get dHash during scan)
                fpath = img.get("path", "").replace("/", os.sep)
                if not fpath or not os.path.isfile(fpath) or img.get("media_type") == "video":
                    return None
                dh = compute_dhash(fpath)
                if h and dh is not None:
                    _dhash_cache[h] = dh
                return dh

            # Load CLIP vectors for diversity scoring
            _clip_vec_data = {}
            if clip_engine is not None and os.path.isfile(CLIP_VECTORS_PATH):
                try:
                    _clip_vec_data = clip_engine.load_vectors(CLIP_VECTORS_PATH)
                except Exception:
                    pass

            def _get_clip_vector(img):
                h = img.get("hash")
                if h and h in _clip_vec_data:
                    cv = _clip_vec_data[h]
                    return np.array(cv, dtype=np.float32) if isinstance(cv, list) else cv
                return None

            # Pre-compute vectors for rated images (for visual preference learning)
            _update_task("Building preference model...")
            vector_lookup = {}
            rated_imgs = [i for i in images if i.get("preference") in ("like", "dislike")]
            for img in rated_imgs:
                v = _get_vector(img)
                if v is not None:
                    vector_lookup[img.get("hash")] = v

            ranker.learn_from_feedback(images, vector_lookup)
            if ranker._n_likes + ranker._n_dislikes > 0:
                _update_task(f"Preference model: {ranker._n_likes} likes, {ranker._n_dislikes} dislikes"
                             + (", visual similarity active" if ranker._session_vector is not None else ""))

            dedup_skipped = [0]

            total_selected = 0
            total_videos = 0
            total_cats = len(categories)
            for cat_idx, cat in enumerate(categories):
                if _is_cancelled():
                    _update_task("Stopped by user.")
                    _finish_task("Cancelled")
                    return
                _update_task_percent(int(cat_idx * 100 / total_cats) if total_cats else 0)

                cid = cat["id"]
                img_target = cat.get("target", target_per_cat)
                vid_target = cat.get("video_target", 0)
                pool = by_cat.get(cid, [])

                img_pool = [i for i in pool if i.get("media_type") != "video"]
                vid_pool = [i for i in pool if i.get("media_type") == "video"]

                def _score_and_select(candidates, target_count, media_label):
                    # Hard filters first
                    if face_names:
                        age_days_to = cat.get("age_days_to", 99999)
                        max_dist = 0.45 if age_days_to <= 365 else 0.50
                        if is_manual_cat:
                            for img in candidates:
                                if not img.get("has_target_face"):
                                    img["_no_face_penalty"] = True
                        else:
                            if media_label == "images":
                                candidates = [i for i in candidates
                                              if i.get("face_distance") is not None and i.get("face_distance") <= max_dist]

                    already = [i for i in candidates if i.get("status") == "selected"]
                    unselected = [i for i in candidates if i.get("status") != "selected"]

                    # Register already-selected in ranker state
                    for img in already:
                        v = _get_vector(img)
                        dh = _get_dhash(img)
                        cv = _get_clip_vector(img)
                        ranker.register_selected(img, v, dh, clip_vector=cv)

                    # Score all unselected with ranking engine
                    for img in unselected:
                        v = _get_vector(img)
                        dh = _get_dhash(img)
                        cv = _get_clip_vector(img)
                        ranker.score(img, vector=v, dhash=dh, clip_vector=cv, quality_weights=weights)

                    unselected.sort(key=lambda x: x.get("_score", 0), reverse=True)

                    if unlimited:
                        remaining = len(unselected)
                    else:
                        remaining = max(0, target_count - len(already))

                    # Select top candidates, re-checking duplicates incrementally
                    # (batch scoring can't catch intra-round duplicates)
                    picked = 0
                    for img in unselected:
                        if picked >= remaining:
                            break
                        v = _get_vector(img)
                        dh = _get_dhash(img)
                        # Re-compute dedup against everything selected so far
                        dup = ranker._compute_duplicate_penalty(v, dh, path=img.get("path"))
                        if dup >= 100:
                            dedup_skipped[0] += 1
                            continue
                        img["status"] = "selected"
                        picked += 1
                        cv = _get_clip_vector(img)
                        ranker.register_selected(img, v, dh, clip_vector=cv)

                    return len(already) + picked, picked

                img_total, img_picked = _score_and_select(img_pool, img_target, "images")
                vid_total, vid_picked = _score_and_select(vid_pool, vid_target, "videos")

                total_selected += img_total
                total_videos += vid_total

                parts = []
                if unlimited:
                    parts.append(f"{img_picked} images, {vid_picked} videos")
                else:
                    parts.append(f"{img_total}/{img_target} images")
                    if vid_target > 0 or vid_picked > 0:
                        parts.append(f"{vid_total}/{vid_target} videos")
                _update_task(f"{cat.get('display', cid)}: {', '.join(parts)}")

            # Clean temp scores
            for img in images:
                img.pop("_score", None)
                img.pop("_score_breakdown", None)
                img.pop("_no_face_penalty", None)

            save_scan_db(db)
            msg = f"Done! {total_selected} images"
            if total_videos > 0:
                msg += f" + {total_videos} videos"
            msg += " selected."
            if dedup_skipped[0] > 0:
                msg += f" ({dedup_skipped[0]} near-duplicates skipped)"
            _update_task(msg)
            _finish_task()
        except Exception as e:
            _finish_task(str(e))

    threading.Thread(target=run_quick, daemon=True).start()
    return jsonify({"ok": True, "job_id": job_id})


@app.route("/api/images")
def api_images():
    """Get images for gallery. Supports pagination and filters."""
    db = load_scan_db()
    if not db:
        return jsonify({"images": [], "total": 0})

    images = db["images"]
    category = request.args.get("category")
    status = request.args.get("status")
    source = request.args.get("source")
    offset = int(request.args.get("offset", 0))
    limit = int(request.args.get("limit", 200))

    filtered = images
    if category:
        filtered = [i for i in filtered if i.get("category") == category]
    if status:
        filtered = [i for i in filtered if i.get("status") == status]
    if source:
        filtered = [i for i in filtered if i.get("source_label") == source]
    location = request.args.get("location")
    if location:
        filtered = [i for i in filtered if i.get("location") == location]
    # By default hide suppressed duplicates; pass show_duplicates=1 to include them
    show_dups = request.args.get("show_duplicates", "0") == "1"
    if not show_dups:
        filtered = [i for i in filtered if not i.get("suppressed_by")]

    total = len(filtered)
    page = filtered[offset:offset + limit]

    # Strip thumbnails for lighter response if requested
    compact = request.args.get("compact") == "1"
    if compact:
        page = [{k: v for k, v in img.items() if k != "thumb"} for img in page]

    return jsonify({"images": page, "total": total, "offset": offset})


@app.route("/api/images/move", methods=["POST"])
def api_images_move():
    """Move images between categories or to pool."""
    data = request.json
    hashes = set(data.get("hashes", []))
    to_category = data.get("to_category")
    to_status = data.get("to_status", "qualified")

    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data"}), 404

    moved = 0
    for img in db["images"]:
        if img["hash"] in hashes:
            if to_category:
                img["category"] = to_category
            img["status"] = to_status
            if to_status == "pool":
                img["reject_reason"] = "manual_reject"
            else:
                img["reject_reason"] = None
            moved += 1

    save_scan_db(db)
    return jsonify({"ok": True, "moved": moved})


@app.route("/api/curate/save", methods=["POST"])
def api_curate_save():
    """Bulk-save gallery changes. Expects {changes: [{hash, category, status}, ...]}."""
    data = request.json or {}
    items = data.get("changes", [])
    if not items:
        return jsonify({"ok": True, "updated": 0})

    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data"}), 404

    lookup = {}
    for item in items:
        lookup[item["hash"]] = item

    updated = 0
    for img in db["images"]:
        change = lookup.get(img["hash"])
        if not change:
            continue
        img["category"] = change.get("category") or img.get("category")
        img["status"] = change.get("status", img.get("status", "pool"))
        if img["status"] == "pool":
            img["reject_reason"] = img.get("reject_reason") or "manual_reject"
        elif img["status"] == "qualified":
            img["reject_reason"] = None
        updated += 1

    save_scan_db(db)

    # Also persist to project dir if a project is loaded
    config = load_config()
    if config:
        proj_name = config.get("project_name") or config.get("event_name") or ""
        if proj_name:
            pdir = os.path.join(PROJECTS_DIR, proj_name)
            if os.path.isdir(pdir):
                import shutil
                shutil.copy2(SCAN_DB_PATH, os.path.join(pdir, "scan_db.json"))

    return jsonify({"ok": True, "updated": updated})


@app.route("/api/images/preference", methods=["POST"])
def api_images_preference():
    """Set like/dislike preference for an image. Accepts {hash, preference} where
    preference is 'like', 'dislike', or null to clear."""
    data = request.json or {}
    img_hash = data.get("hash")
    preference = data.get("preference")  # 'like', 'dislike', or None

    if not img_hash:
        return jsonify({"error": "Missing hash"}), 400
    if preference not in ("like", "dislike", None):
        return jsonify({"error": "Invalid preference, must be 'like', 'dislike', or null"}), 400

    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data"}), 404

    found = False
    for img in db["images"]:
        if img["hash"] == img_hash:
            if preference is None:
                img.pop("preference", None)
            else:
                img["preference"] = preference
            found = True
            break

    if not found:
        return jsonify({"error": "Image not found"}), 404

    save_scan_db(db)
    return jsonify({"ok": True, "hash": img_hash, "preference": preference})


@app.route("/api/preferences/summary")
def api_preferences_summary():
    """Return counts of liked/disliked/unrated images."""
    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data"}), 404

    images = db.get("images", [])
    liked = sum(1 for i in images if i.get("preference") == "like")
    disliked = sum(1 for i in images if i.get("preference") == "dislike")
    unrated = len(images) - liked - disliked

    # Also collect per-image preference data for future model training
    preferences = []
    for img in images:
        pref = img.get("preference")
        if pref:
            preferences.append({
                "hash": img["hash"],
                "preference": pref,
                "category": img.get("category"),
                "face_count": img.get("face_count", 0),
                "has_target_face": img.get("has_target_face", False),
                "size_kb": img.get("size_kb", 0),
                "width": img.get("width", 0),
                "height": img.get("height", 0),
                "device": img.get("device"),
                "source_label": img.get("source_label"),
            })

    return jsonify({
        "liked": liked,
        "disliked": disliked,
        "unrated": unrated,
        "total": len(images),
        "preferences": preferences,
    })


@app.route("/api/preferences/quiz", methods=["POST"])
def api_preferences_quiz():
    """Save user taste quiz answers to the project config."""
    data = request.json or {}
    config = load_config()
    if not config:
        return jsonify({"error": "No config"}), 404
    config["taste_quiz"] = data
    save_config(config)
    return jsonify({"ok": True})


@app.route("/api/images/serve/<img_hash>")
def api_images_serve(img_hash):
    """Serve a full-size image or video by hash."""
    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data"}), 404
    # MIME types for video formats Flask may not auto-detect
    _VIDEO_MIMES = {
        ".mp4": "video/mp4", ".m4v": "video/mp4",
        ".mov": "video/quicktime", ".avi": "video/x-msvideo",
        ".mkv": "video/x-matroska", ".webm": "video/webm",
        ".wmv": "video/x-ms-wmv", ".mpg": "video/mpeg",
        ".mpeg": "video/mpeg", ".3gp": "video/3gpp",
    }
    for img in db["images"]:
        if img["hash"] == img_hash:
            fpath = img["path"].replace("/", os.sep)
            if os.path.isfile(fpath):
                ext = os.path.splitext(fpath)[1].lower()
                mime = _VIDEO_MIMES.get(ext)
                # conditional=True enables Range requests (required for video seeking/playback)
                if mime:
                    return send_file(fpath, mimetype=mime, conditional=True)
                return send_file(fpath, conditional=True)
            break
    return jsonify({"error": "Not found"}), 404


@app.route("/api/images/select", methods=["POST"])
def api_images_select():
    """Select or deselect individual images."""
    data = request.json
    hashes = set(data.get("hashes", []))
    action = data.get("action", "select")  # "select" or "deselect"

    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data"}), 404

    changed = 0
    for img in db["images"]:
        if img["hash"] in hashes:
            if action == "select":
                img["status"] = "selected"
            else:
                img["status"] = "qualified"
            changed += 1

    save_scan_db(db)
    return jsonify({"ok": True, "changed": changed})


@app.route("/api/selections/reset", methods=["POST"])
def api_selections_reset():
    """Reset all selected images back to qualified."""
    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data"}), 404
    count = 0
    for img in db["images"]:
        if img.get("status") == "selected":
            img["status"] = "qualified"
            count += 1
    save_scan_db(db)
    return jsonify({"ok": True, "reset": count})


@app.route("/api/locations/summary")
@login_required
def api_locations_summary():
    """Get distinct location values from scan_db."""
    db = load_scan_db()
    if not db:
        return jsonify({"locations": []})
    locs = sorted(set(img.get("location") for img in db.get("images", []) if img.get("location")))
    counts = {}
    for img in db.get("images", []):
        loc = img.get("location")
        if loc:
            counts[loc] = counts.get(loc, 0) + 1
    return jsonify({"locations": [{"name": loc, "count": counts.get(loc, 0)} for loc in locs]})


@app.route("/api/categories/summary")
def api_categories_summary():
    """Get per-category counts for selection UI."""
    db = load_scan_db()
    config = load_config()
    if not db:
        return jsonify([])

    images = db["images"]
    config_cats = config.get("categories", []) if config else []
    valid_ids = {c["id"] for c in config_cats} if config_cats else set()

    cat_data = {}
    for c in config_cats:
        cat_data[c["id"]] = {
            "id": c["id"],
            "display": c.get("display", c["id"]),
            "target": c.get("target", config.get("target_per_category", 75)),
            "qualified": 0,
            "selected": 0,
            "total": 0,
        }

    for img in images:
        cat = img.get("category")
        if cat and cat in cat_data:
            cat_data[cat]["total"] += 1
            if img.get("status") == "selected":
                cat_data[cat]["selected"] += 1
            elif img.get("status") == "qualified":
                cat_data[cat]["qualified"] += 1

    return jsonify(list(cat_data.values()))


@app.route("/api/categories/update-target", methods=["POST"])
def api_categories_update_target():
    """Update target for a specific category."""
    data = request.json
    cat_id = data.get("id")
    new_target = data.get("target")

    config = load_config()
    if not config or not config.get("categories"):
        return jsonify({"error": "No config"}), 400

    for c in config["categories"]:
        if c["id"] == cat_id:
            c["target"] = new_target
            break

    save_config(config)
    return jsonify({"ok": True})


@app.route("/api/categories/substitutes")
def api_categories_substitutes():
    """Get substitute candidates for under-filled categories.
    Returns qualified+suppressed images with target face, sorted by quality.
    For categories below target, these are images a user can manually add."""
    db = load_scan_db()
    config = load_config()
    if not db:
        return jsonify({"categories": []})

    images = db["images"]
    config_cats = config.get("categories", []) if config else []
    face_names = config.get("face_names", [])
    target = int(request.args.get("target", config.get("target_per_category", 75)))
    cat_id = request.args.get("category")  # optional: filter to one category
    limit = int(request.args.get("limit", 200))

    results = []
    for c in config_cats:
        cid = c["id"]
        if cat_id and cid != cat_id:
            continue
        cat_target = c.get("target", target)
        cat_imgs = [i for i in images if i.get("category") == cid]
        selected = [i for i in cat_imgs if i.get("status") == "selected"]
        gap = cat_target - len(selected)
        if gap <= 0 and not cat_id:
            continue  # fully filled, skip unless specifically requested

        # Substitutes: qualified or suppressed, with target face
        subs = []
        for i in cat_imgs:
            if i.get("status") in ("qualified", "pool"):
                has_face = (not face_names or
                            any(n in i.get("faces_found", []) for n in face_names))
                if has_face:
                    subs.append(i)

        # Sort by photo_grade composite (best first)
        def _grade(img):
            g = img.get("photo_grade")
            return g.get("composite", 0) if isinstance(g, dict) else 0
        subs.sort(key=_grade, reverse=True)

        results.append({
            "id": cid,
            "display": c.get("display", cid),
            "target": cat_target,
            "selected": len(selected),
            "gap": gap,
            "substitutes": subs[:limit],
        })

    return jsonify({"categories": results})


@app.route("/api/export", methods=["POST"])
def api_export():
    """Export selected/qualified images to output folder."""
    if _any_job_running():
        return jsonify({"error": "A task is already running"}), 409

    data = request.json or {}
    default_downloads = os.path.join(os.path.expanduser("~"), "Downloads", "E-z Photo Collection")
    output_dir = data.get("output_dir", default_downloads)
    status_filter = data.get("status", "selected")

    config = load_config()
    proj_name = (config.get("project_name") or config.get("event_name") or "") if config else ""
    job_id = _create_job("export", proj_name)

    def run_export():
        try:
            import shutil
            _reset_task("export")
            db = load_scan_db()
            if not db:
                _finish_task("No scan data.")
                return

            os.makedirs(output_dir, exist_ok=True)
            exported = 0

            images = [i for i in db["images"] if i.get("status") == status_filter]
            if not images:
                _update_task("No selected media found. Use the Select step first.")
                _finish_task("No media to export")
                return

            total = len(images)
            vid_count = sum(1 for i in images if i.get("media_type") == "video")
            img_count = total - vid_count
            label = f"{img_count} images"
            if vid_count:
                label += f" + {vid_count} videos"
            _update_task(f"Exporting {label}...")

            skipped = 0
            for i, img in enumerate(images):
                if _is_cancelled():
                    _update_task("Stopped by user.")
                    _finish_task("Cancelled")
                    return

                src = img["path"].replace("/", os.sep)
                if not os.path.isfile(src):
                    skipped += 1
                    continue

                cat = img.get("category", "uncategorized")
                dest_dir = os.path.join(output_dir, cat)
                os.makedirs(dest_dir, exist_ok=True)

                dest = os.path.join(dest_dir, img["filename"])
                if os.path.exists(dest):
                    base, ext = os.path.splitext(img["filename"])
                    j = 1
                    while os.path.exists(dest):
                        dest = os.path.join(dest_dir, f"{base}_{j}{ext}")
                        j += 1

                shutil.copy2(src, dest)
                exported += 1

                if (i + 1) % 50 == 0:
                    _update_task_percent(int((i + 1) * 100 / total) if total else 0)
                    _update_task(f"Exported {exported}/{total}...")

            msg = f"Done! {exported} files exported to {output_dir}"
            if skipped:
                msg += f" ({skipped} skipped — source files not found, check if all drives are connected)"
            _update_task(msg)
            _finish_task()

        except Exception as e:
            _finish_task(str(e))

    threading.Thread(target=run_export, daemon=True).start()
    return jsonify({"ok": True, "job_id": job_id})


@app.route("/api/export/pptx", methods=["POST"])
def api_export_pptx():
    """Generate a PowerPoint presentation from selected images."""
    if _any_job_running():
        return jsonify({"error": "A task is already running"}), 409

    data = request.json or {}
    opt_captions = data.get("captions", True)
    opt_dividers = data.get("dividers", True)
    opt_sort = data.get("sort", "date")  # "date" or "score"

    config = load_config()
    proj_name = (config.get("project_name") or config.get("event_name") or "Photo Collection") if config else "Photo Collection"
    job_id = _create_job("pptx_export", proj_name)

    def run_pptx_export():
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from io import BytesIO
            from PIL import Image

            _reset_task("pptx_export")
            _update_task("Preparing presentation...")

            db = load_scan_db()
            if not db:
                _finish_task("No scan data.")
                return

            cfg = load_config() or {}
            categories = cfg.get("categories", [])
            cat_display = {}
            for cat in categories:
                cat_display[cat["id"]] = cat.get("display", cat["id"])

            # Filter to selected images, exclude videos (can't embed in slides)
            _IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif",
                           ".gif", ".webp", ".heic", ".heif"}
            all_selected = [i for i in db["images"] if i.get("status") == "selected"]
            images = []
            skipped_videos = 0
            for i in all_selected:
                ext = os.path.splitext(i["filename"])[1].lower()
                if i.get("media_type") == "video" or ext not in _IMAGE_EXTS:
                    skipped_videos += 1
                else:
                    images.append(i)

            if not images and not skipped_videos:
                _update_task("No selected images. Use the Select step first.")
                _finish_task("No images to export")
                return
            if not images:
                _update_task(f"All {skipped_videos} selected items are videos (cannot embed in PowerPoint).")
                _finish_task("No images to export")
                return

            # Group by category, preserving category order from config
            from collections import OrderedDict
            cat_order = [c["id"] for c in categories]
            by_cat = OrderedDict()
            for cid in cat_order:
                by_cat[cid] = []
            by_cat["uncategorized"] = []
            for img in images:
                c = img.get("category", "uncategorized")
                if c not in by_cat:
                    by_cat[c] = []
                by_cat[c].append(img)
            # Remove empty categories
            by_cat = OrderedDict((k, v) for k, v in by_cat.items() if v)

            total = len(images)
            label = f"{total} photo{'s' if total != 1 else ''}"
            if skipped_videos:
                label += f" ({skipped_videos} video{'s' if skipped_videos != 1 else ''} skipped)"
            _update_task(f"Building presentation with {label}...")

            # ── Create presentation (16:9 widescreen) ──
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # ── Color palette ──
            BG_DARK = RGBColor(0x1A, 0x20, 0x2C)
            TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
            TEXT_SOFT = RGBColor(0xA0, 0xAE, 0xC0)

            def _set_slide_bg(slide, color):
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = color

            def _add_textbox(slide, left, top, width, height, text, font_size=18,
                             color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(font_size)
                p.font.color.rgb = color
                p.font.bold = bold
                p.alignment = alignment
                return txBox

            def _pluralize(n, word):
                return f"{n} {word}" if n == 1 else f"{n} {word}s"

            # ── Title Slide ──
            title_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_bg(title_slide, BG_DARK)
            _add_textbox(title_slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5),
                         proj_name, font_size=44, color=TEXT_WHITE, bold=True,
                         alignment=PP_ALIGN.CENTER)
            subtitle = cfg.get("template", "")
            if subtitle:
                _add_textbox(title_slide, Inches(1), Inches(3.6), Inches(11.3), Inches(0.8),
                             subtitle, font_size=22, color=TEXT_SOFT,
                             alignment=PP_ALIGN.CENTER)
            _add_textbox(title_slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.6),
                         f"{_pluralize(total, 'photo')} across {_pluralize(len(by_cat), 'chapter')}",
                         font_size=16, color=TEXT_SOFT, alignment=PP_ALIGN.CENTER)

            # ── Per-category slides ──
            done = 0
            skipped_missing = 0
            for cat_id, cat_images in by_cat.items():
                if _is_cancelled():
                    _update_task("Stopped by user.")
                    _finish_task("Cancelled")
                    return

                display_name = cat_display.get(cat_id, cat_id.replace("_", " ").title())

                # Sort images within category
                if opt_sort == "score":
                    cat_images.sort(key=lambda x: (x.get("photo_grade") or {}).get("composite", 0), reverse=True)
                else:
                    cat_images.sort(key=lambda x: x.get("date") or "")

                # Section divider slide
                if opt_dividers:
                    sec_slide = prs.slides.add_slide(prs.slide_layouts[6])
                    _set_slide_bg(sec_slide, BG_DARK)
                    _add_textbox(sec_slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
                                 display_name, font_size=40, color=TEXT_WHITE, bold=True,
                                 alignment=PP_ALIGN.CENTER)
                    _add_textbox(sec_slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.6),
                                 _pluralize(len(cat_images), "photo"),
                                 font_size=18, color=TEXT_SOFT, alignment=PP_ALIGN.CENTER)

                # Image slides
                for img in cat_images:
                    if _is_cancelled():
                        _update_task("Stopped by user.")
                        _finish_task("Cancelled")
                        return

                    src_path = img["path"].replace("/", os.sep)
                    if not os.path.isfile(src_path):
                        done += 1
                        skipped_missing += 1
                        continue

                    img_slide = prs.slides.add_slide(prs.slide_layouts[6])
                    _set_slide_bg(img_slide, BG_DARK)

                    # Load image to get dimensions for proper scaling
                    try:
                        pil_img = Image.open(src_path)
                        img_w, img_h = pil_img.size
                        pil_img.close()
                    except Exception:
                        img_w = img.get("width", 1920)
                        img_h = img.get("height", 1080)

                    if img_w <= 0: img_w = 1920
                    if img_h <= 0: img_h = 1080

                    # Scale image to fit slide, reserving space for caption
                    aspect = img_w / img_h
                    caption_h = 0.5 if opt_captions else 0
                    max_w_in = 12.733   # 13.333 - 0.6 (0.3" margin each side)
                    max_h_in = 7.5 - 0.4 - caption_h  # top margin + caption reserve

                    if aspect >= max_w_in / max_h_in:
                        pic_w = max_w_in
                        pic_h = pic_w / aspect
                    else:
                        pic_h = max_h_in
                        pic_w = pic_h * aspect

                    left = (13.333 - pic_w) / 2
                    avail_h = 7.5 - caption_h
                    top = (avail_h - pic_h) / 2

                    try:
                        img_slide.shapes.add_picture(
                            src_path,
                            Inches(left), Inches(top),
                            Inches(pic_w), Inches(pic_h)
                        )
                    except Exception:
                        # Fallback: load via PIL, convert to JPEG in memory
                        try:
                            pil_img = Image.open(src_path)
                            if pil_img.mode != "RGB":
                                pil_img = pil_img.convert("RGB")
                            buf = BytesIO()
                            pil_img.save(buf, format="JPEG", quality=92)
                            buf.seek(0)
                            pil_img.close()
                            img_slide.shapes.add_picture(
                                buf, Inches(left), Inches(top),
                                Inches(pic_w), Inches(pic_h)
                            )
                        except Exception:
                            _add_textbox(img_slide, Inches(2), Inches(3), Inches(9), Inches(1),
                                         f"Could not load: {img['filename']}",
                                         font_size=14, color=TEXT_SOFT, alignment=PP_ALIGN.CENTER)

                    # Caption
                    if opt_captions:
                        caption = img.get("date", "")
                        if caption and display_name:
                            caption = f"{display_name}  |  {caption}"
                        elif display_name:
                            caption = display_name
                        _add_textbox(img_slide, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.4),
                                     caption, font_size=11, color=TEXT_SOFT,
                                     alignment=PP_ALIGN.CENTER)

                    done += 1
                    if done % 10 == 0:
                        pct = int(done * 100 / total) if total else 0
                        _update_task_percent(pct)
                        _update_task(f"Added {done}/{total} images...")

            # ── End slide ──
            end_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_bg(end_slide, BG_DARK)
            _add_textbox(end_slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1),
                         proj_name, font_size=36, color=TEXT_WHITE, bold=True,
                         alignment=PP_ALIGN.CENTER)
            _add_textbox(end_slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.6),
                         f"{_pluralize(total, 'photo')}  |  {_pluralize(len(by_cat), 'chapter')}",
                         font_size=16, color=TEXT_SOFT, alignment=PP_ALIGN.CENTER)

            # ── Save ──
            safe_name = "".join(c if c.isalnum() or c in " _-" else "_" for c in proj_name).strip()[:80]
            out_name = f"{safe_name}.pptx" if safe_name else "presentation.pptx"
            out_path = os.path.join(PROJECT_DIR, out_name)
            _update_task("Saving presentation...")
            prs.save(out_path)

            size_mb = os.path.getsize(out_path) / 1e6
            summary = f"Done! {out_name} ({size_mb:.1f} MB) — {_pluralize(total, 'photo')}, {len(prs.slides)} slides"
            if skipped_videos:
                summary += f", {_pluralize(skipped_videos, 'video')} skipped"
            if skipped_missing:
                summary += f", {skipped_missing} files not found"
            _update_task(summary)
            # Store path for download
            with _jobs_lock:
                job = _jobs.get(job_id)
                if job:
                    job["result_file"] = out_path
            _finish_task()

        except Exception as e:
            import traceback
            traceback.print_exc()
            _finish_task(str(e))

    threading.Thread(target=run_pptx_export, daemon=True).start()
    return jsonify({"ok": True, "job_id": job_id})


@app.route("/api/export/pptx/download")
def api_export_pptx_download():
    """Download the last generated PowerPoint file."""
    # Find the most recent pptx job
    with _jobs_lock:
        pptx_jobs = [j for j in _jobs.values() if j.get("type") == "pptx_export" and j.get("result_file")]
    if not pptx_jobs:
        return jsonify({"error": "No presentation file available. Generate one first."}), 404
    latest = max(pptx_jobs, key=lambda j: j.get("started_at", ""))
    fpath = latest["result_file"]
    if not os.path.isfile(fpath):
        return jsonify({"error": "Presentation file not found."}), 404
    return send_file(fpath, as_attachment=True, download_name=os.path.basename(fpath))


@app.route("/api/ref-faces")
def api_ref_faces():
    """List reference face folders."""
    ref_dir = os.path.join(PROJECT_DIR, "ref_faces")
    if not os.path.isdir(ref_dir):
        return jsonify([])
    # Check if encodings cache exists
    cache_path = os.path.join(ref_dir, "_encodings_cache.json")
    cached_persons = set()
    if os.path.isfile(cache_path):
        try:
            import json as _json
            with open(cache_path, "r", encoding="utf-8") as f:
                cache = _json.load(f)
            cached_persons = set(cache.keys())
        except Exception:
            pass

    # Also check library encoding files for diversity/verification info
    lib_info = {}
    for person_dir_name in os.listdir(ref_dir):
        enc_file = os.path.join(ref_dir, person_dir_name, "_face_encodings.json")
        if not os.path.isfile(enc_file):
            # Check library
            lib_enc = os.path.join(FACE_LIBRARY_DIR, person_dir_name, "_face_encodings.json")
            if os.path.isfile(lib_enc):
                enc_file = lib_enc
            else:
                continue
        try:
            import json as _json
            with open(enc_file, "r", encoding="utf-8") as f:
                edata = _json.load(f)
            lib_info[person_dir_name] = {
                "diversity_score": edata.get("diversity_score"),
                "verified_photos": edata.get("verified_photos", []),
                "encoding_count": len(edata.get("encodings", [])),
            }
        except Exception:
            pass

    result = []
    for person in sorted(os.listdir(ref_dir)):
        pdir = os.path.join(ref_dir, person)
        if os.path.isdir(pdir):
            photos = [f for f in os.listdir(pdir) if os.path.splitext(f)[1].lower() in {".jpg", ".jpeg", ".png"}]
            entry = {"name": person, "photo_count": len(photos), "has_encodings": person in cached_persons}
            if person in lib_info:
                entry["diversity_score"] = lib_info[person].get("diversity_score")
                entry["verified_photos"] = lib_info[person].get("verified_photos", [])
                entry["encoding_count"] = lib_info[person].get("encoding_count", 0)
            result.append(entry)
    return jsonify(result)


@app.route("/api/ref-faces/upload", methods=["POST"])
def api_ref_faces_upload():
    """Upload reference face photos for a person.
    Accepts multipart form with 'person' field and 'photos' files,
    or JSON with 'person' and base64-encoded 'photos' array.
    """
    ref_dir = os.path.join(PROJECT_DIR, "ref_faces")
    os.makedirs(ref_dir, exist_ok=True)

    if request.content_type and "multipart" in request.content_type:
        person = request.form.get("person", "").strip().lower()
        if not person:
            return jsonify({"error": "Person name is required"}), 400
        pdir = os.path.join(ref_dir, person)
        os.makedirs(pdir, exist_ok=True)
        saved = []
        for f in request.files.getlist("photos"):
            if f.filename:
                safe_name = os.path.basename(f.filename)
                dest = os.path.join(pdir, safe_name)
                f.save(dest)
                auto_rotate_image(dest)
                saved.append(safe_name)
        return jsonify({"ok": True, "person": person, "saved": saved})
    else:
        data = request.json or {}
        person = data.get("person", "").strip().lower()
        if not person:
            return jsonify({"error": "Person name is required"}), 400
        pdir = os.path.join(ref_dir, person)
        os.makedirs(pdir, exist_ok=True)
        saved = []
        for i, photo_b64 in enumerate(data.get("photos", [])):
            img_data = base64.b64decode(photo_b64)
            fname = f"ref_{i+1:03d}.jpg"
            with open(os.path.join(pdir, fname), "wb") as f:
                f.write(img_data)
            saved.append(fname)
        return jsonify({"ok": True, "person": person, "saved": saved})


@app.route("/api/ref-faces/<person>/replace", methods=["POST"])
def api_ref_faces_replace(person):
    """Replace a single reference photo. Expects multipart form with 'filename' and 'photo'."""
    ref_dir = os.path.join(PROJECT_DIR, "ref_faces")
    pdir = os.path.join(ref_dir, person)
    if not os.path.isdir(pdir):
        return jsonify({"error": "Person not found"}), 404

    old_name = request.form.get("filename", "").strip()
    photo = request.files.get("photo")
    if not old_name or not photo:
        return jsonify({"error": "filename and photo are required"}), 400

    # Remove old file
    old_path = os.path.join(pdir, old_name)
    if os.path.isfile(old_path):
        os.remove(old_path)

    # Save new file with same name, auto-rotate
    new_path = os.path.join(pdir, old_name)
    photo.save(new_path)
    auto_rotate_image(new_path)
    return jsonify({"ok": True, "filename": old_name})


@app.route("/api/ref-faces/<person>/verify-photos", methods=["POST"])
def api_ref_faces_verify_photos(person):
    """Verify specific photos only. Expects JSON {filenames: [...]}."""
    data = request.json or {}
    filenames = data.get("filenames", [])
    pdir = os.path.join(PROJECT_DIR, "ref_faces", person)
    if not os.path.isdir(pdir):
        return jsonify({"error": "Person not found"}), 404
    results = {}
    for fn in filenames:
        fpath = os.path.join(pdir, fn)
        if os.path.isfile(fpath):
            results[fn] = _verify_single_photo(fpath)
        else:
            results[fn] = {"status": "error", "message": "File not found"}
    return jsonify({"ok": True, "results": results})


@app.route("/api/ref-faces/<person>/rotate", methods=["POST"])
def api_ref_faces_rotate(person):
    """Rotate a reference photo 90 degrees clockwise. Expects JSON {filename, direction}."""
    data = request.json or {}
    filename = data.get("filename", "")
    direction = data.get("direction", "cw")  # cw or ccw

    pdir = os.path.join(PROJECT_DIR, "ref_faces", person)
    fpath = os.path.join(pdir, filename)
    if not os.path.isfile(fpath):
        return jsonify({"error": "File not found"}), 404

    from PIL import Image
    try:
        img = Image.open(fpath)
        angle = -90 if direction == "cw" else 90
        rotated = img.rotate(angle, expand=True)
        rotated.save(fpath, quality=95)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/ref-faces/<person>/photo/<filename>", methods=["DELETE"])
def api_ref_faces_delete_photo(person, filename):
    """Delete a single reference photo."""
    pdir = os.path.join(PROJECT_DIR, "ref_faces", person)
    fpath = os.path.join(pdir, filename)
    if os.path.isfile(fpath):
        os.remove(fpath)
    return jsonify({"ok": True})


@app.route("/api/ref-faces/<person>", methods=["DELETE"])
def api_ref_faces_delete(person):
    """Delete all reference photos for a person."""
    pdir = os.path.join(PROJECT_DIR, "ref_faces", person)
    if os.path.isdir(pdir):
        shutil.rmtree(pdir)
    return jsonify({"ok": True})


# ── Face Library (global, persists across projects) ──────────────────────────

FACE_LIBRARY_DIR = os.path.join(PROJECT_DIR, "face_library")


@app.route("/api/face-library")
def api_face_library():
    """List all people in the global face library."""
    if not os.path.isdir(FACE_LIBRARY_DIR):
        return jsonify([])
    result = []
    for person in sorted(os.listdir(FACE_LIBRARY_DIR)):
        pdir = os.path.join(FACE_LIBRARY_DIR, person)
        if not os.path.isdir(pdir):
            continue
        photos = [f for f in os.listdir(pdir) if os.path.splitext(f)[1].lower() in {".jpg", ".jpeg", ".png"}]
        result.append({"name": person, "photo_count": len(photos)})
    return jsonify(result)


@app.route("/api/face-library/save", methods=["POST"])
def api_face_library_save():
    """Save a person from the current project's ref_faces to the global library."""
    data = request.json or {}
    person = data.get("person", "").strip().lower()
    if not person:
        return jsonify({"error": "Person name required"}), 400

    src_dir = os.path.join(PROJECT_DIR, "ref_faces", person)
    if not os.path.isdir(src_dir):
        return jsonify({"error": f"No reference photos for {person}"}), 404

    # Verify each photo and only save ones with valid face encodings
    try:
        import face_recognition as fr
        import numpy as np
    except ImportError:
        return jsonify({"error": "face_recognition library not installed"}), 500

    dst_dir = os.path.join(FACE_LIBRARY_DIR, person)
    # Clear old library photos for this person
    if os.path.isdir(dst_dir):
        shutil.rmtree(dst_dir)
    os.makedirs(dst_dir, exist_ok=True)

    image_exts = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"}
    saved = 0
    skipped = 0
    encodings_list = []
    saved_photo_names = []

    for fname in sorted(os.listdir(src_dir)):
        if os.path.splitext(fname)[1].lower() not in image_exts:
            continue
        fpath = os.path.join(src_dir, fname)
        try:
            from PIL import Image
            pil_img = Image.open(fpath).convert("RGB")
            max_dim = 1600
            w, h = pil_img.size
            if w > max_dim or h > max_dim:
                pil_img.thumbnail((max_dim, max_dim), Image.LANCZOS)
            arr = np.array(pil_img)
            locations = fr.face_locations(arr, model="hog")
            if not locations:
                skipped += 1
                continue
            if len(locations) > 1:
                areas = [(b-t)*(r-l) for t, r, b, l in locations]
                best = areas.index(max(areas))
                enc = fr.face_encodings(arr, [locations[best]])
            else:
                enc = fr.face_encodings(arr, locations)
            if not enc:
                skipped += 1
                continue
            # Photo is validated — save it
            shutil.copy2(fpath, os.path.join(dst_dir, fname))
            encodings_list.append(enc[0].tolist())
            saved_photo_names.append(fname)
            saved += 1
        except Exception:
            skipped += 1

    if not encodings_list:
        # Clean up empty dir
        if os.path.isdir(dst_dir):
            shutil.rmtree(dst_dir)
        return jsonify({"error": f"No valid face photos found for {person}. Verify faces first."}), 400

    # Compute diversity score
    import numpy as np
    diversity_score = 0.0
    if len(encodings_list) >= 3:
        encs_np = [np.array(e) for e in encodings_list]
        dists = []
        for i in range(len(encs_np)):
            for j in range(i+1, len(encs_np)):
                dists.append(float(np.linalg.norm(encs_np[i] - encs_np[j])))
        diversity_score = min(np.mean(dists) / 0.6, 1.0) if dists else 0.0
    elif len(encodings_list) >= 1:
        diversity_score = 0.3  # minimal with few photos

    # Save encodings alongside the photos
    import json as _json
    enc_path = os.path.join(dst_dir, "_face_encodings.json")
    with open(enc_path, "w", encoding="utf-8") as f:
        _json.dump({
            "person": person,
            "encodings": encodings_list,
            "photo_count": saved,
            "diversity_score": round(diversity_score, 2),
            "verified_photos": saved_photo_names,
        }, f)

    # Also update the global library cache
    cache_path = os.path.join(FACE_LIBRARY_DIR, "_encodings_cache.json")
    try:
        cache = {}
        if os.path.isfile(cache_path):
            with open(cache_path, "r", encoding="utf-8") as f:
                cache = _json.load(f)
        from curate import _get_cache_fingerprint
        cache[person] = {
            "fingerprint": _get_cache_fingerprint(dst_dir),
            "encodings": encodings_list,
        }
        with open(cache_path, "w", encoding="utf-8") as f:
            _json.dump(cache, f)
    except Exception:
        pass

    return jsonify({"ok": True, "person": person, "photos_saved": saved, "photos_skipped": skipped,
                     "encodings": len(encodings_list)})


@app.route("/api/face-library/import", methods=["POST"])
def api_face_library_import():
    """Import a person from the global library into the current project's ref_faces.
    Only copies validated photos (no re-detection needed) and their pre-computed encodings."""
    data = request.json or {}
    person = data.get("person", "").strip().lower()
    if not person:
        return jsonify({"error": "Person name required"}), 400

    src_dir = os.path.join(FACE_LIBRARY_DIR, person)
    if not os.path.isdir(src_dir):
        return jsonify({"error": f"{person} not found in library"}), 404

    dst_dir = os.path.join(PROJECT_DIR, "ref_faces", person)
    os.makedirs(dst_dir, exist_ok=True)

    # Copy only image files (not cache/encoding files)
    image_exts = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"}
    copied = 0
    for fname in os.listdir(src_dir):
        if os.path.splitext(fname)[1].lower() not in image_exts:
            continue
        src = os.path.join(src_dir, fname)
        dst = os.path.join(dst_dir, fname)
        if os.path.isfile(src):
            shutil.copy2(src, dst)
            copied += 1

    # Copy _face_encodings.json into the project person dir
    enc_src = os.path.join(src_dir, "_face_encodings.json")
    if os.path.isfile(enc_src):
        shutil.copy2(enc_src, os.path.join(dst_dir, "_face_encodings.json"))

    # Copy encodings from library's stored encodings + cache
    import json as _json

    # Copy from per-person encoding file
    enc_src = os.path.join(src_dir, "_face_encodings.json")
    cache_dst = os.path.join(PROJECT_DIR, "ref_faces", "_encodings_cache.json")

    encodings_data = None
    # Try per-person encoding file first
    if os.path.isfile(enc_src):
        try:
            with open(enc_src, "r", encoding="utf-8") as f:
                enc_data = _json.load(f)
            encodings_data = enc_data.get("encodings", [])
        except Exception:
            pass

    # Fall back to global library cache
    if not encodings_data:
        cache_src = os.path.join(FACE_LIBRARY_DIR, "_encodings_cache.json")
        if os.path.isfile(cache_src):
            try:
                with open(cache_src, "r", encoding="utf-8") as f:
                    src_cache = _json.load(f)
                if person in src_cache:
                    encodings_data = src_cache[person].get("encodings", [])
            except Exception:
                pass

    # Write encodings to project's cache
    if encodings_data:
        try:
            from curate import _get_cache_fingerprint
            dst_cache = {}
            if os.path.isfile(cache_dst):
                with open(cache_dst, "r", encoding="utf-8") as f:
                    dst_cache = _json.load(f)
            dst_cache[person] = {
                "fingerprint": _get_cache_fingerprint(dst_dir),
                "encodings": encodings_data,
            }
            with open(cache_dst, "w", encoding="utf-8") as f:
                _json.dump(dst_cache, f)
        except Exception:
            pass

    return jsonify({"ok": True, "person": person, "photos_imported": copied})


@app.route("/api/face-library/<person>", methods=["DELETE"])
def api_face_library_delete(person):
    """Remove a person from the global library."""
    pdir = os.path.join(FACE_LIBRARY_DIR, person)
    if os.path.isdir(pdir):
        shutil.rmtree(pdir)
    return jsonify({"ok": True})


@app.route("/api/ref-faces/verify", methods=["POST"])
def api_ref_faces_verify():
    """Verify face reference quality. Checks:
    - Each photo has a detectable face
    - Enough photos for reliable multi-age recognition
    - Encoding diversity (different angles/ages)
    Returns per-photo results and overall readiness score.
    """
    data = request.json or {}
    person = data.get("person")
    ref_dir = os.path.join(PROJECT_DIR, "ref_faces")

    # If no person specified, verify all
    persons_to_check = []
    if person:
        persons_to_check = [person]
    else:
        if os.path.isdir(ref_dir):
            persons_to_check = [d for d in sorted(os.listdir(ref_dir))
                                if os.path.isdir(os.path.join(ref_dir, d))]

    if not persons_to_check:
        return jsonify({"persons": [], "ready": False,
                        "message": "No reference faces found. Add photos first."})

    try:
        import face_recognition as fr
        import numpy as np
    except ImportError:
        return jsonify({"error": "face_recognition library not installed. "
                        "Install with: pip install face_recognition"}), 500

    results = []
    for pname in persons_to_check:
        pdir = os.path.join(ref_dir, pname)
        if not os.path.isdir(pdir):
            results.append({"person": pname, "photos": [], "encodings": 0,
                            "ready": False, "issues": ["Folder not found"]})
            continue

        photos = []
        encodings = []
        image_exts = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"}

        for fname in sorted(os.listdir(pdir)):
            if os.path.splitext(fname)[1].lower() not in image_exts:
                continue
            fpath = os.path.join(pdir, fname)
            photo_result = {"filename": fname, "status": "unknown", "message": ""}

            try:
                from PIL import Image
                pil_img = Image.open(fpath).convert("RGB")
                w, h = pil_img.size
                photo_result["dimensions"] = f"{w}x{h}"

                # Resize large images to prevent dlib segfaults
                max_dim = 1600
                if w > max_dim or h > max_dim:
                    pil_img.thumbnail((max_dim, max_dim), Image.LANCZOS)

                arr = np.array(pil_img)
                locations = fr.face_locations(arr, model="hog")
                if not locations:
                    photo_result["status"] = "no_face"
                    photo_result["message"] = "No face detected in this photo"
                elif len(locations) > 1:
                    # Multiple faces — use the largest
                    areas = [(b-t)*(r-l) for t, r, b, l in locations]
                    best = areas.index(max(areas))
                    enc = fr.face_encodings(arr, [locations[best]])
                    if enc:
                        encodings.append(enc[0])
                        photo_result["status"] = "ok_multi"
                        photo_result["message"] = f"{len(locations)} faces found, using largest"
                    else:
                        photo_result["status"] = "encode_fail"
                        photo_result["message"] = "Face found but encoding failed"
                else:
                    enc = fr.face_encodings(arr, locations)
                    if enc:
                        encodings.append(enc[0])
                        photo_result["status"] = "ok"
                        photo_result["message"] = "Face detected and encoded"
                    else:
                        photo_result["status"] = "encode_fail"
                        photo_result["message"] = "Face found but encoding failed"
            except Exception as e:
                photo_result["status"] = "error"
                photo_result["message"] = str(e)

            photos.append(photo_result)

        # Analyze encoding diversity
        n_enc = len(encodings)
        issues = []
        tips = []
        diversity_score = 0.0

        if n_enc == 0:
            issues.append("No faces could be encoded from any photo")
        elif n_enc < 3:
            issues.append(f"Only {n_enc} encoding(s). Need at least 3-5 for reliable recognition")
            tips.append("Add more photos with clear, front-facing views")
        else:
            # Compute pairwise distances to measure diversity
            dists = []
            for i in range(n_enc):
                for j in range(i+1, n_enc):
                    dists.append(np.linalg.norm(encodings[i] - encodings[j]))
            avg_dist = np.mean(dists) if dists else 0
            min_dist = np.min(dists) if dists else 0
            max_dist = np.max(dists) if dists else 0
            diversity_score = min(avg_dist / 0.6, 1.0)  # 0.6 is typical tolerance

            if avg_dist < 0.3:
                issues.append("Photos are very similar — all seem to be the same angle/age")
                tips.append("Add photos from different ages and angles for better multi-age recognition")
            elif avg_dist < 0.45:
                tips.append("Consider adding photos from more diverse ages/angles")

            if n_enc < 5:
                tips.append(f"You have {n_enc} good encodings. 5-10 is ideal for multi-age recognition")

        # Age coverage assessment
        if n_enc >= 3:
            tips.append("For best results across ages: include baby, toddler, child, and recent photos")

        ok_count = sum(1 for p in photos if p["status"] in ("ok", "ok_multi"))
        fail_count = len(photos) - ok_count
        ready = n_enc >= 3 and len(issues) == 0

        person_result = {
            "person": pname,
            "photos": photos,
            "total_photos": len(photos),
            "encodings": n_enc,
            "ok_count": ok_count,
            "fail_count": fail_count,
            "diversity_score": round(diversity_score, 2),
            "ready": ready,
            "issues": issues,
            "tips": tips,
        }
        results.append(person_result)

    all_ready = all(r["ready"] for r in results)
    return jsonify({
        "persons": results,
        "ready": all_ready,
        "message": "All faces verified and ready!" if all_ready else
                   "Some issues found — see details below"
    })


@app.route("/api/ref-faces/<person>/photos")
def api_ref_faces_photos(person):
    """Get thumbnails of reference photos for a person."""
    from PIL import Image, ImageOps
    import io

    pdir = os.path.join(PROJECT_DIR, "ref_faces", person)
    if not os.path.isdir(pdir):
        return jsonify([])

    image_exts = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"}
    photos = []
    for fname in sorted(os.listdir(pdir)):
        if os.path.splitext(fname)[1].lower() not in image_exts:
            continue
        fpath = os.path.join(pdir, fname)
        try:
            img = Image.open(fpath)
            img = ImageOps.exif_transpose(img)
            img.thumbnail((120, 120))
            if img.mode in ("RGBA", "P", "LA"):
                img = img.convert("RGB")
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=70)
            thumb = base64.b64encode(buf.getvalue()).decode()
            photos.append({"filename": fname, "thumb": thumb})
        except Exception:
            photos.append({"filename": fname, "thumb": ""})

    return jsonify(photos)


@app.route("/api/ref-faces/<person>/photo/<filename>")
def api_ref_faces_photo_full(person, filename):
    """Serve a full-size reference photo."""
    fpath = os.path.join(PROJECT_DIR, "ref_faces", person, filename)
    if not os.path.isfile(fpath):
        return jsonify({"error": "Not found"}), 404
    return send_file(fpath)


@app.route("/api/stats")
def api_stats():
    """Quick stats for the dashboard."""
    config = load_config()
    db = load_scan_db()

    has_config = config is not None
    has_scan = db is not None
    has_sources = bool(config.get("sources", [])) if config else False

    stats = {
        "has_config": has_config,
        "has_scan": has_scan,
        "has_sources": has_sources,
        "event_type": config.get("event_type") if config else None,
        "template_name": config.get("template") if config else None,
    }

    if db:
        images = db["images"]
        stats["total_images"] = sum(1 for i in images if i.get("media_type") != "video")
        stats["total_videos"] = sum(1 for i in images if i.get("media_type") == "video")
        stats["total_media"] = len(images)
        stats["qualified"] = sum(1 for i in images if i.get("status") == "qualified")
        stats["selected"] = sum(1 for i in images if i.get("status") == "selected")
        stats["selected_videos"] = sum(1 for i in images if i.get("status") == "selected" and i.get("media_type") == "video")
        stats["pool"] = sum(1 for i in images if i.get("status") == "pool")
        stats["rejected_blurry"] = sum(1 for i in images if i.get("reject_reason") == "blurry")
        graded = [i.get("photo_grade", {}).get("composite", 0) for i in images if i.get("photo_grade")]
        if graded:
            stats["grade_avg"] = round(sum(graded) / len(graded), 1)
            stats["grade_high"] = sum(1 for g in graded if g >= 70)
            stats["grade_medium"] = sum(1 for g in graded if 40 <= g < 70)
            stats["grade_low"] = sum(1 for g in graded if g < 40)
        stats["sources"] = list(set(i.get("source_label", "") for i in images))
        db_stats = db.get("stats", {})
        stats["partial"] = bool(db_stats.get("partial"))
        stats["total_scanned"] = db_stats.get("total_scanned", len(images))
        stats["scan_date"] = db.get("scan_date")

    stats["default_export_dir"] = os.path.join(os.path.expanduser("~"), "Downloads", "E-z Photo Collection")
    return jsonify(stats)


# ── Report generation ─────────────────────────────────────────────────────────

@app.route("/api/report")
def api_report():
    """Generate the interactive HTML gallery report and serve it."""
    if not os.path.isfile(SCAN_DB_PATH):
        return jsonify({"error": "No scan data. Run scan first."}), 400

    # Import report generation from curate.py
    sys.path.insert(0, PROJECT_DIR)
    import importlib
    import curate
    importlib.reload(curate)  # Ensure fresh module state

    class FakeArgs:
        output = os.path.join(PROJECT_DIR, "curate_report.html")
        no_open = True

    try:
        curate.cmd_report(FakeArgs())
    except SystemExit:
        pass

    report_path = os.path.join(PROJECT_DIR, "curate_report.html")
    if os.path.isfile(report_path):
        return send_file(report_path, mimetype="text/html")
    return jsonify({"error": "Report generation failed"}), 500


# ── Project save/load ─────────────────────────────────────────────────────────

def _safe_project_name(name):
    """Sanitize project name for use as directory name."""
    import re
    safe = re.sub(r'[<>:"/\\|?*]', '_', name.strip())
    return safe[:100] or "untitled"


def _project_meta_path(pdir):
    return os.path.join(pdir, "project_meta.json")


def _current_project_state():
    """Gather current wizard step from config, for saving."""
    config = load_config()
    return {
        "has_config": config is not None,
        "event_type": config.get("event_type") if config else None,
        "template": config.get("template") if config else None,
    }


@app.route("/api/projects")
def api_projects_list():
    """List all saved projects."""
    os.makedirs(PROJECTS_DIR, exist_ok=True)
    projects = []
    for name in sorted(os.listdir(PROJECTS_DIR)):
        pdir = os.path.join(PROJECTS_DIR, name)
        if not os.path.isdir(pdir):
            continue
        meta_path = _project_meta_path(pdir)
        if os.path.isfile(meta_path):
            with open(meta_path, "r", encoding="utf-8") as f:
                meta = json.load(f)
        else:
            meta = {"name": name}
        meta["dir_name"] = name
        projects.append(meta)
    return jsonify(projects)


@app.route("/api/projects/save", methods=["POST"])
def api_projects_save():
    """Save current state as a project. Expects {name, step, overwrite?}."""
    data = request.json or {}
    raw_name = data.get("name", "").strip()
    if not raw_name:
        return jsonify({"error": "Project name is required"}), 400

    step = data.get("step", 0)
    allow_overwrite = data.get("overwrite", False)
    dir_name = _safe_project_name(raw_name)
    pdir = os.path.join(PROJECTS_DIR, dir_name)

    # Check for duplicate name
    if os.path.isdir(pdir) and not allow_overwrite:
        return jsonify({"error": "A project with this name already exists. Choose a different name."}), 409

    os.makedirs(pdir, exist_ok=True)

    # Save meta
    meta_path = _project_meta_path(pdir)
    existing_meta = {}
    if os.path.isfile(meta_path):
        with open(meta_path, "r", encoding="utf-8") as f:
            existing_meta = json.load(f)

    config = load_config()
    meta = {
        "name": raw_name,
        "dir_name": dir_name,
        "created": existing_meta.get("created", datetime.now().isoformat()),
        "modified": datetime.now().isoformat(),
        "step": step,
        "event_type": config.get("event_type") if config else None,
        "template": config.get("template") if config else None,
    }
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

    # Copy config
    if os.path.isfile(CONFIG_PATH):
        shutil.copy2(CONFIG_PATH, os.path.join(pdir, "curate_config.json"))

    # Copy sidecars BEFORE scan_db — on crash, project has newer
    # sidecars + older scan_db (safe: extra vectors are harmless)
    if os.path.isfile(IMAGE_VECTORS_PATH):
        shutil.copy2(IMAGE_VECTORS_PATH,
                     os.path.join(pdir, "image_vectors.npz"))
    if os.path.isfile(CLIP_VECTORS_PATH):
        shutil.copy2(CLIP_VECTORS_PATH,
                     os.path.join(pdir, "clip_vectors.npz"))
    if os.path.isfile(FACE_ENCODINGS_PATH):
        shutil.copy2(FACE_ENCODINGS_PATH,
                     os.path.join(pdir, "face_encodings.npz"))

    # Copy scan_db last
    if os.path.isfile(SCAN_DB_PATH):
        with _db_lock:
            shutil.copy2(SCAN_DB_PATH, os.path.join(pdir, "scan_db.json"))

    # Copy ref_faces
    ref_src = os.path.join(PROJECT_DIR, "ref_faces")
    ref_dst = os.path.join(pdir, "ref_faces")
    if os.path.isdir(ref_src):
        if os.path.isdir(ref_dst):
            shutil.rmtree(ref_dst)
        shutil.copytree(ref_src, ref_dst)

    return jsonify({"ok": True, "project": meta})


@app.route("/api/projects/load", methods=["POST"])
def api_projects_load():
    """Load a saved project. Expects {dir_name}."""
    data = request.json or {}
    dir_name = data.get("dir_name", "")
    pdir = os.path.join(PROJECTS_DIR, dir_name)
    if not os.path.isdir(pdir):
        return jsonify({"error": "Project not found"}), 404

    meta_path = _project_meta_path(pdir)
    meta = {}
    if os.path.isfile(meta_path):
        with open(meta_path, "r", encoding="utf-8") as f:
            meta = json.load(f)

    # Restore config
    src_config = os.path.join(pdir, "curate_config.json")
    if os.path.isfile(src_config):
        shutil.copy2(src_config, CONFIG_PATH)
    elif os.path.isfile(CONFIG_PATH):
        os.remove(CONFIG_PATH)

    # Restore scan_db
    src_db = os.path.join(pdir, "scan_db.json")
    with _db_lock:
        if os.path.isfile(src_db):
            shutil.copy2(src_db, SCAN_DB_PATH)
        elif os.path.isfile(SCAN_DB_PATH):
            os.remove(SCAN_DB_PATH)

    # Restore CLIP vectors sidecar
    src_clip = os.path.join(pdir, "clip_vectors.npz")
    if os.path.isfile(src_clip):
        shutil.copy2(src_clip, CLIP_VECTORS_PATH)
    elif os.path.isfile(CLIP_VECTORS_PATH):
        os.remove(CLIP_VECTORS_PATH)

    # Restore image vectors sidecar
    src_ivec = os.path.join(pdir, "image_vectors.npz")
    if os.path.isfile(src_ivec):
        shutil.copy2(src_ivec, IMAGE_VECTORS_PATH)
    elif os.path.isfile(IMAGE_VECTORS_PATH):
        os.remove(IMAGE_VECTORS_PATH)

    # Restore face encodings sidecar
    src_face = os.path.join(pdir, "face_encodings.npz")
    if os.path.isfile(src_face):
        shutil.copy2(src_face, FACE_ENCODINGS_PATH)
    elif os.path.isfile(FACE_ENCODINGS_PATH):
        os.remove(FACE_ENCODINGS_PATH)

    # Restore ref_faces
    ref_dst = os.path.join(PROJECT_DIR, "ref_faces")
    ref_src = os.path.join(pdir, "ref_faces")
    if os.path.isdir(ref_dst):
        shutil.rmtree(ref_dst)
    if os.path.isdir(ref_src):
        shutil.copytree(ref_src, ref_dst)
    else:
        os.makedirs(ref_dst, exist_ok=True)

    # Activate auto-save for this project
    _set_active_project(load_config())

    return jsonify({"ok": True, "step": meta.get("step", 0), "project": meta})


@app.route("/api/projects/new", methods=["POST"])
def api_projects_new():
    """Start a fresh project. Clears current config, scan_db, sidecars, and ref_faces."""
    if os.path.isfile(CONFIG_PATH):
        os.remove(CONFIG_PATH)
    with _db_lock:
        if os.path.isfile(SCAN_DB_PATH):
            os.remove(SCAN_DB_PATH)
    for sidecar in (IMAGE_VECTORS_PATH, CLIP_VECTORS_PATH, FACE_ENCODINGS_PATH):
        if os.path.isfile(sidecar):
            os.remove(sidecar)
    ref_dir = os.path.join(PROJECT_DIR, "ref_faces")
    if os.path.isdir(ref_dir):
        shutil.rmtree(ref_dir)
    os.makedirs(ref_dir, exist_ok=True)
    return jsonify({"ok": True})


@app.route("/api/projects/<dir_name>/rename", methods=["POST"])
def api_projects_rename(dir_name):
    """Rename a saved project. Expects {name}."""
    data = request.json or {}
    new_name = data.get("name", "").strip()
    if not new_name:
        return jsonify({"error": "Name is required"}), 400

    pdir = os.path.join(PROJECTS_DIR, _safe_project_name(dir_name))
    if not os.path.isdir(pdir):
        return jsonify({"error": "Project not found"}), 404

    meta_path = _project_meta_path(pdir)
    meta = {}
    if os.path.isfile(meta_path):
        with open(meta_path, "r", encoding="utf-8") as f:
            meta = json.load(f)

    meta["name"] = new_name
    meta["modified"] = datetime.now().isoformat()
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

    return jsonify({"ok": True, "name": new_name})


@app.route("/api/projects/<dir_name>", methods=["DELETE"])
def api_projects_delete(dir_name):
    """Delete a saved project."""
    pdir = os.path.join(PROJECTS_DIR, _safe_project_name(dir_name))
    if not os.path.isdir(pdir):
        return jsonify({"error": "Project not found"}), 404
    shutil.rmtree(pdir)
    return jsonify({"ok": True})


# ── Cleanup / Trash ──────────────────────────────────────────────────────────

@app.route("/api/cleanup/images")
def api_cleanup_images():
    """Get all images for cleanup review, grouped by category."""
    db = load_scan_db()
    if not db:
        return jsonify({"images": [], "total": 0})

    images = db["images"]
    status_filter = request.args.get("status")  # e.g. "candidate", "pool"
    category = request.args.get("category")
    media_type = request.args.get("media_type")
    trash_only = request.args.get("trash") == "1"

    filtered = images
    if trash_only:
        filtered = [i for i in filtered if i.get("trash")]
    if status_filter:
        filtered = [i for i in filtered if i.get("status") == status_filter]
    if category:
        filtered = [i for i in filtered if i.get("category") == category]
    if media_type:
        filtered = [i for i in filtered if i.get("media_type", "image") == media_type]
    location = request.args.get("location")
    if location:
        filtered = [i for i in filtered if i.get("location") == location]

    offset = int(request.args.get("offset", 0))
    limit = int(request.args.get("limit", 200))
    total = len(filtered)
    page = filtered[offset:offset + limit]
    return jsonify({"images": page, "total": total, "offset": offset})


@app.route("/api/cleanup/mark-trash", methods=["POST"])
def api_cleanup_mark_trash():
    """Mark images for trash (soft mark, not deleted yet)."""
    data = request.json or {}
    hashes = set(data.get("hashes", []))
    if not hashes:
        return jsonify({"error": "No hashes provided"}), 400

    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data"}), 404

    marked = 0
    for img in db["images"]:
        if img["hash"] in hashes:
            img["trash"] = True
            marked += 1

    save_scan_db(db)
    return jsonify({"ok": True, "marked": marked})


@app.route("/api/cleanup/unmark-trash", methods=["POST"])
def api_cleanup_unmark_trash():
    """Remove trash mark from images."""
    data = request.json or {}
    hashes = set(data.get("hashes", []))

    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data"}), 404

    unmarked = 0
    for img in db["images"]:
        if img["hash"] in hashes and img.get("trash"):
            img["trash"] = False
            unmarked += 1

    save_scan_db(db)
    return jsonify({"ok": True, "unmarked": unmarked})


@app.route("/api/cleanup/trash-count")
def api_cleanup_trash_count():
    """Get count of images marked for trash."""
    db = load_scan_db()
    if not db:
        return jsonify({"count": 0, "size_mb": 0})
    count = 0
    size_kb = 0
    for img in db["images"]:
        if img.get("trash"):
            count += 1
            size_kb += img.get("size_kb", 0)
    return jsonify({"count": count, "size_mb": round(size_kb / 1024, 1)})


@app.route("/api/cleanup/confirm-trash", methods=["POST"])
def api_cleanup_confirm_trash():
    """Send all trash-marked images to the OS recycle bin."""
    try:
        from send2trash import send2trash
    except ImportError:
        return jsonify({"error": "send2trash not installed. Run: pip install send2trash"}), 500

    db = load_scan_db()
    if not db:
        return jsonify({"error": "No scan data"}), 404

    trashed = []
    failed = []
    for img in db["images"]:
        if img.get("trash"):
            fpath = img["path"].replace("/", os.sep)
            if os.path.isfile(fpath):
                try:
                    send2trash(fpath)
                    trashed.append(img["hash"])
                except Exception as e:
                    failed.append({"hash": img["hash"], "error": str(e)})
            else:
                trashed.append(img["hash"])  # Already gone

    # Remove trashed entries from scan_db
    db["images"] = [i for i in db["images"] if i["hash"] not in set(trashed)]
    save_scan_db(db)

    return jsonify({
        "ok": True,
        "recycled": len(trashed),
        "failed": len(failed),
        "failures": failed[:10],
    })


# ── Age Assessment (standalone) ──────────────────────────────────────────────

@app.route("/api/age-assess/start", methods=["POST"])
def api_age_assess_start():
    """Run age assessment as a standalone background task on selected folders."""
    if _any_job_running():
        return jsonify({"error": "A task is already running"}), 409

    data = request.json or {}
    folders = data.get("folders", [])
    if not folders:
        return jsonify({"error": "No folders selected"}), 400
    face_mode = data.get("face_mode", "all")  # "all" or "specific"
    person_name = data.get("person_name", "")

    config = load_config()
    proj_name = (config.get("project_name") or config.get("event_name") or "") if config else ""
    job_id = _create_job("age_assess", proj_name)

    def run_age_assess():
        try:
            _reset_task("age_assess")
            _update_task("Loading age estimation model (first run may download ~500 MB)...")

            from deepface import DeepFace
            import numpy as np

            # Warm up
            _dummy = np.zeros((100, 100, 3), dtype=np.uint8)
            try:
                DeepFace.analyze(_dummy, actions=["age"], enforce_detection=False, silent=True)
            except Exception:
                pass
            _update_task("Model ready.")

            if _is_cancelled():
                _finish_task("Cancelled")
                return

            # Load face references if specific person mode
            ref_encodings = {}
            use_face_filter = False
            if face_mode == "specific" and person_name:
                from curate import load_reference_faces
                face_dir = os.path.join(PROJECT_DIR, "ref_faces")
                if os.path.isdir(face_dir):
                    _update_task(f"Loading face references for {person_name}...")
                    ref_encodings = load_reference_faces(face_dir)
                    if person_name in ref_encodings:
                        use_face_filter = True
                        _update_task(f"Face references loaded for {person_name}.")
                    else:
                        _update_task(f"No face reference found for '{person_name}'. Running on all faces.")

            from curate import IMAGE_EXTS, file_hash, make_thumbnail_b64

            results = []
            total_files = 0
            processed = 0
            face_matched = 0

            # Count files first
            for folder in folders:
                if os.path.isdir(folder):
                    for dp, dn, fns in os.walk(folder):
                        for fn in fns:
                            if os.path.splitext(fn)[1].lower() in IMAGE_EXTS:
                                total_files += 1

            _update_task(f"Found {total_files} images to analyze...")

            for folder in folders:
                if not os.path.isdir(folder):
                    _update_task(f"Folder not found: {folder}")
                    continue

                folder_label = os.path.basename(folder) or folder

                for dp, dn, fns in os.walk(folder):
                    for fn in fns:
                        if _is_cancelled():
                            _finish_task("Cancelled")
                            return

                        ext = os.path.splitext(fn)[1].lower()
                        if ext not in IMAGE_EXTS:
                            continue

                        fpath = os.path.join(dp, fn)
                        processed += 1

                        if processed % 10 == 0:
                            _update_task_percent(int(processed * 100 / total_files) if total_files else 0)
                            _update_task(f"Analyzing {processed}/{total_files} — {len(results)} faces aged...")

                        # If specific person, check face first
                        if use_face_filter:
                            fc, ff, ok, dist, _enc = _fast_face_detect(
                                fpath, ref_encodings, [person_name], 0.6)
                            if person_name not in ff:
                                continue
                            face_matched += 1

                        est = _estimate_age(fpath)
                        if est is not None:
                            thumb = make_thumbnail_b64(fpath, 120)
                            fh = file_hash(fpath)
                            results.append({
                                "path": fpath.replace("\\", "/"),
                                "filename": fn,
                                "folder": folder_label,
                                "hash": fh,
                                "estimated_age": est,
                                "person": person_name if use_face_filter else None,
                                "thumb": thumb,
                            })

            # Also store results in scan_db entries if they exist
            db = load_scan_db()
            if db:
                result_map = {r["hash"]: r["estimated_age"] for r in results}
                updated = 0
                for img in db.get("images", []):
                    if img["hash"] in result_map and "estimated_age" not in img:
                        img["estimated_age"] = result_map[img["hash"]]
                        updated += 1
                if updated > 0:
                    save_scan_db(db)
                    _update_task(f"Updated {updated} entries in scan database.")

            # Save results to a temp file for the UI to fetch
            results_path = os.path.join(PROJECT_DIR, "age_results.json")
            with open(results_path, "w", encoding="utf-8") as f:
                json.dump(results, f, ensure_ascii=False, cls=NumpyEncoder)

            msg = f"Done! Estimated age for {len(results)} images out of {processed} analyzed."
            if use_face_filter:
                msg += f" ({face_matched} photos matched {person_name})"
            _update_task(msg)
            _finish_task()

        except Exception as e:
            _finish_task(str(e))

    threading.Thread(target=run_age_assess, daemon=True).start()
    return jsonify({"ok": True, "job_id": job_id})


@app.route("/api/age-assess/results")
def api_age_assess_results():
    """Return saved age assessment results."""
    results_path = os.path.join(PROJECT_DIR, "age_results.json")
    if not os.path.isfile(results_path):
        return jsonify([])
    with open(results_path, "r", encoding="utf-8") as f:
        return jsonify(json.load(f))


# ── Serve UI ──────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return WIZARD_HTML


# ── Wizard HTML (single page) ────────────────────────────────────────────────

WIZARD_HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>E-z Photo Organizer</title>
<style>
* { margin:0; padding:0; box-sizing:border-box; }
body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; background:#f0f4f8; color:#2d3748; min-height:100vh; }

/* ── Layout ── */
.app { max-width:1200px; margin:0 auto; padding:20px; padding-left:72px; }
.header { text-align:center; padding:36px 0 22px; position:relative; }
.header h1 { font-size:2.2em; color:#1a365d; margin-bottom:4px; letter-spacing:-.02em; font-weight:800; }
.header .version-badge { display:inline-block; font-size:.6em; font-weight:600; color:#3182ce; background:#dbeafe; padding:2px 10px; border-radius:99px; vertical-align:middle; margin-left:8px; letter-spacing:.03em; }
.header p { color:#718096; font-size:.92em; margin-top:2px; }

/* ── Steps nav ── */
.steps-wrap { display:flex; align-items:center; gap:12px; margin:25px 0; }
.steps { display:flex; justify-content:center; gap:0; margin:0; flex:1; overflow:hidden; border-radius:10px; box-shadow:0 1px 4px rgba(0,0,0,.06); }
.step-dot {
    display:flex; align-items:center; gap:8px; padding:11px 18px;
    background:#fff; cursor:pointer; font-size:.82em; color:#a0aec0;
    border:1px solid #e2e8f0; border-left:none; transition:all .25s;
    position:relative;
}
.step-dot:first-child { border-left:1px solid #e2e8f0; border-radius:10px 0 0 10px; }
.step-dot:last-child { border-radius:0 10px 10px 0; }
.step-dot:hover { background:#f7fafc; color:#4a5568; }
.step-dot.done { background:#f0fff4; color:#276749; border-color:#9ae6b4; }
.step-dot.done + .step-dot { border-left-color:#9ae6b4; }
.step-dot.active { background:linear-gradient(135deg, #2b6cb0, #3182ce); color:white; border-color:#2b6cb0; font-weight:600; box-shadow:0 2px 8px rgba(43,108,176,.3); }
.step-dot.active + .step-dot { border-left-color:#2b6cb0; }
.step-dot .num {
    width:22px; height:22px; border-radius:50%; background:#edf2f7; color:#a0aec0;
    display:flex; align-items:center; justify-content:center; font-size:.72em; font-weight:bold;
    transition:all .25s;
}
.step-dot.done .num { background:#38a169; color:white; }
.step-dot.active .num { background:rgba(255,255,255,.95); color:#2b6cb0; }

/* ── Panels ── */
.panel {
    display:none; background:#fff; border-radius:14px; padding:32px 34px; margin-top:14px;
    box-shadow:0 2px 12px rgba(0,0,0,.06); border:1px solid #e2e8f0;
    position:relative;
}
.panel.active { display:block; animation:panelIn .3s ease-out; }
@keyframes panelIn { from { opacity:0; transform:translateY(8px); } to { opacity:1; transform:translateY(0); } }
.panel h2 {
    color:#1a365d; margin-bottom:16px; font-size:1.35em; font-weight:700;
    padding-bottom:12px; border-bottom:2px solid #ebf4ff;
    display:flex; align-items:center; gap:10px;
}
.panel p { color:#4a5568; line-height:1.65; margin-bottom:16px; }

/* ── Footer ── */
.app-footer { text-align:center; padding:20px 0 12px; color:#a0aec0; font-size:.75em; }
.app-footer a { color:#718096; text-decoration:none; }
.app-footer a:hover { color:#2b6cb0; }

/* ── Shared design system ── */

/* Stat cards (summary numbers) */
.stat-row { display:flex; gap:14px; margin-bottom:18px; flex-wrap:wrap; }
.stat-card {
    border-radius:12px; padding:18px 22px; flex:1; min-width:140px; text-align:center;
    background:#ebf8ff; border:1px solid #bee3f8; transition:transform .15s, box-shadow .15s;
}
.stat-card:hover { transform:translateY(-1px); box-shadow:0 3px 10px rgba(0,0,0,.06); }
.stat-card .stat-value { font-size:1.9em; font-weight:800; color:#2b6cb0; line-height:1.2; letter-spacing:-.02em; }
.stat-card .stat-label { color:#718096; font-size:.8em; margin-top:4px; font-weight:500; text-transform:uppercase; letter-spacing:.04em; }
.stat-card.green { background:#f0fff4; border-color:#c6f6d5; }
.stat-card.green .stat-value { color:#38a169; }
.stat-card.red { background:#fff5f5; border-color:#fed7d7; }
.stat-card.red .stat-value { color:#e53e3e; }

/* Notice/alert boxes */
.notice { padding:13px 16px; border-radius:10px; font-size:.84em; margin-top:12px; line-height:1.6; }
.notice strong { font-weight:600; }
.notice-error { background:#fff5f5; border:1px solid #feb2b2; color:#9b2c2c; }
.notice-warn { background:#fffff0; border:1px solid #fefcbf; color:#744210; }
.notice-info { background:#ebf8ff; border:1px solid #bee3f8; color:#2b6cb0; }
.notice-tip { background:#f0fff4; border:1px solid #c6f6d5; color:#276749; }

/* Option checkbox row */
.option-check {
    display:flex; align-items:flex-start; gap:7px; padding:7px 0;
    font-size:.82em; cursor:pointer;
}
.option-check input[type=checkbox],
.option-check input[type=radio] {
    width:16px; height:16px; accent-color:#3b82f6; cursor:pointer; margin:0; margin-top:1px; flex-shrink:0;
}
.option-check .opt-text { color:#4a5568; font-weight:600; }
.option-check .opt-hint { font-weight:400; color:#a0aec0; font-size:.92em; }

/* Radio group */
.radio-group { margin:6px 0 4px 23px; }
.radio-row {
    display:flex; align-items:flex-start; gap:7px; padding:4px 0; font-size:.82em; cursor:pointer;
}
.radio-row input[type=radio] { width:16px; height:16px; accent-color:#3b82f6; cursor:pointer; margin:0; margin-top:1px; flex-shrink:0; }
.radio-row strong { color:#4a5568; font-weight:600; }
.radio-row .radio-hint { color:#a0aec0; }

/* Section title (subheading within panel) */
.section-title { font-weight:600; color:#2d3748; font-size:.95em; margin-bottom:8px; }

/* Panel section grouping */
.panel-section { margin-top:20px; }
.section-break { margin-top:20px; padding-top:16px; border-top:1px solid #edf2f7; }

/* Help/hint text */
.help-text { font-size:.8em; color:#718096; margin-top:5px; line-height:1.5; }

/* Info box (highlighted section with border) */
.info-box { padding:14px 18px; background:#ebf8ff; border:1px solid #bee3f8; border-radius:8px; margin-top:12px; }
.info-box .info-box-title { font-weight:600; font-size:.9em; color:#2d3748; margin-bottom:6px; display:flex; align-items:center; gap:6px; }

/* Export cards */
.export-cards { display:flex; gap:18px; margin-top:22px; flex-wrap:wrap; }
.export-card { flex:1; min-width:260px; border:2px solid #e2e8f0; border-radius:14px; padding:24px; background:#f7fafc; transition:all .2s; }
.export-card:hover { box-shadow:0 4px 16px rgba(0,0,0,.06); transform:translateY(-2px); }
.export-card.primary { border-color:#3182ce; background:linear-gradient(135deg, #ebf8ff 0%, #dbeafe 100%); min-width:280px; }
.export-card.primary:hover { box-shadow:0 4px 16px rgba(49,130,206,.15); }
.export-card .card-title { font-size:1.1em; font-weight:700; color:#2d3748; margin-bottom:8px; }
.export-card.primary .card-title { color:#1a365d; }
.export-card .card-desc { font-size:.85em; color:#4a5568; margin:0 0 14px; line-height:1.55; }

/* Category card list */
.cat-card-list { display:flex; flex-direction:column; gap:8px; }
.cat-card {
    display:flex; align-items:center; gap:14px; padding:14px 16px;
    background:#fff; border:1px solid #e2e8f0; border-radius:12px;
    transition:border-color .2s, box-shadow .2s;
}
.cat-card:hover { border-color:#bfdbfe; box-shadow:0 2px 8px rgba(0,0,0,.04); }
.cat-card .cat-num {
    width:26px; height:26px; border-radius:50%; background:#f1f5f9;
    color:#94a3b8; font-size:.72em; font-weight:700;
    display:flex; align-items:center; justify-content:center; flex-shrink:0;
}
.cat-card .cat-name-wrap { flex:1; min-width:0; }
.cat-card .cat-name-input {
    width:100%; border:1px solid transparent; background:transparent;
    padding:5px 8px; border-radius:6px; font-size:.92em; font-weight:600;
    color:#1e293b; transition:all .15s;
}
.cat-card .cat-name-input:hover { border-color:#e2e8f0; background:#f8fafc; }
.cat-card .cat-name-input:focus { border-color:#93c5fd; background:#fff; box-shadow:0 0 0 2px rgba(59,130,246,.15); outline:none; }
.cat-card .cat-targets {
    display:flex; gap:10px; flex-shrink:0; align-items:center;
}
.cat-card .cat-target-group {
    display:flex; flex-direction:column; align-items:center; gap:2px;
}
.cat-card .cat-target-label {
    font-size:.62em; font-weight:600; text-transform:uppercase; letter-spacing:.05em; color:#94a3b8;
}
.cat-card .cat-target-input {
    width:62px; text-align:center; border:1px solid #e2e8f0; border-radius:8px;
    padding:5px 4px; font-size:.85em; font-weight:500; color:#334155;
    background:#f8fafc; transition:all .15s;
}
.cat-card .cat-target-input:hover { border-color:#cbd5e1; }
.cat-card .cat-target-input:focus { border-color:#93c5fd; background:#fff; box-shadow:0 0 0 2px rgba(59,130,246,.15); outline:none; }
.cat-card .cat-delete {
    width:28px; height:28px; border:none; background:none; border-radius:6px;
    color:#cbd5e1; cursor:pointer; display:flex; align-items:center; justify-content:center;
    font-size:1.15em; transition:all .15s; flex-shrink:0;
}
.cat-card .cat-delete:hover { background:#fef2f2; color:#ef4444; }

/* Inline form group */
.form-row { display:flex; gap:10px; align-items:center; flex-wrap:wrap; }
.form-row input, .form-row select { width:auto; }

/* ── Forms ── */
label { display:block; font-size:.85em; color:#718096; margin-bottom:4px; margin-top:12px; font-weight:500; }
input, select { width:100%; padding:10px 14px; border-radius:8px; border:1px solid #cbd5e0; background:#fff; color:#2d3748; font-size:.9em; transition:border-color .2s, box-shadow .2s; }
input:focus, select:focus { outline:none; border-color:#63b3ed; box-shadow:0 0 0 3px rgba(66,153,225,.15); }
.row { display:flex; gap:12px; }
.row > * { flex:1; }

/* ── Buttons ── */
.btn {
    padding:10px 24px; border:none; border-radius:8px; cursor:pointer;
    font-size:.88em; font-weight:600; transition:all .2s; letter-spacing:.01em;
}
.btn-primary { background:linear-gradient(135deg, #3182ce, #2b6cb0); color:white; box-shadow:0 2px 6px rgba(49,130,206,.25); }
.btn-primary:hover { background:linear-gradient(135deg, #2b6cb0, #2c5282); box-shadow:0 3px 10px rgba(49,130,206,.35); transform:translateY(-1px); }
.btn-primary:active { transform:translateY(0); box-shadow:0 1px 3px rgba(49,130,206,.2); }
.btn-primary:disabled { background:#cbd5e0; color:#a0aec0; cursor:not-allowed; box-shadow:none; transform:none; }
.btn-secondary { background:#edf2f7; color:#2b6cb0; border:1px solid #e2e8f0; }
.btn-secondary:hover { background:#e2e8f0; border-color:#cbd5e0; }
.btn-danger { background:#fff5f5; color:#c53030; border:1px solid #feb2b2; }
.btn-danger:hover { background:#fed7d7; }
.btn-group { display:flex; gap:10px; margin-top:24px; }

/* ── Cards ── */
.template-grid { display:grid; grid-template-columns:repeat(auto-fill, minmax(280px, 1fr)); gap:14px; margin-top:18px; }
.template-card {
    background:#fff; border:2px solid #e2e8f0; border-radius:12px; padding:18px 20px;
    cursor:pointer; transition:all .2s;
}
.template-card:hover { border-color:#90cdf4; box-shadow:0 4px 14px rgba(66,153,225,.12); transform:translateY(-2px); }
.template-card.selected { border-color:#3182ce; background:#ebf8ff; box-shadow:0 0 0 3px rgba(49,130,206,.15); }
.template-card h3 { color:#2b6cb0; margin-bottom:6px; font-size:1.05em; }
.template-card .desc { color:#718096; font-size:.82em; line-height:1.5; }
.template-card .meta { color:#a0aec0; font-size:.75em; margin-top:10px; font-weight:500; }

/* ── Sources list ── */
.source-list { display:flex; flex-direction:column; gap:10px; margin-top:16px; }

/* Source card — each source is a self-contained block */
.src-card {
    background:#fff; border:1px solid #e2e8f0; border-radius:12px;
    padding:16px 18px 14px; transition:border-color .2s, box-shadow .2s;
}
.src-card:hover { border-color:#bfdbfe; box-shadow:0 2px 8px rgba(0,0,0,.04); }

/* Top row: number pill, label input, remove button */
.src-header {
    display:flex; align-items:center; gap:10px; margin-bottom:8px;
}
.src-num {
    width:24px; height:24px; border-radius:50%; background:#f1f5f9;
    color:#94a3b8; font-size:.68em; font-weight:700;
    display:flex; align-items:center; justify-content:center; flex-shrink:0;
}
.src-label {
    border:1px solid transparent; background:transparent;
    padding:3px 7px; border-radius:6px; font-size:.82em; font-weight:600;
    color:#64748b; transition:all .15s; max-width:200px;
}
.src-label:hover { border-color:#e2e8f0; background:#f8fafc; }
.src-label:focus { border-color:#93c5fd; background:#fff; box-shadow:0 0 0 2px rgba(59,130,246,.15); outline:none; }
.src-remove {
    margin-left:auto; width:26px; height:26px; border:none; background:none;
    border-radius:6px; color:#cbd5e1; cursor:pointer; font-size:1.1em;
    display:flex; align-items:center; justify-content:center; transition:all .15s; flex-shrink:0;
}
.src-remove:hover { background:#fef2f2; color:#ef4444; }

/* Path input — visually dominant */
.src-path {
    width:100%; padding:9px 12px; border:1px solid #e2e8f0; border-radius:8px;
    font-size:.9em; color:#1e293b; background:#f8fafc; transition:all .15s;
    box-sizing:border-box;
}
.src-path:hover { border-color:#cbd5e1; background:#fff; }
.src-path:focus { border-color:#93c5fd; background:#fff; box-shadow:0 0 0 2px rgba(59,130,246,.15); outline:none; }
.src-path.path-valid { border-color:#68d391; }
.src-path.path-invalid { border-color:#fc8181; }
.src-path.path-warn { border-color:#ecc94b; }

/* Meta row: badge + coverage hint + browse action */
.src-meta {
    display:flex; align-items:center; gap:8px; margin-top:8px;
}
.source-badge {
    font-size:.72em; padding:3px 9px; border-radius:10px; white-space:nowrap;
    display:inline-flex; align-items:center; gap:4px; font-weight:500;
}
.src-coverage {
    font-size:.72em; color:#94a3b8; font-weight:400; font-style:italic;
}
.src-actions { margin-left:auto; display:flex; gap:6px; align-items:center; }
.src-browse {
    background:#edf2f7; border:1px solid #cbd5e0; border-radius:6px; padding:5px 12px;
    cursor:pointer; font-size:.78em; color:#4a5568; white-space:nowrap; transition:all .15s;
    display:inline-flex; align-items:center; gap:4px;
}
.src-browse:hover { background:#e2e8f0; border-color:#a0aec0; }
.src-browse svg { flex-shrink:0; }
.source-badge.valid { background:#c6f6d5; color:#276749; }
.source-badge.empty { background:#fefcbf; color:#744210; }
.source-badge.invalid { background:#fed7d7; color:#9b2c2c; }
.source-badge.checking { background:#e2e8f0; color:#718096; }

/* Source summary row */
.src-summary-row { display:flex; gap:12px; margin-bottom:16px; flex-wrap:wrap; }
.src-summary-card {
    flex:1; min-width:120px; padding:14px 16px; text-align:center;
    border-radius:10px; background:#f8fafc; border:1px solid #e2e8f0;
}
.src-summary-val { font-size:1.5em; font-weight:800; color:#2b6cb0; line-height:1.2; }
.src-summary-lbl { color:#94a3b8; font-size:.7em; font-weight:500; text-transform:uppercase; letter-spacing:.04em; margin-top:2px; }

/* Coverage guidance badge */
.src-readiness {
    display:inline-flex; align-items:center; gap:6px; padding:6px 14px;
    border-radius:8px; font-size:.78em; font-weight:500; margin-top:4px;
}
.src-readiness.good { background:#f0fff4; border:1px solid #c6f6d5; color:#276749; }
.src-readiness.okay { background:#fffff0; border:1px solid #fefcbf; color:#744210; }
.src-readiness.low { background:#fff5f5; border:1px solid #fed7d7; color:#9b2c2c; }

/* Add-source inline button */
.src-add-btn {
    display:flex; align-items:center; justify-content:center; gap:6px;
    width:100%; padding:12px; border:2px dashed #e2e8f0; border-radius:12px;
    background:transparent; color:#94a3b8; font-size:.85em; font-weight:500;
    cursor:pointer; transition:all .2s;
}
.src-add-btn:hover { border-color:#93c5fd; color:#3b82f6; background:#f8fbff; }

/* Subtle notice (collapsed) */
.notice-subtle {
    padding:0; margin-top:16px; border:none; background:none;
}
.notice-subtle-toggle {
    font-size:.78em; color:#94a3b8; cursor:pointer; display:inline-flex; align-items:center; gap:4px;
    transition:color .15s; border:none; background:none; padding:0;
}
.notice-subtle-toggle:hover { color:#64748b; }
.notice-subtle-toggle svg { transition:transform .2s; }
.notice-subtle-toggle.open svg { transform:rotate(90deg); }
.notice-subtle-body {
    display:none; margin-top:8px; padding:10px 14px; border-radius:8px;
    font-size:.8em; line-height:1.55;
}
.notice-subtle-body.show { display:block; }
.notice-subtle-body.warn { background:#fffdf5; border:1px solid #fef3c7; color:#92400e; }
.notice-subtle-body.error { background:#fef7f7; border:1px solid #fecaca; color:#991b1b; }

/* ── Folder picker: multi-select checkbox ── */
.fp-item input[type=checkbox] {
    width:16px; height:16px; flex-shrink:0; accent-color:#3182ce; cursor:pointer;
}
.fp-checked-count {
    font-size:.82em; color:#3182ce; font-weight:600;
}

/* ── PPTX export option controls ── */
.pptx-options {
    display:flex; flex-wrap:wrap; gap:6px; margin-bottom:16px;
}
.pptx-opt {
    display:inline-flex; align-items:center; gap:7px;
    padding:6px 12px; border-radius:8px; background:#dbeafe;
    border:1px solid #bfdbfe; cursor:pointer; position:relative;
    transition:background .15s, border-color .15s;
    user-select:none;
}
.pptx-opt:hover { background:#c7d7f5; border-color:#93c5fd; }
.pptx-opt svg {
    width:16px; height:16px; flex-shrink:0; stroke:#4573b0;
    fill:none; stroke-width:1.8; stroke-linecap:round; stroke-linejoin:round;
}
.pptx-opt .pptx-opt-label {
    font-size:.82em; font-weight:500; color:#2d4a7a; white-space:nowrap;
}
.pptx-opt input[type=checkbox] {
    width:15px; height:15px; accent-color:#3b82f6; cursor:pointer;
    flex-shrink:0; margin:0;
}
.pptx-opt select {
    font-size:.82em; padding:2px 6px; border:1px solid #93c5fd;
    border-radius:5px; background:#fff; color:#2d4a7a; cursor:pointer;
    outline:none;
}
.pptx-opt select:focus { border-color:#3b82f6; box-shadow:0 0 0 2px rgba(59,130,246,.15); }
.pptx-opt .pptx-hint {
    display:none; position:absolute; bottom:calc(100% + 8px); left:50%;
    transform:translateX(-50%); width:230px; padding:8px 11px;
    background:#1e293b; color:#e2e8f0; font-size:11.5px; font-weight:400;
    border-radius:7px; line-height:1.45; z-index:100; pointer-events:none;
    text-align:left; box-shadow:0 4px 12px rgba(0,0,0,.2);
}
.pptx-opt .pptx-hint::after {
    content:''; position:absolute; top:100%; left:50%; transform:translateX(-50%);
    border:5px solid transparent; border-top-color:#1e293b;
}
.pptx-opt:hover .pptx-hint { display:block; }

/* ── Folder picker: recent folders ── */
.fp-recent { padding:6px 18px; border-bottom:1px solid #edf2f7; }
.fp-recent-title { font-size:.75em; color:#a0aec0; text-transform:uppercase; letter-spacing:.5px; margin-bottom:4px; }
.fp-recent-item {
    display:inline-flex; align-items:center; gap:4px; padding:3px 10px; margin:2px 4px 2px 0;
    background:#edf2f7; border-radius:12px; font-size:.78em; color:#4a5568; cursor:pointer;
}
.fp-recent-item:hover { background:#e2e8f0; }

/* ── Folder picker modal ── */
.fp-overlay {
    display:none; position:fixed; inset:0; background:rgba(0,0,0,.45);
    z-index:9000; align-items:center; justify-content:center;
}
.fp-overlay.active { display:flex; }
.fp-modal {
    background:#fff; border-radius:10px; width:560px; max-width:92vw;
    max-height:80vh; display:flex; flex-direction:column;
    box-shadow:0 8px 30px rgba(0,0,0,.25);
}
.fp-header {
    padding:14px 18px; border-bottom:1px solid #e2e8f0;
    display:flex; align-items:center; justify-content:space-between;
}
.fp-header h3 { margin:0; font-size:1em; color:#2d3748; }
.fp-breadcrumb {
    padding:8px 18px; background:#f7fafc; border-bottom:1px solid #edf2f7;
    font-size:.82em; color:#718096; display:flex; flex-wrap:wrap; gap:2px; align-items:center;
}
.fp-breadcrumb span { cursor:pointer; color:#3182ce; }
.fp-breadcrumb span:hover { text-decoration:underline; }
.fp-breadcrumb .sep { color:#a0aec0; cursor:default; margin:0 2px; }
.fp-list {
    flex:1; overflow-y:auto; padding:6px 0; min-height:200px;
}
.fp-item {
    display:flex; align-items:center; padding:7px 18px; cursor:pointer;
    font-size:.88em; color:#2d3748; gap:8px;
}
.fp-item:hover { background:#ebf8ff; }
.fp-item.selected { background:#bee3f8; }
.fp-item .fp-icon { color:#ecc94b; font-size:1.1em; flex-shrink:0; }
.fp-item .fp-arrow { color:#a0aec0; margin-left:auto; font-size:.9em; }
.fp-empty { padding:30px 18px; text-align:center; color:#a0aec0; font-size:.88em; }
.fp-footer {
    padding:12px 18px; border-top:1px solid #e2e8f0;
    display:flex; align-items:center; justify-content:space-between; gap:10px;
}
.fp-footer .fp-path { flex:1; font-size:.78em; color:#718096; overflow:hidden; text-overflow:ellipsis; white-space:nowrap; }
.fp-footer button { padding:6px 16px; border-radius:5px; font-size:.85em; cursor:pointer; border:1px solid #cbd5e0; }
.fp-footer .fp-select { background:#3182ce; color:#fff; border-color:#3182ce; }
.fp-footer .fp-select:hover { background:#2b6cb0; }
.fp-footer .fp-select:disabled { background:#a0aec0; border-color:#a0aec0; cursor:not-allowed; }
.fp-footer .fp-cancel { background:#fff; color:#4a5568; }
.fp-footer .fp-cancel:hover { background:#f7fafc; }

/* ── Progress ── */
.progress-box {
    background:#1a202c; border-radius:10px; padding:16px 18px; margin-top:16px;
    font-family:'SF Mono', 'Fira Code', Consolas, monospace; font-size:.8em; max-height:300px; overflow-y:auto;
    box-shadow:inset 0 2px 4px rgba(0,0,0,.2);
}
.progress-box .line { padding:2px 0; color:#63b3ed; }
.progress-box .current { color:#fc8181; font-weight:bold; }

/* ── Analysis ── */
.analysis-table { width:100%; border-collapse:collapse; margin-top:12px; }
.analysis-table th, .analysis-table td { padding:10px 14px; text-align:left; border-bottom:1px solid #edf2f7; font-size:.85em; }
.analysis-table th { color:#718096; font-weight:600; text-transform:uppercase; font-size:.72em; letter-spacing:.05em; }
.analysis-table tbody tr { transition:background .15s; }
.analysis-table tbody tr:hover { background:#f7fafc; }
.status-ok { color:#38a169; }
.status-close { color:#d69e2e; }
.status-low { color:#dd6b20; }
.status-critical { color:#e53e3e; }
.status-empty { color:#e53e3e; font-weight:bold; }
.status-overflow { color:#3182ce; }

.rec-card {
    background:#f7fafc; border-radius:10px; padding:14px 16px; margin-bottom:10px;
    border-left:4px solid #cbd5e0; transition:transform .15s;
}
.rec-card:hover { transform:translateX(2px); }
.rec-card.critical { border-left-color:#e53e3e; background:#fff5f5; }
.rec-card.warning { border-left-color:#dd6b20; background:#fffff0; }
.rec-card.info { border-left-color:#3182ce; background:#ebf8ff; }
.rec-card.tip { border-left-color:#38a169; background:#f0fff4; }
.rec-card .title { font-weight:600; margin-bottom:4px; color:#2d3748; font-size:.92em; }
.rec-card .detail { color:#718096; font-size:.84em; line-height:1.5; }

/* ── Gallery (embedded) ── */
.gallery-frame { width:100%; height:80vh; border:none; border-radius:8px; margin-top:15px; background:#f7fafc; }

/* ── Export ── */
.export-summary { background:#ebf8ff; border-radius:8px; padding:20px; margin-top:15px; border:1px solid #bee3f8; }
.export-summary .big-num { font-size:2.5em; color:#2b6cb0; font-weight:bold; }
@keyframes spin { to { transform: translate(-50%,-50%) rotate(360deg); } }

/* ── Lightbox ── */
.lightbox { display:none; position:fixed; top:0; left:0; width:100%; height:100%; background:rgba(0,0,0,.8); backdrop-filter:blur(6px); -webkit-backdrop-filter:blur(6px); z-index:9999; align-items:center; justify-content:center; }
.lightbox.open { display:flex; }
.lightbox img { max-width:90vw; max-height:90vh; border-radius:10px; box-shadow:0 8px 40px rgba(0,0,0,.5); }
.lightbox .close-btn { position:absolute; top:20px; right:30px; font-size:2em; color:white; cursor:pointer; background:rgba(255,255,255,.15); border:none; border-radius:50%; width:44px; height:44px; display:flex; align-items:center; justify-content:center; backdrop-filter:blur(4px); transition:background .2s; }
.lightbox .close-btn:hover { background:rgba(255,255,255,.3); }

/* ── Inline spinner ── */
.inline-loader { display:inline-flex; align-items:center; gap:10px; padding:12px 16px; background:#ebf8ff; border:1px solid #bee3f8; border-radius:8px; color:#2b6cb0; font-size:.9em; }
.inline-loader .spin { width:18px; height:18px; border:3px solid #bee3f8; border-top:3px solid #3182ce; border-radius:50%; animation:spinA .7s linear infinite; }
@keyframes spinA { to { transform:rotate(360deg); } }

/* ── Task overlay ── */
#task-overlay { display:none; position:fixed; top:0; left:0; width:100%; height:100%; background:rgba(255,255,255,.92); backdrop-filter:blur(4px); -webkit-backdrop-filter:blur(4px); z-index:9990; align-items:center; justify-content:center; flex-direction:column; }
#task-overlay.active { display:flex; }
#task-overlay .task-box { background:#fff; border-radius:16px; box-shadow:0 8px 32px rgba(0,0,0,.1); padding:36px 44px; min-width:400px; max-width:600px; text-align:center; border:1px solid #e2e8f0; }
#task-overlay .task-title { font-size:1.2em; font-weight:700; color:#1a365d; margin-bottom:12px; }
#task-overlay .task-status { font-size:.9em; color:#718096; margin-bottom:16px; min-height:1.2em; }
#task-overlay .task-bar-bg { height:8px; background:#e2e8f0; border-radius:4px; overflow:hidden; margin-bottom:16px; }
#task-overlay .task-bar { height:100%; background:linear-gradient(90deg, #3182ce, #63b3ed); border-radius:4px; transition:width .5s ease; }
@keyframes taskPulse { 0%,100% { opacity:.6; } 50% { opacity:1; } }
#task-overlay .task-bar.indeterminate { width:100% !important; animation: taskPulse 1.5s ease-in-out infinite; }
#task-overlay .task-lines { text-align:left; max-height:120px; overflow-y:auto; font-size:.8em; color:#718096; background:#f7fafc; border-radius:6px; padding:8px 12px; margin-bottom:16px; }
#task-overlay .task-lines .line { margin:2px 0; }
#task-overlay .btn-stop { background:#e53e3e; color:#fff; border:none; padding:8px 24px; border-radius:6px; font-size:.9em; cursor:pointer; }
#task-overlay .btn-stop:hover { background:#c53030; }

/* ── Side menu (projects) ── */
/* ── Icon rail (always visible) ── */
#icon-rail {
    position:fixed; top:0; left:0; width:52px; height:100%;
    background:linear-gradient(180deg, #f8fafc 0%, #fff 100%);
    z-index:10000; display:flex; flex-direction:column; align-items:center;
    padding:14px 0 10px; border-right:1px solid #e2e8f0;
    box-shadow:1px 0 8px rgba(0,0,0,.04);
}
#icon-rail .rail-btn {
    width:36px; height:36px; border:none; background:none; border-radius:10px;
    cursor:pointer; display:flex; align-items:center; justify-content:center;
    color:#94a3b8; transition:all .2s; position:relative; margin-bottom:2px;
}
#icon-rail .rail-btn:hover { background:#dbeafe; color:#2563eb; }
#icon-rail .rail-btn.active { background:#dbeafe; color:#2563eb; }
#icon-rail .rail-btn svg { width:19px; height:19px; }
#icon-rail .rail-btn .rail-tip {
    display:none; position:absolute; left:50px; top:50%; transform:translateY(-50%);
    background:#1e293b; color:#f1f5f9; font-size:.72em; font-weight:500;
    padding:5px 12px; border-radius:6px; white-space:nowrap; pointer-events:none; z-index:10;
    box-shadow:0 2px 8px rgba(0,0,0,.15);
}
#icon-rail .rail-btn .rail-tip::before {
    content:''; position:absolute; left:-4px; top:50%; transform:translateY(-50%) rotate(45deg);
    width:8px; height:8px; background:#1e293b;
}
#icon-rail .rail-btn:hover .rail-tip { display:block; }
#icon-rail .rail-spacer { flex:1; }
#icon-rail .rail-divider { width:22px; height:1px; background:#e2e8f0; margin:8px 0; }
#icon-rail .rail-btn .job-badge {
    position:absolute; top:-3px; right:-3px; background:#ef4444; color:#fff;
    font-size:9px; font-weight:700; width:16px; height:16px; border-radius:50%;
    display:flex; align-items:center; justify-content:center; line-height:1;
    box-shadow:0 1px 3px rgba(239,68,68,.4);
}
.rail-submenu {
    position:absolute; left:50px; top:0; background:#fff; border:1px solid #e2e8f0;
    border-radius:10px; box-shadow:0 6px 20px rgba(0,0,0,.1); padding:6px; z-index:20;
    white-space:nowrap; min-width:160px;
}
.rail-submenu-btn {
    display:flex; align-items:center; gap:8px; width:100%; padding:9px 14px;
    border:none; background:none; cursor:pointer; border-radius:8px;
    font-size:.82em; color:#475569; text-align:left; transition:all .15s;
}
.rail-submenu-btn:hover { background:#eff6ff; color:#1d4ed8; }
.rail-submenu-btn svg { color:#94a3b8; }
.rail-submenu-btn:hover svg { color:#3b82f6; }
/* ── Jobs panel ── */
#jobs-panel {
    display:none; position:fixed; bottom:60px; left:60px; width:360px;
    background:#fff; border-radius:14px; box-shadow:0 8px 30px rgba(0,0,0,.12);
    z-index:10001; max-height:400px; overflow:hidden; border:1px solid #e2e8f0;
}
#jobs-panel.open { display:flex; flex-direction:column; }
#jobs-panel .jp-header {
    padding:14px 18px; border-bottom:1px solid #edf2f7; font-weight:700;
    font-size:.88em; color:#1e293b; display:flex; justify-content:space-between; align-items:center;
}
#jobs-panel .jp-list { overflow-y:auto; max-height:340px; padding:4px 0; }
#jobs-panel .jp-empty { padding:28px; text-align:center; color:#a0aec0; font-size:.85em; }
#jobs-panel .jp-item { padding:12px 18px; border-bottom:1px solid #f7fafc; transition:background .15s; }
#jobs-panel .jp-item:hover { background:#f8fafc; }
#jobs-panel .jp-item:last-child { border-bottom:none; }
#jobs-panel .jp-row { display:flex; align-items:center; justify-content:space-between; margin-bottom:4px; }
#jobs-panel .jp-type { font-weight:600; font-size:.82em; color:#1e293b; }
#jobs-panel .jp-proj { font-size:.72em; color:#a0aec0; margin-left:6px; }
#jobs-panel .jp-status { font-size:.72em; color:#718096; }
#jobs-panel .jp-bar-bg { height:5px; background:#e2e8f0; border-radius:3px; overflow:hidden; margin-bottom:4px; }
#jobs-panel .jp-bar { height:100%; border-radius:3px; transition:width .5s ease; }
#jobs-panel .jp-bar.running { background:linear-gradient(90deg, #3182ce, #63b3ed); animation:taskPulse 1.5s ease-in-out infinite; }
#jobs-panel .jp-bar.done { background:#22c55e; }
#jobs-panel .jp-bar.error { background:#ef4444; }
#jobs-panel .jp-actions { display:flex; gap:6px; }
#jobs-panel .jp-actions button { border:none; background:none; cursor:pointer; font-size:.75em; padding:3px 10px; border-radius:6px; transition:background .15s; }
#jobs-panel .jp-actions .jp-stop { color:#ef4444; }
#jobs-panel .jp-actions .jp-stop:hover { background:#fef2f2; }
#jobs-panel .jp-actions .jp-goto { color:#3182ce; }
#jobs-panel .jp-actions .jp-goto:hover { background:#eff6ff; }
#jobs-panel .jp-actions .jp-dismiss { color:#a0aec0; }
#jobs-panel .jp-actions .jp-dismiss:hover { background:#f8fafc; color:#718096; }
.menu-btn { display:none; }
/* ── Cleanup overlay ── */
#cleanup-overlay { display:none; position:fixed; top:0; left:48px; width:calc(100% - 48px); height:100%; background:#f5f7fa; z-index:1300; overflow-y:auto; }
#cleanup-overlay.active { display:block; }
#cleanup-overlay .cleanup-header { position:sticky; top:0; background:#fff; z-index:10; padding:16px 24px; border-bottom:1px solid #e2e8f0; display:flex; justify-content:space-between; align-items:center; box-shadow:0 2px 8px rgba(0,0,0,.05); }
#cleanup-overlay .cleanup-header h2 { margin:0; color:#2d3748; font-size:1.3em; }
#cleanup-overlay .cleanup-body { padding:20px 24px; max-width:1400px; margin:0 auto; }
#cleanup-overlay .cleanup-filters { display:flex; gap:12px; flex-wrap:wrap; align-items:center; margin-bottom:16px; }
#cleanup-overlay .cleanup-filters select, #cleanup-overlay .cleanup-filters input { padding:6px 10px; border:1px solid #e2e8f0; border-radius:6px; font-size:.85em; }
#age-overlay { display:none; position:fixed; top:0; left:48px; width:calc(100% - 48px); height:100%; background:#f5f7fa; z-index:1300; overflow-y:auto; }
#age-overlay.active { display:block; }
#age-overlay .age-header { position:sticky; top:0; background:#fff; z-index:10; padding:16px 24px; border-bottom:1px solid #e2e8f0; display:flex; justify-content:space-between; align-items:center; box-shadow:0 2px 8px rgba(0,0,0,.05); }
#age-overlay .age-header h2 { margin:0; color:#2d3748; font-size:1.3em; }
#age-overlay .age-body { padding:20px 24px; max-width:1200px; margin:0 auto; }
#age-overlay .age-grid { display:grid; grid-template-columns:repeat(auto-fill, minmax(140px, 1fr)); gap:12px; }
#age-overlay .age-card { background:#fff; border:1px solid #e2e8f0; border-radius:8px; overflow:hidden; text-align:center; transition:box-shadow .2s; }
#age-overlay .age-card:hover { box-shadow:0 4px 12px rgba(0,0,0,.1); }
#age-overlay .age-card img { width:100%; height:120px; object-fit:cover; }
#age-overlay .age-card .age-info { padding:6px 8px; font-size:.78em; color:#4a5568; }
#age-overlay .age-card .age-badge { display:inline-block; background:#ebf8ff; color:#2b6cb0; padding:2px 8px; border-radius:10px; font-weight:600; font-size:.85em; }
.cleanup-grid { display:grid; grid-template-columns:repeat(auto-fill, minmax(150px, 1fr)); gap:10px; }
.cleanup-card { position:relative; border-radius:8px; overflow:hidden; cursor:pointer; border:2px solid transparent; transition:all .15s; background:#fff; box-shadow:0 1px 3px rgba(0,0,0,.08); }
.cleanup-card:hover { box-shadow:0 2px 8px rgba(0,0,0,.15); }
.cleanup-card.trashed { border-color:#e53e3e; opacity:.7; }
.cleanup-card.trashed::after { content:'TRASH'; position:absolute; top:50%; left:50%; transform:translate(-50%,-50%) rotate(-20deg); font-size:1.4em; font-weight:900; color:#e53e3e; background:rgba(255,255,255,.85); padding:4px 16px; border-radius:6px; pointer-events:none; }
.cleanup-card img { width:100%; height:130px; object-fit:cover; display:block; }
.cleanup-card .card-info { padding:6px 8px; font-size:.72em; color:#718096; white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
.cleanup-trash-bar { position:sticky; bottom:0; background:#fff; border-top:1px solid #e2e8f0; padding:12px 24px; display:flex; justify-content:space-between; align-items:center; box-shadow:0 -2px 8px rgba(0,0,0,.05); z-index:10; }
.cleanup-trash-bar .trash-count { font-size:.95em; color:#4a5568; font-weight:600; }
.cleanup-trash-bar .trash-size { font-size:.8em; color:#a0aec0; margin-left:8px; }

#side-drawer {
    position:fixed; top:0; left:-324px; width:304px; height:100%;
    background:#fff; z-index:1150;
    box-shadow:6px 0 24px rgba(0,0,0,.08);
    transition:left .28s cubic-bezier(.4,0,.2,1);
    display:flex; flex-direction:column; padding-left:52px;
}
#side-drawer.open { left:0; }
#side-drawer .drawer-header {
    padding:18px 20px 14px; border-bottom:1px solid #edf2f7;
    display:flex; justify-content:space-between; align-items:center;
}
#side-drawer .drawer-header h3 { color:#1e293b; margin:0; font-size:1.05em; font-weight:700; }
#side-drawer .drawer-close {
    background:none; border:none; font-size:1.4em; cursor:pointer;
    color:#94a3b8; padding:2px 6px; border-radius:6px; transition:all .15s;
}
#side-drawer .drawer-close:hover { color:#1e293b; background:#f1f5f9; }
#side-drawer .drawer-section {
    border-bottom:1px solid #edf2f7;
}
#side-drawer .drawer-section-header {
    display:flex; align-items:center; justify-content:space-between;
    padding:12px 20px; cursor:pointer; user-select:none; transition:background .15s;
}
#side-drawer .drawer-section-header:hover { background:#f8fafc; }
#side-drawer .drawer-section-title {
    font-weight:600; font-size:.88em; color:#334155; display:flex; align-items:center; gap:8px;
}
#side-drawer .drawer-section-title svg { color:#94a3b8; }
#side-drawer .drawer-section-body { padding:0 20px 14px; }
#side-drawer .drawer-btn {
    width:100%; padding:10px 14px; border:1px solid #e2e8f0; border-radius:10px;
    background:#f8fafc; color:#475569; font-size:.84em; cursor:pointer;
    display:flex; align-items:center; gap:9px; text-align:left;
    transition:all .15s; margin-bottom:6px;
}
#side-drawer .drawer-btn:hover { background:#eff6ff; border-color:#bfdbfe; color:#1d4ed8; }
#side-drawer .drawer-btn:last-child { margin-bottom:0; }
#side-drawer .drawer-btn svg { color:#94a3b8; flex-shrink:0; }
#side-drawer .drawer-btn:hover svg { color:#3b82f6; }
#side-drawer .drawer-btn.primary {
    border-color:#bfdbfe; background:#eff6ff; color:#1d4ed8; font-weight:600;
}
#side-drawer .drawer-btn.primary:hover { background:#dbeafe; border-color:#93c5fd; }
#side-drawer .drawer-btn.primary svg { color:#3b82f6; }
#side-drawer .drawer-actions { padding:12px 20px; display:flex; gap:8px; border-bottom:1px solid #edf2f7; }
#side-drawer .drawer-actions button {
    flex:1; padding:7px 0; border-radius:8px; font-size:.8em; cursor:pointer;
    border:1px solid #e2e8f0; background:#f8fafc; color:#475569; transition:all .15s;
}
#side-drawer .drawer-actions button:hover { background:#eff6ff; border-color:#93c5fd; color:#2563eb; }
#side-drawer .project-list { flex:1; overflow-y:auto; padding:6px 0; }
#side-drawer .project-item {
    padding:11px 20px; cursor:pointer; border-bottom:1px solid #f8fafc; transition:background .15s;
}
#side-drawer .project-item:hover { background:#f0f9ff; }
#side-drawer .project-item .p-name { font-weight:600; color:#1e293b; font-size:.88em; }
#side-drawer .project-item .p-meta { color:#94a3b8; font-size:.72em; margin-top:3px; }
#side-drawer .project-item .p-actions { display:flex; gap:6px; margin-top:6px; }
#side-drawer .project-item .p-del {
    background:none; border:none; color:#ef4444; cursor:pointer;
    font-size:.75em; padding:2px 8px; border-radius:5px; transition:background .15s;
}
#side-drawer .project-item .p-del:hover { background:#fef2f2; }
#side-drawer .drawer-footer {
    padding:14px 20px; border-top:1px solid #edf2f7; margin-top:auto;
}
#side-drawer .drawer-footer .user-info {
    font-size:.78em; color:#94a3b8; margin-bottom:10px;
}
#drawer-backdrop { display:none; position:fixed; top:0; left:0; width:100%; height:100%; background:rgba(0,0,0,.25); backdrop-filter:blur(2px); -webkit-backdrop-filter:blur(2px); z-index:1150; }
#drawer-backdrop.open { display:block; }

/* ── Tutorial overlay ── */
#tour-backdrop { display:none; position:fixed; top:0; left:0; width:100%; height:100%; z-index:2000; }
#tour-backdrop.active { display:block; }
#tour-highlight {
    position:absolute; z-index:2001; border-radius:8px;
    box-shadow: 0 0 0 4000px rgba(0,0,0,.55);
    transition: all .4s cubic-bezier(.4,0,.2,1);
    pointer-events:none;
}
#tour-tooltip {
    position:absolute; z-index:2002; background:white; border-radius:12px; padding:20px 24px;
    box-shadow: 0 8px 30px rgba(0,0,0,.2); max-width:360px; min-width:260px;
    opacity:0; transform:translateY(12px); transition: opacity .35s ease, transform .35s ease;
}
#tour-tooltip.show { opacity:1; transform:translateY(0); }
#tour-tooltip h3 { font-size:1em; color:#2b6cb0; margin-bottom:6px; }
#tour-tooltip p { font-size:.88em; color:#4a5568; line-height:1.5; margin:0 0 14px; }
#tour-tooltip .tour-actions { display:flex; justify-content:space-between; align-items:center; }
#tour-tooltip .tour-dots { display:flex; gap:5px; }
#tour-tooltip .tour-dot { width:7px; height:7px; border-radius:50%; background:#e2e8f0; }
#tour-tooltip .tour-dot.active { background:#667eea; }
#tour-tooltip .tour-skip { font-size:.8em; color:#a0aec0; cursor:pointer; border:none; background:none; }
#tour-tooltip .tour-skip:hover { color:#718096; }
#tour-tooltip .tour-next {
    padding:7px 18px; border:none; border-radius:6px; background:#667eea; color:white;
    font-size:.85em; font-weight:600; cursor:pointer; transition: background .2s;
}
#tour-tooltip .tour-next:hover { background:#5a67d8; }
@keyframes tour-pulse { 0%,100% { box-shadow: 0 0 0 4000px rgba(0,0,0,.55); } 50% { box-shadow: 0 0 0 4000px rgba(0,0,0,.5), 0 0 0 8px rgba(102,126,234,.4); } }
#tour-highlight.pulse { animation: tour-pulse 1.5s ease infinite; }

/* ── Select page: hover like/dislike overlay ── */
.sel-thumb .sel-hover-overlay {
    position:absolute; top:0; left:0; width:100%; height:100%;
    background:rgba(0,0,0,.45); display:flex; align-items:center;
    justify-content:center; gap:12px; opacity:0; transition:opacity .18s;
    pointer-events:none;
}
.sel-thumb:hover .sel-hover-overlay { opacity:1; pointer-events:auto; }
.sel-pref-btn {
    width:32px; height:32px; border-radius:50%; border:2px solid rgba(255,255,255,.5);
    background:rgba(0,0,0,.5); font-size:16px; cursor:pointer; display:flex;
    align-items:center; justify-content:center; transition:all .15s; padding:0;
}
.sel-pref-btn:hover { transform:scale(1.25); border-color:#fff; background:rgba(0,0,0,.7); }
.sel-pref-active-like { background:#48bb78!important; border-color:#48bb78!important; }
.sel-pref-active-dislike { background:#e53e3e!important; border-color:#e53e3e!important; }
.sel-pred-badge { pointer-events:none; }

/* ── Questionnaire modal ── */
.pref-quiz-overlay {
    position:fixed; top:0; left:0; width:100%; height:100%;
    background:rgba(0,0,0,.6); z-index:9998; display:flex;
    align-items:center; justify-content:center;
}
.pref-quiz {
    background:#1e1e30; border-radius:12px; padding:28px 32px; max-width:560px;
    width:90%; color:#e2e8f0; box-shadow:0 8px 40px rgba(0,0,0,.5);
    max-height:85vh; overflow-y:auto;
}
.pref-quiz h3 { margin:0 0 16px; color:#667eea; font-size:1.1em; }
.pref-quiz .q-group { margin-bottom:14px; }
.pref-quiz label { display:block; font-size:.85em; color:#a0aec0; margin-bottom:4px; }
.pref-quiz select, .pref-quiz input[type=range] { width:100%; }
.pref-quiz .q-row { display:flex; gap:8px; flex-wrap:wrap; margin-bottom:6px; }
.pref-quiz .q-chip {
    padding:4px 12px; border-radius:16px; border:1px solid #4a5568;
    background:#2d3748; color:#cbd5e0; cursor:pointer; font-size:.8em;
    transition:all .15s;
}
.pref-quiz .q-chip.active { background:#667eea; border-color:#667eea; color:white; }
.pref-quiz .q-btns { display:flex; gap:10px; margin-top:18px; justify-content:flex-end; }
.pref-quiz .q-btn {
    padding:8px 20px; border:none; border-radius:6px; cursor:pointer;
    font-size:.85em; font-weight:600;
}
.pref-quiz .q-btn-skip { background:#4a5568; color:#e2e8f0; }
.pref-quiz .q-btn-save { background:#667eea; color:white; }
</style>
</head>
<body>

<!-- Icon rail (always visible) -->
<div id="icon-rail">
    <button class="rail-btn" onclick="toggleDrawer()" id="rail-toggle" title="Menu">
        <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><line x1="3" y1="6" x2="21" y2="6"/><line x1="3" y1="12" x2="21" y2="12"/><line x1="3" y1="18" x2="21" y2="18"/></svg>
        <span class="rail-tip">Menu</span>
    </button>
    <div class="rail-divider"></div>
    <button class="rail-btn" onclick="newProject()" title="New Project">
        <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><line x1="12" y1="5" x2="12" y2="19"/><line x1="5" y1="12" x2="19" y2="12"/></svg>
        <span class="rail-tip">New Project</span>
    </button>
    <button class="rail-btn" onclick="openDrawerToProjects()" title="Saved Projects">
        <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><path d="M22 19a2 2 0 01-2 2H4a2 2 0 01-2-2V5a2 2 0 012-2h5l2 3h9a2 2 0 012 2z"/></svg>
        <span class="rail-tip">Saved Projects</span>
    </button>
    <div class="rail-divider"></div>
    <div style="position:relative;" id="rail-cleanup-group">
        <button class="rail-btn" onclick="toggleRailCleanup()" title="Cleanup Assistant" id="rail-cleanup-btn">
            <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><polyline points="3 6 5 6 21 6"/><path d="M19 6v14a2 2 0 01-2 2H7a2 2 0 01-2-2V6m3 0V4a2 2 0 012-2h4a2 2 0 012 2v2"/></svg>
            <span class="rail-tip">Cleanup Assistant</span>
        </button>
        <div id="rail-cleanup-sub" class="rail-submenu" style="display:none;">
            <button class="rail-submenu-btn" onclick="openCleanup(); closeRailCleanup();">
                <svg width="15" height="15" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><rect x="2" y="3" width="20" height="14" rx="2" ry="2"/><line x1="8" y1="21" x2="16" y2="21"/><line x1="12" y1="17" x2="12" y2="21"/></svg>
                Hard Drives
            </button>
            <button class="rail-submenu-btn" onclick="openPhoneImages(); closeRailCleanup();">
                <svg width="15" height="15" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><rect x="5" y="2" width="14" height="20" rx="2" ry="2"/><line x1="12" y1="18" x2="12.01" y2="18"/></svg>
                Mobile Images
            </button>
        </div>
    </div>
    <button class="rail-btn" onclick="openAgeAssessment()" title="Age Assessment">
        <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><circle cx="12" cy="8" r="4"/><path d="M20 21v-2a4 4 0 00-4-4H8a4 4 0 00-4 4v2"/><path d="M16 3.13a4 4 0 010 7.75"/></svg>
        <span class="rail-tip">Age Assessment</span>
    </button>
    <div class="rail-spacer"></div>
    <div class="rail-divider"></div>
    <button class="rail-btn" onclick="toggleJobsPanel()" title="Running Jobs" id="rail-jobs-btn">
        <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><rect x="2" y="7" width="20" height="14" rx="2" ry="2"/><path d="M16 7V5a2 2 0 00-2-2h-4a2 2 0 00-2 2v2"/><line x1="12" y1="12" x2="12" y2="12.01"/></svg>
        <span class="rail-tip">Running Jobs</span>
    </button>
    <button class="rail-btn" onclick="startTutorial()" title="Tutorial">
        <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><circle cx="12" cy="12" r="10"/><path d="M9.09 9a3 3 0 015.83 1c0 2-3 3-3 3"/><line x1="12" y1="17" x2="12.01" y2="17"/></svg>
        <span class="rail-tip">Tutorial</span>
    </button>
    <button class="rail-btn" onclick="logout()" title="Sign Out">
        <svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><path d="M9 21H5a2 2 0 01-2-2V5a2 2 0 012-2h4"/><polyline points="16 17 21 12 16 7"/><line x1="21" y1="12" x2="9" y2="12"/></svg>
        <span class="rail-tip">Sign Out</span>
    </button>
</div>

<!-- Jobs panel (popup from icon rail) -->
<div id="jobs-panel">
    <div class="jp-header">
        <span>Running Jobs</span>
        <button onclick="toggleJobsPanel()" style="background:none; border:none; cursor:pointer; color:#94a3b8; font-size:1.2em; padding:2px 6px; border-radius:6px; transition:all .15s;" onmouseover="this.style.background='#f1f5f9';this.style.color='#334155'" onmouseout="this.style.background='none';this.style.color='#94a3b8'">&times;</button>
    </div>
    <div class="jp-list" id="jp-list"></div>
</div>

<!-- Side drawer (slides from behind icon rail) -->
<div id="drawer-backdrop" onclick="toggleDrawer()"></div>
<div id="side-drawer">
    <div class="drawer-header">
        <h3 id="drawer-title">Menu</h3>
        <button class="drawer-close" onclick="toggleDrawer()">&times;</button>
    </div>

    <!-- New Project button -->
    <div class="drawer-section">
        <div style="padding:14px 20px;">
            <button class="drawer-btn primary" onclick="newProject()">
                <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><line x1="12" y1="5" x2="12" y2="19"/><line x1="5" y1="12" x2="19" y2="12"/></svg>
                New Project
            </button>
        </div>
    </div>

    <!-- Saved Projects section (collapsible) -->
    <div class="drawer-section">
        <div class="drawer-section-header" onclick="toggleProjectsSection()" id="projects-toggle">
            <span class="drawer-section-title">
                <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><path d="M22 19a2 2 0 01-2 2H4a2 2 0 01-2-2V5a2 2 0 012-2h5l2 3h9a2 2 0 012 2z"/></svg>
                Saved Projects
            </span>
            <svg id="projects-arrow" width="14" height="14" viewBox="0 0 24 24" fill="none" stroke="#94a3b8" stroke-width="2.5" stroke-linecap="round" style="transition:transform .2s;"><polyline points="6 9 12 15 18 9"/></svg>
        </div>
        <div id="projects-section" style="display:none;">
            <div class="project-list" id="project-list">
                <div style="padding:20px; color:#94a3b8; font-size:.85em; text-align:center;">Loading...</div>
            </div>
        </div>
    </div>

    <!-- Cleanup section -->
    <div class="drawer-section">
        <div class="drawer-section-header" onclick="toggleCleanupSection()">
            <span class="drawer-section-title">
                <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><polyline points="3 6 5 6 21 6"/><path d="M19 6v14a2 2 0 01-2 2H7a2 2 0 01-2-2V6m3 0V4a2 2 0 012-2h4a2 2 0 012 2v2"/></svg>
                Cleanup
            </span>
            <svg id="cleanup-arrow" width="14" height="14" viewBox="0 0 24 24" fill="none" stroke="#94a3b8" stroke-width="2.5" stroke-linecap="round" style="transition:transform .2s;"><polyline points="6 9 12 15 18 9"/></svg>
        </div>
        <div id="cleanup-section" style="display:none;">
            <div class="drawer-section-body">
                <button class="drawer-btn" onclick="toggleDrawer(); openCleanup();">
                    <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><rect x="2" y="3" width="20" height="14" rx="2" ry="2"/><line x1="8" y1="21" x2="16" y2="21"/><line x1="12" y1="17" x2="12" y2="21"/></svg>
                    Desktop / Laptop / Portable Memory
                </button>
                <button class="drawer-btn" onclick="toggleDrawer(); openPhoneImages();">
                    <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><rect x="5" y="2" width="14" height="20" rx="2" ry="2"/><line x1="12" y1="18" x2="12.01" y2="18"/></svg>
                    Mobile
                </button>
            </div>
        </div>
    </div>

    <!-- Age Assessment -->
    <div class="drawer-section">
        <div class="drawer-section-header" onclick="toggleDrawer(); openAgeAssessment();" style="cursor:pointer;">
            <span class="drawer-section-title">
                <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><circle cx="12" cy="8" r="4"/><path d="M20 21v-2a4 4 0 00-4-4H8a4 4 0 00-4 4v2"/><path d="M16 3.13a4 4 0 010 7.75"/></svg>
                Age Assessment
            </span>
        </div>
    </div>

    <!-- Bottom actions -->
    <div class="drawer-footer">
        <div class="user-info" id="user-info"></div>
        <button class="drawer-btn" onclick="toggleDrawer(); startTutorial();">
            <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><circle cx="12" cy="12" r="10"/><path d="M9.09 9a3 3 0 015.83 1c0 2-3 3-3 3"/><line x1="12" y1="17" x2="12.01" y2="17"/></svg>
            Tutorial
        </button>
        <button class="drawer-btn" onclick="logout()" style="border-color:#fecaca; color:#dc2626; background:#fef2f2;">
            <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><path d="M9 21H5a2 2 0 01-2-2V5a2 2 0 012-2h4"/><polyline points="16 17 21 12 16 7"/><line x1="21" y1="12" x2="9" y2="12"/></svg>
            Sign Out
        </button>
    </div>
</div>

<!-- Cleanup overlay -->
<div id="cleanup-overlay">
    <div class="cleanup-header">
        <h2>Cleanup Tool</h2>
        <div style="display:flex; gap:10px; align-items:center;">
            <span id="cleanup-trash-badge" style="background:#fed7d7; color:#c53030; padding:4px 12px; border-radius:12px; font-size:.85em; font-weight:600; display:none;">0 in trash</span>
            <button class="btn btn-secondary" onclick="showCleanupTrash()" id="btn-review-trash" style="display:none;">Review Trash</button>
            <button class="btn btn-secondary" onclick="closeCleanup()" style="font-size:.85em;">Close</button>
        </div>
    </div>
    <div class="cleanup-body">
        <div class="cleanup-filters">
            <select id="cleanup-filter-cat" onchange="loadCleanupImages()">
                <option value="">All Categories</option>
            </select>
            <select id="cleanup-filter-status" onchange="loadCleanupImages()">
                <option value="">All Statuses</option>
                <option value="candidate">Candidate</option>
                <option value="qualified">Qualified</option>
                <option value="selected">Selected</option>
                <option value="pool">Rejected / Pool</option>
            </select>
            <select id="cleanup-filter-media" onchange="loadCleanupImages()">
                <option value="">All Media</option>
                <option value="image">Images Only</option>
                <option value="video">Videos Only</option>
            </select>
            <label style="display:flex; align-items:center; gap:4px; font-size:.85em; color:#4a5568; cursor:pointer;">
                <input type="checkbox" id="cleanup-show-trash" onchange="loadCleanupImages()"> Show trash only
            </label>
            <div style="flex:1;"></div>
            <button class="btn btn-secondary" onclick="cleanupSelectAll()" style="font-size:.8em; padding:6px 14px;">Select All Visible</button>
            <button class="btn btn-secondary" onclick="cleanupDeselectAll()" style="font-size:.8em; padding:6px 14px;">Deselect All</button>
        </div>
        <div id="cleanup-grid" class="cleanup-grid">
            <div style="grid-column:1/-1; text-align:center; color:#a0aec0; padding:40px;">Open the cleanup tool to review and trash unwanted images.</div>
        </div>
        <div id="cleanup-load-more" style="text-align:center; padding:16px; display:none;">
            <button class="btn btn-secondary" onclick="loadCleanupMore()">Load More</button>
        </div>
    </div>
    <div class="cleanup-trash-bar" id="cleanup-bottom-bar" style="display:none;">
        <div>
            <span class="trash-count" id="cleanup-bar-count">0 items in trash</span>
            <span class="trash-size" id="cleanup-bar-size"></span>
        </div>
        <div style="display:flex; gap:10px;">
            <button class="btn btn-secondary" onclick="cleanupClearTrash()">Clear All Trash Marks</button>
            <button class="btn" style="background:#e53e3e; color:#fff;" onclick="cleanupConfirmRecycle()">Move to Recycle Bin</button>
        </div>
    </div>
</div>

<!-- Age Assessment overlay -->
<div id="age-overlay">
    <div class="age-header">
        <h2>Age Assessment</h2>
        <div style="display:flex; gap:10px; align-items:center;">
            <span id="age-status" style="font-size:.85em; color:#718096;"></span>
            <button class="btn btn-secondary" onclick="closeAgeAssessment()" style="font-size:.85em;">Close</button>
        </div>
    </div>
    <div class="age-body">
        <p style="color:#4a5568; font-size:.9em; margin-bottom:16px;">Estimate the age of people in your photos using AI face analysis.</p>

        <div id="age-setup" style="margin-bottom:20px;">
            <!-- Step 1: Face mode -->
            <div style="font-weight:600; font-size:.9em; color:#2d3748; margin-bottom:6px;">1. Who to analyze</div>
            <div style="margin-bottom:12px;">
                <table style="border-collapse:collapse;">
                    <tr>
                        <td style="padding:2px 6px 2px 0; vertical-align:middle;"><input type="radio" name="age-face-mode" value="all" checked style="margin:0; cursor:pointer;" onchange="toggleAgeFaceMode()"></td>
                        <td style="padding:2px 0; font-size:.82em; cursor:pointer;" onclick="this.previousElementSibling.querySelector('input').checked=true; toggleAgeFaceMode()"><strong style="color:#4a5568;">All faces</strong> <span style="color:#a0aec0;">— Estimate age for every face found in each image</span></td>
                    </tr>
                    <tr>
                        <td style="padding:2px 6px 2px 0; vertical-align:middle;"><input type="radio" name="age-face-mode" value="specific" style="margin:0; cursor:pointer;" onchange="toggleAgeFaceMode()"></td>
                        <td style="padding:2px 0; font-size:.82em; cursor:pointer;" onclick="this.previousElementSibling.querySelector('input').checked=true; toggleAgeFaceMode()"><strong style="color:#4a5568;">Specific person</strong> <span style="color:#a0aec0;">— Only estimate age for a recognized person (requires reference photos)</span></td>
                    </tr>
                </table>
            </div>

            <!-- Face reference upload (shown only for specific person) -->
            <div id="age-face-ref" style="display:none; margin-bottom:14px; padding:12px 16px; background:#f7fafc; border:1px solid #e2e8f0; border-radius:8px;">
                <div style="font-weight:600; font-size:.85em; color:#2d3748; margin-bottom:6px;">Reference photos for recognition</div>
                <div style="display:flex; gap:8px; align-items:center; margin-bottom:10px;">
                    <input type="text" id="age-ref-person-name" placeholder="Person name (e.g. Reef)" style="padding:6px 10px; border:1px solid #e2e8f0; border-radius:6px; font-size:.82em; max-width:180px;">
                    <label class="btn btn-secondary" style="font-size:.8em; padding:6px 12px; margin:0; cursor:pointer;" for="age-ref-upload">+ Add Photos</label>
                    <input type="file" id="age-ref-upload" multiple accept="image/*" style="display:none" onchange="uploadAgeRefPhotos(this.files)">
                    <button class="btn btn-secondary" id="btn-age-verify" onclick="verifyAgeRefFaces()" style="font-size:.8em; padding:6px 12px; margin:0; display:none;">Verify Faces</button>
                </div>
                <div id="age-ref-status" style="font-size:.82em; color:#718096; margin-bottom:6px;"></div>
                <div id="age-ref-thumbs" style="display:flex; gap:8px; flex-wrap:wrap; margin-bottom:8px;"></div>
                <div id="age-verify-result" style="display:none; font-size:.82em; margin-bottom:8px; padding:8px 12px; border-radius:6px;"></div>
                <div style="font-size:.75em; color:#a0aec0;">Upload 3-5 clear face photos from different angles. Verify before running.</div>
                <div id="age-use-existing" style="margin-top:8px;"></div>
            </div>

            <!-- Step 2: Folders -->
            <div style="font-weight:600; font-size:.9em; color:#2d3748; margin-bottom:6px;">2. Select folders to analyze</div>
            <div id="age-folder-list" style="margin-bottom:12px;"></div>

            <!-- Run -->
            <div style="display:flex; gap:10px; align-items:center; margin-bottom:12px;">
                <button class="btn btn-primary" id="btn-run-age" onclick="runAgeAssessment()">Run Age Assessment</button>
                <button class="btn btn-secondary" id="btn-stop-age" onclick="stopAgeAssessment()" style="display:none;">Stop</button>
            </div>
            <div id="age-progress" style="display:none; padding:10px 14px; background:#f7fafc; border:1px solid #e2e8f0; border-radius:6px; font-size:.85em; color:#4a5568;"></div>
        </div>

        <div id="age-results" style="display:none;">
            <div style="display:flex; justify-content:space-between; align-items:center; margin-bottom:12px;">
                <div style="font-weight:600; font-size:.95em; color:#2d3748;">Results <span id="age-result-count" style="font-weight:400; color:#718096;"></span></div>
                <div style="display:flex; gap:8px; align-items:center;">
                    <label style="font-size:.8em; color:#4a5568;">Sort:</label>
                    <select id="age-sort" onchange="renderAgeResults()" style="padding:4px 8px; border:1px solid #e2e8f0; border-radius:5px; font-size:.8em;">
                        <option value="age-asc">Age (youngest first)</option>
                        <option value="age-desc">Age (oldest first)</option>
                        <option value="name">Person name</option>
                        <option value="date">Photo date</option>
                    </select>
                    <label style="font-size:.8em; color:#4a5568;">Filter person:</label>
                    <select id="age-filter-person" onchange="renderAgeResults()" style="padding:4px 8px; border:1px solid #e2e8f0; border-radius:5px; font-size:.8em;">
                        <option value="">All</option>
                    </select>
                </div>
            </div>
            <div id="age-grid" class="age-grid"></div>
        </div>
    </div>
</div>

<!-- Task overlay -->
<div id="task-overlay">
    <div class="task-box">
        <div class="task-title" id="task-overlay-title">Working...</div>
        <div class="task-status" id="task-overlay-status"></div>
        <div class="task-bar-bg"><div class="task-bar indeterminate" id="task-overlay-bar"></div></div>
        <div class="task-lines" id="task-overlay-lines"></div>
        <div style="display:flex; gap:10px; justify-content:center;">
            <button class="btn-stop" onclick="sendToBackground()" style="background:#3182ce;">Run in Background</button>
            <button class="btn-stop" onclick="stopTask()">Stop</button>
        </div>
    </div>
</div>

<div class="app">

<div class="header">
    <h1>E-z Photo Organizer <span class="version-badge">v2.0</span></h1>
    <div id="project-name-bar" style="display:none; font-size:1.6em; color:#2b6cb0; font-weight:700; margin-top:2px; margin-bottom:4px;">Project: <span id="project-name-text"></span></div>
    <p id="header-greeting">Build a meaningful photo collection for your special event</p>
</div>

<div class="steps-wrap">
    <div class="steps" id="steps-nav">
        <div class="step-dot active" onclick="goStep(0)"><span class="num">1</span> Event</div>
        <div class="step-dot" onclick="goStep(1)"><span class="num">2</span> Categories</div>
        <div class="step-dot" onclick="goStep(2)"><span class="num">3</span> Sources</div>
        <div class="step-dot" onclick="goStep(3)"><span class="num">4</span> Faces</div>
        <div class="step-dot" onclick="goStep(4)"><span class="num">5</span> Scan</div>
        <div class="step-dot" onclick="goStep(5)"><span class="num">6</span> Analyze</div>
        <div class="step-dot" onclick="goStep(6)"><span class="num">7</span> Select</div>
        <div class="step-dot" onclick="goStep(7)"><span class="num">8</span> Review</div>
        <div class="step-dot" onclick="goStep(8)"><span class="num">9</span> Export</div>
    </div>
    <button id="btn-save-project" onclick="saveCurrentProject()" title="Save Project" style="flex-shrink:0; width:38px; height:38px; border:1px solid #e2e8f0; background:#fff; border-radius:8px; cursor:pointer; display:flex; align-items:center; justify-content:center; color:#4a5568; transition:all .2s;" onmouseover="this.style.background='#ebf8ff';this.style.borderColor='#90cdf4';this.style.color='#2b6cb0'" onmouseout="this.style.background='#fff';this.style.borderColor='#e2e8f0';this.style.color='#4a5568'">
        <svg viewBox="0 0 24 24" width="20" height="20" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><path d="M19 21H5a2 2 0 01-2-2V5a2 2 0 012-2h11l5 5v11a2 2 0 01-2 2z"/><polyline points="17 21 17 13 7 13 7 21"/><polyline points="7 3 7 8 15 8"/></svg>
    </button>
</div>

<!-- ── STEP 0: Choose Event ── -->
<div class="panel active" id="panel-0">
    <h2>What are you creating?</h2>
    <p>Choose the type of event. This sets up the right categories and targets for your photo collection.</p>
    <div class="template-grid" id="template-grid"></div>
    <div id="more-events" style="margin-top:18px; display:none;">
        <div style="color:#718096; font-size:.85em; margin-bottom:6px; cursor:pointer;" onclick="document.getElementById('more-events-list').style.display = document.getElementById('more-events-list').style.display === 'none' ? 'block' : 'none'; this.querySelector('span').textContent = document.getElementById('more-events-list').style.display === 'none' ? '&#9654;' : '&#9660;'">
            <span>&#9654;</span> More events...
        </div>
        <div id="more-events-list" style="display:none;"></div>
    </div>

    <div id="event-fields" style="margin-top:20px; display:none;">
        <div class="row">
            <div id="field-birthday" style="display:none">
                <label>Child's Birthday</label>
                <div style="display:flex; align-items:center; gap:10px;">
                    <input type="date" id="inp-birthday" style="width:160px;" max="">
                    <label style="font-size:.8em; color:#718096; display:flex; align-items:center; gap:4px; white-space:nowrap;"><input type="checkbox" id="skip-birthday" onchange="toggleSkipDate('birthday')"> Not needed</label>
                </div>
            </div>
            <div id="field-event-date" style="display:none">
                <label>Event Date</label>
                <div style="display:flex; align-items:center; gap:10px;">
                    <input type="date" id="inp-event-date" style="width:160px;">
                    <label style="font-size:.8em; color:#718096; display:flex; align-items:center; gap:4px; white-space:nowrap;"><input type="checkbox" id="skip-event-date" onchange="toggleSkipDate('event-date')"> Not needed</label>
                </div>
            </div>
            <div id="field-end-date" style="display:none">
                <label>End Date</label>
                <div style="display:flex; align-items:center; gap:10px;">
                    <input type="date" id="inp-end-date" style="width:160px;">
                    <label style="font-size:.8em; color:#718096; display:flex; align-items:center; gap:4px; white-space:nowrap;"><input type="checkbox" id="skip-end-date" onchange="toggleSkipDate('end-date')"> Not needed</label>
                </div>
            </div>
            <div id="field-year" style="display:none">
                <label>Year</label>
                <div style="display:flex; align-items:center; gap:10px;">
                    <input type="number" id="inp-year" placeholder="2024" min="2000" max="2030" style="width:100px;">
                    <label style="font-size:.8em; color:#718096; display:flex; align-items:center; gap:4px; white-space:nowrap;"><input type="checkbox" id="skip-year" onchange="toggleSkipDate('year')"> Not needed</label>
                </div>
            </div>
        </div>
    </div>

    <div id="event-tips" style="margin-top:15px; display:none;"></div>

    <div class="btn-group">
        <button class="btn btn-primary" id="btn-next-0" disabled onclick="completeStep0()">Next: Categories</button>
    </div>
</div>

<!-- ── STEP 1: Categories ── -->
<div class="panel" id="panel-1">
    <h2>Customize Your Categories</h2>
    <p>These categories come from your event template. Adjust names, targets, or add/remove as needed.</p>

    <div style="margin-bottom:20px;">
        <label class="section-title">Project Name <span style="color:#ef4444;">*</span></label>
        <input type="text" id="inp-project-name" placeholder="e.g. Reef's Bar Mitzva Photos" style="max-width:420px; margin-top:4px;" oninput="validateProjectName()">
        <div id="project-name-error" style="color:#ef4444; font-size:.8em; margin-top:4px; display:none;">Project name is required</div>
    </div>

    <div class="stat-row">
        <div class="stat-card">
            <div class="stat-value" id="cat-total-count">0</div>
            <div class="stat-label">Categories</div>
        </div>
        <div class="stat-card green">
            <div class="stat-value" id="cat-total-target">0</div>
            <div class="stat-label">Total target</div>
        </div>
        <div class="stat-card" style="background:#f8fafc; border-color:#e2e8f0;">
            <div class="stat-value" id="cat-avg-target" style="color:#64748b;">0</div>
            <div class="stat-label">Per category avg</div>
        </div>
    </div>

    <label class="option-check" style="margin-bottom:16px;">
        <input type="checkbox" id="chk-unlimited" onchange="toggleUnlimited(this.checked)">
        <span><span class="opt-text">Select all matching media (no count limit)</span> <span class="opt-hint">— Selects every image/video that matches your face reference</span></span>
    </label>

    <div class="section-title" style="margin-bottom:10px;">Categories</div>
    <div id="cat-card-list" class="cat-card-list"></div>

    <div style="margin-top:12px;">
        <button class="btn btn-secondary" onclick="addCategory()">+ Add Category</button>
    </div>

    <div class="btn-group">
        <button class="btn btn-secondary" onclick="goStep(0)">Back</button>
        <button class="btn btn-primary" onclick="completeCategoriesStep()">Next: Add Sources</button>
    </div>
</div>

<!-- ── STEP 2: Sources ── -->
<div class="panel" id="panel-2">
    <h2>Where are your photos?</h2>
    <p>Point to folders where your photos live — phone backups, cloud exports, USB drives, anything.</p>
    <div class="help-text" style="margin-top:-10px; margin-bottom:16px;">Adding multiple sources gives better coverage across years and events.</div>

    <div id="src-summary" class="src-summary-row" style="display:none;"></div>

    <div class="source-list" id="source-list"></div>
    <button class="src-add-btn" onclick="addSource()">
        <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2.5" stroke-linecap="round"><line x1="12" y1="5" x2="12" y2="19"/><line x1="5" y1="12" x2="19" y2="12"/></svg>
        Add source folder
    </button>

    <div id="src-readiness" style="margin-top:12px;"></div>

    <div class="notice-subtle">
        <button class="notice-subtle-toggle" onclick="toggleSourceTips(this)">
            <svg width="12" height="12" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2.5" stroke-linecap="round"><polyline points="9 18 15 12 9 6"/></svg>
            Tips &amp; warnings
        </button>
        <div id="src-tips-body" class="notice-subtle-body warn">
            <strong>ZIP files:</strong> Windows Explorer can browse ZIP contents, but the scanner cannot read them. Extract the ZIP first, then add the extracted folder.<br><br>
            <strong>Rotated images:</strong> If your images have wrong orientation, face detection may miss faces and scanning retries with rotation correction (slower). Make sure source images have correct EXIF orientation.
        </div>
    </div>

    <div class="btn-group">
        <button class="btn btn-secondary" onclick="goStep(1)">Back</button>
        <button class="btn btn-primary" onclick="completeStep2()">Next: Face References</button>
    </div>
</div>

<!-- ── STEP 3: Face References ── -->
<div class="panel" id="panel-3">
    <h2>Face Recognition Setup</h2>
    <p>Add reference photos of the people you want to find in your collection. More photos from different ages and angles means better recognition.</p>

    <div style="margin-bottom:20px;">
        <label>Person Name</label>
        <div class="form-row" style="margin-top:4px;">
            <input type="text" id="inp-face-name" placeholder="e.g. daniel" style="max-width:250px">
            <button class="btn btn-secondary" onclick="addPerson()" style="margin:0">Add Person</button>
            <button class="btn btn-secondary" onclick="showFaceLibrary()" style="margin:0; background:#ebf8ff; color:#2b6cb0; border-color:#bee3f8;">Pick from Library</button>
        </div>
        <div id="face-library-panel" style="display:none; margin-top:12px; padding:14px; background:#f7fafc; border:1px solid #e2e8f0; border-radius:8px;">
            <div style="display:flex; justify-content:space-between; align-items:center; margin-bottom:10px;">
                <strong style="color:#2b6cb0; font-size:.9em;">Saved People</strong>
                <span style="cursor:pointer; color:#a0aec0; font-size:1.2em;" onclick="document.getElementById('face-library-panel').style.display='none'">&times;</span>
            </div>
            <div id="face-library-list" class="help-text">Loading...</div>
        </div>
    </div>

    <div id="face-persons"></div>

    <div id="face-match-mode" class="info-box" style="display:none; margin:16px 0;">
        <div class="info-box-title">Face matching mode <span style="color:#e53e3e;">*</span>
            <span style="position:relative; display:inline-flex; cursor:help;" id="face-mode-info">
                <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="#718096" stroke-width="2" stroke-linecap="round"><circle cx="12" cy="12" r="10"/><line x1="12" y1="16" x2="12" y2="12"/><line x1="12" y1="8" x2="12.01" y2="8"/></svg>
                <span style="display:none; position:absolute; left:24px; top:-8px; width:320px; padding:10px 14px; background:#2d3748; color:#fff; font-size:.8em; font-weight:400; border-radius:6px; line-height:1.5; z-index:10; box-shadow:0 4px 12px rgba(0,0,0,.2); pointer-events:none;" id="face-mode-tooltip">
                    <strong>Any person:</strong> A photo is included if it contains at least one of the people you added. Good for collecting all photos of each person separately.<br><br>
                    <strong>All people together:</strong> A photo is only included if every person appears in it. Great for finding group shots where everyone is together.
                </span>
            </span>
        </div>
        <label class="radio-row" style="margin-top:4px;">
            <input type="radio" name="face-match" value="any" checked>
            <span><strong>Any person</strong> <span class="radio-hint">— At least one appears</span></span>
        </label>
        <label class="radio-row">
            <input type="radio" name="face-match" value="all">
            <span><strong>All together</strong> <span class="radio-hint">— Everyone must appear</span></span>
        </label>
        <div id="face-match-people" class="help-text"></div>
    </div>

    <div class="btn-group" id="face-verify-btn-group" style="display:none;">
        <button class="btn btn-primary" onclick="runVerifyAll()">Verify All Faces</button>
    </div>

    <div id="face-verify-results" style="display:none; margin-top:20px;"></div>

    <div class="btn-group">
        <button class="btn btn-secondary" onclick="goStep(2)">Back</button>
        <button class="btn btn-secondary" id="btn-skip-faces" onclick="skipFaces()">Skip (no face recognition)</button>
        <button class="btn btn-primary" id="btn-next-3" disabled onclick="completeFacesStep()">Next: Start Scan</button>
    </div>
</div>

<!-- ── STEP 4: Scan ── -->
<div class="panel" id="panel-4">
    <h2>Scanning your photos</h2>
    <p>This scans all your sources, extracts dates, detects duplicates, and creates thumbnails. This runs once — future scans are incremental.</p>

    <div id="scan-summary-stats" style="display:none;"></div>

    <div class="section-title" style="margin-top:16px;">Scan options</div>
    <label class="option-check" style="margin-bottom:6px;">
        <input type="checkbox" id="chk-nsfw-filter" onchange="toggleNsfwFilter(this.checked)">
        <span><span class="opt-text">Filter out nudity / inappropriate content</span> <span class="opt-hint">(AI detection, ~25 MB model)</span></span>
    </label>

    <label class="option-check" style="margin-bottom:6px;">
        <input type="checkbox" id="chk-age-estimation" onchange="toggleAgeEstimation(this.checked)">
        <span><span class="opt-text">Estimate age from faces</span> <span class="opt-hint">(AI age detection — helps sort undated photos, first run downloads ~500 MB model)</span></span>
    </label>
    <div id="age-est-options" style="display:none;" class="radio-group">
        <label class="radio-row" onclick="event.stopPropagation()">
            <input type="radio" name="age-est-scope" value="all" checked>
            <span><strong>All scanned images</strong> <span class="radio-hint">— Estimate age on every photo with a recognized face</span></span>
        </label>
        <label class="radio-row" onclick="event.stopPropagation(); showAgeEstFolders()">
            <input type="radio" name="age-est-scope" value="folders">
            <span><strong>Specific folders only</strong> <span class="radio-hint">— Choose which source folders to run age estimation on</span></span>
        </label>
        <div id="age-est-folders" style="display:none; margin-top:6px; text-align:left;"></div>
    </div>

    <div class="btn-group">
        <button class="btn btn-primary" id="btn-start-scan" onclick="startScan(false)">Start Scan</button>
        <button class="btn btn-secondary" onclick="startScan(true)">Full Rescan</button>
    </div>

    <div class="progress-box" id="scan-progress" style="display:none"></div>

    <div class="btn-group">
        <button class="btn btn-secondary" onclick="goStep(3)">Back</button>
        <button class="btn btn-primary" id="btn-next-4" disabled onclick="goStep(5)">Next: Analyze</button>
    </div>
</div>

<!-- ── STEP 5: Analyze ── -->
<div class="panel" id="panel-5">
    <h2>Collection Analysis</h2>
    <p>Understanding what you have and what's missing.</p>

    <!-- Quick stats (loads instantly) -->
    <div id="a-quick" style="margin-top:12px;"></div>

    <div id="a-next-hint" style="display:none;"></div>

    <!-- Actions -->
    <div class="section-break" style="display:flex; gap:12px; flex-wrap:wrap; align-items:center;">
        <button class="btn btn-secondary" onclick="runFullAnalysis()" id="btn-full-analysis">Run Full Analysis</button>
        <button class="btn btn-danger" onclick="resetAllSelections()">Reset All Selections</button>
        <span class="help-text" style="margin:0;">Full analysis gives recommendations on which categories need attention.</span>
    </div>
    <div id="a-full-results" style="display:none; margin-top:16px;"></div>

    <div class="btn-group">
        <button class="btn btn-secondary" onclick="goStep(4)">Back</button>
        <button class="btn btn-primary" id="btn-next-5" onclick="goStep(6)">Next: Select Photos</button>
    </div>
</div>

<!-- ── STEP 6: Select ── -->
<div class="panel" id="panel-6" style="max-width:1200px">
    <h2>Select Photos</h2>
    <div id="sel-overall-bar" style="display:flex; align-items:center; gap:12px; margin-bottom:14px;">
        <div style="font-size:.85em; font-weight:600; color:#334155; min-width:60px;" id="sel-overall-label">0 / 0</div>
        <div style="flex:1; height:8px; background:#e2e8f0; border-radius:4px; overflow:hidden;">
            <div id="sel-overall-fill" style="height:100%; width:0%; background:linear-gradient(90deg, #3b82f6, #22c55e); border-radius:4px; transition:width .4s;"></div>
        </div>
        <div style="font-size:.78em; font-weight:600; color:#64748b;" id="sel-overall-pct">0%</div>
    </div>

    <div style="display:flex; gap:16px; margin-top:0;">
        <!-- Category sidebar -->
        <div id="sel-cat-list" style="min-width:240px; max-width:280px; border-right:1px solid #e2e8f0; padding-right:12px; max-height:70vh; overflow-y:auto;">
            <div class="section-title">Categories</div>
        </div>

        <!-- Image grid area -->
        <div style="flex:1; min-width:0;">
            <div id="sel-cat-header" style="display:flex; align-items:center; gap:12px; margin-bottom:10px; flex-wrap:wrap;">
                <span id="sel-cat-title" style="font-size:1.1em; font-weight:700; color:#1e293b;">Select a category</span>
                <span id="sel-cat-counter" style="font-size:.85em; color:#64748b; background:#f1f5f9; padding:2px 10px; border-radius:99px; font-weight:500;"></span>
                <div id="sel-target-edit" style="display:none; margin-left:auto; align-items:center; gap:6px; font-size:.85em;">
                    <label style="color:#4a5568; margin:0;">Target:</label>
                    <input type="number" id="sel-target-input" style="width:60px; padding:2px 6px; border:1px solid #cbd5e0; border-radius:4px;" min="0" max="999">
                    <button class="btn btn-secondary" style="padding:2px 10px; font-size:.85em;" onclick="saveTarget()">Set</button>
                </div>
            </div>
            <div id="sel-filter-bar" style="display:none; margin-bottom:8px; font-size:.85em; align-items:center; gap:10px;">
                <label style="display:inline-flex; align-items:center; gap:4px; margin:0; cursor:pointer; font-size:1em; color:#4a5568;">
                    <input type="checkbox" id="sel-show-selected" checked onchange="renderSelGrid()" style="margin:0;"> Show selected
                </label>
                <button class="btn btn-secondary" style="padding:2px 10px; font-size:.85em;" onclick="selectAllVisible()">Select All</button>
                <button class="btn btn-secondary" style="padding:2px 10px; font-size:.85em;" onclick="deselectAllVisible()">Deselect All</button>
                <label style="margin-left:12px; font-size:.9em; color:#4a5568;">Sort:</label>
                <select id="sel-sort" onchange="sortSelImages()" style="font-size:.85em; padding:2px 4px;">
                    <option value="date">Date</option>
                    <option value="grade">Best grade first</option>
                    <option value="pref">Liked first</option>
                    <option value="predicted">AI recommended</option>
                </select>
                <button class="btn btn-secondary" style="padding:2px 10px; font-size:.85em; margin-left:8px; background:#667eea; color:white; border-color:#667eea;" onclick="showPrefQuiz()">&#x2753; Taste Quiz</button>
                <span id="sel-pref-stats" style="font-size:.8em; color:#718096; margin-left:8px;"></span>
            </div>
            <div id="sel-grid" style="display:flex; flex-wrap:wrap; gap:6px; max-height:55vh; overflow-y:auto; padding:4px;"></div>
            <div id="sel-grid-paging" style="margin-top:8px; display:none; font-size:.85em; color:#718096;">
                <button class="btn btn-secondary" style="padding:2px 10px; font-size:.85em;" id="sel-prev" onclick="selPage(-1)">Prev</button>
                <span id="sel-page-info"></span>
                <button class="btn btn-secondary" style="padding:2px 10px; font-size:.85em;" id="sel-next" onclick="selPage(1)">Next</button>
            </div>

            <!-- Preference Library -->
            <div id="pref-library" style="margin-top:16px; display:none;">
                <div style="display:flex; gap:12px; align-items:center; margin-bottom:10px;">
                    <button class="btn btn-secondary" id="pref-lib-toggle" style="padding:4px 14px; font-size:.85em;" onclick="togglePrefLibrary()">&#x2764; Show Liked &amp; Disliked Library</button>
                </div>
                <div id="pref-lib-content" style="display:none;">
                    <div style="margin-bottom:12px;">
                        <h4 style="color:#48bb78; margin:0 0 6px; font-size:.95em;">&#x1F44D; Liked Images (<span id="pref-lib-like-count">0</span>)</h4>
                        <div id="pref-lib-liked" style="display:flex; flex-wrap:wrap; gap:6px; max-height:200px; overflow-y:auto; padding:4px; background:rgba(72,187,120,.08); border-radius:8px; min-height:40px;"></div>
                    </div>
                    <div>
                        <h4 style="color:#e53e3e; margin:0 0 6px; font-size:.95em;">&#x1F44E; Disliked Images (<span id="pref-lib-dislike-count">0</span>)</h4>
                        <div id="pref-lib-disliked" style="display:flex; flex-wrap:wrap; gap:6px; max-height:200px; overflow-y:auto; padding:4px; background:rgba(229,62,62,.08); border-radius:8px; min-height:40px;"></div>
                    </div>
                </div>
            </div>
        </div>
    </div>

    <!-- Auto-fill section -->
    <div class="section-break">
        <div class="section-title">Auto-fill remaining slots</div>
        <div style="display:flex; gap:12px; flex-wrap:wrap;">
            <button class="btn btn-primary" onclick="runQuickFill()">Quick Fill (fast)</button>
            <button class="btn btn-secondary" onclick="runAutoSelect()">Smart Fill (slow, dedup)</button>
        </div>
        <div class="help-text">
            Quick Fill picks top-quality images instantly. Smart Fill also removes visual duplicates (takes several minutes).
        </div>
    </div>

    <div class="btn-group">
        <button class="btn btn-secondary" onclick="goStep(5)">Back</button>
        <button class="btn btn-primary" id="btn-next-6" onclick="goStep(7)">Next: Review</button>
    </div>
</div>

<!-- ── STEP 7: Review ── -->
<div class="panel" id="panel-7">
    <h2>Review Your Selection</h2>
    <p>Check your collection is ready, then open the gallery to fine-tune.</p>

    <div id="review-stats" style="margin-top:12px;"></div>

    <div id="review-readiness" style="margin-top:16px;"></div>

    <div id="review-cat-breakdown" style="margin-top:16px; display:none;"></div>

    <div style="margin-top:20px; display:flex; gap:12px; align-items:center; flex-wrap:wrap;">
        <button class="btn btn-primary" onclick="openGallery()">Open Gallery in New Tab</button>
        <button class="btn btn-secondary" onclick="loadReviewStats()">Refresh</button>
    </div>

    <div class="notice notice-info" style="margin-top:14px;">
        <strong>Tip:</strong> In the gallery you can drag images between categories, mark favorites, and remove unwanted photos. Changes are saved automatically — return here when you're done.
    </div>

    <div class="btn-group">
        <button class="btn btn-secondary" onclick="goStep(6)">Back</button>
        <button class="btn btn-primary" id="btn-next-7" onclick="goStep(8)">Next: Export</button>
    </div>
</div>

<!-- ── STEP 8: Export ── -->
<div class="panel" id="panel-8">
    <h2>Export Your Collection</h2>
    <p>Export your curated photos as a ready-to-use presentation or organized folder.</p>

    <div id="export-stats" style="margin-top:15px;"></div>

    <!-- Export options as cards -->
    <div class="export-cards">

        <!-- PowerPoint card -->
        <div class="export-card primary">
            <div class="card-title">PowerPoint Presentation</div>
            <p class="card-desc">
                One slide per photo, organized by category. Ready to present or customize.
            </p>
            <div class="pptx-options">
                <label class="pptx-opt">
                    <svg viewBox="0 0 24 24"><path d="M4 6h16M4 12h10M4 18h14"/></svg>
                    <span class="pptx-opt-label">Captions</span>
                    <input type="checkbox" id="pptx-captions" checked>
                    <span class="pptx-hint">Show category name and date below each photo. Turn off for clean, image-only slides.</span>
                </label>
                <label class="pptx-opt">
                    <svg viewBox="0 0 24 24"><rect x="3" y="3" width="18" height="7" rx="1.5"/><rect x="3" y="14" width="18" height="7" rx="1.5"/></svg>
                    <span class="pptx-opt-label">Section dividers</span>
                    <input type="checkbox" id="pptx-dividers" checked>
                    <span class="pptx-hint">Add a title slide before each category group. Turn off for a continuous photo flow.</span>
                </label>
                <label class="pptx-opt">
                    <svg viewBox="0 0 24 24"><path d="M3 6h7M3 12h10M3 18h5"/><path d="M17 6l3 3-3 3"/><path d="M20 9H14"/></svg>
                    <span class="pptx-opt-label">Sort</span>
                    <select id="pptx-sort">
                        <option value="date">By date</option>
                        <option value="score">By quality</option>
                    </select>
                    <span class="pptx-hint">By date: chronological within each category. By quality: highest-graded photos first.</span>
                </label>
            </div>
            <button class="btn btn-primary" id="btn-export-pptx" onclick="runPptxExport()">Create Presentation</button>
            <div id="pptx-download" style="margin-top:10px; display:none;">
                <a id="pptx-download-link" href="/api/export/pptx/download" style="color:#2b6cb0; font-weight:600; font-size:.9em;">Download .pptx</a>
            </div>
        </div>

        <!-- Folder export card -->
        <div class="export-card">
            <div class="card-title">Export to Folder</div>
            <p class="card-desc">
                Copy photos to a folder organized by category subfolders.
            </p>
            <label style="font-size:.82em; color:#718096;">Output Folder</label>
            <input type="text" id="inp-export-dir" style="margin-bottom:10px;">
            <button class="btn btn-secondary" id="btn-export" onclick="runExport()">Export Photos</button>
        </div>

    </div>

    <div class="progress-box" id="export-progress" style="display:none"></div>

    <div class="btn-group" style="margin-top:20px;">
        <button class="btn btn-secondary" onclick="goStep(7)">Back</button>
    </div>
</div>

<div class="app-footer">E-z Photo Organizer v2.0</div>

</div><!-- .app -->

<div class="lightbox" id="lightbox" onclick="closeLightbox()">
    <button class="close-btn" onclick="closeLightbox()">&times;</button>
    <img id="lightbox-img" src="">
    <video id="lightbox-video" controls style="display:none; max-width:90vw; max-height:80vh;" onclick="event.stopPropagation()"></video>
    <div id="lightbox-info" style="color:#ccc; font-size:.85em; text-align:center; margin-top:8px;"></div>
</div>

<script>
let currentStep = 0;
let selectedTemplate = null;
let templates = [];
let config = null;
let taskPoll = null;

function showLoader(el, msg) {
    if (typeof el === 'string') el = document.getElementById(el);
    el.innerHTML = '<div class=inline-loader><div class=spin></div><span>' + esc(msg || 'Loading...') + '</span></div>';
    el.style.display = 'block';
}

// ── Task overlay ──
const TASK_TITLES = {
    'scan': 'Scanning Images...',
    'auto-select': 'Auto-Selecting Photos...',
    'export': 'Exporting Collection...',
};

function showTaskOverlay(taskType) {
    const overlay = document.getElementById('task-overlay');
    document.getElementById('task-overlay-title').textContent = TASK_TITLES[taskType] || 'Working...';
    document.getElementById('task-overlay-status').textContent = 'Starting...';
    document.getElementById('task-overlay-lines').innerHTML = '';
    const bar = document.getElementById('task-overlay-bar');
    bar.className = 'task-bar indeterminate';
    bar.style.width = '100%';
    overlay.classList.add('active');

    if (taskPoll) clearInterval(taskPoll);
    taskPoll = setInterval(async () => {
        try {
            const res = await fetch('/api/scan/status');
            const st = await res.json();
            document.getElementById('task-overlay-status').textContent = st.progress || '';
            const linesEl = document.getElementById('task-overlay-lines');
            const clean = st.lines.filter(l => !l.trimStart().startsWith('File "') && !l.startsWith('Traceback') && !l.trimStart().startsWith('json.decoder'));
            linesEl.innerHTML = clean.map(l => '<div class="line">' + esc(l) + '</div>').join('');
            linesEl.scrollTop = linesEl.scrollHeight;

            if (st.done || st.error || st.cancelled) {
                clearInterval(taskPoll);
                taskPoll = null;
                overlay.classList.remove('active');
                if (st.cancelled || st.error === 'Cancelled') {
                    // User stopped — stay on current step
                } else if (st.error) {
                    alert('Error: ' + st.error);
                }
            }
        } catch(e) {}
    }, 1500);
}

async function stopTask() {
    document.getElementById('task-overlay-status').textContent = 'Stopping...';
    document.querySelector('#task-overlay .btn-stop').disabled = true;
    try {
        await fetch('/api/task/stop', { method: 'POST' });
    } catch(e) {}
    // Wait briefly for backend to acknowledge, then dismiss
    setTimeout(() => {
        if (taskPoll) { clearInterval(taskPoll); taskPoll = null; }
        document.getElementById('task-overlay').classList.remove('active');
        document.querySelector('#task-overlay .btn-stop').disabled = false;
    }, 2000);
}

function sendToBackground() {
    // Dismiss overlay but keep task running
    if (taskPoll) { clearInterval(taskPoll); taskPoll = null; }
    document.getElementById('task-overlay').classList.remove('active');
    // Open jobs panel so user sees progress there
    if (!_jobsPanelOpen) toggleJobsPanel();
}

// ── Jobs panel ──
var _jobsPanelOpen = false;
var _jobsPoll = null;

var JOB_TYPE_LABELS = {
    scan: 'Scan',
    'auto-select': 'Auto Select',
    export: 'Export',
    age_assess: 'Age Assessment'
};

var JOB_NEXT_STEP = {
    scan: 5,
    'auto-select': 6,
    export: 8,
    age_assess: null
};

function toggleJobsPanel() {
    var panel = document.getElementById('jobs-panel');
    _jobsPanelOpen = !_jobsPanelOpen;
    if (_jobsPanelOpen) {
        panel.classList.add('open');
        refreshJobsPanel();
        if (!_jobsPoll) _jobsPoll = setInterval(refreshJobsPanel, 2000);
    } else {
        panel.classList.remove('open');
        if (_jobsPoll) { clearInterval(_jobsPoll); _jobsPoll = null; }
        dismissFinishedJobs();
    }
}

async function dismissFinishedJobs() {
    try {
        var res = await fetch('/api/jobs');
        var data = await res.json();
        var jobs = data.jobs || [];
        for (var i = 0; i < jobs.length; i++) {
            if (!jobs[i].running) {
                await fetch('/api/jobs/' + jobs[i].id + '/dismiss', { method: 'POST' });
            }
        }
        updateJobsBadge(jobs.filter(function(j) { return j.running; }).length);
    } catch(e) {}
}

function updateJobsBadge(count) {
    var btn = document.getElementById('rail-jobs-btn');
    var badge = btn.querySelector('.job-badge');
    if (count > 0) {
        if (!badge) {
            badge = document.createElement('span');
            badge.className = 'job-badge';
            btn.appendChild(badge);
        }
        badge.textContent = count;
    } else if (badge) {
        badge.remove();
    }
}

async function refreshJobsPanel() {
    try {
        var res = await fetch('/api/jobs');
        var data = await res.json();
        var jobs = data.jobs || [];
        var runCount = jobs.filter(function(j) { return j.running; }).length;
        updateJobsBadge(runCount);

        var list = document.getElementById('jp-list');
        if (!jobs.length) {
            list.innerHTML = '<div class="jp-empty">No jobs running</div>';
            return;
        }
        list.innerHTML = '';
        jobs.forEach(function(j) {
            var item = document.createElement('div');
            item.className = 'jp-item';

            var row = document.createElement('div');
            row.className = 'jp-row';
            var typeSpan = document.createElement('span');
            typeSpan.className = 'jp-type';
            typeSpan.textContent = JOB_TYPE_LABELS[j.type] || j.type;
            if (j.project_name) {
                var projSpan = document.createElement('span');
                projSpan.className = 'jp-proj';
                projSpan.textContent = j.project_name;
                typeSpan.appendChild(projSpan);
            }
            row.appendChild(typeSpan);

            var statusSpan = document.createElement('span');
            statusSpan.className = 'jp-status';
            if (j.running) {
                var pct = j.percent || 0;
                statusSpan.textContent = pct > 0 ? (pct + '%') : 'Starting...';
            }
            else if (j.cancelled) statusSpan.textContent = 'Cancelled';
            else if (j.error && j.error !== 'Cancelled') statusSpan.textContent = 'Error';
            else statusSpan.textContent = '100% - Done';
            row.appendChild(statusSpan);
            item.appendChild(row);

            var barBg = document.createElement('div');
            var barRow = document.createElement('div');
            barRow.style.cssText = 'display:flex; align-items:center; gap:6px; margin-bottom:4px;';
            var barBg = document.createElement('div');
            barBg.className = 'jp-bar-bg';
            barBg.style.flex = '1';
            var bar = document.createElement('div');
            bar.className = 'jp-bar';
            if (j.running) {
                var pctW = j.percent || 0;
                bar.classList.add('running');
                bar.style.width = pctW > 0 ? (pctW + '%') : '5%';
                bar.style.animation = 'none';
            } else if (j.error && j.error !== 'Cancelled') {
                bar.classList.add('error');
                bar.style.width = '100%';
            } else {
                bar.classList.add('done');
                bar.style.width = '100%';
            }
            barBg.appendChild(bar);
            barRow.appendChild(barBg);

            if (j.running) {
                var expandBtn = document.createElement('button');
                expandBtn.style.cssText = 'background:none; border:none; cursor:pointer; color:#3182ce; padding:0; display:flex; align-items:center;';
                expandBtn.title = 'Open task view';
                expandBtn.innerHTML = '<svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><polyline points="15 3 21 3 21 9"/><polyline points="9 21 3 21 3 15"/><line x1="21" y1="3" x2="14" y2="10"/><line x1="3" y1="21" x2="10" y2="14"/></svg>';
                (function(jobType) {
                    expandBtn.onclick = function() {
                        toggleJobsPanel();
                        showTaskOverlay(jobType);
                    };
                })(j.type);
                barRow.appendChild(expandBtn);
            }

            item.appendChild(barRow);

            var actions = document.createElement('div');
            actions.className = 'jp-actions';
            if (j.running) {
                var stopBtn = document.createElement('button');
                stopBtn.className = 'jp-stop';
                stopBtn.textContent = 'Stop';
                (function(jobId) {
                    stopBtn.onclick = function() {
                        fetch('/api/jobs/' + jobId + '/stop', { method: 'POST' });
                    };
                })(j.id);
                actions.appendChild(stopBtn);
            } else {
                if (!j.error || j.error === 'Cancelled') {
                    var nextStep = JOB_NEXT_STEP[j.type];
                    if (nextStep !== null && nextStep !== undefined) {
                        var gotoBtn = document.createElement('button');
                        gotoBtn.className = 'jp-goto';
                        gotoBtn.textContent = 'Go to next step';
                        (function(step, jobId) {
                            gotoBtn.onclick = function() {
                                goStep(step);
                                fetch('/api/jobs/' + jobId + '/dismiss', { method: 'POST' });
                                refreshJobsPanel();
                            };
                        })(nextStep, j.id);
                        actions.appendChild(gotoBtn);
                    }
                }
                var dismissBtn = document.createElement('button');
                dismissBtn.className = 'jp-dismiss';
                dismissBtn.textContent = 'Dismiss';
                (function(jobId) {
                    dismissBtn.onclick = function() {
                        fetch('/api/jobs/' + jobId + '/dismiss', { method: 'POST' }).then(function() {
                            refreshJobsPanel();
                        });
                    };
                })(j.id);
                actions.appendChild(dismissBtn);
            }
            item.appendChild(actions);
            list.appendChild(item);
        });
    } catch(e) {}
}

// Poll jobs badge even when panel is closed
setInterval(async function() {
    if (_jobsPanelOpen) return; // panel polling handles it
    try {
        var res = await fetch('/api/jobs');
        var data = await res.json();
        var runCount = (data.jobs || []).filter(function(j) { return j.running; }).length;
        updateJobsBadge(runCount);
    } catch(e) {}
}, 5000);

// ── Step navigation ──
function goStep(n) {
    currentStep = n;
    document.querySelectorAll('.panel').forEach((p, i) => p.classList.toggle('active', i === n));
    document.querySelectorAll('.step-dot').forEach((d, i) => {
        d.classList.toggle('active', i === n);
    });

    if (n === 1) loadCategoriesStep();
    if (n === 2) renderSources();
    if (n === 3) loadFaceStep();
    if (n === 4) checkExistingScan();
    if (n === 5) runAnalysis();
    if (n === 6) loadSelCategories();
    if (n === 7) loadReviewStats();
    if (n === 8) loadExportStats();
}

// ── Step 0: Choose event ──
async function loadTemplates() {
    const res = await fetch('/api/templates');
    templates = await res.json();
    var main = templates.filter(t => !t.extra);
    var extras = templates.filter(t => t.extra);
    main.sort((a, b) => (a.event_type === 'custom') - (b.event_type === 'custom'));
    extras.sort((a, b) => a.display_name.localeCompare(b.display_name));
    // Reorder: main first, then extras
    templates = main.concat(extras);
    var grid = document.getElementById('template-grid');
    grid.innerHTML = main.map(t => `
        <div class="template-card" onclick="selectTemplate('${t.event_type}')" id="tpl-${t.event_type}">
            <h3>${t.display_name}</h3>
            <div class="desc">${t.description}</div>
            <div class="meta">${t.num_categories} categories | ${t.categorization}</div>
        </div>
    `).join('');
    var moreEl = document.getElementById('more-events');
    var listEl = document.getElementById('more-events-list');
    if (extras.length) {
        moreEl.style.display = 'block';
        listEl.innerHTML = '';
        extras.forEach(function(t) {
            var row = document.createElement('div');
            row.id = 'tpl-' + t.event_type;
            row.style.cssText = 'display:flex; align-items:center; gap:12px; padding:8px 12px; cursor:pointer; border-radius:6px; border:1px solid #e2e8f0; margin-bottom:4px; transition:all .15s;';
            row.onmouseenter = function() { row.style.borderColor = '#63b3ed'; row.style.background = '#f0f9ff'; };
            row.onmouseleave = function() { if (!row.classList.contains('selected')) { row.style.borderColor = '#e2e8f0'; row.style.background = ''; } };
            row.onclick = function() { selectTemplate(t.event_type); };
            var name = document.createElement('span');
            name.style.cssText = 'font-weight:500; color:#2b6cb0; min-width:160px;';
            name.textContent = t.display_name;
            var desc = document.createElement('span');
            desc.style.cssText = 'color:#718096; font-size:.82em; flex:1;';
            desc.textContent = t.description;
            var meta = document.createElement('span');
            meta.style.cssText = 'color:#a0aec0; font-size:.75em; white-space:nowrap;';
            meta.textContent = t.num_categories + ' cats';
            row.appendChild(name);
            row.appendChild(desc);
            row.appendChild(meta);
            listEl.appendChild(row);
        });
    }

    // Check existing config
    const cfgRes = await fetch('/api/config');
    config = await cfgRes.json();
    if (config && config.event_type) {
        selectTemplate(config.event_type);
        // Restore fields
        if (config.subject_birthday) document.getElementById('inp-birthday').value = config.subject_birthday;
        if (config.event_date) document.getElementById('inp-event-date').value = config.event_date;
        if (config.end_date) document.getElementById('inp-end-date').value = config.end_date;
        if (config.year) document.getElementById('inp-year').value = config.year;
    }
}

function selectTemplate(type) {
    selectedTemplate = templates.find(t => t.event_type === type);
    // Clear all selections (cards + extra rows)
    document.querySelectorAll('.template-card').forEach(c => c.classList.remove('selected'));
    document.querySelectorAll('#more-events-list > div').forEach(function(r) {
        r.classList.remove('selected');
        r.style.borderColor = '#e2e8f0';
        r.style.background = '';
    });
    var el = document.getElementById('tpl-' + type);
    if (el) {
        el.classList.add('selected');
        if (el.parentElement && el.parentElement.id === 'more-events-list') {
            el.style.borderColor = '#3182ce';
            el.style.background = '#ebf8ff';
            // Auto-expand the more events list
            document.getElementById('more-events-list').style.display = 'block';
        }
    }
    document.getElementById('btn-next-0').disabled = false;

    // Show required fields
    const fields = selectedTemplate.required_fields || {};
    document.getElementById('event-fields').style.display = 'block';
    document.getElementById('field-birthday').style.display = 'subject_birthday' in fields ? 'block' : 'none';
    document.getElementById('field-event-date').style.display = 'event_date' in fields ? 'block' : 'none';
    document.getElementById('field-end-date').style.display = 'end_date' in fields ? 'block' : 'none';
    document.getElementById('field-year').style.display = 'year' in fields ? 'block' : 'none';

    // Tips
    const tips = selectedTemplate.tips || [];
    const tipsEl = document.getElementById('event-tips');
    if (tips.length) {
        tipsEl.style.display = 'block';
        tipsEl.innerHTML = '<h3 style="color:#4caf50; font-size:.9em; margin-bottom:8px;">Tips</h3>' +
            tips.map(t => '<div style="color:#718096; font-size:.85em; padding:3px 0;">• ' + t + '</div>').join('');
    }
}

function toggleSkipDate(field) {
    var inp = document.getElementById('inp-' + field);
    var skip = document.getElementById('skip-' + field).checked;
    inp.disabled = skip;
    inp.style.opacity = skip ? '0.4' : '1';
    if (skip) inp.value = '';
}

async function completeStep0() {
    if (!selectedTemplate) return;

    const data = {
        event_type: selectedTemplate.event_type,
        sources: config?.sources || [],
        face_names: config?.face_names || [],
    };
    const bday = document.getElementById('inp-birthday').value;
    const edate = document.getElementById('inp-event-date').value;
    const enddate = document.getElementById('inp-end-date').value;
    const year = document.getElementById('inp-year').value;
    if (bday) data.subject_birthday = bday;
    if (edate) data.event_date = edate;
    if (enddate) data.end_date = enddate;
    if (year) data.year = year;

    await fetch('/api/init', { method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(data) });

    const cfgRes = await fetch('/api/config');
    config = await cfgRes.json();

    document.querySelectorAll('.step-dot')[0].classList.add('done');
    goStep(1);
}

// ── Step 1: Categories ──
function loadCategoriesStep() {
    // Populate project name from config
    var nameInp = document.getElementById('inp-project-name');
    if (config && config.project_name) {
        nameInp.value = config.project_name;
    }
    const cats = config?.categories || [];
    renderCategories(cats);
}

function renderCategories(cats) {
    const unlimited = config?.unlimited_mode || false;
    document.getElementById('chk-unlimited').checked = unlimited;
    var el = document.getElementById('cat-card-list');
    el.innerHTML = cats.map(function(c, i) {
        var nameVal = esc(c.display || c.id || '');
        var photoVal = c.target || config?.target_per_category || 75;
        var videoVal = c.video_target || 0;
        var dis = unlimited ? ' disabled' : '';
        return '<div class="cat-card">' +
            '<div class="cat-num">' + (i + 1) + '</div>' +
            '<div class="cat-name-wrap">' +
                '<input class="cat-name-input" type="text" value="' + nameVal + '" onchange="updateCatField(' + i + ',\\'display\\',this.value)">' +
            '</div>' +
            '<div class="cat-targets">' +
                '<div class="cat-target-group">' +
                    '<span class="cat-target-label">Photos</span>' +
                    '<input class="cat-target-input" type="number" value="' + photoVal + '" min="0" max="500" onchange="updateCatField(' + i + ',\\'target\\',parseInt(this.value))"' + dis + '>' +
                '</div>' +
                '<div class="cat-target-group">' +
                    '<span class="cat-target-label">Videos</span>' +
                    '<input class="cat-target-input" type="number" value="' + videoVal + '" min="0" max="500" onchange="updateCatField(' + i + ',\\'video_target\\',parseInt(this.value))"' + dis + '>' +
                '</div>' +
            '</div>' +
            '<button class="cat-delete" onclick="removeCategory(' + i + ')" title="Remove category">&times;</button>' +
        '</div>';
    }).join('');
    updateCatSummary();
}

function updateCatField(i, field, value) {
    if (!config?.categories) return;
    config.categories[i][field] = value;
    updateCatSummary();
}

function toggleUnlimited(checked) {
    config = config || {};
    config.unlimited_mode = checked;
    renderCategories(config.categories || []);
}

function updateCatSummary() {
    const cats = config?.categories || [];
    const defaultTarget = config?.target_per_category || 75;
    document.getElementById('cat-total-count').textContent = cats.length;
    var avgEl = document.getElementById('cat-avg-target');
    if (config?.unlimited_mode) {
        document.getElementById('cat-total-target').textContent = 'All matching';
        if (avgEl) avgEl.textContent = '--';
    } else {
        const imgTotal = cats.reduce((sum, c) => sum + (c.target || defaultTarget), 0);
        const vidTotal = cats.reduce((sum, c) => sum + (c.video_target || 0), 0);
        document.getElementById('cat-total-target').textContent = imgTotal + (vidTotal ? ' + ' + vidTotal + ' vid' : '');
        if (avgEl) avgEl.textContent = cats.length > 0 ? Math.round(imgTotal / cats.length) : 0;
    }
}

function addCategory() {
    if (!config) return;
    if (!config.categories) config.categories = [];
    const n = config.categories.length + 1;
    config.categories.push({
        id: 'custom_' + n,
        display: 'New Category ' + n,
        target: config.target_per_category || 75,
        video_target: 0
    });
    renderCategories(config.categories);
}

function removeCategory(i) {
    if (!config?.categories) return;
    config.categories.splice(i, 1);
    renderCategories(config.categories);
}

async function completeCategoriesStep() {
    // Validate project name
    var projName = document.getElementById('inp-project-name').value.trim();
    if (!projName) {
        validateProjectName();
        document.getElementById('inp-project-name').focus();
        return;
    }
    if (!config) config = {};
    config.project_name = projName;

    // Ensure all categories have an id
    if (config?.categories) {
        config.categories.forEach((c, i) => {
            if (!c.id) c.id = 'cat_' + i;
        });
    }
    await fetch('/api/config', { method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(config) });
    document.querySelectorAll('.step-dot')[1].classList.add('done');
    goStep(2);
    renderSources();
}

// ── Step 2: Sources ──

// --- Validation ---
var _validateTimers = {};
var _validateCache = {};  // path → result

function validateSourcePath(idx) {
    var path = (config.sources[idx] && config.sources[idx].path) || '';
    var badge = document.getElementById('source-badge-' + idx);
    var pathInput = document.getElementById('source-path-' + idx);
    if (!badge) return;
    if (!path.trim()) {
        badge.className = 'source-badge';
        badge.textContent = '';
        if (pathInput) { pathInput.style.borderColor = ''; }
        return;
    }
    // Check cache first
    if (_validateCache[path]) {
        applyValidationBadge(idx, _validateCache[path]);
        return;
    }
    badge.className = 'source-badge checking';
    badge.textContent = 'checking...';
    if (pathInput) pathInput.style.borderColor = '';
    // Debounce per-row
    if (_validateTimers[idx]) clearTimeout(_validateTimers[idx]);
    _validateTimers[idx] = setTimeout(async function() {
        try {
            var res = await fetch('/api/browse/validate?path=' + encodeURIComponent(path));
            var data = await res.json();
            _validateCache[path] = data;
            applyValidationBadge(idx, data);
        } catch (e) {
            badge.className = 'source-badge invalid';
            badge.textContent = 'error';
        }
    }, 300);
}

function _fmtCount(n) {
    return n >= 1000 ? n.toLocaleString() : String(n);
}

function applyValidationBadge(idx, data) {
    var badge = document.getElementById('source-badge-' + idx);
    var hint = document.getElementById('source-hint-' + idx);
    var pathInput = document.getElementById('source-path-' + idx);
    if (!badge) return;
    if (!data.valid) {
        badge.className = 'source-badge invalid';
        badge.textContent = 'not found';
        if (pathInput) { pathInput.className = 'src-path path-invalid'; }
        if (hint) { hint.textContent = ''; }
    } else if (!data.accessible) {
        badge.className = 'source-badge invalid';
        badge.textContent = 'no access';
        if (pathInput) { pathInput.className = 'src-path path-invalid'; }
        if (hint) { hint.textContent = ''; }
    } else if (data.media_count === 0) {
        badge.className = 'source-badge empty';
        badge.textContent = 'no media';
        if (pathInput) { pathInput.className = 'src-path path-warn'; }
        if (hint) { hint.textContent = 'empty folder'; }
    } else {
        var cnt = data.media_count;
        var label = data.capped ? _fmtCount(cnt) + '+ photos' : _fmtCount(cnt) + ' photos';
        badge.className = 'source-badge valid';
        badge.textContent = label;
        if (pathInput) { pathInput.className = 'src-path path-valid'; }
        if (hint) {
            if (cnt >= 2000) hint.textContent = 'great coverage';
            else if (cnt >= 500) hint.textContent = 'good coverage';
            else if (cnt >= 50) hint.textContent = '';
            else hint.textContent = 'very small source';
        }
    }
    updateSourceSummary();
}

// --- Source row rendering ---
function renderSources() {
    const list = document.getElementById('source-list');
    if (!config) config = { sources: [] };
    if (!config.sources) config.sources = [];
    config.sources = config.sources.map(function(s, i) {
        if (typeof s === 'string') return { path: s, label: 'Source ' + (i + 1) };
        return s;
    });
    var sources = config.sources;
    list.innerHTML = '';
    sources.forEach(function(s, i) {
        var lbl = s.label || ('Source ' + (i + 1));
        // Card container
        var card = document.createElement('div');
        card.className = 'src-card';

        // Header row: num + label + remove
        var header = document.createElement('div');
        header.className = 'src-header';
        var num = document.createElement('span');
        num.className = 'src-num';
        num.textContent = i + 1;
        var lblInput = document.createElement('input');
        lblInput.type = 'text';
        lblInput.className = 'src-label';
        lblInput.value = lbl;
        lblInput.placeholder = 'Source ' + (i + 1);
        lblInput.onchange = function() { updateSource(i, 'label', this.value); };
        var removeBtn = document.createElement('button');
        removeBtn.type = 'button';
        removeBtn.className = 'src-remove';
        removeBtn.innerHTML = '&times;';
        removeBtn.title = 'Remove source';
        removeBtn.onclick = function() { removeSource(i); };
        header.appendChild(num);
        header.appendChild(lblInput);
        header.appendChild(removeBtn);

        // Path input — visually dominant
        var pathInput = document.createElement('input');
        pathInput.type = 'text';
        pathInput.id = 'source-path-' + i;
        pathInput.className = 'src-path';
        pathInput.value = s.path || '';
        pathInput.placeholder = 'Paste folder path or click Browse';
        pathInput.onchange = function() { updateSource(i, 'path', this.value); validateSourcePath(i); };
        pathInput.onblur = function() { validateSourcePath(i); };

        // Meta row: badge + coverage hint + browse
        var meta = document.createElement('div');
        meta.className = 'src-meta';
        var badge = document.createElement('span');
        badge.id = 'source-badge-' + i;
        badge.className = 'source-badge';
        var hint = document.createElement('span');
        hint.id = 'source-hint-' + i;
        hint.className = 'src-coverage';
        var actions = document.createElement('span');
        actions.className = 'src-actions';
        var browseBtn = document.createElement('button');
        browseBtn.type = 'button';
        browseBtn.className = 'src-browse';
        browseBtn.innerHTML = '<svg width="14" height="14" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M22 19a2 2 0 01-2 2H4a2 2 0 01-2-2V5a2 2 0 012-2h5l2 3h9a2 2 0 012 2z"/></svg> Browse';
        (function(idx) {
            browseBtn.onclick = async function() {
                var startPath = config.sources[idx].path || '';
                var paths = await nativePick(startPath);
                if (!paths.length) return;
                config.sources[idx].path = paths[0];
                var curLabel = config.sources[idx].label || '';
                if (!curLabel || curLabel.match(/^Source [0-9]+$/)) {
                    config.sources[idx].label = _folderLabel(paths[0]);
                }
                for (var p = 1; p < paths.length; p++) {
                    config.sources.push({ path: paths[p], label: _folderLabel(paths[p]) });
                }
                fpAddRecent(paths);
                renderSources();
                saveSourcesConfig();
            };
        })(i);
        actions.appendChild(browseBtn);
        meta.appendChild(badge);
        meta.appendChild(hint);
        meta.appendChild(actions);

        // Assemble card
        card.appendChild(header);
        card.appendChild(pathInput);
        card.appendChild(meta);
        list.appendChild(card);

        // Trigger validation for existing paths
        if (s.path) validateSourcePath(i);
    });
    updateSourceSummary();
}

function toggleSourceTips(btn) {
    btn.classList.toggle('open');
    var body = document.getElementById('src-tips-body');
    body.classList.toggle('show');
}

function updateSourceSummary() {
    var summaryEl = document.getElementById('src-summary');
    var readyEl = document.getElementById('src-readiness');
    var sources = (config && config.sources) || [];
    var validCount = 0;
    var totalPhotos = 0;
    sources.forEach(function(s) {
        var cached = _validateCache[s.path];
        if (cached && cached.valid && cached.accessible && cached.media_count > 0) {
            validCount++;
            totalPhotos += cached.media_count;
        }
    });
    // Summary row
    if (sources.length > 0) {
        summaryEl.style.display = 'flex';
        summaryEl.innerHTML =
            '<div class="src-summary-card"><div class="src-summary-val">' + sources.length + '</div><div class="src-summary-lbl">Sources</div></div>' +
            '<div class="src-summary-card"><div class="src-summary-val">' + validCount + '</div><div class="src-summary-lbl">Valid</div></div>' +
            '<div class="src-summary-card"><div class="src-summary-val">' + (totalPhotos > 0 ? _fmtCount(totalPhotos) : '--') + '</div><div class="src-summary-lbl">Total Photos</div></div>';
    } else {
        summaryEl.style.display = 'none';
    }
    // Readiness hint
    if (totalPhotos >= 2000) {
        readyEl.innerHTML = '<span class="src-readiness good">Great coverage — ready to scan</span>';
    } else if (totalPhotos >= 500) {
        readyEl.innerHTML = '<span class="src-readiness good">Good coverage</span>';
    } else if (totalPhotos > 0 && totalPhotos < 500) {
        readyEl.innerHTML = '<span class="src-readiness okay">Consider adding more sources for better results</span>';
    } else if (sources.length > 0 && totalPhotos === 0) {
        readyEl.innerHTML = '<span class="src-readiness low">No photos found yet — check your paths</span>';
    } else {
        readyEl.innerHTML = '';
    }
}

function _folderLabel(path) {
    var parts = (path || '').replace(/[/\\\\]+$/, '').split(/[/\\\\]/);
    return parts[parts.length - 1] || path;
}

var _pickerOpen = false;
async function nativePick(initialdir) {
    if (_pickerOpen) return [];
    _pickerOpen = true;
    try {
        // Launch the dialog (non-blocking server-side)
        await fetch('/api/browse/native', {
            method: 'POST',
            headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({ initialdir: initialdir || '' })
        });
        // Poll until dialog closes
        while (true) {
            await new Promise(function(r) { setTimeout(r, 300); });
            var res = await fetch('/api/browse/native', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({ poll: true })
            });
            var data = await res.json();
            if (data.status === 'done') return data.paths || [];
        }
    } catch(e) {
        return [];
    } finally {
        _pickerOpen = false;
    }
}

async function addSource() {
    if (!config) config = { sources: [] };
    if (!config.sources) config.sources = [];
    var paths = await nativePick('');
    if (!paths.length) return;
    paths.forEach(function(p) {
        config.sources.push({ path: p, label: _folderLabel(p) });
    });
    fpAddRecent(paths);
    renderSources();
    saveSourcesConfig();
}

function removeSource(i) {
    config.sources.splice(i, 1);
    renderSources();
    saveSourcesConfig();
}

function updateSource(i, field, value) {
    config.sources[i][field] = value;
    saveSourcesConfig();
}

var _sourcesSaveTimer = null;
function saveSourcesConfig() {
    if (_sourcesSaveTimer) clearTimeout(_sourcesSaveTimer);
    _sourcesSaveTimer = setTimeout(function() {
        if (!config || !config.sources) return;
        config.sources.forEach(function(s, i) {
            if (!s.label || !s.label.trim()) s.label = 'Source ' + (i + 1);
        });
        fetch('/api/config/sources', { method: 'POST', headers: {'Content-Type': 'application/json'}, body: JSON.stringify({ sources: config.sources }) });
    }, 500);
}

async function completeStep2() {
    // Fill in empty labels before saving
    if (config && config.sources) {
        config.sources.forEach(function(s, i) {
            if (!s.label || !s.label.trim()) s.label = 'Source ' + (i + 1);
        });
    }
    await fetch('/api/config', { method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(config) });

    document.querySelectorAll('.step-dot')[2].classList.add('done');
    goStep(3);
}

// ── Step 2: Face References ──
let facesVerified = false;

async function loadFaceStep() {
    const res = await fetch('/api/ref-faces');
    const faces = await res.json();
    renderFacePersons(faces);
    // Show verify button if there are persons with photos
    const hasPhotos = faces.some(f => f.photo_count > 0);
    document.getElementById('face-verify-btn-group').style.display = hasPhotos ? 'flex' : 'none';

    // Auto-verify if all people have cached encodings and at least 1 photo each
    var allHaveEncodings = hasPhotos && faces.every(f => f.has_encodings && f.photo_count > 0);
    var anyEmpty = faces.some(f => f.photo_count === 0);
    if (allHaveEncodings && !anyEmpty) {
        facesVerified = true;
    }
    if (anyEmpty) {
        facesVerified = false;
    }

    if (!facesVerified) document.getElementById('btn-next-3').disabled = true;
    else document.getElementById('btn-next-3').disabled = false;
    document.getElementById('btn-skip-faces').style.display = facesVerified ? 'none' : '';
    // Restore saved face match mode
    var savedMode = config && config.face_match_mode ? config.face_match_mode : 'any';
    var radio = document.querySelector('input[name="face-match"][value="' + savedMode + '"]');
    if (radio) radio.checked = true;
}

function renderFacePersons(faces) {
    const container = document.getElementById('face-persons');
    if (!faces.length) {
        container.innerHTML = '<p style="color:#a0aec0; font-style:italic;">No reference faces added yet. Add a person above, or skip this step.</p>';
        return;
    }

    // Pre-populate faceVerifyCache for people with verified_photos from library
    faces.forEach(function(f) {
        if (f.verified_photos && f.verified_photos.length && !faceVerifyCache[f.name]) {
            var statuses = {};
            f.verified_photos.forEach(function(fn) { statuses[fn] = 'ok'; });
            faceVerifyCache[f.name] = statuses;
        }
    });

    container.innerHTML = faces.map(f => {
        var divScore = f.diversity_score != null ? Math.round(f.diversity_score * 100) : null;
        var divColor = divScore !== null ? (divScore >= 60 ? '#4caf50' : divScore >= 40 ? '#ff9800' : '#f44336') : '#a0aec0';
        var divBadge = divScore !== null ? '<span style="font-size:.8em; font-weight:600; color:' + divColor + '; margin-left:8px;">Diversity: ' + divScore + '%</span>' : '';
        var encBadge = f.has_encodings ? '<span style="font-size:.75em; color:#38a169; margin-left:6px;" title="Face encodings cached — no re-detection needed">&#10003; Verified</span>' : '';
        return '<div style="background:#f7fafc; border:1px solid ' + (f.has_encodings ? '#9ae6b4' : '#e2e8f0') + '; border-radius:8px; padding:15px; margin-bottom:12px;">' +
            '<div style="display:flex; justify-content:space-between; align-items:center; margin-bottom:10px;">' +
                '<h3 style="color:#2b6cb0; margin:0; text-transform:capitalize;">' + esc(f.name) + encBadge + '</h3>' +
                '<div style="display:flex; gap:8px; align-items:center;">' +
                    '<span style="color:#718096;">' + f.photo_count + ' photo(s)</span>' + divBadge +
                    '<button class="btn btn-secondary" onclick="saveToLibrary(\\\'' + esc(f.name) + '\\\')" style="margin:0; padding:4px 10px; font-size:.8em; background:#ebf8ff; color:#2b6cb0;" title="Save to global face library for reuse in other projects">Save to Library</button>' +
                    '<button class="btn btn-secondary" onclick="removePerson(\\\'' + esc(f.name) + '\\\')" style="margin:0; padding:4px 10px; font-size:.8em; background:#fed7d7; color:#c53030;">Remove</button>' +
                '</div>' +
            '</div>' +
            '<div id="face-thumbs-' + f.name + '" style="display:flex; gap:6px; flex-wrap:wrap; margin-bottom:10px;"></div>' +
            '<div style="display:flex; gap:8px; align-items:center;">' +
                '<label style="margin:0; cursor:pointer;" class="btn btn-secondary" for="upload-' + f.name + '">+ Add Photos</label>' +
                '<input type="file" id="upload-' + f.name + '" multiple accept="image/*" style="display:none" onchange="uploadFacePhotos(\\\'' + esc(f.name) + '\\\', this.files)">' +
                '<button class="btn btn-secondary" onclick="verifyPerson(\\\'' + esc(f.name) + '\\\')" style="margin:0">Verify Face</button>' +
            '</div>' +
            '<div id="face-status-' + f.name + '" style="margin-top:8px;"></div>' +
        '</div>';
    }).join('');

    // Load thumbnails for each person
    faces.forEach(f => loadFaceThumbs(f.name));

    // Show face match mode selector when multiple persons have photos
    var personsWithPhotos = faces.filter(f => f.photo_count > 0);
    var matchArea = document.getElementById('face-match-mode');
    if (personsWithPhotos.length > 1) {
        matchArea.style.display = 'block';
        var names = personsWithPhotos.map(function(f) { return f.name.charAt(0).toUpperCase() + f.name.slice(1); });
        document.getElementById('face-match-people').textContent = 'People: ' + names.join(', ');
        updateMatchModeStyle();
    } else if (personsWithPhotos.length === 1) {
        matchArea.style.display = 'none';
    } else {
        matchArea.style.display = 'none';
    }
}

let faceVerifyCache = {};
let replacedPhotos = {};  // {person: Set of filenames}

async function loadFaceThumbs(person) {
    const res = await fetch('/api/ref-faces/' + encodeURIComponent(person) + '/photos');
    const photos = await res.json();
    const container = document.getElementById('face-thumbs-' + person);
    if (!container) return;
    const statuses = faceVerifyCache[person] || {};
    const replaced = replacedPhotos[person] || new Set();
    container.innerHTML = '';
    photos.forEach(function(p) {
        var st = statuses[p.filename];
        var isReplaced = replaced.has(p.filename) && !st;
        var border = st === 'ok' || st === 'ok_multi' ? '3px solid #4caf50' : st === 'no_face' || st === 'encode_fail' || st === 'error' ? '3px solid #f44336' : isReplaced ? '3px dashed #dd6b20' : '2px solid #cbd5e0';
        var hasAnyVerified = Object.keys(statuses).length > 0;

        var wrap = document.createElement('div');
        wrap.style.cssText = 'display:inline-flex; flex-direction:column; align-items:center; gap:2px; margin-right:8px; margin-bottom:6px;';

        // Verify or Replace button
        if (isReplaced && hasAnyVerified) {
            var verifySpan = document.createElement('span');
            verifySpan.textContent = 'Verify';
            verifySpan.style.cssText = 'font-size:.65em; color:#e53e3e; cursor:pointer; padding:1px 4px; background:#fff5f5; border:1px solid #fed7d7; border-radius:3px; white-space:nowrap;';
            (function(pn, fn) { verifySpan.onclick = function() { verifySinglePhoto(pn, fn); }; })(person, p.filename);
            wrap.appendChild(verifySpan);
        } else {
            var replaceLabel = document.createElement('label');
            replaceLabel.textContent = 'Replace';
            replaceLabel.style.cssText = 'font-size:.65em; color:#3182ce; cursor:pointer; padding:1px 4px; background:#ebf8ff; border-radius:3px; white-space:nowrap;';
            replaceLabel.setAttribute('for', 'replace-' + person + '-' + p.filename);
            wrap.appendChild(replaceLabel);
        }

        // Hidden file input for replace
        var fileInput = document.createElement('input');
        fileInput.type = 'file';
        fileInput.id = 'replace-' + person + '-' + p.filename;
        fileInput.accept = 'image/*';
        fileInput.style.display = 'none';
        (function(pn, fn) { fileInput.onchange = function() { replaceFacePhoto(pn, fn, this.files); }; })(person, p.filename);
        wrap.appendChild(fileInput);

        // Image container
        var imgDiv = document.createElement('div');
        imgDiv.style.cssText = 'position:relative;';
        imgDiv.id = 'face-img-' + person + '-' + p.filename.replace(/[^a-zA-Z0-9]/g, '-');
        imgDiv.onmouseenter = function() { var x = this.querySelector('.face-x'); if (x) x.style.display = 'flex'; };
        imgDiv.onmouseleave = function() { var x = this.querySelector('.face-x'); if (x) x.style.display = 'none'; };

        if (p.thumb) {
            var img = document.createElement('img');
            img.src = 'data:image/jpeg;base64,' + p.thumb;
            img.style.cssText = 'width:70px; height:70px; object-fit:cover; border-radius:4px; border:' + border + '; cursor:pointer;';
            img.title = 'Double-click to enlarge';
            (function(pn, fn) { img.ondblclick = function() { openLightbox(pn, fn); }; })(person, p.filename);
            imgDiv.appendChild(img);
        } else {
            var placeholder = document.createElement('div');
            placeholder.style.cssText = 'width:70px; height:70px; background:#e2e8f0; border-radius:4px; border:' + border + '; display:flex; align-items:center; justify-content:center; font-size:.7em; color:#718096;';
            placeholder.textContent = p.filename;
            imgDiv.appendChild(placeholder);
        }

        // Remove X button
        var xBtn = document.createElement('span');
        xBtn.className = 'face-x remove-face-btn';
        xBtn.setAttribute('data-person', person);
        xBtn.setAttribute('data-filename', p.filename);
        xBtn.innerHTML = '&#x2715;';
        xBtn.style.cssText = 'display:none; position:absolute; top:-4px; right:-4px; width:18px; height:18px; background:#e53e3e; color:white; border-radius:50%; font-size:11px; align-items:center; justify-content:center; cursor:pointer; line-height:1; box-shadow:0 1px 3px rgba(0,0,0,.3);';
        imgDiv.appendChild(xBtn);

        // Status label
        if (st === 'no_face') {
            var lbl = document.createElement('span');
            lbl.innerHTML = '&#x2718;';
            lbl.style.cssText = 'position:absolute; bottom:1px; right:3px; font-size:14px; color:#f44336; text-shadow:0 0 3px #fff;';
            imgDiv.appendChild(lbl);
        } else if (st === 'ok' || st === 'ok_multi') {
            var lbl = document.createElement('span');
            lbl.innerHTML = '&#x2714;';
            lbl.style.cssText = 'position:absolute; bottom:1px; right:3px; font-size:14px; color:#4caf50; text-shadow:0 0 3px #fff;';
            imgDiv.appendChild(lbl);
        }

        wrap.appendChild(imgDiv);

        // Rotate buttons
        var rotDiv = document.createElement('div');
        rotDiv.style.cssText = 'display:flex; gap:2px;';
        var rotL = document.createElement('button');
        rotL.innerHTML = '&#x21BA;';
        rotL.title = 'Rotate left';
        rotL.style.cssText = 'font-size:.7em; padding:1px 5px; cursor:pointer; background:#edf2f7; border:1px solid #cbd5e0; border-radius:3px;';
        (function(pn, fn) { rotL.onclick = function() { rotateFace(pn, fn, 'ccw'); }; })(person, p.filename);
        var rotR = document.createElement('button');
        rotR.innerHTML = '&#x21BB;';
        rotR.title = 'Rotate right';
        rotR.style.cssText = 'font-size:.7em; padding:1px 5px; cursor:pointer; background:#edf2f7; border:1px solid #cbd5e0; border-radius:3px;';
        (function(pn, fn) { rotR.onclick = function() { rotateFace(pn, fn, 'cw'); }; })(person, p.filename);
        rotDiv.appendChild(rotL);
        rotDiv.appendChild(rotR);
        wrap.appendChild(rotDiv);

        container.appendChild(wrap);
    });

    // Add "Verify All Replaced" button if multiple unverified replacements
    const unverified = [...replaced].filter(fn => !statuses[fn]);
    if (unverified.length > 1) {
        var btn = document.createElement('div');
        btn.style.marginTop = '6px';
        var b = document.createElement('button');
        b.className = 'btn btn-secondary';
        b.style.cssText = 'font-size:.8em; padding:4px 12px; color:#e53e3e; border-color:#fed7d7;';
        b.textContent = 'Verify ' + unverified.length + ' Replaced Photos';
        b.onclick = function() { verifyReplacedPhotos(person); };
        btn.appendChild(b);
        container.appendChild(btn);
    }
}

// Event delegation for remove-face buttons (avoids inline onclick quote issues)
document.addEventListener('click', function(e) {
    if (e.target.classList.contains('remove-face-btn')) {
        var person = e.target.getAttribute('data-person');
        var filename = e.target.getAttribute('data-filename');
        removeFacePhoto(person, filename);
    }
});

async function removeFacePhoto(person, filename) {
    if (!confirm('Remove "' + filename + '" from ' + person + "'s reference photos?")) return;
    await fetch('/api/ref-faces/' + encodeURIComponent(person) + '/photo/' + encodeURIComponent(filename), {method: 'DELETE'});
    // Clean caches
    if (faceVerifyCache[person]) delete faceVerifyCache[person][filename];
    if (replacedPhotos[person]) replacedPhotos[person].delete(filename);
    await loadFaceThumbs(person);
    await loadFaces();
}

async function verifySinglePhoto(person, filename) {
    // Show spinner on the image
    showFaceLoader(person, filename);
    const res = await fetch('/api/ref-faces/' + encodeURIComponent(person) + '/verify-photos', {
        method: 'POST', headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ filenames: [filename] })
    });
    const data = await res.json();
    if (data.results && data.results[filename]) {
        if (!faceVerifyCache[person]) faceVerifyCache[person] = {};
        faceVerifyCache[person][filename] = data.results[filename].status;
    }
    // Remove from replaced set since it's now verified
    if (replacedPhotos[person]) replacedPhotos[person].delete(filename);
    await loadFaceThumbs(person);
    // Show result
    const statusEl = document.getElementById('face-status-' + person);
    if (statusEl && data.results && data.results[filename]) {
        const v = data.results[filename];
        const color = v.status === 'ok' ? '#38a169' : v.status === 'ok_multi' ? '#dd6b20' : '#e53e3e';
        statusEl.innerHTML = '<span style="color:' + color + ';">' + esc(filename) + ': ' + esc(v.message) + '</span>';
    }
}

async function verifyReplacedPhotos(person) {
    const replaced = replacedPhotos[person];
    if (!replaced || !replaced.size) return;
    const filenames = [...replaced].filter(fn => !(faceVerifyCache[person] || {})[fn]);
    if (!filenames.length) return;

    const statusEl = document.getElementById('face-status-' + person);
    showLoader(statusEl, 'Verifying ' + filenames.length + ' replaced photos...');

    const res = await fetch('/api/ref-faces/' + encodeURIComponent(person) + '/verify-photos', {
        method: 'POST', headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ filenames })
    });
    const data = await res.json();
    if (data.results) {
        if (!faceVerifyCache[person]) faceVerifyCache[person] = {};
        for (const [fn, result] of Object.entries(data.results)) {
            faceVerifyCache[person][fn] = result.status;
            if (replacedPhotos[person]) replacedPhotos[person].delete(fn);
        }
    }
    await loadFaceThumbs(person);

    // Show summary
    if (statusEl && data.results) {
        const entries = Object.entries(data.results);
        const ok = entries.filter(([,v]) => v.status === 'ok' || v.status === 'ok_multi').length;
        const fail = entries.length - ok;
        statusEl.innerHTML = '<span style="color:' + (fail ? '#e53e3e' : '#38a169') + ';">' + ok + '/' + entries.length + ' faces detected' + (fail ? ' (' + fail + ' failed)' : '') + '</span>';
    }
}

async function addPerson() {
    const inp = document.getElementById('inp-face-name');
    const name = inp.value.trim().toLowerCase();
    if (!name) return;

    // Create folder via upload with no files (just creates the dir)
    await fetch('/api/ref-faces/upload', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({person: name, photos: []})
    });
    inp.value = '';
    // New person has no encodings — force re-verification
    facesVerified = false;
    document.getElementById('btn-next-3').disabled = true;
    loadFaceStep();
}

async function showFaceLibrary() {
    var panel = document.getElementById('face-library-panel');
    panel.style.display = panel.style.display === 'none' ? 'block' : 'none';
    if (panel.style.display === 'none') return;
    var listEl = document.getElementById('face-library-list');
    listEl.innerHTML = 'Loading...';
    try {
        var res = await fetch('/api/face-library');
        var people = await res.json();
        if (!people.length) {
            listEl.innerHTML = '<p style="color:#a0aec0; font-style:italic;">No saved people yet. Add reference photos and click "Save to Library" to build your library.</p>';
            return;
        }
        listEl.innerHTML = '';
        people.forEach(function(p) {
            var row = document.createElement('div');
            row.style.cssText = 'display:flex; align-items:center; justify-content:space-between; padding:8px 12px; border:1px solid #e2e8f0; border-radius:6px; margin-bottom:6px; background:white;';
            var info = document.createElement('div');
            info.style.cssText = 'display:flex; align-items:center; gap:10px;';
            var nameEl = document.createElement('span');
            nameEl.style.cssText = 'font-weight:600; color:#2d3748; text-transform:capitalize;';
            nameEl.textContent = p.name;
            var countEl = document.createElement('span');
            countEl.style.cssText = 'font-size:.8em; color:#a0aec0;';
            countEl.textContent = p.photo_count + ' photo(s)';
            info.appendChild(nameEl);
            info.appendChild(countEl);
            var btns = document.createElement('div');
            btns.style.cssText = 'display:flex; gap:6px;';
            var useBtn = document.createElement('button');
            useBtn.className = 'btn btn-secondary';
            useBtn.style.cssText = 'margin:0; padding:4px 12px; font-size:.8em; background:#f0fff4; color:#276749; border-color:#9ae6b4;';
            useBtn.textContent = 'Use in Project';
            (function(name) {
                useBtn.onclick = function() { importFromLibrary(name); };
            })(p.name);
            var delBtn = document.createElement('button');
            delBtn.className = 'btn btn-secondary';
            delBtn.style.cssText = 'margin:0; padding:4px 8px; font-size:.8em; background:#fed7d7; color:#c53030;';
            delBtn.textContent = 'Delete';
            (function(name) {
                delBtn.onclick = function() { deleteFromLibrary(name); };
            })(p.name);
            btns.appendChild(useBtn);
            btns.appendChild(delBtn);
            row.appendChild(info);
            row.appendChild(btns);
            listEl.appendChild(row);
        });
    } catch(e) { listEl.innerHTML = 'Error loading library'; }
}

async function importFromLibrary(name) {
    var res = await fetch('/api/face-library/import', {
        method: 'POST', headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ person: name })
    });
    var data = await res.json();
    if (data.error) { alert(data.error); return; }
    document.getElementById('face-library-panel').style.display = 'none';
    // Library faces are pre-verified — enable Next button
    facesVerified = true;
    document.getElementById('btn-next-3').disabled = false;
    document.getElementById('btn-skip-faces').style.display = 'none';
    loadFaceStep();
}

async function saveToLibrary(name) {
    var statusEl = document.getElementById('face-status-' + name);
    if (statusEl) statusEl.innerHTML = '<span style="color:#dd6b20; font-size:.85em;">Verifying all photos before saving...</span>';

    // First verify all photos have detectable faces
    var vRes = await fetch('/api/ref-faces/verify', {
        method: 'POST', headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ person: name })
    });
    var vData = await vRes.json();
    var personResult = vData.persons && vData.persons[0];
    if (!personResult) { alert('Verification failed'); return; }

    if (personResult.fail_count > 0) {
        var badPhotos = personResult.photos.filter(function(p) { return p.status !== 'ok' && p.status !== 'ok_multi'; });
        var names = badPhotos.map(function(p) { return p.filename; }).join(', ');
        if (statusEl) statusEl.innerHTML = '<span style="color:#e53e3e; font-size:.85em;">Cannot save: ' + personResult.fail_count + ' photo(s) have no detectable face (' + names + '). Remove or replace them first.</span>';
        return;
    }

    // All photos verified — now save
    if (statusEl) statusEl.innerHTML = '<span style="color:#dd6b20; font-size:.85em;">Saving to library...</span>';
    var res = await fetch('/api/face-library/save', {
        method: 'POST', headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ person: name })
    });
    var data = await res.json();
    if (data.error) { alert(data.error); return; }
    if (statusEl) statusEl.innerHTML = '<span style="color:#38a169; font-size:.85em;">Saved to library (' + data.photos_saved + ' photos, ' + data.encodings + ' face encodings)</span>';
}

async function deleteFromLibrary(name) {
    if (!confirm('Remove ' + name + ' from the global library?')) return;
    await fetch('/api/face-library/' + encodeURIComponent(name), { method: 'DELETE' });
    showFaceLibrary();
}

async function removePerson(name) {
    if (!confirm('Remove all reference photos for ' + name + '?')) return;
    await fetch('/api/ref-faces/' + encodeURIComponent(name), {method: 'DELETE'});
    // Reset verification — loadFaceStep will re-check if all remaining have encodings
    facesVerified = false;
    loadFaceStep();
}

function showFaceLoader(person, filename) {
    const safeId = 'face-img-' + person + '-' + filename.replace(/[^a-zA-Z0-9]/g, '-');
    const wrapper = document.getElementById(safeId);
    if (!wrapper) return;
    const imgEl = wrapper.querySelector('img') || wrapper.querySelector('div[style*="70px"]');
    if (imgEl) imgEl.style.opacity = '0.3';
    const old = wrapper.querySelector('.face-spinner');
    if (old) old.remove();
    const spinner = document.createElement('div');
    spinner.className = 'face-spinner';
    spinner.style.cssText = 'position:absolute;top:50%;left:50%;transform:translate(-50%,-50%);width:20px;height:20px;border:3px solid #e2e8f0;border-top:3px solid #3182ce;border-radius:50%;animation:spin 0.8s linear infinite;';
    wrapper.appendChild(spinner);
}

function markFacesDirty() {
    faceVerifyCache = {};
    facesVerified = false;
    document.getElementById('btn-next-3').disabled = true;
    document.getElementById('face-verify-btn-group').style.display = 'flex';
}

async function rotateFace(person, filename, direction) {
    showFaceLoader(person, filename);
    await fetch('/api/ref-faces/' + encodeURIComponent(person) + '/rotate', {
        method: 'POST', headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({filename, direction})
    });
    markFacesDirty();
    await loadFaceThumbs(person);
}

async function replaceFacePhoto(person, filename, files) {
    if (!files || !files.length) return;
    showFaceLoader(person, filename);

    const formData = new FormData();
    formData.append('filename', filename);
    formData.append('photo', files[0]);

    await fetch('/api/ref-faces/' + encodeURIComponent(person) + '/replace', {method: 'POST', body: formData});

    // Track as replaced, clear old verify status
    if (!replacedPhotos[person]) replacedPhotos[person] = new Set();
    replacedPhotos[person].add(filename);
    if (faceVerifyCache[person]) delete faceVerifyCache[person][filename];

    markFacesDirty();
    await loadFaceThumbs(person);
}

async function uploadFacePhotos(person, files) {
    if (!files.length) return;
    const formData = new FormData();
    formData.append('person', person);
    for (const f of files) formData.append('photos', f);

    const statusEl = document.getElementById('face-status-' + person);
    statusEl.innerHTML = '<span style="color:#3182ce;">Uploading ' + files.length + ' photo(s)...</span>';

    await fetch('/api/ref-faces/upload', {method: 'POST', body: formData});
    statusEl.innerHTML = '<span style="color:#4caf50;">Uploaded! Click <strong>Verify All Faces</strong> to check.</span>';
    facesVerified = false;
    document.getElementById('btn-next-3').disabled = true;
    loadFaceStep();
}

async function verifyPerson(person) {
    const statusEl = document.getElementById('face-status-' + person);
    showLoader(statusEl, 'Verifying face encodings...');

    const res = await fetch('/api/ref-faces/verify', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({person})
    });
    const data = await res.json();

    if (data.error) {
        statusEl.innerHTML = '<div style="color:#f44336;">' + esc(data.error) + '</div>';
        return;
    }

    const p = data.persons[0];
    if (!p) return;

    // Cache statuses and re-render thumbnails with color borders
    const statuses = {};
    p.photos.forEach(ph => { statuses[ph.filename] = ph.status; });
    faceVerifyCache[person] = statuses;
    loadFaceThumbs(person);

    let html = '<div style="margin-top:8px;">';

    // Per-photo results
    html += '<table style="width:100%; font-size:.85em; margin-bottom:10px;"><thead><tr style="color:#718096;"><th style="text-align:left;">Photo</th><th>Size</th><th>Status</th></tr></thead><tbody>';
    for (const photo of p.photos) {
        const color = photo.status === 'ok' ? '#4caf50' : photo.status === 'ok_multi' ? '#ff9800' : '#f44336';
        html += '<tr><td>' + esc(photo.filename) + '</td><td style="color:#718096;">' + (photo.dimensions || '') + '</td><td style="color:' + color + ';">' + esc(photo.message) + '</td></tr>';
    }
    html += '</tbody></table>';

    // Summary
    html += '<div style="display:flex; gap:20px; flex-wrap:wrap; margin-bottom:8px;">';
    html += '<div><strong style="color:#2b6cb0;">' + p.encodings + '</strong> / ' + p.total_photos + ' faces encoded</div>';
    html += '<div>Diversity: <strong style="color:' + (p.diversity_score > 0.6 ? '#4caf50' : p.diversity_score > 0.3 ? '#ff9800' : '#f44336') + ';">' + Math.round(p.diversity_score * 100) + '%</strong></div>';
    html += '<div>Status: <strong style="color:' + (p.ready ? '#4caf50' : '#ff9800') + ';">' + (p.ready ? 'READY' : 'NEEDS WORK') + '</strong></div>';
    html += '</div>';

    // Issues
    if (p.issues.length) {
        html += '<div style="margin-bottom:6px;">';
        p.issues.forEach(i => { html += '<div style="color:#f44336; font-size:.85em;">&#9888; ' + esc(i) + '</div>'; });
        html += '</div>';
    }

    // Tips
    if (p.tips.length) {
        html += '<div>';
        p.tips.forEach(t => { html += '<div style="color:#718096; font-size:.85em;">&#128161; ' + esc(t) + '</div>'; });
        html += '</div>';
    }

    html += '</div>';
    statusEl.innerHTML = html;
}

async function verifyAllFaces() {
    const resultsEl = document.getElementById('face-verify-results');
    showLoader(resultsEl, 'Verifying all faces... this may take a moment');

    const res = await fetch('/api/ref-faces/verify', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({})
    });
    const data = await res.json();

    if (data.error) {
        resultsEl.innerHTML = '<div style="color:#f44336;">' + esc(data.error) + '</div>';
        return;
    }

    // Cache statuses, re-render thumbnails, and update inline per-person details
    for (const p of (data.persons || [])) {
        const statuses = {};
        p.photos.forEach(ph => { statuses[ph.filename] = ph.status; });
        faceVerifyCache[p.person] = statuses;
        loadFaceThumbs(p.person);

        // Update inline status for each person
        const statusEl = document.getElementById('face-status-' + p.person);
        if (statusEl) {
            let h = '<div style="margin-top:8px;">';
            h += '<table style="width:100%; font-size:.85em; margin-bottom:10px;"><thead><tr style="color:#718096;"><th style="text-align:left;">Photo</th><th>Size</th><th>Status</th></tr></thead><tbody>';
            for (const photo of p.photos) {
                const color = photo.status === 'ok' ? '#4caf50' : photo.status === 'ok_multi' ? '#ff9800' : '#f44336';
                h += '<tr><td>' + esc(photo.filename) + '</td><td style="color:#718096;">' + (photo.dimensions || '') + '</td><td style="color:' + color + ';">' + esc(photo.message) + '</td></tr>';
            }
            h += '</tbody></table>';
            h += '<div style="display:flex; gap:20px; flex-wrap:wrap; margin-bottom:8px;">';
            h += '<div><strong style="color:#2b6cb0;">' + p.encodings + '</strong> / ' + p.total_photos + ' faces encoded</div>';
            h += '<div>Diversity: <strong style="color:' + (p.diversity_score > 0.6 ? '#4caf50' : p.diversity_score > 0.3 ? '#ff9800' : '#f44336') + ';">' + Math.round(p.diversity_score * 100) + '%</strong></div>';
            h += '<div>Status: <strong style="color:' + (p.ready ? '#4caf50' : '#ff9800') + ';">' + (p.ready ? 'READY' : 'NEEDS WORK') + '</strong></div>';
            h += '</div>';
            if (p.issues.length) p.issues.forEach(i => { h += '<div style="color:#f44336; font-size:.85em;">&#9888; ' + esc(i) + '</div>'; });
            if (p.tips.length) p.tips.forEach(t => { h += '<div style="color:#718096; font-size:.85em;">&#128161; ' + esc(t) + '</div>'; });
            h += '</div>';
            statusEl.innerHTML = h;
        }
    }

    if (!data.persons.length) {
        resultsEl.innerHTML = '<div style="color:#718096;">No reference faces to verify.</div>';
        return;
    }

    let allReady = data.ready;
    resultsEl.innerHTML = '<div style="padding:15px; border-radius:8px; background:' + (allReady ? '#f0fff4' : '#fff5f5') + '; border:1px solid ' + (allReady ? '#4caf50' : '#f44336') + ';">' +
        '<strong style="color:' + (allReady ? '#4caf50' : '#f44336') + ';">' + (allReady ? '&#10003; All faces verified and ready!' : '&#9888; ' + esc(data.message)) + '</strong></div>';

    return allReady;
}

async function runVerifyAll() {
    const allReady = await verifyAllFaces();
    document.getElementById('btn-skip-faces').style.display = 'none';
    if (allReady) {
        facesVerified = true;
        document.getElementById('btn-next-3').disabled = false;
    } else {
        // Check specific issues — block if any person has no encodings, failed photos, or low diversity
        var res = await fetch('/api/ref-faces/verify', {
            method: 'POST', headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({})
        });
        var data = await res.json();
        var canProceed = true;
        var blockReasons = [];
        for (var p of (data.persons || [])) {
            if (p.encodings === 0) {
                canProceed = false;
                blockReasons.push(p.person + ': no faces detected');
            } else if (p.fail_count > 0) {
                canProceed = false;
                blockReasons.push(p.person + ': ' + p.fail_count + ' photo(s) have no detectable face — remove or replace them');
            } else if (p.diversity_score < 0.5) {
                canProceed = false;
                blockReasons.push(p.person + ': diversity too low (' + Math.round(p.diversity_score * 100) + '%) — add more varied photos');
            }
        }
        if (canProceed) {
            facesVerified = true;
            document.getElementById('btn-next-3').disabled = false;
        } else {
            facesVerified = false;
            document.getElementById('btn-next-3').disabled = true;
            var resultsEl = document.getElementById('face-verify-results');
            resultsEl.innerHTML += '<div style="margin-top:10px; padding:12px; background:#fff5f5; border:1px solid #feb2b2; border-radius:6px; font-size:.85em; color:#9b2c2c;">' +
                '<strong>Cannot proceed:</strong><ul style="margin:6px 0 0 16px;">' +
                blockReasons.map(function(r) { return '<li>' + r + '</li>'; }).join('') +
                '</ul></div>';
        }
    }
}

// Info tooltip hover
(function() {
    var info = document.getElementById('face-mode-info');
    var tip = document.getElementById('face-mode-tooltip');
    if (info && tip) {
        info.addEventListener('mouseenter', function() { tip.style.display = 'block'; });
        info.addEventListener('mouseleave', function() { tip.style.display = 'none'; });
    }
})();

function updateMatchModeStyle() {
    // No visual highlight — just the radio button selection is enough
}

function getFaceMatchMode() {
    var radio = document.querySelector('input[name="face-match"]:checked');
    return radio ? radio.value : 'any';
}

function skipFaces() {
    config.face_names = [];
    fetch('/api/config', { method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(config) });
    document.querySelectorAll('.step-dot')[3].classList.add('done');
    goStep(4);
}

async function completeFacesStep() {
    var btn = document.getElementById('btn-next-3');
    btn.disabled = true;
    btn.innerHTML = '<span style="display:inline-block;width:14px;height:14px;border:2px solid #bee3f8;border-top:2px solid #fff;border-radius:50%;animation:spinA .7s linear infinite;vertical-align:middle;margin-right:6px;"></span>Loading...';

    try {
        // Get list of persons with faces
        const res = await fetch('/api/ref-faces');
        const faces = await res.json();
        const personsWithPhotos = faces.filter(f => f.photo_count > 0);

        var personNames = personsWithPhotos.map(f => f.name);

        if (personNames.length > 0) {
            if (!facesVerified) {
                btn.disabled = false;
                btn.textContent = 'Next: Start Scan';
                await runVerifyAll();
                return;
            }
            const verifyRes = await fetch('/api/ref-faces/verify', {
                method: 'POST', headers: {'Content-Type': 'application/json'}, body: JSON.stringify({})
            });
            const verifyData = await verifyRes.json();
            if (!verifyData.ready) {
                if (!confirm('Some face references need more photos for reliable recognition. Proceed anyway?')) {
                    btn.disabled = false;
                    btn.textContent = 'Next: Start Scan';
                    return;
                }
            }
            config.face_names = personNames;
            config.face_match_mode = personNames.length > 1 ? getFaceMatchMode() : 'any';
        } else {
            config.face_names = [];
            config.face_match_mode = 'any';
        }

        await fetch('/api/config', { method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(config) });
        document.querySelectorAll('.step-dot')[3].classList.add('done');
        goStep(4);
    } finally {
        btn.disabled = false;
        btn.textContent = 'Next: Start Scan';
    }
}

// ── Step 4: Scan ──
async function checkExistingScan() {
    // Restore NSFW checkbox from config
    const chk = document.getElementById('chk-nsfw-filter');
    if (chk && config?.nsfw_filter) chk.checked = true;

    const res = await fetch('/api/stats');
    const st = await res.json();
    var scanBtn = document.getElementById('btn-start-scan');
    var summaryEl = document.getElementById('scan-summary-stats');
    if (st.has_scan && st.total_media > 0 && st.has_sources) {
        document.getElementById('btn-next-4').disabled = false;

        // Render stat cards
        summaryEl.style.display = 'block';
        summaryEl.innerHTML = '<div class="stat-row">' +
            '<div class="stat-card"><div class="stat-value">' + st.total_media + '</div><div class="stat-label">Images found</div></div>' +
            '<div class="stat-card green"><div class="stat-value">' + (st.qualified || 0) + '</div><div class="stat-label">Qualified</div></div>' +
            '<div class="stat-card"><div class="stat-value">' + (st.sources?.length || 0) + '</div><div class="stat-label">Sources</div></div>' +
            '</div>';

        document.getElementById('scan-progress').style.display = 'block';
        if (st.partial) {
            document.getElementById('scan-progress').innerHTML =
                '<div class="line" style="color:#dd6b20;">' +
                '<strong>Partial scan saved.</strong> Click <strong>Continue Scanning</strong> to resume from where you left off.</div>';
            scanBtn.textContent = 'Continue Scanning';
            scanBtn.onclick = function() { startScan(false); };
        } else {
            document.getElementById('scan-progress').innerHTML =
                '<div class="line" style="color:#38a169;">' +
                '<strong>Scan complete.</strong> You can proceed to analysis or rescan.</div>';
            scanBtn.textContent = 'Continue Scanning';
            scanBtn.onclick = function() { startScan(false); };
        }
    } else {
        summaryEl.style.display = 'none';
        scanBtn.textContent = 'Start Scan';
        scanBtn.onclick = function() { startScan(false); };
    }
}

function toggleNsfwFilter(checked) {
    config = config || {};
    config.nsfw_filter = checked;
    fetch('/api/config', { method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify(config) });
}

function toggleAgeEstimation(checked) {
    document.getElementById('age-est-options').style.display = checked ? 'block' : 'none';
    if (checked) buildAgeEstFolderList();
}

function showAgeEstFolders() {
    document.querySelector('input[name="age-est-scope"][value="folders"]').checked = true;
    document.getElementById('age-est-folders').style.display = 'block';
    buildAgeEstFolderList();
}

function buildAgeEstFolderList() {
    var container = document.getElementById('age-est-folders');
    var validSources = (config && config.sources) ? config.sources.filter(function(s) { return s.path && s.path.trim(); }) : [];
    if (!validSources.length) {
        container.innerHTML = '<div style="font-size:.8em; color:#a0aec0;">No sources configured yet.</div>';
        return;
    }
    var tbl = document.createElement('table');
    tbl.style.cssText = 'border-collapse:collapse;';
    validSources.forEach(function(src) {
        var tr = document.createElement('tr');
        var td1 = document.createElement('td');
        td1.style.cssText = 'padding:2px 6px 2px 0; vertical-align:middle;';
        var cb = document.createElement('input');
        cb.type = 'checkbox';
        cb.className = 'age-est-folder-cb';
        cb.value = src.path;
        cb.checked = true;
        cb.style.cssText = 'margin:0;';
        td1.appendChild(cb);
        var td2 = document.createElement('td');
        td2.style.cssText = 'padding:2px 0; font-size:.8em; color:#4a5568;';
        td2.textContent = src.path;
        tr.appendChild(td1);
        tr.appendChild(td2);
        tbl.appendChild(tr);
    });
    container.innerHTML = '';
    container.appendChild(tbl);
    // Wire up radio buttons to toggle folder list visibility
    document.querySelectorAll('input[name="age-est-scope"]').forEach(function(r) {
        r.onchange = function() {
            document.getElementById('age-est-folders').style.display = this.value === 'folders' ? 'block' : 'none';
        };
    });
}

function getAgeEstConfig() {
    var chk = document.getElementById('chk-age-estimation');
    if (!chk || !chk.checked) return null;
    var scope = document.querySelector('input[name="age-est-scope"]:checked');
    var result = { enabled: true, scope: scope ? scope.value : 'all' };
    if (result.scope === 'folders') {
        result.folders = [...document.querySelectorAll('.age-est-folder-cb:checked')].map(function(cb) { return cb.value; });
    }
    return result;
}

async function startScan(full) {
    var btn = document.getElementById('btn-start-scan');
    btn.disabled = true;
    btn.innerHTML = '<span style="display:inline-block;width:14px;height:14px;border:2px solid #bee3f8;border-top:2px solid #fff;border-radius:50%;animation:spinA .7s linear infinite;vertical-align:middle;margin-right:6px;"></span>Starting...';
    const nsfwFilter = document.getElementById('chk-nsfw-filter')?.checked || false;
    const ageEst = getAgeEstConfig();
    await fetch('/api/scan/start', { method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({full, nsfw_filter: nsfwFilter, age_estimation: ageEst}) });
    btn.disabled = false;
    btn.textContent = 'Continue Scanning';
    showTaskOverlay('scan');
    // Override completion handler
    const origPoll = taskPoll;
    const check = setInterval(async () => {
        try {
            const res = await fetch('/api/scan/status');
            const st = await res.json();
            if (st.done || st.error || st.cancelled) {
                clearInterval(check);
                if (!st.error && !st.cancelled) {
                    document.getElementById('btn-next-4').disabled = false;
                    document.querySelectorAll('.step-dot')[4].classList.add('done');
                }
                // Refresh scan status display
                checkExistingScan();
            }
        } catch(e) {}
    }, 2000);
}

// ── Step 5: Analyze ──
function renderCategoryBars(container, cats, showSelected) {
    let totalAvail = 0, totalTarget = 0, totalSel = 0;
    const rows = cats.map(c => {
        const avail = c.qualified + (c.selected || 0);
        const sel = c.selected || 0;
        const target = c.target || 75;
        totalAvail += avail; totalTarget += target; totalSel += sel;
        const count = showSelected ? sel : avail;
        const pct = Math.min(100, Math.round(count / target * 100));
        const full = count >= target;
        const low = count < target * 0.5;
        const barColor = full ? '#38a169' : low ? '#e53e3e' : '#dd6b20';
        const statusText = full ? 'FULL' : (count === 0 ? 'EMPTY' : 'FILLING');
        const statusColor = full ? '#38a169' : (count === 0 ? '#a0aec0' : '#dd6b20');
        return `
        <div style="display:flex; align-items:center; gap:12px; padding:8px 0; border-bottom:1px solid #f0f0f0;">
            <div style="min-width:140px; font-weight:500; color:#2d3748; font-size:.9em;">${esc(c.display)}</div>
            <div style="flex:1; height:8px; background:#e2e8f0; border-radius:4px; overflow:hidden;">
                <div style="height:100%; width:${pct}%; background:${barColor}; border-radius:4px; transition:width .5s;"></div>
            </div>
            <div style="min-width:80px; font-size:.85em; color:#4a5568; text-align:right; font-weight:500;">${count} / ${target}</div>
            <div style="min-width:55px; font-size:.75em; font-weight:700; color:${statusColor}; text-align:center;">${statusText}</div>
        </div>`;
    }).join('');

    const overallCount = showSelected ? totalSel : totalAvail;
    const overallPct = Math.min(100, Math.round(overallCount / (totalTarget || 1) * 100));
    const overallFull = overallCount >= totalTarget;

    container.innerHTML = `
        <div class="stat-row">
            <div class="stat-card">
                <div class="stat-value">${totalAvail}</div>
                <div class="stat-label">Available</div>
            </div>
            ${showSelected ? `<div class="stat-card green">
                <div class="stat-value">${totalSel}</div>
                <div class="stat-label">Selected</div>
            </div>` : ''}
            <div class="stat-card red">
                <div class="stat-value">${totalTarget}</div>
                <div class="stat-label">Target</div>
            </div>
        </div>
        <div style="display:flex; align-items:center; gap:12px; margin-bottom:16px;">
            <div style="font-weight:600; color:#2d3748; min-width:70px;">Overall:</div>
            <div style="flex:1; height:12px; background:#e2e8f0; border-radius:6px; overflow:hidden;">
                <div style="height:100%; width:${overallPct}%; background:${overallFull ? '#38a169' : '#3182ce'}; border-radius:6px;"></div>
            </div>
            <div style="font-weight:600; color:${overallFull ? '#38a169' : '#2d3748'}; min-width:110px; text-align:right;">${overallCount} / ${totalTarget} (${overallPct}%)</div>
        </div>
        <div class="section-title">Category Fill Status</div>
        ${rows}`;
}

async function runAnalysis() {
    const el = document.getElementById('a-quick');
    showLoader(el, 'Loading category stats...');
    const res = await fetch('/api/categories/summary');
    const cats = await res.json();
    if (!cats.length) { el.innerHTML = '<div class="notice notice-error"><strong>No scan data found.</strong> Go back and run a scan first.</div>'; return; }
    renderCategoryBars(el, cats, true);
    document.querySelectorAll('.step-dot')[5].classList.add('done');

    // Show contextual hint
    var hintEl = document.getElementById('a-next-hint');
    if (hintEl) {
        var totalSel = 0, totalTarget = 0;
        cats.forEach(function(c) { totalSel += (c.selected || 0); totalTarget += (c.target || 75); });
        if (totalSel >= totalTarget) {
            hintEl.innerHTML = '<div class="notice notice-tip" style="margin-top:12px;"><strong>All targets met!</strong> Proceed to Select to review your choices, or continue to Export.</div>';
        } else if (totalSel > 0) {
            hintEl.innerHTML = '<div class="notice notice-info" style="margin-top:12px;"><strong>' + totalSel + ' of ' + totalTarget + ' selected.</strong> Go to Select to pick more images or use Auto-fill.</div>';
        } else {
            hintEl.innerHTML = '<div class="notice notice-info" style="margin-top:12px;"><strong>Ready to select images.</strong> Proceed to the next step to start picking photos for each category.</div>';
        }
        hintEl.style.display = 'block';
    }
}

async function runFullAnalysis() {
    const el = document.getElementById('a-full-results');
    showLoader(el, 'Running full analysis (this may take a moment)...');
    document.getElementById('btn-full-analysis').disabled = true;
    const res = await fetch('/api/analyze');
    document.getElementById('btn-full-analysis').disabled = false;
    if (!res.ok) { el.innerHTML = '<div style="color:#e53e3e;">Analysis failed.</div>'; return; }
    const a = await res.json();
    el.style.display = 'block';
    let html = '<h3 style="color:#718096; margin:0 0 8px;">Recommendations</h3>';
    html += (a.recommendations || []).map(r => `
        <div class="rec-card ${r.type}">
            <div class="title">${esc(r.title)}</div>
            <div class="detail">${esc(r.detail)}</div>
        </div>`).join('');
    html += '<h3 style="color:#718096; margin:16px 0 8px;">What makes a great collection</h3>';
    html += (a.priorities || []).map((p, i) => `<div style="color:#718096; padding:3px 0;">${i+1}. ${esc(p)}</div>`).join('');
    el.innerHTML = html;
}

async function resetAllSelections() {
    if (!confirm('Reset all selections? All 900 selected images will go back to the available pool.')) return;
    showLoader('a-quick', 'Resetting selections...');
    await fetch('/api/selections/reset', { method: 'POST' });
    await runAnalysis();
}

// ── Step 6: Select ──
let selCats = [];
let selActiveCat = null;
let selImages = [];
let selOffset = 0;
const SEL_PAGE = 100;
async function loadSelCategories() {
    const res = await fetch('/api/categories/summary');
    selCats = await res.json();
    renderSelCatList();
    // Load saved quiz preferences
    try {
        const cfgRes = await fetch('/api/config');
        const cfg = await cfgRes.json();
        if (cfg.taste_quiz) {
            _prefPredictor._quizPrefs = cfg.taste_quiz;
        }
    } catch(e) {}
}

function renderSelCatList() {
    const el = document.getElementById('sel-cat-list');
    el.innerHTML = '<div class="section-title">Categories</div>';
    let totalSel = 0, totalTarget = 0;
    selCats.forEach(c => {
        totalSel += c.selected;
        totalTarget += c.target;
        const pct = c.target > 0 ? Math.min(100, Math.round(c.selected / c.target * 100)) : 0;
        const full = c.selected >= c.target;
        const empty = c.selected === 0;
        const active = selActiveCat === c.id;
        el.innerHTML += `
            <div onclick="selectCategory('${c.id}')" style="padding:8px 10px; cursor:pointer; border-radius:8px; margin-bottom:4px;
                background:${active ? '#eff6ff' : '#fff'}; border:1px solid ${active ? '#3b82f6' : '#e2e8f0'};
                ${full ? 'border-left:3px solid #22c55e;' : ''} transition:all .15s;">
                <div style="display:flex; justify-content:space-between; align-items:center;">
                    <span style="font-size:.88em; font-weight:${active ? '600' : '500'}; color:#1e293b;">${esc(c.display)}</span>
                    <span style="font-size:.72em; font-weight:600; color:${full ? '#22c55e' : empty ? '#94a3b8' : '#f59e0b'};">${full ? 'FULL' : empty ? 'EMPTY' : pct + '%'}</span>
                </div>
                <div style="font-size:.78em; color:#64748b; margin-top:2px;">${c.selected} / ${c.target} selected</div>
                <div style="height:3px; background:#e2e8f0; border-radius:2px; margin-top:4px;">
                    <div style="height:100%; width:${pct}%; background:${full ? '#22c55e' : '#3b82f6'}; border-radius:2px; transition:width .3s;"></div>
                </div>
            </div>`;
    });
    // Overall progress bar
    var overallPct = totalTarget > 0 ? Math.min(100, Math.round(totalSel / totalTarget * 100)) : 0;
    var lbl = document.getElementById('sel-overall-label');
    var fill = document.getElementById('sel-overall-fill');
    var pctEl = document.getElementById('sel-overall-pct');
    if (lbl) lbl.textContent = totalSel + ' / ' + totalTarget + ' selected';
    if (fill) fill.style.width = overallPct + '%';
    if (pctEl) pctEl.textContent = overallPct + '%';
}

async function selectCategory(catId) {
    if (catId === selActiveCat) return;
    selActiveCat = catId;
    selOffset = 0;
    renderSelCatList();
    const cat = selCats.find(c => c.id === catId);
    document.getElementById('sel-cat-title').textContent = cat ? cat.display : catId;
    if (cat) {
        var ctr = document.getElementById('sel-cat-counter');
        ctr.textContent = (cat.selected || 0) + ' / ' + cat.target + ' selected';
    }
    document.getElementById('sel-target-edit').style.display = 'flex';
    document.getElementById('sel-target-input').value = cat ? cat.target : 75;
    document.getElementById('sel-filter-bar').style.display = 'flex';
    document.getElementById('pref-library').style.display = 'block';
    await loadSelImages();
}

async function loadSelImages() {
    const grid = document.getElementById('sel-grid');
    showLoader(grid, 'Loading images...');
    // Load both selected and qualified for this category
    const [rSel, rQual] = await Promise.all([
        fetch(`/api/images?category=${selActiveCat}&status=selected&limit=500`),
        fetch(`/api/images?category=${selActiveCat}&status=qualified&limit=500`)
    ]);
    const dSel = await rSel.json();
    const dQual = await rQual.json();
    // Mark them so we know which are selected
    selImages = [
        ...dSel.images.map(i => ({...i, _sel: true})),
        ...dQual.images.map(i => ({...i, _sel: false}))
    ];
    // Train preference predictor from all rated images
    _prefPredictor.train(selImages);
    renderSelGrid();
    updateSelCounter();
    updatePrefStatsDisplay();
}

function updateSelCounter() {
    const cat = selCats.find(c => c.id === selActiveCat);
    const nSel = selImages.filter(i => i._sel).length;
    const target = cat ? cat.target : 0;
    var ctr = document.getElementById('sel-cat-counter');
    ctr.textContent = nSel + ' / ' + target + ' selected';
    ctr.style.background = nSel >= target ? '#dcfce7' : '#f1f5f9';
    ctr.style.color = nSel >= target ? '#16a34a' : '#64748b';
}

// ── Preference Predictor (learns from user likes/dislikes) ──
const _prefPredictor = {
    // Feature extraction: turn image metadata into a numeric vector
    _features(img) {
        const g = img.photo_grade || {};
        return [
            (g.resolution || 50) / 100,
            (g.sharpness || 50) / 100,
            (g.noise || 50) / 100,
            (g.compression || 50) / 100,
            (g.color || 50) / 100,
            (g.exposure || 50) / 100,
            (g.focus || 50) / 100,
            (g.distortion || 50) / 100,
            (g.composite || 50) / 100,
            img.face_count ? Math.min(img.face_count / 5, 1) : 0,
            img.has_target_face ? 1 : 0,
            img.face_distance != null ? (1 - img.face_distance) : 0.5,
            (img.size_kb || 500) / 5000,
            ((img.width || 1000) * (img.height || 1000)) / 20000000,
            img.media_type === 'video' ? 1 : 0,
        ];
    },

    // Learned weights (one per feature + bias), initialized to 0
    _weights: null,
    _bias: 0,
    _trained: false,
    _quizPrefs: null,  // from questionnaire

    // Sigmoid
    _sig(x) { return 1 / (1 + Math.exp(-Math.max(-10, Math.min(10, x)))); },

    // Train from all rated images (logistic regression via gradient descent)
    train(images) {
        const rated = images.filter(i => i.preference === 'like' || i.preference === 'dislike');
        if (rated.length < 5) { this._trained = false; return; }

        const X = rated.map(i => this._features(i));
        const Y = rated.map(i => i.preference === 'like' ? 1 : 0);
        const nFeat = X[0].length;

        // Initialize weights
        let w = new Array(nFeat).fill(0);
        let b = 0;
        const lr = 0.5;
        const epochs = 80;

        for (let ep = 0; ep < epochs; ep++) {
            let dw = new Array(nFeat).fill(0);
            let db = 0;
            for (let i = 0; i < X.length; i++) {
                let z = b;
                for (let j = 0; j < nFeat; j++) z += w[j] * X[i][j];
                const pred = this._sig(z);
                const err = pred - Y[i];
                for (let j = 0; j < nFeat; j++) dw[j] += err * X[i][j];
                db += err;
            }
            for (let j = 0; j < nFeat; j++) w[j] -= lr * dw[j] / X.length;
            b -= lr * db / X.length;
        }

        this._weights = w;
        this._bias = b;
        this._trained = true;

        // Apply quiz preferences as weight adjustments
        if (this._quizPrefs) {
            this._applyQuizBoosts();
        }
    },

    _applyQuizBoosts() {
        if (!this._weights || !this._quizPrefs) return;
        const q = this._quizPrefs;
        // Boost/penalize features based on quiz answers
        // Index mapping: 0=resolution, 1=sharpness, 2=noise, 3=compression,
        //   4=color, 5=exposure, 6=focus, 7=distortion, 8=composite,
        //   9=face_count_norm, 10=has_target, 11=face_closeness, 12=size, 13=megapixels, 14=is_video
        if (q.sharpness === 'high') { this._weights[1] += 0.3; this._weights[6] += 0.3; }
        if (q.colorful === 'yes') { this._weights[4] += 0.4; }
        if (q.faces === 'many') { this._weights[9] += 0.5; }
        if (q.faces === 'few') { this._weights[9] -= 0.3; }
        if (q.closeups === 'yes') { this._weights[11] += 0.3; }
        if (q.resolution === 'high') { this._weights[0] += 0.3; this._weights[13] += 0.3; }
        if (q.style === 'candid') { this._weights[10] -= 0.2; }
        if (q.style === 'posed') { this._weights[10] += 0.3; }
    },

    // Predict preference for an unrated image
    predict(img) {
        if (!this._trained || !this._weights) return null;
        const x = this._features(img);
        let z = this._bias;
        for (let j = 0; j < x.length; j++) z += this._weights[j] * x[j];
        const p = this._sig(z);
        if (p > 0.65) return 'like';
        if (p < 0.35) return 'dislike';
        return null;  // uncertain
    },

    // Get confidence score 0-100
    confidence(img) {
        if (!this._trained || !this._weights) return 0;
        const x = this._features(img);
        let z = this._bias;
        for (let j = 0; j < x.length; j++) z += this._weights[j] * x[j];
        const p = this._sig(z);
        return Math.round(Math.abs(p - 0.5) * 200);  // 0=uncertain, 100=very confident
    }
};

function setSelPref(img, pref) {
    img.preference = pref;
    fetch('/api/images/preference', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ hash: img.hash, preference: pref })
    }).catch(err => console.warn('Preference save failed:', err));
    // Retrain predictor with updated preferences
    _prefPredictor.train(selImages);
    renderSelGrid();
    updateSelCounter();
    updatePrefStatsDisplay();
    // Refresh library if open
    const libContent = document.getElementById('pref-lib-content');
    if (libContent && libContent.style.display !== 'none') {
        renderPrefLibrary();
    }
}

// ── Preference Questionnaire ──
let _quizShown = false;
function showPrefQuiz() {
    if (document.querySelector('.pref-quiz-overlay')) return;
    const overlay = document.createElement('div');
    overlay.className = 'pref-quiz-overlay';
    overlay.innerHTML = `
    <div class="pref-quiz">
        <h3>What kind of photos do you prefer? (optional)</h3>
        <p style="font-size:.8em; color:#718096; margin:0 0 14px;">This helps us suggest images you'll like. Skip any question you're unsure about.</p>

        <div class="q-group">
            <label>Photo sharpness preference:</label>
            <div class="q-row" data-field="sharpness">
                <span class="q-chip" data-val="any">Don't mind</span>
                <span class="q-chip" data-val="high">Must be sharp</span>
            </div>
        </div>

        <div class="q-group">
            <label>Color preference:</label>
            <div class="q-row" data-field="colorful">
                <span class="q-chip" data-val="any">No preference</span>
                <span class="q-chip" data-val="yes">Vibrant &amp; colorful</span>
                <span class="q-chip" data-val="muted">Soft / muted tones</span>
            </div>
        </div>

        <div class="q-group">
            <label>People in photos:</label>
            <div class="q-row" data-field="faces">
                <span class="q-chip" data-val="any">Mix is fine</span>
                <span class="q-chip" data-val="many">Prefer group shots</span>
                <span class="q-chip" data-val="few">Prefer 1-2 people</span>
                <span class="q-chip" data-val="none">Landscapes / objects too</span>
            </div>
        </div>

        <div class="q-group">
            <label>Close-ups of the subject:</label>
            <div class="q-row" data-field="closeups">
                <span class="q-chip" data-val="any">No preference</span>
                <span class="q-chip" data-val="yes">Love close-up portraits</span>
                <span class="q-chip" data-val="no">Prefer wider shots</span>
            </div>
        </div>

        <div class="q-group">
            <label>Photo style:</label>
            <div class="q-row" data-field="style">
                <span class="q-chip" data-val="any">Any style</span>
                <span class="q-chip" data-val="candid">Candid / natural moments</span>
                <span class="q-chip" data-val="posed">Posed / formal</span>
            </div>
        </div>

        <div class="q-group">
            <label>Image quality importance:</label>
            <div class="q-row" data-field="resolution">
                <span class="q-chip" data-val="any">Content matters more</span>
                <span class="q-chip" data-val="high">High resolution preferred</span>
            </div>
        </div>

        <hr style="border-color:#2d3748; margin:16px 0;">
        <h3 style="margin:0 0 12px; font-size:1em; color:#a78bfa;">What types of images do you love?</h3>
        <p style="font-size:.78em; color:#718096; margin:0 0 12px;">Pick all that apply — helps us prioritize the right moments.</p>

        <div class="q-group">
            <label>Moments you love (pick multiple):</label>
            <div class="q-row" data-field="moments" data-multi="true">
                <span class="q-chip" data-val="laughing">Laughing &amp; smiling</span>
                <span class="q-chip" data-val="hugging">Hugs &amp; affection</span>
                <span class="q-chip" data-val="playing">Playing &amp; action</span>
                <span class="q-chip" data-val="eating">Meals &amp; celebrations</span>
                <span class="q-chip" data-val="sleeping">Quiet / sleeping</span>
                <span class="q-chip" data-val="milestone">Milestones (first steps, school, etc.)</span>
                <span class="q-chip" data-val="silly">Funny / silly faces</span>
            </div>
        </div>

        <div class="q-group">
            <label>Scene types you prefer (pick multiple):</label>
            <div class="q-row" data-field="scenes" data-multi="true">
                <span class="q-chip" data-val="outdoor">Outdoors / nature</span>
                <span class="q-chip" data-val="beach">Beach / pool</span>
                <span class="q-chip" data-val="home">Home / everyday</span>
                <span class="q-chip" data-val="travel">Travel / vacation</span>
                <span class="q-chip" data-val="event">Events / parties</span>
                <span class="q-chip" data-val="school">School / activities</span>
                <span class="q-chip" data-val="sport">Sports / physical</span>
            </div>
        </div>

        <div class="q-group">
            <label>Who should be in the photos:</label>
            <div class="q-row" data-field="subjects" data-multi="true">
                <span class="q-chip" data-val="subject_alone">Just the subject</span>
                <span class="q-chip" data-val="with_parents">With parents</span>
                <span class="q-chip" data-val="with_siblings">With siblings</span>
                <span class="q-chip" data-val="with_friends">With friends</span>
                <span class="q-chip" data-val="with_grandparents">With grandparents</span>
                <span class="q-chip" data-val="family_group">Whole family</span>
                <span class="q-chip" data-val="pets">With pets</span>
            </div>
        </div>

        <div class="q-group">
            <label>Mood / energy:</label>
            <div class="q-row" data-field="mood">
                <span class="q-chip" data-val="any">Mix of everything</span>
                <span class="q-chip" data-val="happy">Happy &amp; upbeat</span>
                <span class="q-chip" data-val="calm">Calm &amp; intimate</span>
                <span class="q-chip" data-val="dramatic">Dramatic &amp; artistic</span>
            </div>
        </div>

        <div class="q-group">
            <label>Anything to avoid?</label>
            <div class="q-row" data-field="avoid" data-multi="true">
                <span class="q-chip" data-val="messy_bg">Messy backgrounds</span>
                <span class="q-chip" data-val="dark">Too dark images</span>
                <span class="q-chip" data-val="selfies">Selfies</span>
                <span class="q-chip" data-val="screenshots">Screenshots / memes</span>
                <span class="q-chip" data-val="duplicates">Similar / near-duplicates</span>
                <span class="q-chip" data-val="no_face">Photos without faces</span>
            </div>
        </div>

        <div class="q-btns">
            <button class="q-btn q-btn-skip" onclick="closePrefQuiz()">Skip</button>
            <button class="q-btn q-btn-save" onclick="savePrefQuiz()">Apply Preferences</button>
        </div>
    </div>`;
    document.body.appendChild(overlay);

    // Chip toggle logic: single-select or multi-select
    overlay.querySelectorAll('.q-chip').forEach(chip => {
        chip.onclick = () => {
            const row = chip.parentElement;
            if (row.dataset.multi === 'true') {
                // Multi-select: toggle individual chip
                chip.classList.toggle('active');
            } else {
                // Single-select: deselect others, select this
                row.querySelectorAll('.q-chip').forEach(c => c.classList.remove('active'));
                chip.classList.add('active');
            }
        };
    });
}

function closePrefQuiz() {
    document.querySelector('.pref-quiz-overlay')?.remove();
}

function savePrefQuiz() {
    const prefs = {};
    document.querySelectorAll('.pref-quiz .q-row').forEach(row => {
        const field = row.dataset.field;
        if (row.dataset.multi === 'true') {
            // Collect all active chips as array
            const vals = [...row.querySelectorAll('.q-chip.active')].map(c => c.dataset.val);
            if (vals.length) prefs[field] = vals;
        } else {
            const active = row.querySelector('.q-chip.active');
            if (active) prefs[field] = active.dataset.val;
        }
    });
    _prefPredictor._quizPrefs = prefs;
    // Save to server for persistence
    fetch('/api/preferences/quiz', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify(prefs)
    }).catch(() => {});
    // Retrain with quiz applied
    if (_prefPredictor._trained) _prefPredictor._applyQuizBoosts();
    closePrefQuiz();
    renderSelGrid();
}

function renderSelGrid() {
    const grid = document.getElementById('sel-grid');
    const showSelected = document.getElementById('sel-show-selected').checked;
    let visible = showSelected ? selImages : selImages.filter(i => !i._sel);
    // Pagination
    const total = visible.length;
    const paged = visible.slice(selOffset, selOffset + SEL_PAGE);
    const pagingEl = document.getElementById('sel-grid-paging');
    if (total > SEL_PAGE) {
        pagingEl.style.display = 'block';
        document.getElementById('sel-page-info').textContent =
            ` ${selOffset + 1}-${Math.min(selOffset + SEL_PAGE, total)} of ${total} `;
        document.getElementById('sel-prev').disabled = selOffset === 0;
        document.getElementById('sel-next').disabled = selOffset + SEL_PAGE >= total;
    } else {
        pagingEl.style.display = 'none';
    }
    grid.innerHTML = '';
    paged.forEach(img => {
        const div = document.createElement('div');
        div.className = 'sel-thumb';
        div.style.cssText = `width:100px; height:100px; border-radius:6px; overflow:hidden; cursor:pointer; position:relative;
            border:3px solid ${img._sel ? '#3182ce' : img.preference === 'like' ? '#48bb78' : img.preference === 'dislike' ? '#e53e3e' : '#e2e8f0'}; flex-shrink:0;`;
        if (img._sel) div.style.boxShadow = '0 0 0 2px #bee3f8';
        const thumbSrc = img.thumb ? 'data:image/jpeg;base64,' + img.thumb : '';
        const isVid = img.media_type === 'video';
        const vidBadge = isVid ? '<div style="position:absolute; bottom:2px; left:2px; background:rgba(0,0,0,.7); color:#fff; border-radius:4px; padding:1px 5px; font-size:10px; font-weight:600;">&#9654; VID</div>' : '';
        const gc = img.photo_grade?.composite;
        const gradeColor = gc >= 70 ? '#48bb78' : gc >= 40 ? '#ed8936' : gc != null ? '#e53e3e' : '';
        const gradeBadge = gc != null ? `<div style="position:absolute; bottom:2px; right:2px; background:${gradeColor}; color:#fff; border-radius:4px; padding:1px 5px; font-size:10px; font-weight:700;">${Math.round(gc)}</div>` : '';
        // Predicted preference indicator
        const pred = _prefPredictor.predict(img);
        const predBadge = (!img.preference && pred) ? `<div class="sel-pred-badge" style="position:absolute; top:2px; left:2px; font-size:10px; background:rgba(0,0,0,.6); color:${pred==='like'?'#48bb78':'#e53e3e'}; border-radius:3px; padding:1px 4px;">AI: ${pred==='like'?'&#x1F44D;':'&#x1F44E;'}</div>` : '';
        const prefIndicator = img.preference === 'like' ? 'active-like' : img.preference === 'dislike' ? 'active-dislike' : '';
        div.innerHTML = `<img src="${thumbSrc}" style="width:100%; height:100%; object-fit:cover;">
            ${vidBadge}${gradeBadge}${predBadge}
            ${img._sel ? '<div style="position:absolute; top:2px; right:2px; background:#3182ce; color:#fff; border-radius:50%; width:18px; height:18px; font-size:12px; display:flex; align-items:center; justify-content:center;">&#10003;</div>' : ''}
            <div class="sel-hover-overlay">
                <button class="sel-pref-btn sel-like-btn ${prefIndicator === 'active-like' ? 'sel-pref-active-like' : ''}" data-hash="${img.hash}" data-action="like" title="Like">&#x1F44D;</button>
                <button class="sel-pref-btn sel-dislike-btn ${prefIndicator === 'active-dislike' ? 'sel-pref-active-dislike' : ''}" data-hash="${img.hash}" data-action="dislike" title="Dislike">&#x1F44E;</button>
            </div>`;
        // Like/dislike button handlers
        div.querySelectorAll('.sel-pref-btn').forEach(btn => {
            btn.onclick = (e) => {
                e.stopPropagation();
                const action = btn.dataset.action;
                const newPref = img.preference === action ? null : action;
                setSelPref(img, newPref);
            };
        });
        div.onclick = (e) => {
            if (e.target.classList.contains('sel-pref-btn')) return;
            toggleSelImage(img);
        };
        div.ondblclick = (e) => { e.stopPropagation(); showSelLightbox(img); };
        grid.appendChild(div);
    });
    if (paged.length === 0) {
        grid.innerHTML = '<div style="color:#718096; padding:20px;">No photos or videos in this category.</div>';
    }
}

function sortSelImages() {
    const mode = document.getElementById('sel-sort')?.value || 'date';
    if (mode === 'grade') {
        selImages.sort((a, b) => {
            const ga = a.photo_grade?.composite || 0;
            const gb = b.photo_grade?.composite || 0;
            return gb - ga;
        });
    } else if (mode === 'pref') {
        const prefOrder = {like: 0, undefined: 1, null: 1, dislike: 2};
        selImages.sort((a, b) => (prefOrder[a.preference] || 1) - (prefOrder[b.preference] || 1));
    } else if (mode === 'predicted') {
        // Sort by AI predicted preference (liked first, then confidence)
        selImages.sort((a, b) => {
            // Explicit prefs first
            const pa = a.preference === 'like' ? 2 : a.preference === 'dislike' ? -2 : 0;
            const pb = b.preference === 'like' ? 2 : b.preference === 'dislike' ? -2 : 0;
            if (pa !== pb) return pb - pa;
            // Then by prediction confidence
            return _prefPredictor.confidence(b) - _prefPredictor.confidence(a);
        });
    } else {
        selImages.sort((a, b) => (a.date || '').localeCompare(b.date || ''));
    }
    selOffset = 0;
    renderSelGrid();
}

function togglePrefLibrary() {
    const content = document.getElementById('pref-lib-content');
    const btn = document.getElementById('pref-lib-toggle');
    if (content.style.display === 'none') {
        content.style.display = 'block';
        btn.innerHTML = '&#x2764; Hide Library';
        renderPrefLibrary();
    } else {
        content.style.display = 'none';
        btn.innerHTML = '&#x2764; Show Liked &amp; Disliked Library';
    }
}

async function renderPrefLibrary() {
    // Fetch ALL images with preferences (across all categories)
    const res = await fetch('/api/preferences/summary');
    const data = await res.json();

    const likedGrid = document.getElementById('pref-lib-liked');
    const dislikedGrid = document.getElementById('pref-lib-disliked');
    document.getElementById('pref-lib-like-count').textContent = data.liked;
    document.getElementById('pref-lib-dislike-count').textContent = data.disliked;

    // Need thumbnails — fetch from scan db
    const allRes = await fetch('/api/images?limit=10000');
    const allData = await allRes.json();
    const byHash = {};
    allData.images.forEach(i => byHash[i.hash] = i);

    function renderMiniGrid(container, hashes, borderColor) {
        container.innerHTML = '';
        if (hashes.length === 0) {
            container.innerHTML = '<div style="color:#718096; padding:10px; font-size:.85em;">None yet — hover images and click the thumbs up/down buttons.</div>';
            return;
        }
        hashes.forEach(hash => {
            const img = byHash[hash];
            if (!img) return;
            const div = document.createElement('div');
            div.style.cssText = `width:70px; height:70px; border-radius:6px; overflow:hidden; cursor:pointer; position:relative;
                border:2px solid ${borderColor}; flex-shrink:0;`;
            const thumbSrc = img.thumb ? 'data:image/jpeg;base64,' + img.thumb : '';
            const gc = img.photo_grade?.composite;
            const gradeBadge = gc != null ? `<div style="position:absolute; bottom:1px; right:1px; background:rgba(0,0,0,.7); color:#fff; border-radius:3px; padding:0 3px; font-size:9px; font-weight:700;">${Math.round(gc)}</div>` : '';
            div.innerHTML = `<img src="${thumbSrc}" style="width:100%; height:100%; object-fit:cover;">${gradeBadge}`;
            div.title = (img.filename || '') + (img.date ? ' | ' + img.date : '') + (gc != null ? ' | Grade: ' + Math.round(gc) : '');
            div.ondblclick = () => showSelLightbox(img);
            div.onclick = () => {
                // Toggle preference off on click
                const pref = img.preference;
                setSelPref(img, null);
                renderPrefLibrary();
            };
            container.appendChild(div);
        });
    }

    const likedHashes = data.preferences.filter(p => p.preference === 'like').map(p => p.hash);
    const dislikedHashes = data.preferences.filter(p => p.preference === 'dislike').map(p => p.hash);
    renderMiniGrid(likedGrid, likedHashes, '#48bb78');
    renderMiniGrid(dislikedGrid, dislikedHashes, '#e53e3e');
}

function updatePrefStatsDisplay() {
    const el = document.getElementById('sel-pref-stats');
    if (!el) return;
    const liked = selImages.filter(i => i.preference === 'like').length;
    const disliked = selImages.filter(i => i.preference === 'dislike').length;
    const predicted = _prefPredictor._trained ? selImages.filter(i => !i.preference && _prefPredictor.predict(i)).length : 0;
    let html = '';
    if (liked) html += `<span style="color:#48bb78">&#x1F44D; ${liked}</span> `;
    if (disliked) html += `<span style="color:#e53e3e">&#x1F44E; ${disliked}</span> `;
    if (predicted) html += `<span style="color:#667eea">AI: ${predicted} predicted</span>`;
    el.innerHTML = html;
}

function selPage(dir) {
    selOffset += dir * SEL_PAGE;
    if (selOffset < 0) selOffset = 0;
    renderSelGrid();
}

async function toggleSelImage(img) {
    const action = img._sel ? 'deselect' : 'select';
    img._sel = !img._sel;
    renderSelGrid();
    updateSelCounter();
    // Update server
    await fetch('/api/images/select', {
        method: 'POST', headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ hashes: [img.hash], action })
    });
    // Update sidebar count
    const cat = selCats.find(c => c.id === selActiveCat);
    if (cat) {
        cat.selected += (action === 'select' ? 1 : -1);
        renderSelCatList();
    }
}

async function selectAllVisible() {
    const unsel = selImages.filter(i => !i._sel);
    if (unsel.length === 0) return;
    const hashes = unsel.map(i => i.hash);
    unsel.forEach(i => i._sel = true);
    renderSelGrid(); updateSelCounter();
    await fetch('/api/images/select', {
        method: 'POST', headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ hashes, action: 'select' })
    });
    const cat = selCats.find(c => c.id === selActiveCat);
    if (cat) { cat.selected = selImages.filter(i => i._sel).length; renderSelCatList(); }
}

async function deselectAllVisible() {
    const sel = selImages.filter(i => i._sel);
    if (sel.length === 0) return;
    const hashes = sel.map(i => i.hash);
    sel.forEach(i => i._sel = false);
    renderSelGrid(); updateSelCounter();
    await fetch('/api/images/select', {
        method: 'POST', headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ hashes, action: 'deselect' })
    });
    const cat = selCats.find(c => c.id === selActiveCat);
    if (cat) { cat.selected = 0; renderSelCatList(); }
}

async function saveTarget() {
    const val = parseInt(document.getElementById('sel-target-input').value) || 0;
    await fetch('/api/categories/update-target', {
        method: 'POST', headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ id: selActiveCat, target: val })
    });
    const cat = selCats.find(c => c.id === selActiveCat);
    if (cat) cat.target = val;
    renderSelCatList();
    updateSelCounter();
}

function showSelLightbox(img) {
    const overlay = document.createElement('div');
    overlay.style.cssText = 'position:fixed;top:0;left:0;width:100%;height:100%;background:rgba(0,0,0,.85);z-index:9999;display:flex;flex-direction:column;align-items:center;justify-content:center;';
    overlay.onclick = (e) => { if (e.target === overlay) overlay.remove(); };

    const isVid = img.media_type === 'video';
    var mediaEl;
    if (isVid) {
        mediaEl = document.createElement('video');
        mediaEl.src = '/api/images/serve/' + img.hash;
        mediaEl.controls = true;
        mediaEl.autoplay = true;
        mediaEl.style.cssText = 'max-width:90vw;max-height:80vh;border-radius:8px;box-shadow:0 4px 30px rgba(0,0,0,.5);';
        mediaEl.onclick = (e) => e.stopPropagation();
    } else {
        mediaEl = document.createElement('img');
        mediaEl.src = '/api/images/serve/' + img.hash;
        mediaEl.style.cssText = 'max-width:90vw;max-height:85vh;border-radius:8px;box-shadow:0 4px 30px rgba(0,0,0,.5);';
    }

    const close = document.createElement('div');
    close.innerHTML = '&times;';
    close.style.cssText = 'position:absolute;top:20px;right:30px;color:#fff;font-size:36px;cursor:pointer;z-index:10;';
    close.onclick = () => { if (isVid) mediaEl.pause(); overlay.remove(); };

    const info = document.createElement('div');
    info.style.cssText = 'position:absolute;bottom:60px;left:50%;transform:translateX(-50%);color:#fff;font-size:.85em;background:rgba(0,0,0,.6);padding:6px 16px;border-radius:6px;';
    var label = (img.filename || '');
    if (isVid && img.duration) label += ' | ' + Math.round(img.duration) + 's';
    if (img.date || img.date_taken) label += ' | ' + (img.date || img.date_taken);
    if (img.source_label) label += ' | ' + img.source_label;
    const gc = img.photo_grade?.composite;
    if (gc != null) label += ' | Grade: ' + Math.round(gc);
    info.textContent = label;

    // Like/dislike buttons in lightbox
    const prefBar = document.createElement('div');
    prefBar.style.cssText = 'position:absolute;bottom:16px;left:50%;transform:translateX(-50%);display:flex;gap:16px;';
    const likeClass = img.preference === 'like' ? 'sel-pref-active-like' : '';
    const dislikeClass = img.preference === 'dislike' ? 'sel-pref-active-dislike' : '';
    prefBar.innerHTML = `
        <button class="sel-pref-btn ${likeClass}" style="width:42px;height:42px;font-size:22px;" id="lb-sel-like">&#x1F44D;</button>
        <button class="sel-pref-btn ${dislikeClass}" style="width:42px;height:42px;font-size:22px;" id="lb-sel-dislike">&#x1F44E;</button>`;
    prefBar.querySelector('#lb-sel-like').onclick = (e) => {
        e.stopPropagation();
        const newPref = img.preference === 'like' ? null : 'like';
        setSelPref(img, newPref);
        prefBar.querySelector('#lb-sel-like').className = 'sel-pref-btn' + (newPref === 'like' ? ' sel-pref-active-like' : '');
        prefBar.querySelector('#lb-sel-dislike').className = 'sel-pref-btn';
    };
    prefBar.querySelector('#lb-sel-dislike').onclick = (e) => {
        e.stopPropagation();
        const newPref = img.preference === 'dislike' ? null : 'dislike';
        setSelPref(img, newPref);
        prefBar.querySelector('#lb-sel-dislike').className = 'sel-pref-btn' + (newPref === 'dislike' ? ' sel-pref-active-dislike' : '');
        prefBar.querySelector('#lb-sel-like').className = 'sel-pref-btn';
    };

    // Keyboard nav for like/dislike in lightbox
    const keyHandler = (e) => {
        if (e.key === 'ArrowRight' || e.key === 'l') {
            const newPref = img.preference === 'like' ? null : 'like';
            setSelPref(img, newPref);
            prefBar.querySelector('#lb-sel-like').className = 'sel-pref-btn' + (newPref === 'like' ? ' sel-pref-active-like' : '');
            prefBar.querySelector('#lb-sel-dislike').className = 'sel-pref-btn';
        } else if (e.key === 'ArrowLeft' || e.key === 'd') {
            const newPref = img.preference === 'dislike' ? null : 'dislike';
            setSelPref(img, newPref);
            prefBar.querySelector('#lb-sel-dislike').className = 'sel-pref-btn' + (newPref === 'dislike' ? ' sel-pref-active-dislike' : '');
            prefBar.querySelector('#lb-sel-like').className = 'sel-pref-btn';
        } else if (e.key === 'Escape') {
            if (isVid) mediaEl.pause();
            document.removeEventListener('keydown', keyHandler);
            overlay.remove();
        }
    };
    document.addEventListener('keydown', keyHandler);
    overlay.addEventListener('remove', () => document.removeEventListener('keydown', keyHandler));

    overlay.appendChild(mediaEl);
    overlay.appendChild(close);
    overlay.appendChild(info);
    overlay.appendChild(prefBar);
    document.body.appendChild(overlay);
}

async function runQuickFill() {
    await fetch('/api/quick-fill', { method:'POST' });
    showTaskOverlay('auto-select');
    const check = setInterval(async () => {
        try {
            const res = await fetch('/api/scan/status');
            const st = await res.json();
            if (st.done || st.error || st.cancelled) {
                clearInterval(check);
                if (!st.error && !st.cancelled) {
                    document.querySelectorAll('.step-dot')[6].classList.add('done');
                }
                await loadSelCategories();
                if (selActiveCat) await loadSelImages();
            }
        } catch(e) {}
    }, 1500);
}

async function runAutoSelect() {
    await fetch('/api/auto-select', {
        method:'POST', headers:{'Content-Type':'application/json'},
        body:JSON.stringify({ strategy: 'balanced', sim_threshold: 0.85 })
    });
    showTaskOverlay('auto-select');
    const check = setInterval(async () => {
        try {
            const res = await fetch('/api/scan/status');
            const st = await res.json();
            if (st.done || st.error || st.cancelled) {
                clearInterval(check);
                if (!st.error && !st.cancelled) {
                    document.querySelectorAll('.step-dot')[6].classList.add('done');
                }
                await loadSelCategories();
                if (selActiveCat) await loadSelImages();
            }
        } catch(e) {}
    }, 2000);
}

// ── Step 7: Review ──
async function loadReviewStats() {
    const el = document.getElementById('review-stats');
    const readyEl = document.getElementById('review-readiness');
    const breakdownEl = document.getElementById('review-cat-breakdown');
    if (!el) return;
    try {
        const res = await fetch('/api/categories/summary');
        const cats = await res.json();
        if (!cats.length) { el.innerHTML = ''; readyEl.innerHTML = ''; breakdownEl.style.display = 'none'; return; }
        let totalSel = 0, totalTarget = 0, filledCats = 0, emptyCats = 0;
        var underCats = [];
        cats.forEach(c => {
            const sel = c.selected || 0;
            const tgt = c.target || 75;
            totalSel += sel;
            totalTarget += tgt;
            if (sel >= tgt) filledCats++;
            if (sel === 0) emptyCats++;
            if (sel > 0 && sel < tgt) underCats.push(c);
        });
        var overallPct = Math.min(100, Math.round(totalSel / (totalTarget || 1) * 100));
        var ready = filledCats === cats.length;

        // Stat cards
        el.innerHTML =
            '<div class="stat-row">' +
                '<div class="stat-card green"><div class="stat-value">' + totalSel + '</div><div class="stat-label">Images selected</div></div>' +
                '<div class="stat-card"><div class="stat-value">' + cats.length + '</div><div class="stat-label">Categories</div></div>' +
                '<div class="stat-card' + (ready ? ' green' : '') + '"><div class="stat-value">' + filledCats + ' / ' + cats.length + '</div><div class="stat-label">Categories filled</div></div>' +
                '<div class="stat-card' + (ready ? ' green' : '') + '"><div class="stat-value">' + overallPct + '%</div><div class="stat-label">Overall fill</div></div>' +
            '</div>';

        // Readiness notice
        if (ready) {
            readyEl.innerHTML = '<div class="notice notice-tip"><strong>Ready to export!</strong> All categories are filled to target. Open the gallery if you want to fine-tune, or proceed to export.</div>';
        } else if (totalSel === 0) {
            readyEl.innerHTML = '<div class="notice notice-error"><strong>No images selected yet.</strong> Go back to the Select step to pick images or run Auto-fill.</div>';
        } else {
            var msg = '<strong>' + (cats.length - filledCats) + ' categories still need more images.</strong>';
            if (emptyCats > 0) msg += ' ' + emptyCats + ' are completely empty.';
            if (underCats.length > 0 && underCats.length <= 4) {
                msg += ' Under-filled: ' + underCats.map(function(c) { return c.display + ' (' + (c.selected||0) + '/' + (c.target||75) + ')'; }).join(', ') + '.';
            }
            readyEl.innerHTML = '<div class="notice notice-warn">' + msg + '</div>';
        }

        // Per-category breakdown bars
        breakdownEl.style.display = 'block';
        var rows = cats.map(function(c) {
            var sel = c.selected || 0;
            var tgt = c.target || 75;
            var pct = Math.min(100, Math.round(sel / tgt * 100));
            var full = sel >= tgt;
            var barColor = full ? '#22c55e' : (sel === 0 ? '#e2e8f0' : '#3b82f6');
            var statusText = full ? 'FULL' : (sel === 0 ? 'EMPTY' : sel + '/' + tgt);
            var statusColor = full ? '#22c55e' : (sel === 0 ? '#94a3b8' : '#f59e0b');
            return '<div style="display:flex;align-items:center;gap:10px;padding:6px 0;border-bottom:1px solid #f1f5f9;">' +
                '<div style="min-width:130px;font-size:.85em;font-weight:500;color:#334155;">' + esc(c.display) + '</div>' +
                '<div style="flex:1;height:6px;background:#e2e8f0;border-radius:3px;overflow:hidden;">' +
                    '<div style="height:100%;width:' + pct + '%;background:' + barColor + ';border-radius:3px;transition:width .4s;"></div>' +
                '</div>' +
                '<div style="min-width:55px;font-size:.78em;font-weight:600;color:' + statusColor + ';text-align:right;">' + statusText + '</div>' +
            '</div>';
        }).join('');
        breakdownEl.innerHTML = '<div class="section-title">Category breakdown</div>' + rows;

    } catch(e) { el.innerHTML = ''; }
}
function openGallery() {
    window.open('/api/report', '_blank');
}

// ── Step 8: Export ──
async function loadExportStats() {
    const el = document.getElementById('export-stats');
    showLoader(el, 'Loading stats...');
    const res = await fetch('/api/stats');
    const st = await res.json();
    var dirInput = document.getElementById('inp-export-dir');
    if (!dirInput.value && st.default_export_dir) dirInput.value = st.default_export_dir;
    if (st.has_scan) {
        const selected = st.selected || 0;
        const qualified = st.qualified || 0;
        el.innerHTML = `
            <div class="stat-row">
                <div class="stat-card${selected > 0 ? ' green' : ' red'}">
                    <div class="stat-value">${selected}</div>
                    <div class="stat-label">Selected for export</div>
                </div>
                <div class="stat-card">
                    <div class="stat-value">${qualified}</div>
                    <div class="stat-label">Qualified</div>
                </div>
                <div class="stat-card">
                    <div class="stat-value">${st.total_images}</div>
                    <div class="stat-label">Total scanned</div>
                </div>
            </div>
            ${selected === 0 ? '<div class="notice notice-error"><strong>No images selected.</strong> Go back to the Select step to pick images.</div>' :
              '<div class="notice notice-tip"><strong>Ready to export!</strong> Choose PowerPoint for a presentation or export to a folder for organizing by category.</div>'}
        `;
    }
}

async function runExport() {
    const outputDir = document.getElementById('inp-export-dir').value || 'final_collection';

    await fetch('/api/export', {
        method:'POST', headers:{'Content-Type':'application/json'},
        body:JSON.stringify({ output_dir: outputDir })
    });
    showTaskOverlay('export');
    const check = setInterval(async () => {
        try {
            const res = await fetch('/api/scan/status');
            const st = await res.json();
            if (st.done || st.error || st.cancelled) {
                clearInterval(check);
                if (!st.error && !st.cancelled) {
                    document.querySelectorAll('.step-dot')[8].classList.add('done');
                }
            }
        } catch(e) {}
    }, 2000);
}

async function runPptxExport() {
    document.getElementById('btn-export-pptx').disabled = true;
    document.getElementById('btn-export-pptx').textContent = 'Creating...';
    document.getElementById('pptx-download').style.display = 'none';

    var opts = {
        captions: document.getElementById('pptx-captions').checked,
        dividers: document.getElementById('pptx-dividers').checked,
        sort: document.getElementById('pptx-sort').value
    };
    await fetch('/api/export/pptx', {
        method:'POST', headers:{'Content-Type':'application/json'},
        body:JSON.stringify(opts)
    });
    showTaskOverlay('pptx_export');
    var check = setInterval(async function() {
        try {
            var res = await fetch('/api/scan/status');
            var st = await res.json();
            if (st.done || st.error || st.cancelled) {
                clearInterval(check);
                document.getElementById('btn-export-pptx').disabled = false;
                document.getElementById('btn-export-pptx').textContent = 'Create Presentation';
                if (!st.error && !st.cancelled) {
                    document.querySelectorAll('.step-dot')[8].classList.add('done');
                    document.getElementById('pptx-download').style.display = 'block';
                }
            }
        } catch(e) {}
    }, 2000);
}

function esc(s) { const d=document.createElement('div'); d.textContent=s; return d.innerHTML; }

// ── Lightbox ──
function openLightbox(person, filename) {
    document.getElementById('lightbox-img').src = '/api/ref-faces/' + encodeURIComponent(person) + '/photo/' + encodeURIComponent(filename);
    document.getElementById('lightbox').classList.add('open');
}
function closeLightbox() {
    document.getElementById('lightbox').classList.remove('open');
    document.getElementById('lightbox-img').src = '';
}
document.addEventListener('keydown', e => { if (e.key === 'Escape') closeLightbox(); });

// ── Projects drawer ──
function toggleDrawer() {
    var drawer = document.getElementById('side-drawer');
    var backdrop = document.getElementById('drawer-backdrop');
    var toggle = document.getElementById('rail-toggle');
    drawer.classList.toggle('open');
    backdrop.classList.toggle('open');
    if (toggle) toggle.classList.toggle('active', drawer.classList.contains('open'));
    if (drawer.classList.contains('open')) {
        loadProjectList();
        loadUserInfo();
    }
}

function openDrawerToProjects() {
    var drawer = document.getElementById('side-drawer');
    if (!drawer.classList.contains('open')) toggleDrawer();
    // Expand projects section
    document.getElementById('projects-section').style.display = 'block';
    document.getElementById('projects-arrow').style.transform = 'rotate(180deg)';
}

function toggleProjectsSection() {
    var sec = document.getElementById('projects-section');
    var arrow = document.getElementById('projects-arrow');
    if (sec.style.display === 'none') {
        sec.style.display = 'block';
        arrow.style.transform = 'rotate(180deg)';
        loadProjectList();
    } else {
        sec.style.display = 'none';
        arrow.style.transform = '';
    }
}

function toggleCleanupSection() {
    var sec = document.getElementById('cleanup-section');
    var arrow = document.getElementById('cleanup-arrow');
    if (sec.style.display === 'none') {
        sec.style.display = 'block';
        arrow.style.transform = 'rotate(180deg)';
    } else {
        sec.style.display = 'none';
        arrow.style.transform = '';
    }
}

function openPhoneImages() {
    alert('Phone Images feature coming soon!\\nConnect your phone via USB to manage and clean up photos directly.');
}

async function loadUserInfo() {
    try {
        const res = await fetch('/api/auth/me');
        const data = await res.json();
        if (data.authenticated) {
            document.getElementById('user-info').textContent = data.user;
        }
    } catch(e) {}
}

async function logout() {
    if (!confirm('Sign out?')) return;
    await fetch('/api/auth/logout', {method: 'POST'});
    window.location.href = '/login';
}

async function loadProjectList() {
    const el = document.getElementById('project-list');
    try {
        const res = await fetch('/api/projects');
        if (!res.ok) throw new Error('HTTP ' + res.status);
        const projects = await res.json();
        if (!Array.isArray(projects) || !projects.length) {
            el.innerHTML = '<div style="padding:18px; color:#a0aec0; font-size:.85em; text-align:center;">No saved projects yet.</div>';
            return;
        }
        el.innerHTML = '';
        projects.forEach(function(p) {
            var modified = p.modified ? new Date(p.modified).toLocaleDateString() + ' ' + new Date(p.modified).toLocaleTimeString([], {hour:'2-digit', minute:'2-digit'}) : '';
            var info = [p.template || p.event_type || '', 'Step ' + ((p.step || 0) + 1)].filter(Boolean).join(' | ');
            var row = document.createElement('div');
            row.className = 'project-item';
            row.onclick = function() { loadProject(p.dir_name); };
            var nameDiv = document.createElement('div');
            nameDiv.className = 'p-name';
            nameDiv.textContent = p.name;
            var metaDiv = document.createElement('div');
            metaDiv.className = 'p-meta';
            metaDiv.textContent = info + (modified ? ' | ' + modified : '');
            var actDiv = document.createElement('div');
            actDiv.className = 'p-actions';
            var renBtn = document.createElement('button');
            renBtn.className = 'p-del';
            renBtn.style.color = '#2b6cb0';
            renBtn.textContent = 'Rename';
            renBtn.onclick = function(e) { e.stopPropagation(); renameProject(p.dir_name, p.name); };
            actDiv.appendChild(renBtn);
            var delBtn = document.createElement('button');
            delBtn.className = 'p-del';
            delBtn.textContent = 'Delete';
            delBtn.onclick = function(e) { e.stopPropagation(); deleteProject(p.dir_name, p.name); };
            actDiv.appendChild(delBtn);
            row.appendChild(nameDiv);
            row.appendChild(metaDiv);
            row.appendChild(actDiv);
            el.appendChild(row);
        });
    } catch(e) {
        console.error('loadProjectList error:', e);
        el.innerHTML = '<div style="padding:18px; color:#e53e3e; font-size:.85em;">Error loading projects: ' + e.message + '</div>';
    }
}

async function autoSaveDraft() {
    // Auto-save current work as a draft with date-based name
    try {
        var cfg = await (await fetch('/api/config')).json();
        var hasWork = cfg && (cfg.event_type || cfg.sources && cfg.sources.length > 0);
        if (!hasWork) return; // Nothing to save

        var now = new Date();
        var draftName = 'Draft - ' + now.toLocaleDateString() + ' ' + now.toLocaleTimeString([], {hour:'2-digit', minute:'2-digit'});
        var projectName = cfg.project_name || draftName;
        if (!cfg.project_name) projectName = draftName;

        await fetch('/api/projects/save', {
            method: 'POST', headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({ name: projectName, step: currentStep, overwrite: true })
        });
    } catch(e) { /* silent */ }
}

async function saveCurrentProject() {
    var projectName = '';
    if (config && config.project_name) {
        projectName = config.project_name;
    }
    var name = prompt('Project name:', projectName);
    if (!name || !name.trim()) return;
    name = name.trim();
    // Allow overwrite if saving with the same name as current project
    var overwrite = (config && config.project_name && config.project_name === name);
    showLoader('project-list', 'Saving...');
    try {
        var res = await fetch('/api/projects/save', {
            method: 'POST', headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({ name: name, step: currentStep, overwrite: overwrite })
        });
        var data = await res.json();
        if (data.error) { alert(data.error); return; }
        if (config) config.project_name = name;
        updateProjectNameBar(name);
        await loadProjectList();
    } catch(e) { alert('Save failed: ' + e.message); }
}

async function loadProject(dirName) {
    if (!confirm('Load this project? Current unsaved progress will be lost.')) return;
    showLoader('project-list', 'Loading project...');
    try {
        var res = await fetch('/api/projects/load', {
            method: 'POST', headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({ dir_name: dirName })
        });
        var data = await res.json();
        if (data.error) { alert(data.error); return; }
        toggleDrawer();
        // Reload the app at the saved step
        config = null;
        selectedTemplate = null;
        var savedStep = data.step || 0;
        await loadTemplates();
        var cfg = await (await fetch('/api/config')).json();
        if (cfg && cfg.event_type) {
            config = cfg;
            selectedTemplate = cfg.event_type;
        }
        // Reset all step dots, then mark steps before saved step as done
        var dots = document.querySelectorAll('.step-dot');
        dots.forEach(function(d) { d.classList.remove('done'); });
        for (var si = 0; si < savedStep && si < dots.length; si++) {
            dots[si].classList.add('done');
        }
        updateProjectNameBar(config ? config.project_name || config.event_name : '');
        goStep(savedStep);
    } catch(e) { alert('Load failed: ' + e.message); }
}

async function newProject() {
    if (!confirm('Start a new project? Your current work will be saved as a draft.')) return;
    try {
        await autoSaveDraft();
        await fetch('/api/projects/new', { method: 'POST' });
        if (document.getElementById('side-drawer').classList.contains('open')) toggleDrawer();
        config = null;
        selectedTemplate = null;
        faceVerifyCache = {};
        replacedPhotos = {};
        document.getElementById('inp-project-name').value = '';
        document.querySelectorAll('.step-dot').forEach(function(d) { d.classList.remove('done'); });
        updateProjectNameBar('');
        await loadTemplates();
        goStep(0);
    } catch(e) { alert('Error: ' + e.message); }
}

async function deleteProject(dirName, displayName) {
    if (!confirm('Delete project "' + displayName + '"? This cannot be undone.')) return;
    try {
        await fetch('/api/projects/' + encodeURIComponent(dirName), { method: 'DELETE' });
        await loadProjectList();
    } catch(e) { alert('Delete failed: ' + e.message); }
}

async function renameProject(dirName, currentName) {
    var newName = prompt('Rename project:', currentName);
    if (!newName || !newName.trim() || newName.trim() === currentName) return;
    try {
        var res = await fetch('/api/projects/' + encodeURIComponent(dirName) + '/rename', {
            method: 'POST', headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({ name: newName.trim() })
        });
        var data = await res.json();
        if (data.error) { alert(data.error); return; }
        await loadProjectList();
    } catch(e) { alert('Rename failed: ' + e.message); }
}

function validateProjectName() {
    var inp = document.getElementById('inp-project-name');
    var err = document.getElementById('project-name-error');
    if (inp.value.trim()) {
        err.style.display = 'none';
        inp.style.borderColor = '#e2e8f0';
    } else {
        err.style.display = 'block';
        inp.style.borderColor = '#e53e3e';
    }
}

// ── Tutorial system ──

const TOUR_STEPS = [
    {
        target: '.header h1',
        title: 'Welcome!',
        text: 'This is E-z Photo Organizer. It helps you build a curated photo collection for any special event — step by step.',
        position: 'bottom'
    },
    {
        target: '#steps-nav',
        title: 'Step Navigation',
        text: 'These are your 9 steps. Each one guides you through the process — from choosing an event type to exporting your final collection.',
        position: 'bottom'
    },
    {
        target: '.step-dot:first-child',
        title: 'Step 1: Choose Event',
        text: 'Start here. Pick your event type (Bar Mitzva, Wedding, Graduation, etc.) and the app will set up categories and settings for you.',
        position: 'bottom'
    },
    {
        target: '#icon-rail',
        title: 'Side Menu',
        text: 'Quick access to projects, cleanup tools, phone images, tutorial, and more. Click any icon or open the full menu.',
        position: 'right'
    },
    {
        target: '.step-dot:nth-child(3)',
        title: 'Add Your Photo Sources',
        text: 'In the Sources step, add folders from USB drives, cloud backups, phone exports — anywhere your photos live.',
        position: 'bottom'
    },
    {
        target: '.step-dot:nth-child(4)',
        title: 'Face Recognition',
        text: 'Add reference photos of the main person. The app will automatically find them across thousands of photos.',
        position: 'bottom'
    },
    {
        target: '.step-dot:nth-child(5)',
        title: 'Smart Scanning',
        text: 'The scanner extracts dates, detects faces, filters duplicates, and categorizes everything automatically.',
        position: 'bottom'
    },
    {
        target: '.step-dot:nth-child(7)',
        title: 'Auto-Select & Review',
        text: 'Use Quick Fill or Smart Fill to automatically pick the best photos, then manually review and swap as needed.',
        position: 'bottom'
    },
    {
        target: '.step-dot:nth-child(9)',
        title: 'Export',
        text: 'When you are happy with your selection, export everything to a folder — organized and ready for your event!',
        position: 'bottom'
    }
];

let tourIdx = -1;

function startTutorial() {
    tourIdx = 0;
    // Create overlay elements if they don't exist
    if (!document.getElementById('tour-backdrop')) {
        var bd = document.createElement('div');
        bd.id = 'tour-backdrop';
        bd.onclick = function(e) { if (e.target === bd) endTutorial(); };
        document.body.appendChild(bd);

        var hl = document.createElement('div');
        hl.id = 'tour-highlight';
        document.body.appendChild(hl);

        var tt = document.createElement('div');
        tt.id = 'tour-tooltip';
        document.body.appendChild(tt);
    }
    document.getElementById('tour-backdrop').classList.add('active');
    document.getElementById('tour-highlight').style.display = 'block';
    document.getElementById('tour-tooltip').style.display = 'block';
    showTourStep();
}

function showTourStep() {
    var step = TOUR_STEPS[tourIdx];
    var el = document.querySelector(step.target);
    var hl = document.getElementById('tour-highlight');
    var tt = document.getElementById('tour-tooltip');

    tt.classList.remove('show');

    if (!el) { tourNext(); return; }

    // Scroll element into view
    el.scrollIntoView({ behavior: 'smooth', block: 'center' });

    setTimeout(function() {
        var rect = el.getBoundingClientRect();
        var pad = 6;

        // Position highlight
        hl.style.left = (rect.left - pad + window.scrollX) + 'px';
        hl.style.top = (rect.top - pad + window.scrollY) + 'px';
        hl.style.width = (rect.width + pad * 2) + 'px';
        hl.style.height = (rect.height + pad * 2) + 'px';
        hl.classList.add('pulse');

        // Build tooltip content
        var dotsHtml = '';
        for (var i = 0; i < TOUR_STEPS.length; i++) {
            dotsHtml += '<span class="tour-dot' + (i === tourIdx ? ' active' : '') + '"></span>';
        }

        tt.innerHTML = '<h3>' + step.title + '</h3><p>' + step.text + '</p>' +
            '<div class="tour-actions">' +
            '<div class="tour-dots">' + dotsHtml + '</div>' +
            '<div style="display:flex; gap:8px; align-items:center;">' +
            '<button class="tour-skip" onclick="endTutorial()">Skip</button>' +
            '<button class="tour-next" onclick="tourNext()">' + (tourIdx < TOUR_STEPS.length - 1 ? 'Next' : 'Got it!') + '</button>' +
            '</div></div>';

        // Position tooltip
        var pos = step.position || 'bottom';
        var ttW = 340;
        var ttLeft, ttTop;

        if (pos === 'bottom') {
            ttLeft = rect.left + rect.width / 2 - ttW / 2 + window.scrollX;
            ttTop = rect.bottom + 16 + window.scrollY;
        } else if (pos === 'right') {
            ttLeft = rect.right + 16 + window.scrollX;
            ttTop = rect.top + window.scrollY;
        } else if (pos === 'top') {
            ttLeft = rect.left + rect.width / 2 - ttW / 2 + window.scrollX;
            ttTop = rect.top - 140 + window.scrollY;
        }

        // Keep in viewport
        ttLeft = Math.max(12, Math.min(ttLeft, window.innerWidth - ttW - 12));

        tt.style.left = ttLeft + 'px';
        tt.style.top = ttTop + 'px';
        tt.style.width = ttW + 'px';

        // Animate in
        setTimeout(function() { tt.classList.add('show'); }, 50);
    }, 300);
}

function tourNext() {
    tourIdx++;
    if (tourIdx >= TOUR_STEPS.length) {
        endTutorial();
        return;
    }
    showTourStep();
}

function endTutorial() {
    tourIdx = -1;
    var bd = document.getElementById('tour-backdrop');
    var hl = document.getElementById('tour-highlight');
    var tt = document.getElementById('tour-tooltip');
    if (bd) bd.classList.remove('active');
    if (hl) { hl.classList.remove('pulse'); hl.style.display = 'none'; }
    if (tt) { tt.classList.remove('show'); tt.style.display = 'none'; }
    localStorage.setItem('tutorial_seen', '1');
}

// ── Personalized greeting ──
function updateProjectNameBar(name) {
    var bar = document.getElementById('project-name-bar');
    var txt = document.getElementById('project-name-text');
    if (name && name.trim()) {
        txt.textContent = name.trim();
        bar.style.display = '';
        document.title = name.trim() + ' — E-z Photo Organizer';
    } else {
        bar.style.display = 'none';
        document.title = 'E-z Photo Organizer';
    }
}

async function loadGreeting() {
    try {
        var res = await fetch('/api/auth/me');
        var data = await res.json();
        if (data.authenticated && data.name) {
            var firstName = data.name.split(' ')[0];
            document.getElementById('header-greeting').textContent = 'Hi ' + firstName + ', how can I help you today?';

            // First visit — auto-start tutorial
            if (!localStorage.getItem('tutorial_seen')) {
                setTimeout(startTutorial, 800);
            }
        }
    } catch(e) {}
}

// ── Cleanup Tool ──
var cleanupOffset = 0;
var cleanupImages = [];

function toggleRailCleanup() {
    var sub = document.getElementById('rail-cleanup-sub');
    if (sub.style.display === 'none') {
        sub.style.display = 'block';
        document.addEventListener('click', _closeRailCleanupOutside, true);
    } else {
        closeRailCleanup();
    }
}

function closeRailCleanup() {
    document.getElementById('rail-cleanup-sub').style.display = 'none';
    document.removeEventListener('click', _closeRailCleanupOutside, true);
}

function _closeRailCleanupOutside(e) {
    var group = document.getElementById('rail-cleanup-group');
    if (group && !group.contains(e.target)) {
        closeRailCleanup();
    }
}

async function openCleanup() {
    await autoSaveDraft();
    document.getElementById('cleanup-overlay').classList.add('active');
    cleanupOffset = 0;
    cleanupImages = [];
    populateCleanupCategories();
    loadCleanupImages();
    updateTrashBadge();
}

function closeCleanup() {
    document.getElementById('cleanup-overlay').classList.remove('active');
}

async function populateCleanupCategories() {
    try {
        var res = await fetch('/api/categories/summary');
        var data = await res.json();
        var sel = document.getElementById('cleanup-filter-cat');
        sel.innerHTML = '<option value="">All Categories</option>';
        (data.categories || []).forEach(function(c) {
            var opt = document.createElement('option');
            opt.value = c.id;
            opt.textContent = c.display || c.id;
            sel.appendChild(opt);
        });
    } catch(e) {}
}

async function loadCleanupImages() {
    cleanupOffset = 0;
    cleanupImages = [];
    var grid = document.getElementById('cleanup-grid');
    grid.innerHTML = '<div style="grid-column:1/-1; text-align:center; color:#a0aec0; padding:40px;"><span style="display:inline-block;width:18px;height:18px;border:3px solid #bee3f8;border-top:3px solid #3182ce;border-radius:50%;animation:spinA .7s linear infinite;vertical-align:middle;margin-right:8px;"></span>Loading...</div>';
    await fetchCleanupPage();
}

async function fetchCleanupPage() {
    var cat = document.getElementById('cleanup-filter-cat').value;
    var status = document.getElementById('cleanup-filter-status').value;
    var media = document.getElementById('cleanup-filter-media').value;
    var trashOnly = document.getElementById('cleanup-show-trash').checked;

    var params = new URLSearchParams({offset: cleanupOffset, limit: 200});
    if (cat) params.set('category', cat);
    if (status) params.set('status', status);
    if (media) params.set('media_type', media);
    if (trashOnly) params.set('trash', '1');

    try {
        var res = await fetch('/api/cleanup/images?' + params.toString());
        var data = await res.json();
        var grid = document.getElementById('cleanup-grid');

        if (cleanupOffset === 0) grid.innerHTML = '';

        if (data.images.length === 0 && cleanupOffset === 0) {
            grid.innerHTML = '<div style="grid-column:1/-1; text-align:center; color:#a0aec0; padding:40px;">No images found. Run a scan first or adjust filters.</div>';
        }

        data.images.forEach(function(img) {
            cleanupImages.push(img);
            var card = document.createElement('div');
            card.className = 'cleanup-card' + (img.trash ? ' trashed' : '');
            card.dataset.hash = img.hash;

            var isVid = img.media_type === 'video';
            var thumbSrc = img.thumb ? ('data:image/jpeg;base64,' + img.thumb) : '';
            var vidBadge = isVid ? '<div style="position:absolute;top:4px;right:4px;background:rgba(0,0,0,.7);color:#fff;font-size:.65em;padding:2px 6px;border-radius:3px;">VID</div>' : '';

            var imgEl = document.createElement('div');
            imgEl.style.position = 'relative';
            var pic = document.createElement('img');
            pic.src = thumbSrc;
            pic.alt = '';
            pic.style.cssText = 'width:100%;height:130px;object-fit:cover;display:block;';
            pic.onerror = function() { this.style.background = '#edf2f7'; };
            imgEl.appendChild(pic);
            if (isVid) {
                var vb = document.createElement('div');
                vb.style.cssText = 'position:absolute;top:4px;right:4px;background:rgba(0,0,0,.7);color:#fff;font-size:.65em;padding:2px 6px;border-radius:3px;';
                vb.textContent = 'VID';
                imgEl.appendChild(vb);
            }
            var infoEl = document.createElement('div');
            infoEl.className = 'card-info';
            infoEl.textContent = img.filename;
            card.appendChild(imgEl);
            card.appendChild(infoEl);

            card.onclick = function() { toggleCleanupTrash(img.hash, card); };
            grid.appendChild(card);
        });

        var loadMore = document.getElementById('cleanup-load-more');
        if (data.total > cleanupOffset + 200) {
            loadMore.style.display = 'block';
            cleanupOffset += 200;
        } else {
            loadMore.style.display = 'none';
        }
    } catch(e) {
        document.getElementById('cleanup-grid').innerHTML = '<div style="grid-column:1/-1; text-align:center; color:#e53e3e; padding:40px;">Failed to load images.</div>';
    }
}

function loadCleanupMore() {
    fetchCleanupPage();
}

async function toggleCleanupTrash(hash, card) {
    var isTrash = card.classList.contains('trashed');
    var endpoint = isTrash ? '/api/cleanup/unmark-trash' : '/api/cleanup/mark-trash';
    try {
        await fetch(endpoint, {
            method: 'POST',
            headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({hashes: [hash]})
        });
        card.classList.toggle('trashed');
        updateTrashBadge();
    } catch(e) {}
}

async function updateTrashBadge() {
    try {
        var res = await fetch('/api/cleanup/trash-count');
        var data = await res.json();
        var badge = document.getElementById('cleanup-trash-badge');
        var reviewBtn = document.getElementById('btn-review-trash');
        var bar = document.getElementById('cleanup-bottom-bar');
        var barCount = document.getElementById('cleanup-bar-count');
        var barSize = document.getElementById('cleanup-bar-size');

        if (data.count > 0) {
            badge.textContent = data.count + ' in trash';
            badge.style.display = 'inline-block';
            reviewBtn.style.display = 'inline-block';
            bar.style.display = 'flex';
            barCount.textContent = data.count + ' item' + (data.count !== 1 ? 's' : '') + ' in trash';
            barSize.textContent = data.size_mb > 0 ? '(' + data.size_mb + ' MB)' : '';
        } else {
            badge.style.display = 'none';
            reviewBtn.style.display = 'none';
            bar.style.display = 'none';
        }
    } catch(e) {}
}

function showCleanupTrash() {
    document.getElementById('cleanup-show-trash').checked = true;
    loadCleanupImages();
}

function cleanupSelectAll() {
    var cards = document.querySelectorAll('#cleanup-grid .cleanup-card:not(.trashed)');
    var hashes = [];
    cards.forEach(function(c) { hashes.push(c.dataset.hash); });
    if (hashes.length === 0) return;
    if (!confirm('Mark ' + hashes.length + ' visible images for trash?')) return;
    fetch('/api/cleanup/mark-trash', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({hashes: hashes})
    }).then(function() {
        cards.forEach(function(c) { c.classList.add('trashed'); });
        updateTrashBadge();
    });
}

function cleanupDeselectAll() {
    var cards = document.querySelectorAll('#cleanup-grid .cleanup-card.trashed');
    var hashes = [];
    cards.forEach(function(c) { hashes.push(c.dataset.hash); });
    if (hashes.length === 0) return;
    fetch('/api/cleanup/unmark-trash', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({hashes: hashes})
    }).then(function() {
        cards.forEach(function(c) { c.classList.remove('trashed'); });
        updateTrashBadge();
    });
}

function cleanupClearTrash() {
    fetch('/api/cleanup/images?trash=1&limit=99999').then(function(r) { return r.json(); }).then(function(data) {
        var hashes = data.images.map(function(i) { return i.hash; });
        if (hashes.length === 0) return;
        return fetch('/api/cleanup/unmark-trash', {
            method: 'POST',
            headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({hashes: hashes})
        });
    }).then(function() {
        loadCleanupImages();
        updateTrashBadge();
    });
}

async function cleanupConfirmRecycle() {
    var res = await fetch('/api/cleanup/trash-count');
    var data = await res.json();
    if (data.count === 0) { alert('No items in trash.'); return; }
    if (!confirm('Move ' + data.count + ' file' + (data.count !== 1 ? 's' : '') + ' (' + data.size_mb + ' MB) to the Recycle Bin?\\n\\nYou can recover them from the Recycle Bin if needed.')) return;

    var btn = document.querySelector('#cleanup-bottom-bar .btn[style*="e53e3e"]');
    btn.disabled = true;
    btn.textContent = 'Moving...';

    try {
        var r = await fetch('/api/cleanup/confirm-trash', { method: 'POST' });
        var result = await r.json();
        if (result.ok) {
            var msg = result.recycled + ' file' + (result.recycled !== 1 ? 's' : '') + ' moved to Recycle Bin.';
            if (result.failed > 0) msg += '\\n' + result.failed + ' failed.';
            alert(msg);
            loadCleanupImages();
            updateTrashBadge();
        } else {
            alert('Error: ' + (result.error || 'Unknown'));
        }
    } catch(e) {
        alert('Failed to recycle files.');
    }
    btn.disabled = false;
    btn.textContent = 'Move to Recycle Bin';
}

// ── Age Assessment ──
var ageResults = [];

var ageRefPhotos = [];

function openAgeAssessment() {
    document.getElementById('age-overlay').classList.add('active');
    buildAgeFolderList();
    loadAgeResults();
    toggleAgeFaceMode();
    loadExistingAgeRefs();
}

function closeAgeAssessment() {
    document.getElementById('age-overlay').classList.remove('active');
}

function toggleAgeFaceMode() {
    var mode = document.querySelector('input[name="age-face-mode"]:checked');
    document.getElementById('age-face-ref').style.display = (mode && mode.value === 'specific') ? 'block' : 'none';
}

var ageRefVerified = false;

async function loadExistingAgeRefs() {
    try {
        var res = await fetch('/api/ref-faces');
        var faces = await res.json();
        var container = document.getElementById('age-use-existing');
        if (!Array.isArray(faces) || !faces.length) { container.innerHTML = ''; return; }
        var personsWithPhotos = faces.filter(function(f) { return f.photo_count > 0; });
        if (!personsWithPhotos.length) { container.innerHTML = ''; return; }
        container.innerHTML = '<div style="font-size:.8em; color:#718096; margin-bottom:4px;">Or use existing project references:</div>';
        personsWithPhotos.forEach(function(f) {
            var btn = document.createElement('button');
            btn.className = 'btn btn-secondary';
            btn.style.cssText = 'font-size:.78em; padding:4px 10px; margin:0 4px 4px 0;';
            btn.textContent = f.name + ' (' + f.photo_count + ' photos)';
            btn.onclick = function() {
                document.getElementById('age-ref-person-name').value = f.name;
                ageRefVerified = false;
                loadAgeRefThumbs(f.name);
            };
            container.appendChild(btn);
        });
    } catch(e) {}
}

async function uploadAgeRefPhotos(files) {
    var personName = document.getElementById('age-ref-person-name').value.trim();
    if (!personName) { alert('Please enter the person name first.'); return; }

    var status = document.getElementById('age-ref-status');
    status.textContent = 'Uploading ' + files.length + ' photos...';

    var fd = new FormData();
    fd.append('person', personName);
    for (var i = 0; i < files.length; i++) {
        fd.append('photos', files[i]);
    }
    try {
        await fetch('/api/ref-faces/upload', { method: 'POST', body: fd });
    } catch(e) {}
    document.getElementById('age-ref-upload').value = '';
    ageRefVerified = false;
    await loadAgeRefThumbs(personName);
}

async function loadAgeRefThumbs(personName) {
    var thumbs = document.getElementById('age-ref-thumbs');
    var status = document.getElementById('age-ref-status');
    thumbs.innerHTML = '';
    try {
        var res = await fetch('/api/ref-faces/' + encodeURIComponent(personName) + '/photos');
        var photos = await res.json();
        if (!photos.length) {
            status.textContent = 'No photos uploaded yet.';
            document.getElementById('btn-age-verify').style.display = 'none';
            return;
        }
        status.innerHTML = photos.length + ' photo(s) for <strong>' + esc(personName) + '</strong>';
        document.getElementById('btn-age-verify').style.display = '';

        photos.forEach(function(p) {
            var wrap = document.createElement('div');
            wrap.style.cssText = 'position:relative; display:inline-flex; flex-direction:column; align-items:center; gap:2px;';
            wrap.setAttribute('data-filename', p.filename);

            var imgDiv = document.createElement('div');
            imgDiv.style.cssText = 'position:relative;';

            if (p.thumb) {
                var img = document.createElement('img');
                img.src = 'data:image/jpeg;base64,' + p.thumb;
                img.style.cssText = 'width:60px; height:60px; object-fit:cover; border-radius:4px; border:2px solid #cbd5e0;';
                imgDiv.appendChild(img);
            } else {
                var placeholder = document.createElement('div');
                placeholder.style.cssText = 'width:60px; height:60px; background:#e2e8f0; border-radius:4px; display:flex; align-items:center; justify-content:center; font-size:.65em; color:#718096;';
                placeholder.textContent = p.filename;
                imgDiv.appendChild(placeholder);
            }

            // Remove button (X)
            var xBtn = document.createElement('span');
            xBtn.textContent = '\u2715';
            xBtn.style.cssText = 'position:absolute; top:-4px; right:-4px; width:16px; height:16px; background:#e53e3e; color:#fff; border-radius:50%; font-size:10px; display:flex; align-items:center; justify-content:center; cursor:pointer; box-shadow:0 1px 3px rgba(0,0,0,.3);';
            xBtn.title = 'Remove';
            (function(fn) {
                xBtn.onclick = function() { removeAgeRefPhoto(personName, fn); };
            })(p.filename);
            imgDiv.appendChild(xBtn);

            wrap.appendChild(imgDiv);

            // Replace label
            var replaceId = 'age-replace-' + p.filename.replace(/[^a-zA-Z0-9]/g, '-');
            var repLabel = document.createElement('label');
            repLabel.style.cssText = 'font-size:.6em; color:#3182ce; cursor:pointer; padding:1px 4px; background:#ebf8ff; border-radius:3px;';
            repLabel.textContent = 'Replace';
            repLabel.setAttribute('for', replaceId);
            var repInput = document.createElement('input');
            repInput.type = 'file';
            repInput.accept = 'image/*';
            repInput.id = replaceId;
            repInput.style.display = 'none';
            (function(fn) {
                repInput.onchange = function() { replaceAgeRefPhoto(personName, fn, this.files); };
            })(p.filename);
            wrap.appendChild(repLabel);
            wrap.appendChild(repInput);

            thumbs.appendChild(wrap);
        });
    } catch(e) {
        status.textContent = 'Error loading photos.';
    }
}

async function removeAgeRefPhoto(personName, filename) {
    await fetch('/api/ref-faces/' + encodeURIComponent(personName) + '/photo/' + encodeURIComponent(filename), { method: 'DELETE' });
    ageRefVerified = false;
    loadAgeRefThumbs(personName);
}

async function replaceAgeRefPhoto(personName, oldFilename, files) {
    if (!files || !files.length) return;
    // Delete old, upload new
    await fetch('/api/ref-faces/' + encodeURIComponent(personName) + '/photo/' + encodeURIComponent(oldFilename), { method: 'DELETE' });
    var fd = new FormData();
    fd.append('person', personName);
    fd.append('photos', files[0]);
    await fetch('/api/ref-faces/upload', { method: 'POST', body: fd });
    ageRefVerified = false;
    loadAgeRefThumbs(personName);
}

async function verifyAgeRefFaces() {
    var personName = document.getElementById('age-ref-person-name').value.trim();
    if (!personName) { alert('Enter person name first.'); return; }

    var resultDiv = document.getElementById('age-verify-result');
    resultDiv.style.display = 'block';
    resultDiv.style.background = '#f7fafc';
    resultDiv.style.color = '#4a5568';
    resultDiv.textContent = 'Verifying faces...';

    try {
        var res = await fetch('/api/ref-faces/verify', {
            method: 'POST',
            headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({person: personName})
        });
        var data = await res.json();
        var personResult = (data.persons && data.persons.length) ? data.persons.find(function(p) { return p.person === personName; }) : null;

        // Apply green/red borders to thumbnails based on per-photo results
        if (personResult && personResult.photos) {
            personResult.photos.forEach(function(ph) {
                var thumbs = document.getElementById('age-ref-thumbs');
                if (!thumbs) return;
                var wraps = thumbs.children;
                for (var i = 0; i < wraps.length; i++) {
                    var imgEl = wraps[i].querySelector('img');
                    var placeholderEl = wraps[i].querySelector('div[style*="60px"]');
                    var target = imgEl || placeholderEl;
                    if (!target) continue;
                    // Match by data-filename attribute
                    if (wraps[i].getAttribute('data-filename') === ph.filename) {
                        var isOk = (ph.status === 'ok' || ph.status === 'ok_multi');
                        target.style.border = '3px solid ' + (isOk ? '#38a169' : '#e53e3e');
                        break;
                    }
                }
            });
        }

        if (data.ready || (personResult && personResult.ready)) {
            resultDiv.style.background = '#f0fff4';
            resultDiv.style.color = '#276749';
            resultDiv.innerHTML = '<strong>Faces verified successfully.</strong> Ready to run age assessment.';
            ageRefVerified = true;
        } else if (personResult) {
            var msg = personResult.issues && personResult.issues.length ? personResult.issues.join('. ') : 'Some photos may not have detectable faces. Try replacing them.';
            resultDiv.style.background = personResult.ok_count > 0 ? '#fffbeb' : '#fff5f5';
            resultDiv.style.color = personResult.ok_count > 0 ? '#92400e' : '#c53030';
            resultDiv.innerHTML = '<strong>' + personResult.ok_count + '/' + personResult.total_photos + ' photos verified.</strong> ' + esc(msg);
            ageRefVerified = false;
        } else {
            resultDiv.style.background = '#fff5f5';
            resultDiv.style.color = '#c53030';
            resultDiv.innerHTML = '<strong>No face data found.</strong> Upload reference photos and try again.';
            ageRefVerified = false;
        }
    } catch(e) {
        resultDiv.style.background = '#fff5f5';
        resultDiv.style.color = '#c53030';
        resultDiv.textContent = 'Verification failed: ' + e.message;
        ageRefVerified = false;
    }
}

function buildAgeFolderList() {
    var container = document.getElementById('age-folder-list');
    container.innerHTML = '';

    // Checkboxes area (above input)
    var cbArea = document.createElement('div');
    cbArea.id = 'age-folder-cbs';
    cbArea.style.cssText = 'margin-bottom:6px;';
    container.appendChild(cbArea);

    // Add source folders from config if available
    if (config && config.sources && config.sources.length) {
        config.sources.forEach(function(src) {
            if (src.path && src.path.trim()) {
                addAgeFolderCheckbox(src.path, false);
            }
        });
    }

    // Manual folder input row
    var addRow = document.createElement('div');
    addRow.style.cssText = 'display:flex; gap:8px; align-items:center;';
    var inp = document.createElement('input');
    inp.type = 'text';
    inp.id = 'age-folder-input';
    inp.placeholder = 'Enter folder path and click Add';
    inp.style.cssText = 'flex:1; padding:6px 10px; border:1px solid #e2e8f0; border-radius:6px; font-size:.85em;';
    var addBtn = document.createElement('button');
    addBtn.className = 'btn btn-secondary';
    addBtn.style.cssText = 'font-size:.8em; padding:6px 12px; margin:0;';
    addBtn.textContent = 'Add Folder';
    addBtn.onclick = function() {
        var val = inp.value.trim();
        if (val) { addAgeFolderCheckbox(val, true); inp.value = ''; }
    };
    addRow.appendChild(inp);
    addRow.appendChild(addBtn);
    container.appendChild(addRow);
}

function addAgeFolderCheckbox(path, checked) {
    var cbArea = document.getElementById('age-folder-cbs');
    var row = document.createElement('div');
    row.style.cssText = 'display:flex; align-items:center; gap:6px; margin-bottom:3px;';
    var cb = document.createElement('input');
    cb.type = 'checkbox';
    cb.className = 'age-folder-cb';
    cb.value = path;
    cb.checked = checked;
    cb.style.cssText = 'margin:0; flex-shrink:0;';
    row.appendChild(cb);
    var txt = document.createElement('span');
    txt.textContent = path;
    txt.style.cssText = 'font-size:.82em; color:#4a5568;';
    row.appendChild(txt);
    cbArea.appendChild(row);
}

async function runAgeAssessment() {
    var folders = [...document.querySelectorAll('.age-folder-cb:checked')].map(function(cb) { return cb.value; });
    if (!folders.length) { alert('Please select at least one folder.'); return; }

    var faceMode = document.querySelector('input[name="age-face-mode"]:checked');
    var mode = faceMode ? faceMode.value : 'all';
    var personName = '';
    if (mode === 'specific') {
        personName = document.getElementById('age-ref-person-name').value.trim();
        if (!personName) { alert('Please enter a person name and upload reference photos.'); return; }
    }

    document.getElementById('btn-run-age').disabled = true;
    document.getElementById('btn-run-age').textContent = 'Running...';
    document.getElementById('btn-stop-age').style.display = '';
    document.getElementById('age-progress').style.display = 'block';
    document.getElementById('age-progress').textContent = 'Starting age assessment...';

    try {
        var res = await fetch('/api/age-assess/start', {
            method: 'POST',
            headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({ folders: folders, face_mode: mode, person_name: personName })
        });
        var data = await res.json();
        if (data.error) {
            document.getElementById('age-progress').textContent = 'Error: ' + data.error;
            document.getElementById('btn-run-age').disabled = false;
            document.getElementById('btn-run-age').textContent = 'Run Age Assessment';
            document.getElementById('btn-stop-age').style.display = 'none';
            return;
        }
        pollAgeProgress();
    } catch(e) {
        document.getElementById('age-progress').textContent = 'Failed to start: ' + e.message;
        document.getElementById('btn-run-age').disabled = false;
        document.getElementById('btn-run-age').textContent = 'Run Age Assessment';
        document.getElementById('btn-stop-age').style.display = 'none';
    }
}

function pollAgeProgress() {
    var interval = setInterval(async function() {
        try {
            var res = await fetch('/api/scan/status');
            var st = await res.json();
            var prog = document.getElementById('age-progress');
            if (st.lines && st.lines.length) {
                prog.textContent = st.lines[st.lines.length - 1];
            } else if (st.progress) {
                prog.textContent = st.progress;
            }
            if (!st.running || st.done || st.error || st.cancelled) {
                clearInterval(interval);
                document.getElementById('btn-run-age').disabled = false;
                document.getElementById('btn-run-age').textContent = 'Run Age Assessment';
                document.getElementById('btn-stop-age').style.display = 'none';
                if (st.error) {
                    prog.textContent = 'Error: ' + st.error;
                } else if (st.cancelled) {
                    prog.textContent = 'Cancelled.';
                }
                loadAgeResults();
            }
        } catch(e) {
            clearInterval(interval);
        }
    }, 1000);
}

async function stopAgeAssessment() {
    await fetch('/api/task/stop', { method: 'POST' });
    document.getElementById('btn-stop-age').style.display = 'none';
}

async function loadAgeResults() {
    try {
        var res = await fetch('/api/age-assess/results');
        ageResults = await res.json();
        if (ageResults.length) {
            document.getElementById('age-results').style.display = 'block';
            // Populate person filter
            var personSet = new Set();
            ageResults.forEach(function(r) {
                if (r.person) personSet.add(r.person);
            });
            var sel = document.getElementById('age-filter-person');
            sel.innerHTML = '<option value="">All</option>';
            [...personSet].sort().forEach(function(p) {
                var opt = document.createElement('option');
                opt.value = p; opt.textContent = p;
                sel.appendChild(opt);
            });
            renderAgeResults();
        }
    } catch(e) {}
}

function renderAgeResults() {
    var sort = document.getElementById('age-sort').value;
    var filterPerson = document.getElementById('age-filter-person').value;

    var filtered = ageResults;
    if (filterPerson) {
        filtered = filtered.filter(function(r) { return r.person === filterPerson; });
    }

    var sorted = filtered.slice();
    if (sort === 'age-asc') sorted.sort(function(a, b) { return a.estimated_age - b.estimated_age; });
    else if (sort === 'age-desc') sorted.sort(function(a, b) { return b.estimated_age - a.estimated_age; });
    else if (sort === 'name') sorted.sort(function(a, b) { return (a.person || '').localeCompare(b.person || ''); });
    else if (sort === 'date') sorted.sort(function(a, b) { return (a.date || '').localeCompare(b.date || ''); });

    document.getElementById('age-result-count').textContent = '(' + sorted.length + ' images)';

    var grid = document.getElementById('age-grid');
    grid.innerHTML = '';
    sorted.forEach(function(r) {
        var card = document.createElement('div');
        card.className = 'age-card';
        var img = document.createElement('img');
        img.src = r.thumb ? 'data:image/jpeg;base64,' + r.thumb : '';
        img.alt = r.filename;
        img.onerror = function() { this.style.background = '#edf2f7'; this.style.minHeight = '80px'; };
        card.appendChild(img);
        var info = document.createElement('div');
        info.className = 'age-info';
        var badge = document.createElement('span');
        badge.className = 'age-badge';
        badge.textContent = 'Age ~' + r.estimated_age;
        info.appendChild(badge);
        var fn = document.createElement('div');
        fn.style.cssText = 'margin-top:4px; font-size:.9em; color:#718096; white-space:nowrap; overflow:hidden; text-overflow:ellipsis;';
        fn.textContent = r.filename;
        fn.title = r.path || r.filename;
        info.appendChild(fn);
        if (r.folder) {
            var fld = document.createElement('div');
            fld.style.cssText = 'font-size:.85em; color:#a0aec0;';
            fld.textContent = r.folder;
            info.appendChild(fld);
        }
        card.appendChild(info);
        grid.appendChild(card);
    });
}

// ── Init ──
document.getElementById('inp-birthday').max = new Date().toISOString().split('T')[0];
loadTemplates().then(function() {
    // On initial load, show project name and check progress
    if (config && (config.project_name || config.event_name)) {
        updateProjectNameBar(config.project_name || config.event_name);
    }
    if (config && config.event_type) {
        fetch('/api/stats').then(r => r.json()).then(function(st) {
            var dots = document.querySelectorAll('.step-dot');
            // Determine highest completed step based on what exists
            var highest = 0;
            if (config.event_type) highest = 1;  // Event done
            if (config.categories && config.categories.length) highest = 2;  // Categories done
            if (config.sources && config.sources.length) highest = 3;  // Sources done
            if (config.face_names && config.face_names.length) highest = 4;  // Faces done
            if (st.has_scan && st.total_media > 0) highest = 4;  // Scan at least started
            if (st.has_scan && st.total_media > 0 && !st.partial) highest = 5;  // Scan complete
            if (highest >= 5) highest = 6;  // Analyze accessible once scan done
            if (st.selected > 0) highest = 9;  // All steps unlocked when images are selected
            for (var si = 0; si < highest && si < dots.length; si++) {
                dots[si].classList.add('done');
            }
        });
    }
});
// ── Folder Picker ──
var _fpCallback = null;   // function(selectedPath) called on confirm
var _fpCurrent = '';       // path currently being browsed
var _fpSelected = '';      // path user has clicked to highlight
var _fpChecked = {};       // multi-select: { path: true }

function fpOpen(callback, startPath) {
    _fpCallback = callback;
    _fpSelected = '';
    _fpChecked = {};
    document.getElementById('fp-selected-path').textContent = '';
    document.getElementById('fp-select-btn').disabled = true;
    document.getElementById('fp-checked-count').textContent = '';
    document.getElementById('fp-overlay').classList.add('active');
    fpNavigate(startPath || '');
}

// Close on Escape key
document.addEventListener('keydown', function(e) {
    if (e.key === 'Escape' && document.getElementById('fp-overlay').classList.contains('active')) {
        fpClose();
    }
});

function fpClose() {
    document.getElementById('fp-overlay').classList.remove('active');
    _fpCallback = null;
}

function fpConfirm() {
    var checked = Object.keys(_fpChecked);
    if (checked.length > 0 && _fpCallback) {
        _fpCallback(checked);
    } else if (_fpSelected && _fpCallback) {
        _fpCallback([_fpSelected]);
    }
    fpClose();
}

function fpUpdateCheckedCount() {
    var n = Object.keys(_fpChecked).length;
    var el = document.getElementById('fp-checked-count');
    el.textContent = n > 0 ? n + ' folder' + (n > 1 ? 's' : '') + ' selected' : '';
    // Enable select button if anything is checked
    var btn = document.getElementById('fp-select-btn');
    if (n > 0) btn.disabled = false;
    else if (!_fpSelected) btn.disabled = true;
}

// --- Recent folders (localStorage) ---
var _fpRecentKey = 'fp_recent_folders';

function fpAddRecent(paths) {
    if (!paths || paths.length === 0) return;
    try {
        var recent = fpGetRecent();
        paths.forEach(function(p) {
            // Remove if already present, then prepend
            recent = recent.filter(function(r) { return r !== p; });
            recent.unshift(p);
        });
        // Keep max 10
        recent = recent.slice(0, 10);
        localStorage.setItem(_fpRecentKey, JSON.stringify(recent));
    } catch(e) { /* localStorage unavailable */ }
}

function fpGetRecent() {
    try {
        var raw = localStorage.getItem(_fpRecentKey);
        if (raw) return JSON.parse(raw);
    } catch(e) {}
    return [];
}

async function fpNavigate(path) {
    _fpCurrent = path;
    _fpSelected = '';
    document.getElementById('fp-selected-path').textContent = '';
    // Keep select enabled if multi-select has items
    if (Object.keys(_fpChecked).length === 0) {
        document.getElementById('fp-select-btn').disabled = true;
    }
    var list = document.getElementById('fp-list');
    list.innerHTML = '<div class="fp-empty">Loading...</div>';

    // Show/hide "Use this folder" bar
    var useCurrent = document.getElementById('fp-use-current');
    if (path) {
        useCurrent.style.display = 'block';
        document.getElementById('fp-use-current-path').textContent = path;
    } else {
        useCurrent.style.display = 'none';
    }

    try {
        var url = '/api/browse';
        if (path) url += '?path=' + encodeURIComponent(path);
        var res = await fetch(url);
        var data = await res.json();
        if (data.error && !data.folders) {
            list.innerHTML = '<div class="fp-empty">' + data.error + '</div>';
            return;
        }
        fpRenderBreadcrumb(data.path, data.parent);
        fpRenderList(data.folders, data.path, data.error);
        // Show recent folders at drives root
        if (!path) fpRenderRecent(list);
    } catch (e) {
        list.innerHTML = '<div class="fp-empty">Failed to load folders</div>';
    }
}

function fpRenderRecent(listEl) {
    var recent = fpGetRecent();
    if (recent.length === 0) return;
    var section = document.createElement('div');
    section.className = 'fp-recent';
    var hdr = document.createElement('div');
    hdr.className = 'fp-recent-title';
    hdr.textContent = 'Recent Folders';
    section.appendChild(hdr);
    recent.forEach(function(p) {
        var row = document.createElement('div');
        row.className = 'fp-recent-item';
        var icon = document.createElement('span');
        icon.className = 'fp-icon';
        icon.textContent = String.fromCodePoint(0x1F552);
        row.appendChild(icon);
        var name = document.createElement('span');
        name.textContent = p;
        name.style.flex = '1';
        name.style.overflow = 'hidden';
        name.style.textOverflow = 'ellipsis';
        name.style.whiteSpace = 'nowrap';
        name.style.fontSize = '.85em';
        row.appendChild(name);
        row.onclick = function() { fpNavigate(p); };
        section.appendChild(row);
    });
    // Insert before the folder list content
    listEl.insertBefore(section, listEl.firstChild);
}

function fpRenderBreadcrumb(currentPath, parentPath) {
    var bc = document.getElementById('fp-breadcrumb');
    bc.innerHTML = '';
    if (!currentPath) {
        bc.textContent = 'My Computer';
        return;
    }
    // Root link
    var root = document.createElement('span');
    root.textContent = 'Drives';
    root.onclick = function() { fpNavigate(''); };
    bc.appendChild(root);
    // Split path into segments
    var sep = currentPath.indexOf('/') >= 0 && currentPath.indexOf('\\\\') < 0 ? '/' : '\\\\';
    var parts = currentPath.split(/[/\\\\]/).filter(Boolean);
    var built = '';
    parts.forEach(function(part, idx) {
        var s = document.createElement('span');
        s.className = 'sep';
        s.textContent = ' > ';
        bc.appendChild(s);
        built += part + sep;
        var link = document.createElement('span');
        link.textContent = part;
        if (idx < parts.length - 1) {
            (function(p) { link.onclick = function() { fpNavigate(p); }; })(built);
        } else {
            link.style.color = '#2d3748';
            link.style.cursor = 'default';
        }
        bc.appendChild(link);
    });
}

function fpRenderList(folders, currentPath, errMsg) {
    var list = document.getElementById('fp-list');
    list.innerHTML = '';
    if (errMsg) {
        var warn = document.createElement('div');
        warn.className = 'fp-empty';
        warn.textContent = errMsg;
        list.appendChild(warn);
    }
    if (!folders || folders.length === 0) {
        if (!errMsg) {
            var empty = document.createElement('div');
            empty.className = 'fp-empty';
            empty.textContent = currentPath ? 'No subfolders' : 'No drives found';
            list.appendChild(empty);
        }
        return;
    }
    folders.forEach(function(f) {
        var row = document.createElement('div');
        row.className = 'fp-item';
        row.setAttribute('data-path', f.path);
        // Checkbox for multi-select
        var cb = document.createElement('input');
        cb.type = 'checkbox';
        cb.checked = !!_fpChecked[f.path];
        cb.onclick = function(e) { e.stopPropagation(); };
        cb.onchange = function(e) {
            e.stopPropagation();
            if (this.checked) { _fpChecked[f.path] = true; }
            else { delete _fpChecked[f.path]; }
            fpUpdateCheckedCount();
        };
        row.appendChild(cb);
        var icon = document.createElement('span');
        icon.className = 'fp-icon';
        icon.textContent = String.fromCodePoint(0x1F4C1);
        row.appendChild(icon);
        var name = document.createElement('span');
        name.textContent = f.name;
        name.style.flex = '1';
        name.style.overflow = 'hidden';
        name.style.textOverflow = 'ellipsis';
        name.style.whiteSpace = 'nowrap';
        row.appendChild(name);
        if (f.has_children) {
            var arrow = document.createElement('span');
            arrow.className = 'fp-arrow';
            arrow.textContent = '>';
            row.appendChild(arrow);
        }
        // Single click: select this folder (highlight)
        row.onclick = function(e) {
            e.stopPropagation();
            list.querySelectorAll('.fp-item').forEach(function(el) { el.classList.remove('selected'); });
            row.classList.add('selected');
            _fpSelected = f.path;
            document.getElementById('fp-selected-path').textContent = f.path;
            document.getElementById('fp-select-btn').disabled = false;
        };
        // Double click: navigate into
        row.ondblclick = function(e) {
            e.stopPropagation();
            fpNavigate(f.path);
        };
        list.appendChild(row);
    });
}

// Also allow selecting the current directory itself (the one we browsed into)
// by clicking "Select" without clicking a subfolder — means "use this folder"
// This is handled by making the breadcrumb's current folder selectable:
function fpSelectCurrent() {
    if (_fpCurrent) {
        _fpSelected = _fpCurrent;
        document.getElementById('fp-selected-path').textContent = _fpCurrent;
        document.getElementById('fp-select-btn').disabled = false;
        document.getElementById('fp-list').querySelectorAll('.fp-item').forEach(function(el) {
            el.classList.remove('selected');
        });
    }
}

loadGreeting();
</script>

<!-- Folder Picker Modal -->
<div class="fp-overlay" id="fp-overlay">
  <div class="fp-modal">
    <div class="fp-header">
      <h3>Select Folder</h3>
      <span style="cursor:pointer;font-size:1.3em;color:#a0aec0;" onclick="fpClose()">&times;</span>
    </div>
    <div class="fp-breadcrumb" id="fp-breadcrumb"></div>
    <div id="fp-use-current" style="display:none; padding:6px 18px; background:#ebf8ff; border-bottom:1px solid #bee3f8; font-size:.82em;">
      <span style="color:#2b6cb0; cursor:pointer; text-decoration:underline;" onclick="fpSelectCurrent()">Use this folder</span>
      <span id="fp-use-current-path" style="color:#718096; margin-left:6px;"></span>
    </div>
    <div class="fp-list" id="fp-list"></div>
    <div class="fp-footer">
      <div class="fp-path" id="fp-selected-path"></div>
      <span class="fp-checked-count" id="fp-checked-count"></span>
      <button class="fp-cancel" onclick="fpClose()">Cancel</button>
      <button class="fp-select" id="fp-select-btn" disabled onclick="fpConfirm()">Select</button>
    </div>
  </div>
</div>

</body>
</html>"""

# ── Main ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="E-z Photo Organizer Web UI")
    parser.add_argument("--port", type=int, default=5050)
    parser.add_argument("--no-open", action="store_true")
    args = parser.parse_args()

    url = f"http://localhost:{args.port}"
    print(f"E-z Photo Organizer running at {url}")
    _set_active_project(load_config())

    if not args.no_open:
        threading.Timer(1.5, lambda: webbrowser.open(url)).start()

    app.run(host="127.0.0.1", port=args.port, debug=False, threaded=True)
